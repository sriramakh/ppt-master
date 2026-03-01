"""Text rendering — titles, body text, bullets with proper formatting."""

from __future__ import annotations

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Emu

from pptmaster.assets.color_utils import hex_to_rgb
from pptmaster.models import DesignProfile

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def set_title(slide, text: str, profile: DesignProfile) -> None:
    """Set the slide title with template font styling."""
    if not slide.shapes.title:
        return

    title_shape = slide.shapes.title
    title_shape.text = ""
    tf = title_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text

    font = run.font
    font.name = profile.fonts.major
    font.size = Pt(profile.title_style.font_size_pt)
    font.bold = True
    r, g, b = hex_to_rgb(profile.colors.dk2)
    font.color.rgb = RGBColor(r, g, b)


def set_body_text(placeholder, text: str, profile: DesignProfile) -> None:
    """Set body text in a placeholder with template styling."""
    placeholder.text = ""
    tf = placeholder.text_frame
    tf.clear()
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text

    font = run.font
    font.name = profile.fonts.minor
    font.size = Pt(profile.body_style.font_size_pt)
    r, g, b = hex_to_rgb(profile.colors.dk2)
    font.color.rgb = RGBColor(r, g, b)


def set_bullets(
    placeholder,
    bullets: list[str],
    profile: DesignProfile,
    font_size: float | None = None,
    bullet_color: str = "",
) -> None:
    """Set bulleted text in a placeholder with colored bullet markers."""
    if not bullets:
        return

    placeholder.text = ""
    tf = placeholder.text_frame
    tf.clear()
    tf.word_wrap = True

    b_color = bullet_color or profile.colors.accent1
    accent_r, accent_g, accent_b = hex_to_rgb(b_color)
    text_r, text_g, text_b = hex_to_rgb(profile.colors.dk2)
    size = font_size or profile.body_style.font_size_pt

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(3)

        _set_bullet_color(p, accent_r, accent_g, accent_b)
        _set_bullet_char(p, "\u2022")  # Standard bullet

        run = p.add_run()
        run.text = bullet_text
        font = run.font
        font.name = profile.fonts.minor
        font.size = Pt(size)
        font.color.rgb = RGBColor(text_r, text_g, text_b)


def set_metric_text(
    slide,
    value: str,
    label: str,
    left: int,
    top: int,
    width: int,
    height: int,
    profile: DesignProfile,
    color: str = "",
) -> None:
    """Add a key metric (large number + small label) as a text box."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    # Anchor text vertically in middle
    from pptx.enum.text import MSO_ANCHOR
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    try:
        txBody = txBox._element.find(f"{{{_NS_A}}}txBody")
        if txBody is not None:
            bodyPr = txBody.find(f"{{{_NS_A}}}bodyPr")
            if bodyPr is not None:
                bodyPr.set("anchor", "ctr")
    except Exception:
        pass

    # Value line (large, centered)
    p_val = tf.paragraphs[0]
    p_val.alignment = PP_ALIGN.CENTER
    p_val.space_after = Pt(6)
    run_val = p_val.add_run()
    run_val.text = value
    font_val = run_val.font
    font_val.name = profile.fonts.major
    font_val.size = Pt(64)
    font_val.bold = True
    if color:
        r, g, b = hex_to_rgb(color)
    else:
        r, g, b = hex_to_rgb(profile.colors.accent1)
    font_val.color.rgb = RGBColor(r, g, b)

    # Label line (medium, centered, muted color)
    p_label = tf.add_paragraph()
    p_label.alignment = PP_ALIGN.CENTER
    run_label = p_label.add_run()
    run_label.text = label
    font_label = run_label.font
    font_label.name = profile.fonts.minor
    font_label.size = Pt(18)
    text_r, text_g, text_b = hex_to_rgb(profile.colors.dk2)
    font_label.color.rgb = RGBColor(text_r, text_g, text_b)


def set_subtitle_text(
    slide,
    text: str,
    profile: DesignProfile,
    left: int,
    top: int,
    width: int,
    height: int,
) -> None:
    """Add a subtitle/caption text box."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = profile.fonts.minor
    run.font.size = Pt(16)
    r, g, b = hex_to_rgb(profile.colors.dk2)
    run.font.color.rgb = RGBColor(r, g, b)


# ── XML helpers ───────────────────────────────────────────────────────


def _set_bullet_color(paragraph, r: int, g: int, b: int) -> None:
    """Set bullet color via XML manipulation."""
    pPr = paragraph._p.find(f"{{{_NS_A}}}pPr")
    if pPr is None:
        pPr = etree.SubElement(paragraph._p, f"{{{_NS_A}}}pPr")

    # Remove existing bullet color
    for old in pPr.findall(f"{{{_NS_A}}}buClr"):
        pPr.remove(old)

    buClr = etree.SubElement(pPr, f"{{{_NS_A}}}buClr")
    srgbClr = etree.SubElement(buClr, f"{{{_NS_A}}}srgbClr")
    srgbClr.set("val", f"{r:02X}{g:02X}{b:02X}")


def _set_bullet_char(paragraph, char: str) -> None:
    """Set bullet character via XML."""
    pPr = paragraph._p.find(f"{{{_NS_A}}}pPr")
    if pPr is None:
        pPr = etree.SubElement(paragraph._p, f"{{{_NS_A}}}pPr")

    # Remove existing bullet char/font
    for tag in ("buChar", "buAutoNum", "buNone"):
        for old in pPr.findall(f"{{{_NS_A}}}{tag}"):
            pPr.remove(old)

    buChar = etree.SubElement(pPr, f"{{{_NS_A}}}buChar")
    buChar.set("char", char)
