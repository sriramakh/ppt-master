"""Slide 6 — Team Leadership: 4 profile cards with photo circle placeholders."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from pptmaster.builder.design_system import SLIDE_W, SLIDE_H, MARGIN, FONT_SLIDE_TITLE, FONT_BODY, FONT_CARD_HEADING, FONT_CAPTION, CONTENT_TOP, col_span
from pptmaster.builder.helpers import add_textbox, add_card, add_circle, add_gold_accent_line, add_line, add_styled_card, add_slide_title, add_dark_bg, add_background, add_rect
from pptmaster.assets.color_utils import tint, shade


# ── Scholarly variant ─────────────────────────────────────────────────
def _scholarly(slide, t) -> None:
    """Faculty & Leadership — simple text list with thin rules between entries."""
    c = t.content
    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, c.get("team_title", "Faculty & Leadership"), theme=t)

    text_color = t.primary
    sub_color = t.secondary
    m = MARGIN

    # Thin top rule
    rule_y = content_y + 50000
    add_line(slide, m, rule_y, SLIDE_W - m, rule_y, color=t.secondary, width=0.5)

    leaders = c.get("team", [
        ("Jane Smith", "CEO", "20+ years experience"),
        ("John Davis", "CFO", "Fortune 100 veteran"),
        ("Sarah Chen", "CTO", "AI/ML pioneer"),
        ("Michael Brown", "COO", "Operations expert"),
    ])

    y = rule_y + 200000
    content_w = SLIDE_W - 2 * m
    for i, (name, title, bio) in enumerate(leaders[:4]):
        add_textbox(slide, name, m + 100000, y, content_w - 200000, 350000,
                    font_size=FONT_CARD_HEADING, bold=True, color=text_color)
        add_textbox(slide, title, m + 100000, y + 350000, content_w // 3, 300000,
                    font_size=14, bold=True, color=t.accent)
        add_textbox(slide, bio, m + 100000 + content_w // 3, y + 350000,
                    content_w * 2 // 3 - 200000, 300000,
                    font_size=14, color=sub_color)
        y += 750000
        if i < len(leaders[:4]) - 1:
            add_line(slide, m + 100000, y, SLIDE_W - m - 100000, y,
                     color=t.secondary, width=0.5)
            y += 200000


# ── Laboratory variant ────────────────────────────────────────────────
def _laboratory(slide, t) -> None:
    """Research Team — grid of dark cards with subtle borders."""
    c = t.content
    add_background(slide, t.primary)
    content_y = add_slide_title(slide, c.get("team_title", "Research Team"), theme=t)

    text_color = "#FFFFFF"
    sub_color = tint(t.primary, 0.6)
    dark_fill = tint(t.primary, 0.15)
    p = t.palette

    leaders = c.get("team", [
        ("Jane Smith", "CEO", "20+ years experience"),
        ("John Davis", "CFO", "Fortune 100 veteran"),
        ("Sarah Chen", "CTO", "AI/ML pioneer"),
        ("Michael Brown", "COO", "Operations expert"),
    ])

    m = MARGIN
    gap = 150000
    n = min(len(leaders), 4)
    cols = min(n, 2)
    rows = (n + cols - 1) // cols
    card_w = (SLIDE_W - 2 * m - (cols - 1) * gap) // cols
    card_h = (SLIDE_H - content_y - 600000 - (rows - 1) * gap) // rows

    for i, (name, title, bio) in enumerate(leaders[:n]):
        row, col = divmod(i, cols)
        cx = m + col * (card_w + gap)
        cy = content_y + 200000 + row * (card_h + gap)
        accent = p[i % len(p)]

        add_rect(slide, cx, cy, card_w, card_h, fill_color=dark_fill,
                 line_color=tint(t.primary, 0.25), line_width=1)
        add_rect(slide, cx, cy, 50000, card_h, fill_color=accent)  # left accent

        add_textbox(slide, name, cx + 150000, cy + 200000, card_w - 250000, 400000,
                    font_size=FONT_CARD_HEADING, bold=True, color=text_color)
        add_textbox(slide, title, cx + 150000, cy + 600000, card_w - 250000, 300000,
                    font_size=13, bold=True, color=accent)
        add_line(slide, cx + 150000, cy + 950000, cx + card_w - 150000, cy + 950000,
                 color=tint(t.primary, 0.25), width=0.5)
        add_textbox(slide, bio, cx + 150000, cy + 1050000, card_w - 250000, 600000,
                    font_size=12, color=sub_color)


# ── Dashboard variant ─────────────────────────────────────────────────
def _dashboard(slide, t) -> None:
    """Header band + compact 3x2 grid of team tiles."""
    c = t.content
    s = t.ux_style
    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("team_title", "Leadership Team"), theme=t)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    # Accent header band
    add_rect(slide, 0, content_y - 200000, SLIDE_W, 400000, fill_color=t.accent)

    leaders = c.get("team", [
        ("Jane Smith", "CEO", "20+ years experience"),
        ("John Davis", "CFO", "Fortune 100 veteran"),
        ("Sarah Chen", "CTO", "AI/ML pioneer"),
        ("Michael Brown", "COO", "Operations expert"),
    ])

    m = int(MARGIN * 0.7)
    grid_y = content_y + 350000
    gap = 120000
    cols = 3
    rows = 2
    card_w = (SLIDE_W - 2 * m - (cols - 1) * gap) // cols
    card_h = (SLIDE_H - grid_y - 400000 - (rows - 1) * gap) // rows
    p = t.palette

    for i, (name, title, bio) in enumerate(leaders[:6]):
        row, col = divmod(i, cols)
        cx = m + col * (card_w + gap)
        cy = grid_y + row * (card_h + gap)
        accent = p[i % len(p)]

        add_styled_card(slide, cx, cy, card_w, card_h, theme=t, accent_color=accent)
        add_textbox(slide, name, cx + 80000, cy + 150000, card_w - 160000, 350000,
                    font_size=14, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, title, cx + 80000, cy + 500000, card_w - 160000, 250000,
                    font_size=11, bold=True, color=t.accent, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, bio, cx + 80000, cy + 800000, card_w - 160000, card_h - 950000,
                    font_size=10, color=sub_color, alignment=PP_ALIGN.CENTER)


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    s = t.ux_style

    # ── Style dispatch ────────────────────────────────────────────────
    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("team_title", "Leadership Team"), theme=t)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    line_color = tint(t.primary, 0.25) if s.dark_mode else t.light_bg
    photo_bg = tint(t.primary, 0.2) if s.dark_mode else t.light_bg

    leaders = c.get("team", [
        ("Jane Smith", "CEO", "20+ years experience"), ("John Davis", "CFO", "Fortune 100 veteran"),
        ("Sarah Chen", "CTO", "AI/ML pioneer"), ("Michael Brown", "COO", "Operations expert"),
    ])

    card_top = content_y + 100000
    card_height = 4200000
    photo_radius = 300000
    gap = int(200000 * s.gap_factor)

    for i, (name, title, bio) in enumerate(leaders[:4]):
        left, width = col_span(4, i, gap=gap)

        add_styled_card(slide, left, card_top, width, card_height, theme=t)
        center_x = left + width // 2
        photo_y = card_top + 500000

        add_circle(slide, center_x, photo_y, photo_radius, fill_color=photo_bg, line_color=t.accent, line_width=2)
        add_textbox(slide, "Photo", center_x - photo_radius, photo_y - 100000, photo_radius * 2, 200000,
                    font_size=10, color=sub_color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, name, left + 50000, photo_y + photo_radius + 200000, width - 100000, 350000,
                    font_size=FONT_CARD_HEADING, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, title, left + 50000, photo_y + photo_radius + 550000, width - 100000, 300000,
                    font_size=12, bold=True, color=t.accent, alignment=PP_ALIGN.CENTER)

        line_y = photo_y + photo_radius + 900000
        add_line(slide, left + width // 4, line_y, left + 3 * width // 4, line_y, color=line_color, width=0.75)
        add_textbox(slide, bio, left + 80000, line_y + 150000, width - 160000, 800000,
                    font_size=11, color=sub_color, alignment=PP_ALIGN.CENTER)
