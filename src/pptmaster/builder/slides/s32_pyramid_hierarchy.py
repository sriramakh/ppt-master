"""Slide 32 -- Pyramid Hierarchy: 14 unique visual variants dispatched by UX style."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_circle, add_line, add_triangle,
    add_funnel_tier, add_gold_accent_line, add_slide_title,
    add_dark_bg, add_background, add_styled_card, add_color_cell,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    p = t.palette
    s = t.ux_style

    dispatch = {
        "layered": _layered,
        "numbered": _numbered,
        "block-pyramid": _block_pyramid,
        "floating-pyramid": _floating_pyramid,
        "dark-pyramid": _dark_pyramid,
        "split-pyramid": _split_pyramid,
        "angular-pyramid": _angular_pyramid,
        "editorial-pyramid": _editorial_pyramid,
        "gradient-pyramid": _gradient_pyramid,
        "retro-pyramid": _retro_pyramid,
        "creative-pyramid": _creative_pyramid,
        "scholarly-pyramid": _scholarly_pyramid,
        "laboratory-pyramid": _laboratory_pyramid,
        "dashboard-pyramid": _dashboard_pyramid,
    }
    builder = dispatch.get(s.pyramid, _layered)
    builder(slide, t, c)


def _get_layers(c):
    """Extract pyramid layers from content dict with sensible defaults."""
    return c.get("pyramid_layers", [
        ("Vision", "Long-term aspirational goal"),
        ("Strategy", "Multi-year plan to achieve vision"),
        ("Objectives", "Measurable annual targets"),
        ("Tactics", "Quarterly action plans"),
        ("Operations", "Daily execution and processes"),
    ])


# ---------------------------------------------------------------------------
# Geometry helper -- compute centered tier positions for a pyramid.
#
# The pyramid is widest at the bottom and narrowest at the top.
# Each tier is a horizontal band.  We return a list of
# (left, top, width, height) tuples for *n* tiers, stacked vertically
# inside the given bounding box with the narrowest at the top.
# ---------------------------------------------------------------------------

def _pyramid_tiers(
    n: int,
    *,
    area_left: int = MARGIN,
    area_top: int = CONTENT_TOP,
    area_w: int = SLIDE_W - 2 * MARGIN,
    area_h: int = FOOTER_TOP - CONTENT_TOP - 300000,
    gap: int = 40000,
    min_width_frac: float = 0.20,
) -> list[tuple[int, int, int, int]]:
    """Return (left, top, width, height) for each pyramid tier, top to bottom.

    Tier 0 is the narrowest (top).  Tier n-1 is the widest (bottom).
    All tiers are horizontally centred within *area_left..area_left+area_w*.
    """
    tier_h = (area_h - gap * (n - 1)) // n
    center_x = area_left + area_w // 2
    tiers = []
    for i in range(n):
        # fraction goes from min_width_frac (top) to 1.0 (bottom)
        frac = min_width_frac + (1.0 - min_width_frac) * i / max(n - 1, 1)
        w = int(area_w * frac)
        left = center_x - w // 2
        top = area_top + i * (tier_h + gap)
        tiers.append((left, top, w, tier_h))
    return tiers


# =========================================================================
# Variant 1: CLASSIC -- Stacked trapezoid layers (via add_funnel_tier),
# widest at bottom, each tier a different palette color.
# =========================================================================

def _layered(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy"), theme=t)
    layers = _get_layers(c)
    p = t.palette
    s = t.ux_style
    n = len(layers)

    text_color = "#FFFFFF"  # always on colored fill

    area_h = FOOTER_TOP - content_y - 300000
    tier_h = (area_h - 30000 * (n - 1)) // n
    area_w = SLIDE_W - 2 * MARGIN
    center_x = SLIDE_W // 2

    for i, (label, desc) in enumerate(layers):
        # Top tier (i=0) is narrowest; bottom tier is widest
        frac = 0.22 + 0.78 * i / max(n - 1, 1)
        w = int(area_w * frac)
        left = center_x - w // 2
        top = content_y + i * (tier_h + 30000)
        accent = p[i % len(p)]

        add_funnel_tier(slide, left, top, w, tier_h,
                        fill_color=accent, text=label,
                        text_color=text_color, font_size=14)

        # Description to the right
        desc_left = center_x + int(area_w * 0.5 * frac) + 100000
        desc_w = SLIDE_W - MARGIN - desc_left
        if desc_w > 400000:
            sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary
            add_textbox(slide, desc, desc_left, top + tier_h // 4,
                        desc_w, tier_h // 2,
                        font_size=11, color=sub_color,
                        alignment=PP_ALIGN.LEFT)


# =========================================================================
# Variant 2: MINIMAL -- Clean numbered tiers with thin horizontal rules.
# =========================================================================

def _numbered(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy"), theme=t)
    layers = _get_layers(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(layers)
    available_h = FOOTER_TOP - content_y - 400000
    row_h = available_h // n
    content_w = SLIDE_W - 2 * m

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    for i, (label, desc) in enumerate(layers):
        row_top = content_y + 200000 + i * row_h
        # Indent: narrower for earlier (top) rows
        frac = 0.20 + 0.80 * i / max(n - 1, 1)
        usable_w = int(content_w * frac)
        row_left = SLIDE_W // 2 - usable_w // 2

        # Thin horizontal rule
        add_line(slide, row_left, row_top, row_left + usable_w, row_top,
                 color=tint(t.secondary, 0.8), width=0.5)

        # Short accent tick at left
        add_line(slide, row_left, row_top, row_left, row_top + 120000,
                 color=t.accent, width=1.5)

        # Number
        add_textbox(slide, str(i + 1), row_left + 30000, row_top + 30000,
                    350000, 350000,
                    font_size=20, bold=True, color=t.accent,
                    alignment=PP_ALIGN.LEFT)

        # Label
        add_textbox(slide, label, row_left + 400000, row_top + 30000,
                    usable_w // 2, 350000,
                    font_size=14, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, row_left + 400000 + usable_w // 2, row_top + 30000,
                    usable_w // 2 - 400000, 350000,
                    font_size=11, color=sub_color,
                    alignment=PP_ALIGN.LEFT)

    # Bottom rule
    bottom_y = content_y + 200000 + n * row_h
    add_line(slide, m, bottom_y, SLIDE_W - m, bottom_y,
             color=tint(t.secondary, 0.8), width=0.5)


# =========================================================================
# Variant 3: BOLD -- Large colored blocks in pyramid shape, uppercase text.
# =========================================================================

def _block_pyramid(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy").upper(), theme=t)
    layers = _get_layers(c)
    p = t.palette
    n = len(layers)

    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - content_y - 300000
    tier_h = (area_h - 50000 * (n - 1)) // n
    center_x = SLIDE_W // 2

    for i, (label, desc) in enumerate(layers):
        frac = 0.25 + 0.75 * i / max(n - 1, 1)
        w = int(area_w * frac)
        left = center_x - w // 2
        top = content_y + i * (tier_h + 50000)
        accent = p[i % len(p)]

        # Bold colored block
        add_rect(slide, left, top, w, tier_h, fill_color=accent)

        # Large number flush-left inside block
        add_textbox(slide, f"{i + 1:02d}", left + 60000, top + 20000,
                    400000, tier_h - 40000,
                    font_size=28, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # White divider line between number and text
        add_rect(slide, left + 500000, top + tier_h // 4,
                 30000, tier_h // 2, fill_color="#FFFFFF")

        # Label (uppercase)
        add_textbox(slide, label.upper(), left + 600000, top + 20000,
                    w - 700000, tier_h // 2,
                    font_size=15, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.BOTTOM)

        # Description
        add_textbox(slide, desc, left + 600000, top + tier_h // 2,
                    w - 700000, tier_h // 2 - 20000,
                    font_size=11, color=tint(accent, 0.7),
                    alignment=PP_ALIGN.LEFT)


# =========================================================================
# Variant 4: ELEVATED -- Floating card tiers with shadow.
# =========================================================================

def _floating_pyramid(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy"), theme=t)
    layers = _get_layers(c)
    p = t.palette
    s = t.ux_style
    n = len(layers)

    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - content_y - 300000
    tier_h = (area_h - 60000 * (n - 1)) // n
    center_x = SLIDE_W // 2

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    for i, (label, desc) in enumerate(layers):
        frac = 0.25 + 0.75 * i / max(n - 1, 1)
        w = int(area_w * frac)
        left = center_x - w // 2
        top = content_y + i * (tier_h + 60000)
        accent = p[i % len(p)]

        # Floating card with shadow
        add_styled_card(slide, left, top, w, tier_h,
                        theme=t, accent_color=accent)

        # Small colored circle for tier number
        circle_r = 120000
        add_circle(slide, left + 200000, top + tier_h // 2, circle_r,
                   fill_color=accent)
        add_textbox(slide, str(i + 1),
                    left + 200000 - circle_r, top + tier_h // 2 - circle_r,
                    circle_r * 2, circle_r * 2,
                    font_size=14, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Label
        add_textbox(slide, label, left + 400000, top + 40000,
                    w - 500000, tier_h // 2,
                    font_size=14, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.BOTTOM)

        # Description
        add_textbox(slide, desc, left + 400000, top + tier_h // 2 + 10000,
                    w - 500000, tier_h // 2 - 30000,
                    font_size=11, color=sub_color,
                    alignment=PP_ALIGN.LEFT)


# =========================================================================
# Variant 5: DARK -- Dark bg with glowing colored tiers.
# =========================================================================

def _dark_pyramid(slide, t, c):
    add_background(slide, t.primary)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy").upper(), theme=t)
    layers = _get_layers(c)
    p = t.palette
    n = len(layers)

    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - content_y - 300000
    tier_h = (area_h - 40000 * (n - 1)) // n
    center_x = SLIDE_W // 2

    for i, (label, desc) in enumerate(layers):
        frac = 0.22 + 0.78 * i / max(n - 1, 1)
        w = int(area_w * frac)
        left = center_x - w // 2
        top = content_y + i * (tier_h + 40000)
        accent = p[i % len(p)]

        # Outer glow rectangle (wider, slightly taller)
        glow_pad = 30000
        add_rect(slide, left - glow_pad, top - glow_pad,
                 w + glow_pad * 2, tier_h + glow_pad * 2,
                 fill_color=tint(accent, 0.85))

        # Inner rectangle
        add_rect(slide, left, top, w, tier_h, fill_color=accent)

        # Label centered
        add_textbox(slide, label.upper(), left + 30000, top + 10000,
                    w - 60000, tier_h // 2,
                    font_size=14, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.BOTTOM)

        # Description centered below label
        add_textbox(slide, desc, left + 30000, top + tier_h // 2 + 5000,
                    w - 60000, tier_h // 2 - 15000,
                    font_size=11, color=tint(accent, 0.65),
                    alignment=PP_ALIGN.CENTER)


# =========================================================================
# Variant 6: SPLIT -- Pyramid on left, detail list on right.
# =========================================================================

def _split_pyramid(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy"), theme=t)
    layers = _get_layers(c)
    p = t.palette
    s = t.ux_style
    n = len(layers)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    # Left half: pyramid visual
    half_w = (SLIDE_W - 2 * MARGIN) // 2 - 100000
    pyr_left = MARGIN
    area_h = FOOTER_TOP - content_y - 300000
    tier_h = (area_h - 30000 * (n - 1)) // n
    pyr_center = pyr_left + half_w // 2

    for i, (label, _desc) in enumerate(layers):
        frac = 0.30 + 0.70 * i / max(n - 1, 1)
        w = int(half_w * frac)
        left = pyr_center - w // 2
        top = content_y + i * (tier_h + 30000)
        accent = p[i % len(p)]

        add_funnel_tier(slide, left, top, w, tier_h,
                        fill_color=accent, text=label,
                        text_color="#FFFFFF", font_size=12)

    # Vertical divider
    div_x = MARGIN + half_w + 100000
    add_line(slide, div_x, content_y, div_x, FOOTER_TOP - 200000,
             color=tint(t.accent, 0.5), width=1)

    # Right half: numbered detail list
    detail_left = div_x + 150000
    detail_w = SLIDE_W - MARGIN - detail_left
    row_h = area_h // n

    for i, (label, desc) in enumerate(layers):
        accent = p[i % len(p)]
        row_top = content_y + i * row_h

        # Accent dot
        add_circle(slide, detail_left + 60000, row_top + row_h // 2, 40000,
                   fill_color=accent)

        # Label
        add_textbox(slide, f"{i + 1}. {label}",
                    detail_left + 150000, row_top + 30000,
                    detail_w - 150000, row_h // 2,
                    font_size=13, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc,
                    detail_left + 150000, row_top + row_h // 2,
                    detail_w - 150000, row_h // 2 - 30000,
                    font_size=11, color=sub_color,
                    alignment=PP_ALIGN.LEFT)


# =========================================================================
# Variant 7: GEO -- Angular geometric pyramid with borders.
# =========================================================================

def _angular_pyramid(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy"), theme=t)
    layers = _get_layers(c)
    p = t.palette
    s = t.ux_style
    n = len(layers)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - content_y - 300000
    tier_h = (area_h - 50000 * (n - 1)) // n
    center_x = SLIDE_W // 2

    for i, (label, desc) in enumerate(layers):
        frac = 0.20 + 0.80 * i / max(n - 1, 1)
        w = int(area_w * frac)
        left = center_x - w // 2
        top = content_y + i * (tier_h + 50000)
        accent = p[i % len(p)]

        # Angular block with thick border, no fill
        border_color = accent
        fill = tint(accent, 0.15) if s.dark_mode else tint(accent, 0.92)
        add_rect(slide, left, top, w, tier_h,
                 fill_color=fill, line_color=border_color, line_width=2)

        # Small diamond at left edge
        diamond_size = 80000
        from pptmaster.builder.helpers import add_diamond
        add_diamond(slide, left + 200000, top + tier_h // 2, diamond_size,
                    fill_color=accent)

        # Label
        add_textbox(slide, label, left + 380000, top + 20000,
                    w - 500000, tier_h // 2,
                    font_size=14, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.BOTTOM)

        # Description
        add_textbox(slide, desc, left + 380000, top + tier_h // 2 + 5000,
                    w - 500000, tier_h // 2 - 25000,
                    font_size=11, color=sub_color,
                    alignment=PP_ALIGN.LEFT)


# =========================================================================
# Variant 8: EDITORIAL -- Elegant numbered layers with thin accent.
# =========================================================================

def _editorial_pyramid(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy"), theme=t)
    layers = _get_layers(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)
    n = len(layers)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    available_h = FOOTER_TOP - content_y - 400000
    row_h = available_h // n
    content_w = SLIDE_W - 2 * m

    for i, (label, desc) in enumerate(layers):
        row_top = content_y + 200000 + i * row_h
        # Progressive indentation -- narrowest at top
        frac = 0.25 + 0.75 * i / max(n - 1, 1)
        usable_w = int(content_w * frac)
        row_left = m

        # Thin rule
        add_line(slide, row_left, row_top,
                 row_left + usable_w, row_top,
                 color=tint(t.secondary, 0.8), width=0.5)

        # Short accent tick at start of rule
        add_line(slide, row_left, row_top,
                 row_left + 500000, row_top,
                 color=t.accent, width=1.5)

        # Number (large italic)
        add_textbox(slide, f"{i + 1}.", row_left, row_top + 60000,
                    400000, 350000,
                    font_size=22, bold=True, italic=True, color=t.accent,
                    alignment=PP_ALIGN.LEFT)

        # Label
        add_textbox(slide, label, row_left + 450000, row_top + 60000,
                    usable_w * 2 // 5, 350000,
                    font_size=14, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT)

        # Description
        desc_left = row_left + 450000 + usable_w * 2 // 5 + 100000
        add_textbox(slide, desc, desc_left, row_top + 60000,
                    usable_w - (desc_left - row_left), 350000,
                    font_size=12, color=sub_color,
                    alignment=PP_ALIGN.LEFT)

    # Bottom rule
    bottom_y = content_y + 200000 + n * row_h
    add_line(slide, m, bottom_y, SLIDE_W - m, bottom_y,
             color=tint(t.secondary, 0.8), width=0.5)


# =========================================================================
# Variant 9: GRADIENT -- Rounded tiers with gradient-like progressive colors.
# =========================================================================

def _gradient_pyramid(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy"), theme=t)
    layers = _get_layers(c)
    s = t.ux_style
    n = len(layers)

    text_color = "#FFFFFF"  # always on colored fill

    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - content_y - 300000
    tier_h = (area_h - 50000 * (n - 1)) // n
    center_x = SLIDE_W // 2

    for i, (label, desc) in enumerate(layers):
        frac = 0.22 + 0.78 * i / max(n - 1, 1)
        w = int(area_w * frac)
        left = center_x - w // 2
        top = content_y + i * (tier_h + 50000)

        # Progressive tint from accent
        color_factor = i * 0.15
        tier_color = tint(t.accent, color_factor)

        # Outer halo
        add_rect(slide, left - 20000, top - 10000,
                 w + 40000, tier_h + 20000,
                 fill_color=tint(tier_color, 0.80),
                 corner_radius=120000)

        # Rounded tier
        add_rect(slide, left, top, w, tier_h,
                 fill_color=tier_color, corner_radius=100000)

        # Label centred
        add_textbox(slide, label, left + 30000, top + 10000,
                    w - 60000, tier_h // 2,
                    font_size=14, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.BOTTOM)

        # Description centred
        add_textbox(slide, desc, left + 30000, top + tier_h // 2 + 5000,
                    w - 60000, tier_h // 2 - 15000,
                    font_size=11, color=tint(tier_color, 0.65),
                    alignment=PP_ALIGN.CENTER)


# =========================================================================
# Variant 10: RETRO -- Vintage styled with decorative double borders.
# =========================================================================

def _retro_pyramid(slide, t, c):
    add_background(slide, t.light_bg)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy"), theme=t)
    layers = _get_layers(c)
    p = t.palette
    n = len(layers)

    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - content_y - 300000
    tier_h = (area_h - 50000 * (n - 1)) // n
    center_x = SLIDE_W // 2

    for i, (label, desc) in enumerate(layers):
        frac = 0.22 + 0.78 * i / max(n - 1, 1)
        w = int(area_w * frac)
        left = center_x - w // 2
        top = content_y + i * (tier_h + 50000)
        accent = p[i % len(p)]

        # Outer border rectangle
        add_rect(slide, left, top, w, tier_h,
                 fill_color="#FFFFFF",
                 line_color=accent, line_width=2.5,
                 corner_radius=100000)

        # Inner border rectangle (inset)
        inset = 30000
        add_rect(slide, left + inset, top + inset,
                 w - inset * 2, tier_h - inset * 2,
                 line_color=accent, line_width=0.75,
                 corner_radius=80000)

        # Label centred
        add_textbox(slide, label, left + 80000, top + 20000,
                    w - 160000, tier_h // 2,
                    font_size=14, bold=True, color=t.primary,
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.BOTTOM)

        # Description centred
        add_textbox(slide, desc, left + 80000, top + tier_h // 2 + 5000,
                    w - 160000, tier_h // 2 - 25000,
                    font_size=11, color=t.secondary,
                    alignment=PP_ALIGN.CENTER)

        # Decorative dots at the tier bottom-center
        dot_y = top + tier_h - 50000
        for d in range(3):
            add_circle(slide, center_x - 60000 + d * 60000, dot_y,
                       12000, fill_color=accent)


# =========================================================================
# Variant 11: MAGAZINE -- Creative asymmetric pyramid with oversized numbers.
# =========================================================================

def _creative_pyramid(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy"), theme=t)
    layers = _get_layers(c)
    p = t.palette
    s = t.ux_style
    n = len(layers)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - content_y - 300000
    tier_h = (area_h - 40000 * (n - 1)) // n
    m = int(MARGIN * s.margin_factor)

    for i, (label, desc) in enumerate(layers):
        frac = 0.20 + 0.80 * i / max(n - 1, 1)
        w = int(area_w * frac)
        # Asymmetric: left-aligned rather than centred
        left = m
        top = content_y + i * (tier_h + 40000)
        accent = p[i % len(p)]

        # Thin colored bar spanning tier width
        add_rect(slide, left, top, w, tier_h,
                 fill_color=tint(accent, 0.10) if s.dark_mode else tint(accent, 0.92))

        # Accent left edge
        add_rect(slide, left, top, 40000, tier_h, fill_color=accent)

        # Oversized number (ghost, behind text)
        add_textbox(slide, str(i + 1), left + 60000, top - 10000,
                    500000, tier_h + 20000,
                    font_size=42, bold=True,
                    color=tint(accent, 0.60) if s.dark_mode else tint(accent, 0.75),
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Label
        add_textbox(slide, label, left + 500000, top + 15000,
                    w - 550000, tier_h // 2,
                    font_size=15, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.BOTTOM)

        # Description
        add_textbox(slide, desc, left + 500000, top + tier_h // 2 + 5000,
                    w - 550000, tier_h // 2 - 20000,
                    font_size=11, color=sub_color,
                    alignment=PP_ALIGN.LEFT)


# =========================================================================
# Variant 12: SCHOLARLY -- "Figure N:" caption, thin rules, centred,
# numbered sections.  Academic paper aesthetic with generous whitespace.
# =========================================================================

def _scholarly_pyramid(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy"), theme=t)
    layers = _get_layers(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)
    n = len(layers)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary
    rule_color = tint(t.secondary, 0.7)

    content_w = SLIDE_W - 2 * m
    center_x = SLIDE_W // 2

    # "Figure 1: Strategic Hierarchy" caption at top, italic
    caption_text = c.get("pyramid_figure_caption",
                         "Figure 1: Organisational Hierarchy Framework")
    add_textbox(slide, caption_text, m, content_y - 100000,
                content_w, 300000,
                font_size=11, italic=True, color=sub_color,
                alignment=PP_ALIGN.CENTER)

    # Top bounding rule
    fig_top = content_y + 250000
    add_line(slide, m + content_w // 6, fig_top,
             SLIDE_W - m - content_w // 6, fig_top,
             color=rule_color, width=0.75)

    available_h = FOOTER_TOP - fig_top - 600000
    row_h = available_h // n
    section_top = fig_top + 80000

    for i, (label, desc) in enumerate(layers):
        row_top = section_top + i * row_h
        # Progressive width: narrowest at top
        frac = 0.30 + 0.70 * i / max(n - 1, 1)
        usable_w = int(content_w * 0.75 * frac)
        row_left = center_x - usable_w // 2

        # Section number and label on one line, centred
        section_text = f"{i + 1}. {label}"
        add_textbox(slide, section_text, row_left, row_top,
                    usable_w, row_h * 2 // 5,
                    font_size=13, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.BOTTOM)

        # Description centred below
        add_textbox(slide, desc, row_left, row_top + row_h * 2 // 5 + 20000,
                    usable_w, row_h * 2 // 5,
                    font_size=11, color=sub_color,
                    alignment=PP_ALIGN.CENTER)

        # Thin rule after each section (except last)
        if i < n - 1:
            rule_y = row_top + row_h - 20000
            rule_w = int(usable_w * 0.6)
            add_line(slide, center_x - rule_w // 2, rule_y,
                     center_x + rule_w // 2, rule_y,
                     color=rule_color, width=0.5)

    # Bottom bounding rule
    bottom_rule_y = section_top + n * row_h + 30000
    add_line(slide, m + content_w // 6, bottom_rule_y,
             SLIDE_W - m - content_w // 6, bottom_rule_y,
             color=rule_color, width=0.75)

    # Source note at bottom
    add_textbox(slide, "Source: Author's framework", m, bottom_rule_y + 50000,
                content_w, 250000,
                font_size=9, italic=True, color=sub_color,
                alignment=PP_ALIGN.CENTER)


# =========================================================================
# Variant 13: LABORATORY -- Dark bg, grid lines, monospace labels,
# data-coded evidence style.  Looks like a scientific instrument readout.
# =========================================================================

def _laboratory_pyramid(slide, t, c):
    add_background(slide, t.primary)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy"), theme=t)
    layers = _get_layers(c)
    p = t.palette
    n = len(layers)

    m = MARGIN
    area_w = SLIDE_W - 2 * m
    area_h = FOOTER_TOP - content_y - 300000
    center_x = SLIDE_W // 2

    grid_color = tint(t.primary, 0.12)
    label_color = tint(t.primary, 0.60)
    data_color = "#FFFFFF"

    # Draw background grid (horizontal + vertical lines)
    grid_step = 400000
    for gx in range(m, SLIDE_W - m, grid_step):
        add_line(slide, gx, content_y, gx, FOOTER_TOP - 200000,
                 color=grid_color, width=0.25)
    for gy in range(content_y, FOOTER_TOP - 200000, grid_step):
        add_line(slide, m, gy, SLIDE_W - m, gy,
                 color=grid_color, width=0.25)

    # Axis labels along left edge (monospace feel)
    add_textbox(slide, "LEVEL", m + 20000, content_y + 20000,
                600000, 250000,
                font_size=9, bold=True, color=label_color,
                alignment=PP_ALIGN.LEFT)

    tier_h = (area_h - 40000 * (n - 1)) // n

    for i, (label, desc) in enumerate(layers):
        frac = 0.20 + 0.80 * i / max(n - 1, 1)
        w = int(area_w * 0.65 * frac)
        left = center_x - w // 2
        top = content_y + i * (tier_h + 40000)
        accent = p[i % len(p)]

        # Data-coded tier rectangle (colored left border)
        fill = tint(t.primary, 0.08)
        add_rect(slide, left, top, w, tier_h,
                 fill_color=fill, line_color=tint(accent, 0.3), line_width=0.5)

        # Colored left border bar
        add_rect(slide, left, top, 40000, tier_h, fill_color=accent)

        # Level code (e.g., "L1", "L2")
        add_textbox(slide, f"L{i + 1}", left + 60000, top + 10000,
                    300000, tier_h - 20000,
                    font_size=11, bold=True, color=accent,
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Label (monospace feel, uppercase)
        add_textbox(slide, label.upper(), left + 350000, top + 10000,
                    w - 400000, tier_h // 2,
                    font_size=13, bold=True, color=data_color,
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.BOTTOM)

        # Description
        add_textbox(slide, desc, left + 350000, top + tier_h // 2 + 5000,
                    w - 400000, tier_h // 2 - 15000,
                    font_size=10, color=label_color,
                    alignment=PP_ALIGN.LEFT)

        # Right-side evidence box: small metric placeholder
        box_left = center_x + int(area_w * 0.65 * frac) // 2 + 100000
        box_w = SLIDE_W - m - box_left - 50000
        if box_w > 500000:
            add_rect(slide, box_left, top, box_w, tier_h,
                     fill_color=tint(t.primary, 0.06),
                     line_color=tint(accent, 0.2), line_width=0.5)
            add_textbox(slide, f"DATA-{i + 1:02d}", box_left + 30000, top + 10000,
                        box_w - 60000, tier_h // 2,
                        font_size=9, bold=True, color=accent,
                        alignment=PP_ALIGN.LEFT)
            add_textbox(slide, f"n = {(n - i) * 247}", box_left + 30000, top + tier_h // 2,
                        box_w - 60000, tier_h // 2 - 10000,
                        font_size=9, color=label_color,
                        alignment=PP_ALIGN.LEFT)

    # Bottom data timestamp
    add_textbox(slide, "ANALYSIS ID: PYR-001  |  CONFIDENCE: HIGH",
                m, FOOTER_TOP - 280000,
                area_w, 200000,
                font_size=8, color=label_color,
                alignment=PP_ALIGN.LEFT)


# =========================================================================
# Variant 14: DASHBOARD -- Compact pyramid with stats sidebar.
# Dense layout with metric tiles alongside each tier.
# =========================================================================

def _dashboard_pyramid(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("pyramid_title", "Strategic Hierarchy"), theme=t)
    layers = _get_layers(c)
    p = t.palette
    s = t.ux_style
    n = len(layers)

    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary
    card_bg = tint(t.primary, 0.10) if s.dark_mode else "#FFFFFF"

    # Layout: pyramid area (left 55%) + stats sidebar (right 40%), gap 5%
    total_w = SLIDE_W - 2 * m
    pyr_w = int(total_w * 0.55)
    gap = int(total_w * 0.04)
    stats_w = total_w - pyr_w - gap
    stats_left = m + pyr_w + gap

    area_h = FOOTER_TOP - content_y - 300000
    tier_h = (area_h - 30000 * (n - 1)) // n
    pyr_center = m + pyr_w // 2

    # Stats sidebar header
    header_h = 300000
    add_rect(slide, stats_left, content_y, stats_w, header_h,
             fill_color=t.accent)
    add_textbox(slide, "METRICS", stats_left + 50000, content_y + 30000,
                stats_w - 100000, header_h - 60000,
                font_size=11, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Stats tiles and pyramid tiers
    tile_top = content_y + header_h + 30000
    tile_h = (area_h - header_h - 30000 * n) // n

    # Default metric values for the sidebar
    metrics = c.get("pyramid_metrics", [
        ("100%", "Alignment"),
        ("85%", "Execution"),
        ("72%", "Tracking"),
        ("94%", "Activity"),
        ("99%", "Uptime"),
    ])

    for i, (label, desc) in enumerate(layers):
        frac = 0.28 + 0.72 * i / max(n - 1, 1)
        w = int(pyr_w * frac)
        left = pyr_center - w // 2
        top = content_y + i * (tier_h + 30000)
        accent = p[i % len(p)]

        # Compact pyramid tier
        add_rect(slide, left, top, w, tier_h,
                 fill_color=accent, corner_radius=s.card_radius)

        # Label inside tier
        add_textbox(slide, label, left + 30000, top + 5000,
                    w - 60000, tier_h - 10000,
                    font_size=12, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Stats tile on the right
        t_top = tile_top + i * (tile_h + 30000)
        tile_bg = card_bg
        add_rect(slide, stats_left, t_top, stats_w, tile_h,
                 fill_color=tile_bg,
                 line_color=tint(accent, 0.4) if s.dark_mode else tint(accent, 0.7),
                 line_width=0.75,
                 corner_radius=s.card_radius)

        # Colored left accent on tile
        add_rect(slide, stats_left, t_top, 35000, tile_h, fill_color=accent)

        # Metric value
        metric_val, metric_label = (metrics[i] if i < len(metrics)
                                    else (f"{(n - i) * 20}%", desc[:15]))
        add_textbox(slide, metric_val, stats_left + 60000, t_top + 10000,
                    stats_w - 90000, tile_h * 3 // 5,
                    font_size=18, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.BOTTOM)

        # Metric label
        add_textbox(slide, metric_label, stats_left + 60000,
                    t_top + tile_h * 3 // 5 + 5000,
                    stats_w - 90000, tile_h * 2 // 5 - 15000,
                    font_size=9, color=sub_color,
                    alignment=PP_ALIGN.LEFT)
