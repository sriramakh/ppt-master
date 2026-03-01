"""Slide 36 — Kanban Board: 14 unique visual variants dispatched by UX style."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_circle, add_line, add_snip_rect,
    add_gold_accent_line, add_slide_title, add_dark_bg, add_background,
    add_styled_card, add_color_cell,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    s = t.ux_style
    c = t.content

    dispatch = {
        "card-columns": _card_columns,
        "list-minimal": _list_minimal,
        "bold-columns": _bold_columns,
        "floating-cards": _floating_cards,
        "dark-board": _dark_board,
        "split-board": _split_board,
        "geo-board": _geo_board,
        "editorial-board": _editorial_board,
        "gradient-board": _gradient_board,
        "retro-board": _retro_board,
        "creative-board": _creative_board,
        "scholarly-board": _scholarly_board,
        "laboratory-board": _laboratory_board,
        "dashboard-board": _dashboard_board,
    }
    builder = dispatch.get(s.kanban, _card_columns)
    builder(slide, t, c)


def _get_columns(c):
    """Extract kanban column data from content dict with sensible defaults."""
    return c.get("kanban_columns", [
        {"title": "To Do", "cards": ["Define requirements", "Design wireframes", "Set up CI/CD"]},
        {"title": "In Progress", "cards": ["API development", "Frontend build"]},
        {"title": "Done", "cards": ["Project charter", "Team onboarding", "Architecture review"]},
    ])


# ── Variant 1: CLASSIC — 3 columns with card shapes, colored headers ────

def _card_columns(slide, t, c):
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board"), theme=t)
    columns = _get_columns(c)
    p = t.palette

    n = min(len(columns), 4)
    gap = 250000
    total_w = SLIDE_W - 2 * MARGIN
    col_w = (total_w - gap * (n - 1)) // n
    board_top = content_top + 100000
    available_h = FOOTER_TOP - board_top - 200000

    for ci, col_data in enumerate(columns[:n]):
        accent = p[ci % len(p)]
        left = MARGIN + ci * (col_w + gap)

        # Column background
        add_styled_card(slide, left, board_top, col_w, available_h,
                        theme=t, accent_color=accent)

        # Column header
        add_rect(slide, left, board_top, col_w, 400000, fill_color=accent)
        add_textbox(slide, col_data["title"], left + 80000, board_top + 60000,
                    col_w - 160000, 280000,
                    font_size=14, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Card count badge
        count = len(col_data.get("cards", []))
        badge_x = left + col_w - 200000
        add_circle(slide, badge_x, board_top + 200000, 80000,
                   fill_color=shade(accent, 0.2))
        add_textbox(slide, str(count), badge_x - 80000, board_top + 120000,
                    160000, 160000,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Cards
        card_start_y = board_top + 500000
        card_h = 320000
        card_gap = 120000

        for j, card_text in enumerate(col_data.get("cards", [])[:6]):
            card_top = card_start_y + j * (card_h + card_gap)
            if card_top + card_h > board_top + available_h - 80000:
                break

            card_fill = "#FFFFFF" if not t.ux_style.dark_mode else tint(t.primary, 0.2)
            text_color = t.primary if not t.ux_style.dark_mode else "#FFFFFF"

            add_snip_rect(slide, left + 60000, card_top, col_w - 120000, card_h,
                          fill_color=card_fill)
            # Left accent stripe on card
            add_rect(slide, left + 60000, card_top, 40000, card_h, fill_color=accent)
            add_textbox(slide, card_text, left + 140000, card_top + 30000,
                        col_w - 260000, card_h - 60000,
                        font_size=11, color=text_color,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant 2: MINIMAL — Clean list columns with thin borders ────────────

def _list_minimal(slide, t, c):
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board"), theme=t)
    columns = _get_columns(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = min(len(columns), 4)
    total_w = SLIDE_W - 2 * m
    col_w = total_w // n
    board_top = content_top + 200000
    available_h = FOOTER_TOP - board_top - 300000

    for ci, col_data in enumerate(columns[:n]):
        left = m + ci * col_w

        # Thin vertical separator between columns
        if ci > 0:
            add_line(slide, left, board_top, left, board_top + available_h,
                     color=tint(t.secondary, 0.75), width=0.5)

        # Column header — just text with thin underline
        add_textbox(slide, col_data["title"].upper(), left + 60000, board_top,
                    col_w - 120000, 300000,
                    font_size=10, bold=True, color=t.secondary,
                    alignment=PP_ALIGN.LEFT)

        # Short accent line under title
        add_line(slide, left + 60000, board_top + 320000,
                 left + 60000 + 400000, board_top + 320000,
                 color=t.accent, width=1)

        # Card count
        add_textbox(slide, f"({len(col_data.get('cards', []))})",
                    left + col_w - 260000, board_top,
                    200000, 300000,
                    font_size=9, color=tint(t.secondary, 0.5),
                    alignment=PP_ALIGN.RIGHT)

        # List items — no card shapes, just text with thin rule separators
        item_start_y = board_top + 450000
        item_h = 350000

        for j, card_text in enumerate(col_data.get("cards", [])[:7]):
            item_top = item_start_y + j * item_h
            if item_top + item_h > board_top + available_h:
                break

            # Thin rule above each item
            if j > 0:
                add_line(slide, left + 60000, item_top,
                         left + col_w - 60000, item_top,
                         color=tint(t.secondary, 0.85), width=0.25)

            # Bullet marker
            add_textbox(slide, "\u2013", left + 60000, item_top + 40000,
                        80000, 250000,
                        font_size=11, color=t.accent)
            add_textbox(slide, card_text, left + 150000, item_top + 40000,
                        col_w - 220000, 250000,
                        font_size=11, color=t.primary)


# ── Variant 3: BOLD — Bold colored column headers, uppercase, thick accents ─

def _bold_columns(slide, t, c):
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board").upper(), theme=t)
    columns = _get_columns(c)
    p = t.palette

    n = min(len(columns), 4)
    gap = 80000
    total_w = SLIDE_W - 2 * MARGIN
    col_w = (total_w - gap * (n - 1)) // n
    board_top = content_top + 100000
    available_h = FOOTER_TOP - board_top - 200000

    for ci, col_data in enumerate(columns[:n]):
        accent = p[ci % len(p)]
        left = MARGIN + ci * (col_w + gap)

        # Thick left accent bar for entire column
        add_rect(slide, left, board_top, 80000, available_h, fill_color=accent)

        # Column header — bold uppercase on accent bg
        add_rect(slide, left, board_top, col_w, 500000, fill_color=accent)
        add_textbox(slide, col_data["title"].upper(), left + 120000, board_top + 80000,
                    col_w - 240000, 340000,
                    font_size=16, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Thick accent underline below header
        add_rect(slide, left + 120000, board_top + 520000,
                 col_w - 240000, 60000, fill_color=shade(accent, 0.15))

        # Cards
        card_start_y = board_top + 680000
        card_h = 350000
        card_gap = 100000

        for j, card_text in enumerate(col_data.get("cards", [])[:6]):
            card_top = card_start_y + j * (card_h + card_gap)
            if card_top + card_h > board_top + available_h - 60000:
                break

            # Card with thick left border
            card_fill = tint(accent, 0.92)
            add_rect(slide, left + 120000, card_top, col_w - 240000, card_h,
                     fill_color=card_fill)
            add_rect(slide, left + 120000, card_top, 50000, card_h,
                     fill_color=accent)

            add_textbox(slide, card_text.upper(), left + 220000, card_top + 40000,
                        col_w - 400000, card_h - 80000,
                        font_size=11, bold=True, color=t.primary,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant 4: ELEVATED — Floating card shapes with shadows in columns ──

def _floating_cards(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board"), theme=t)
    columns = _get_columns(c)
    p = t.palette

    n = min(len(columns), 4)
    gap = 300000
    total_w = SLIDE_W - 2 * MARGIN
    col_w = (total_w - gap * (n - 1)) // n
    board_top = content_top + 200000
    available_h = FOOTER_TOP - board_top - 300000

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if t.ux_style.dark_mode else t.secondary

    for ci, col_data in enumerate(columns[:n]):
        accent = p[ci % len(p)]
        left = MARGIN + ci * (col_w + gap)

        # Column container card with shadow
        add_styled_card(slide, left, board_top, col_w, available_h,
                        theme=t, accent_color=accent)

        # Rounded header area
        header_h = 380000
        add_rect(slide, left + 40000, board_top + 40000,
                 col_w - 80000, header_h,
                 fill_color=tint(accent, 0.15 if t.ux_style.dark_mode else 0.85),
                 corner_radius=80000)

        # Column title centered in header
        add_textbox(slide, col_data["title"], left + 80000, board_top + 80000,
                    col_w - 160000, 300000,
                    font_size=14, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Floating task cards
        card_start_y = board_top + header_h + 160000
        card_h = 340000
        card_gap = 160000

        for j, card_text in enumerate(col_data.get("cards", [])[:5]):
            card_top = card_start_y + j * (card_h + card_gap)
            if card_top + card_h > board_top + available_h - 100000:
                break

            # Shadow behind card
            shadow_color = shade(t.primary, 0.3) if t.ux_style.dark_mode else shade("#E0E0E0", 0.1)
            add_rect(slide, left + 82000, card_top + 22000,
                     col_w - 160000, card_h,
                     fill_color=shadow_color, corner_radius=100000)

            # Main floating card
            card_fill = tint(t.primary, 0.18) if t.ux_style.dark_mode else "#FFFFFF"
            add_rect(slide, left + 60000, card_top, col_w - 120000, card_h,
                     fill_color=card_fill, corner_radius=100000)

            # Color dot indicator
            add_circle(slide, left + 130000, card_top + card_h // 2, 30000,
                       fill_color=accent)

            add_textbox(slide, card_text, left + 190000, card_top + 40000,
                        col_w - 300000, card_h - 80000,
                        font_size=11, color=text_color,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant 5: DARK — Dark bg, glowing column headers, card borders ──────

def _dark_board(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board").upper(), theme=t)
    columns = _get_columns(c)
    p = t.palette

    n = min(len(columns), 4)
    gap = 200000
    total_w = SLIDE_W - 2 * MARGIN
    col_w = (total_w - gap * (n - 1)) // n
    board_top = content_top + 200000
    available_h = FOOTER_TOP - board_top - 300000

    for ci, col_data in enumerate(columns[:n]):
        accent = p[ci % len(p)]
        left = MARGIN + ci * (col_w + gap)

        # Dark column background with border
        col_fill = tint(t.primary, 0.08)
        add_rect(slide, left, board_top, col_w, available_h,
                 fill_color=col_fill, line_color=tint(t.primary, 0.2),
                 line_width=0.75, corner_radius=60000)

        # Glowing header bar
        add_rect(slide, left + 1, board_top, col_w - 2, 50000,
                 fill_color=accent)

        # Header text
        add_textbox(slide, col_data["title"].upper(), left + 80000, board_top + 100000,
                    col_w - 160000, 350000,
                    font_size=12, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER)

        # Item count
        add_textbox(slide, f"{len(col_data.get('cards', []))} items",
                    left + 80000, board_top + 380000,
                    col_w - 160000, 200000,
                    font_size=9, color=tint(t.primary, 0.4),
                    alignment=PP_ALIGN.CENTER)

        # Dark cards with glowing borders
        card_start_y = board_top + 650000
        card_h = 300000
        card_gap = 120000

        for j, card_text in enumerate(col_data.get("cards", [])[:6]):
            card_top = card_start_y + j * (card_h + card_gap)
            if card_top + card_h > board_top + available_h - 80000:
                break

            card_fill = tint(t.primary, 0.12)
            add_rect(slide, left + 50000, card_top, col_w - 100000, card_h,
                     fill_color=card_fill, line_color=tint(accent, 0.3),
                     line_width=0.5, corner_radius=40000)

            # Thin left glow stripe
            add_rect(slide, left + 50000, card_top + 40000, 30000, card_h - 80000,
                     fill_color=accent)

            add_textbox(slide, card_text, left + 130000, card_top + 30000,
                        col_w - 230000, card_h - 60000,
                        font_size=11, color=tint(t.primary, 0.6),
                        vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Separator line near bottom
        add_line(slide, left + 100000, board_top + available_h - 120000,
                 left + col_w - 100000, board_top + available_h - 120000,
                 color=tint(t.primary, 0.15), width=0.5)


# ── Variant 6: SPLIT — Board with summary stats on side ─────────────────

def _split_board(slide, t, c):
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board"), theme=t)
    columns = _get_columns(c)
    p = t.palette

    n = min(len(columns), 4)
    board_top = content_top + 100000
    available_h = FOOTER_TOP - board_top - 200000

    # Left side: stats summary panel
    stats_w = 2200000
    stats_left = MARGIN
    add_styled_card(slide, stats_left, board_top, stats_w, available_h,
                    theme=t, accent_color=t.accent)

    # Summary stats
    add_textbox(slide, "Summary", stats_left + 100000, board_top + 120000,
                stats_w - 200000, 300000,
                font_size=14, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT)

    add_line(slide, stats_left + 100000, board_top + 440000,
             stats_left + stats_w - 100000, board_top + 440000,
             color=t.accent, width=1.5)

    total_cards = sum(len(col.get("cards", [])) for col in columns[:n])
    stat_items = []
    for col in columns[:n]:
        count = len(col.get("cards", []))
        stat_items.append((col["title"], str(count)))
    stat_items.append(("Total", str(total_cards)))

    stat_y = board_top + 550000
    for label, value in stat_items:
        accent_c = t.accent if label == "Total" else t.secondary
        add_textbox(slide, label, stats_left + 120000, stat_y,
                    stats_w - 400000, 280000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, value, stats_left + stats_w - 400000, stat_y,
                    280000, 280000,
                    font_size=16, bold=True, color=accent_c, alignment=PP_ALIGN.RIGHT)
        stat_y += 350000

    # Right side: kanban columns
    board_left = stats_left + stats_w + 200000
    board_w = SLIDE_W - MARGIN - board_left
    col_gap = 150000
    col_w = (board_w - col_gap * (n - 1)) // n

    for ci, col_data in enumerate(columns[:n]):
        accent = p[ci % len(p)]
        left = board_left + ci * (col_w + col_gap)

        # Column header
        add_rect(slide, left, board_top, col_w, 350000, fill_color=accent,
                 corner_radius=40000)
        add_textbox(slide, col_data["title"], left + 60000, board_top + 50000,
                    col_w - 120000, 250000,
                    font_size=11, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Cards
        card_start_y = board_top + 450000
        card_h = 280000
        card_gap_v = 100000

        for j, card_text in enumerate(col_data.get("cards", [])[:6]):
            card_top = card_start_y + j * (card_h + card_gap_v)
            if card_top + card_h > board_top + available_h - 60000:
                break

            add_styled_card(slide, left, card_top, col_w, card_h,
                            theme=t, accent_color=accent)
            add_textbox(slide, card_text, left + 60000, card_top + 20000,
                        col_w - 120000, card_h - 40000,
                        font_size=10, color=t.primary,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant 7: GEO — Angular geometric cards with borders ───────────────

def _geo_board(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board"), theme=t)
    columns = _get_columns(c)
    p = t.palette

    n = min(len(columns), 4)
    gap = 180000
    total_w = SLIDE_W - 2 * MARGIN
    col_w = (total_w - gap * (n - 1)) // n
    board_top = content_top + 200000
    available_h = FOOTER_TOP - board_top - 300000

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if t.ux_style.dark_mode else t.secondary

    for ci, col_data in enumerate(columns[:n]):
        accent = p[ci % len(p)]
        left = MARGIN + ci * (col_w + gap)

        # Column container with thick border and bottom accent
        add_styled_card(slide, left, board_top, col_w, available_h,
                        theme=t, accent_color=accent)

        # Angular header — sharp-cornered color block
        add_rect(slide, left, board_top, col_w, 420000, fill_color=accent)

        # Diagonal decorative line across header
        add_line(slide, left + col_w - 300000, board_top,
                 left + col_w, board_top + 420000,
                 color=shade(accent, 0.15), width=2)

        # Column title
        add_textbox(slide, col_data["title"], left + 80000, board_top + 60000,
                    col_w - 260000, 300000,
                    font_size=13, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Angular cards — sharp corners with thick accent border
        card_start_y = board_top + 540000
        card_h = 300000
        card_gap_v = 120000

        for j, card_text in enumerate(col_data.get("cards", [])[:6]):
            card_top = card_start_y + j * (card_h + card_gap_v)
            if card_top + card_h > board_top + available_h - 80000:
                break

            card_fill = tint(accent, 0.1 if t.ux_style.dark_mode else 0.92)
            add_rect(slide, left + 50000, card_top, col_w - 100000, card_h,
                     fill_color=card_fill, line_color=accent, line_width=1.5)

            # Bottom accent strip
            add_rect(slide, left + 50000, card_top + card_h - 30000,
                     col_w - 100000, 30000, fill_color=accent)

            add_textbox(slide, card_text, left + 100000, card_top + 30000,
                        col_w - 200000, card_h - 70000,
                        font_size=11, color=text_color,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Decorative bottom line
        add_line(slide, left + 80000, board_top + available_h - 80000,
                 left + col_w - 80000, board_top + available_h - 80000,
                 color=accent, width=0.5)


# ── Variant 8: EDITORIAL — Elegant columns with thin rules ──────────────

def _editorial_board(slide, t, c):
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board"), theme=t)
    columns = _get_columns(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = min(len(columns), 4)
    total_w = SLIDE_W - 2 * m
    col_w = total_w // n
    board_top = content_top + 300000
    available_h = FOOTER_TOP - board_top - 400000

    for ci, col_data in enumerate(columns[:n]):
        left = m + ci * col_w

        # Thin vertical divider between columns
        if ci > 0:
            div_x = left - 1
            add_line(slide, div_x, board_top - 100000,
                     div_x, board_top + available_h,
                     color=tint(t.secondary, 0.8), width=0.5)

        # Full-width thin rule above header
        add_line(slide, left + 40000, board_top,
                 left + col_w - 40000, board_top,
                 color=tint(t.secondary, 0.7), width=0.5)

        # Short accent mark at left of rule
        add_line(slide, left + 40000, board_top,
                 left + 40000 + 300000, board_top,
                 color=t.accent, width=1.5)

        # Column title — elegant, left-aligned
        add_textbox(slide, col_data["title"], left + 40000, board_top + 80000,
                    col_w - 80000, 300000,
                    font_size=13, bold=True, color=t.primary,
                    alignment=PP_ALIGN.LEFT)

        # Italic item count
        add_textbox(slide, f"{len(col_data.get('cards', []))} tasks",
                    left + col_w - 400000, board_top + 80000,
                    360000, 300000,
                    font_size=9, italic=True, color=t.secondary,
                    alignment=PP_ALIGN.RIGHT)

        # Items as elegant list entries with thin separators
        item_start_y = board_top + 480000
        item_h = 380000

        for j, card_text in enumerate(col_data.get("cards", [])[:6]):
            item_top = item_start_y + j * item_h
            if item_top + item_h > board_top + available_h:
                break

            # Thin rule above each item
            add_line(slide, left + 40000, item_top,
                     left + col_w - 40000, item_top,
                     color=tint(t.secondary, 0.85), width=0.25)

            # Small accent dash
            add_line(slide, left + 40000, item_top + item_h // 2,
                     left + 40000 + 60000, item_top + item_h // 2,
                     color=t.accent, width=1)

            add_textbox(slide, card_text, left + 130000, item_top + 60000,
                        col_w - 180000, item_h - 120000,
                        font_size=11, color=t.primary,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Bottom rule
    add_line(slide, m, FOOTER_TOP - 300000, SLIDE_W - m, FOOTER_TOP - 300000,
             color=tint(t.secondary, 0.8), width=0.5)


# ── Variant 9: GRADIENT — Gradient-colored column headers ───────────────

def _gradient_board(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board"), theme=t)
    columns = _get_columns(c)
    p = t.palette

    n = min(len(columns), 4)
    gap = 280000
    total_w = SLIDE_W - 2 * MARGIN
    col_w = (total_w - gap * (n - 1)) // n
    board_top = content_top + 200000
    available_h = FOOTER_TOP - board_top - 300000

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if t.ux_style.dark_mode else t.secondary

    for ci, col_data in enumerate(columns[:n]):
        accent = p[ci % len(p)]
        left = MARGIN + ci * (col_w + gap)

        # Column container
        add_styled_card(slide, left, board_top, col_w, available_h, theme=t)

        # Gradient-like header: two color layers fading from accent
        add_rect(slide, left, board_top, col_w, 450000,
                 fill_color=accent, corner_radius=80000)
        # Lighter overlay on lower half of header to simulate gradient
        add_rect(slide, left + 10000, board_top + 220000,
                 col_w - 20000, 230000,
                 fill_color=tint(accent, 0.3))

        # Halo circle behind title for visual depth
        add_circle(slide, left + col_w // 2, board_top + 225000, 150000,
                   fill_color=tint(accent, 0.2))

        # Column title
        add_textbox(slide, col_data["title"], left + 80000, board_top + 70000,
                    col_w - 160000, 300000,
                    font_size=14, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Soft cards with circle indicator
        card_start_y = board_top + 560000
        card_h = 320000
        card_gap_v = 140000

        for j, card_text in enumerate(col_data.get("cards", [])[:5]):
            card_top = card_start_y + j * (card_h + card_gap_v)
            if card_top + card_h > board_top + available_h - 100000:
                break

            card_fill = tint(accent, 0.1 if t.ux_style.dark_mode else 0.9)
            add_rect(slide, left + 50000, card_top, col_w - 100000, card_h,
                     fill_color=card_fill, corner_radius=80000)

            # Soft color circle indicator
            add_circle(slide, left + 120000, card_top + card_h // 2, 25000,
                       fill_color=accent)

            add_textbox(slide, card_text, left + 180000, card_top + 30000,
                        col_w - 280000, card_h - 60000,
                        font_size=11, color=text_color,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Bottom percentage bar showing column fullness
        pct = len(col_data.get("cards", [])) / max(total_cards_count(columns), 1)
        bar_y = board_top + available_h - 120000
        bar_w = col_w - 120000
        add_rect(slide, left + 60000, bar_y, bar_w, 40000,
                 fill_color=tint(accent, 0.8 if not t.ux_style.dark_mode else 0.15),
                 corner_radius=30000)
        fill_w = int(bar_w * pct)
        if fill_w > 0:
            add_rect(slide, left + 60000, bar_y, fill_w, 40000,
                     fill_color=accent, corner_radius=30000)


def total_cards_count(columns):
    """Helper to count total cards across all columns."""
    return sum(len(col.get("cards", [])) for col in columns)


# ── Variant 10: RETRO — Vintage styled cards with borders ───────────────

def _retro_board(slide, t, c):
    add_background(slide, t.light_bg)
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board"), theme=t)
    columns = _get_columns(c)
    p = t.palette

    n = min(len(columns), 4)
    gap = 220000
    total_w = SLIDE_W - 2 * MARGIN
    col_w = (total_w - gap * (n - 1)) // n
    board_top = content_top + 200000
    available_h = FOOTER_TOP - board_top - 300000

    for ci, col_data in enumerate(columns[:n]):
        accent = p[ci % len(p)]
        left = MARGIN + ci * (col_w + gap)

        # Column with double border (retro frame)
        add_rect(slide, left, board_top, col_w, available_h,
                 fill_color="#FFFFFF", line_color=accent, line_width=2.0,
                 corner_radius=100000)
        add_rect(slide, left + 30000, board_top + 30000,
                 col_w - 60000, available_h - 60000,
                 line_color=accent, line_width=0.75, corner_radius=80000)

        # Decorative circles in header corners
        add_circle(slide, left + 100000, board_top + 100000, 30000,
                   fill_color=accent)
        add_circle(slide, left + col_w - 100000, board_top + 100000, 30000,
                   fill_color=accent)

        # Column header with decorative dashes
        add_textbox(slide, f"\u2014 {col_data['title']} \u2014",
                    left + 60000, board_top + 170000,
                    col_w - 120000, 350000,
                    font_size=13, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER)

        # Decorative rule under header
        rule_left = left + col_w // 4
        add_line(slide, rule_left, board_top + 520000,
                 rule_left + col_w // 2, board_top + 520000,
                 color=accent, width=1.5)

        # Cards — rounded, bordered, vintage feel
        card_start_y = board_top + 640000
        card_h = 300000
        card_gap_v = 130000

        for j, card_text in enumerate(col_data.get("cards", [])[:5]):
            card_top = card_start_y + j * (card_h + card_gap_v)
            if card_top + card_h > board_top + available_h - 80000:
                break

            add_rect(slide, left + 60000, card_top, col_w - 120000, card_h,
                     fill_color=tint(accent, 0.9), line_color=accent,
                     line_width=1.0, corner_radius=60000)

            # Bullet marker
            add_circle(slide, left + 120000, card_top + card_h // 2, 20000,
                       fill_color=accent)

            add_textbox(slide, card_text, left + 170000, card_top + 30000,
                        col_w - 280000, card_h - 60000,
                        font_size=11, color=t.primary,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Decorative triple dots at column bottom
        dot_y = board_top + available_h - 100000
        for d in range(3):
            add_circle(slide, left + col_w // 2 - 60000 + d * 60000, dot_y, 15000,
                       fill_color=accent)


# ── Variant 11: MAGAZINE — Creative staggered cards with color coding ────

def _creative_board(slide, t, c):
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board"), theme=t)
    columns = _get_columns(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = min(len(columns), 4)
    gap = int(60000 * s.gap_factor)
    total_w = SLIDE_W - 2 * m
    col_w = (total_w - gap * (n - 1)) // n
    board_top = content_top + 100000
    available_h = FOOTER_TOP - board_top - 200000

    for ci, col_data in enumerate(columns[:n]):
        accent = p[ci % len(p)]
        left = m + ci * (col_w + gap)

        # Full-color column background block
        add_rect(slide, left, board_top, col_w, available_h,
                 fill_color=accent)

        # Large faded column number
        add_textbox(slide, f"{ci + 1:02d}", left + 40000, board_top + 40000,
                    col_w - 80000, 500000,
                    font_size=36, bold=True, color=tint(accent, 0.3),
                    alignment=PP_ALIGN.LEFT)

        # Column title — white, overlapping number area
        add_textbox(slide, col_data["title"].upper(), left + 40000, board_top + 400000,
                    col_w - 80000, 350000,
                    font_size=14, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT)

        # White accent bar
        add_rect(slide, left + 40000, board_top + 750000,
                 col_w // 3, 40000, fill_color="#FFFFFF")

        # Staggered cards — white cards within colored column
        card_start_y = board_top + 880000
        card_gap_v = 100000

        for j, card_text in enumerate(col_data.get("cards", [])[:5]):
            card_h = 280000
            # Alternate slight left/right offset for stagger effect
            offset = 30000 if j % 2 == 0 else 70000
            card_left = left + offset
            card_w_actual = col_w - offset - 40000
            card_top = card_start_y + j * (card_h + card_gap_v)

            if card_top + card_h > board_top + available_h - 60000:
                break

            # White card on colored background
            add_rect(slide, card_left, card_top, card_w_actual, card_h,
                     fill_color="#FFFFFF")

            add_textbox(slide, card_text, card_left + 50000, card_top + 30000,
                        card_w_actual - 100000, card_h - 60000,
                        font_size=11, bold=True, color=shade(accent, 0.2),
                        vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant 12: SCHOLARLY — Clean numbered columns, thin rules ───────────

def _scholarly_board(slide, t, c):
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board"), theme=t)
    columns = _get_columns(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = min(len(columns), 4)
    total_w = SLIDE_W - 2 * m
    col_w = total_w // n
    board_top = content_top + 300000
    available_h = FOOTER_TOP - board_top - 500000

    # Full-width thin rule above columns
    add_line(slide, m, board_top - 100000, SLIDE_W - m, board_top - 100000,
             color=tint(t.secondary, 0.7), width=0.5)

    # Short accent mark
    add_line(slide, m, board_top - 100000, m + 300000, board_top - 100000,
             color=t.accent, width=1.5)

    for ci, col_data in enumerate(columns[:n]):
        left = m + ci * col_w

        # Thin vertical divider
        if ci > 0:
            add_line(slide, left, board_top - 100000,
                     left, board_top + available_h,
                     color=tint(t.secondary, 0.75), width=0.5)

        # Column number and title — scholarly numbered format
        add_textbox(slide, f"Fig. {ci + 1}.", left + 40000, board_top,
                    200000, 280000,
                    font_size=9, italic=True, color=t.accent,
                    alignment=PP_ALIGN.LEFT)

        add_textbox(slide, col_data["title"], left + 240000, board_top,
                    col_w - 300000, 280000,
                    font_size=12, bold=True, color=t.primary,
                    alignment=PP_ALIGN.LEFT)

        # Thin underline below column header
        add_line(slide, left + 40000, board_top + 320000,
                 left + col_w - 40000, board_top + 320000,
                 color=tint(t.secondary, 0.7), width=0.5)

        # Items as numbered scholarly list
        item_start_y = board_top + 420000
        item_h = 340000

        for j, card_text in enumerate(col_data.get("cards", [])[:6]):
            item_top = item_start_y + j * item_h
            if item_top + item_h > board_top + available_h:
                break

            # Thin separator between items
            if j > 0:
                add_line(slide, left + 40000, item_top,
                         left + col_w - 40000, item_top,
                         color=tint(t.secondary, 0.88), width=0.25)

            # Numbered marker
            add_textbox(slide, f"{ci + 1}.{j + 1}", left + 40000, item_top + 50000,
                        200000, 250000,
                        font_size=9, color=t.accent, alignment=PP_ALIGN.LEFT)

            add_textbox(slide, card_text, left + 250000, item_top + 50000,
                        col_w - 300000, 250000,
                        font_size=11, color=t.primary)

    # Bottom rule with figure caption
    add_line(slide, m, board_top + available_h + 50000,
             SLIDE_W - m, board_top + available_h + 50000,
             color=tint(t.secondary, 0.7), width=0.5)

    add_textbox(slide, f"Figure: Project status board \u2014 {n} workflow stages",
                m, board_top + available_h + 80000,
                total_w, 250000,
                font_size=9, italic=True, color=t.secondary,
                alignment=PP_ALIGN.LEFT)


# ── Variant 13: LABORATORY — Dark bg, data-coded cards ──────────────────

def _laboratory_board(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board"), theme=t)
    columns = _get_columns(c)
    p = t.palette

    n = min(len(columns), 4)
    gap = 160000
    total_w = SLIDE_W - 2 * MARGIN
    col_w = (total_w - gap * (n - 1)) // n
    board_top = content_top + 200000
    available_h = FOOTER_TOP - board_top - 300000

    # Status indicators for columns
    status_icons = ["\u25cb", "\u25d4", "\u25cf"]  # empty, half, full circles

    for ci, col_data in enumerate(columns[:n]):
        accent = p[ci % len(p)]
        left = MARGIN + ci * (col_w + gap)

        # Column container — dark with subtle grid overlay
        col_fill = tint(t.primary, 0.06)
        add_rect(slide, left, board_top, col_w, available_h,
                 fill_color=col_fill, line_color=tint(t.primary, 0.2),
                 line_width=0.75)

        # Color-coded accent bar on left
        add_rect(slide, left, board_top, 40000, available_h, fill_color=accent)

        # Header area with grid-line effect
        add_rect(slide, left, board_top, col_w, 450000,
                 fill_color=tint(t.primary, 0.1))

        # Grid lines in header
        for g in range(3):
            gy = board_top + 120000 + g * 110000
            add_line(slide, left + 50000, gy, left + col_w - 50000, gy,
                     color=tint(t.primary, 0.12), width=0.25)

        # Status icon + column title
        icon = status_icons[ci % len(status_icons)]
        add_textbox(slide, icon, left + 60000, board_top + 60000,
                    120000, 120000,
                    font_size=10, color=accent, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, col_data["title"], left + 180000, board_top + 50000,
                    col_w - 240000, 180000,
                    font_size=11, bold=True, color=accent,
                    alignment=PP_ALIGN.LEFT)

        # Item count label
        count = len(col_data.get("cards", []))
        add_textbox(slide, f"n={count}", left + 60000, board_top + 260000,
                    col_w - 120000, 150000,
                    font_size=8, color=tint(t.primary, 0.4),
                    alignment=PP_ALIGN.LEFT)

        # Data-coded task cards
        card_start_y = board_top + 540000
        card_h = 280000
        card_gap_v = 100000

        for j, card_text in enumerate(col_data.get("cards", [])[:6]):
            card_top = card_start_y + j * (card_h + card_gap_v)
            if card_top + card_h > board_top + available_h - 80000:
                break

            # Card with left color code
            card_fill = tint(t.primary, 0.1)
            add_rect(slide, left + 50000, card_top, col_w - 100000, card_h,
                     fill_color=card_fill, line_color=tint(t.primary, 0.18),
                     line_width=0.5)

            # Data code stripe
            add_rect(slide, left + 50000, card_top, 25000, card_h,
                     fill_color=accent)

            # Item ID label
            add_textbox(slide, f"[{ci + 1}.{j + 1}]", left + 90000, card_top + 20000,
                        200000, 120000,
                        font_size=7, color=tint(t.primary, 0.35))

            add_textbox(slide, card_text, left + 90000, card_top + 100000,
                        col_w - 200000, card_h - 130000,
                        font_size=10, color=tint(t.primary, 0.6),
                        vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Evidence count at column bottom
        add_textbox(slide, f"\u2211 = {count} items", left + 60000,
                    board_top + available_h - 150000,
                    col_w - 120000, 120000,
                    font_size=8, color=tint(t.primary, 0.3),
                    alignment=PP_ALIGN.RIGHT)


# ── Variant 14: DASHBOARD — Dense compact board with status indicators ──

def _dashboard_board(slide, t, c):
    content_top = add_slide_title(slide, c.get("kanban_title", "Project Board"), theme=t)
    columns = _get_columns(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = min(len(columns), 4)
    gap = int(100000 * s.gap_factor)
    total_w = SLIDE_W - 2 * m
    col_w = (total_w - gap * (n - 1)) // n
    board_top = content_top + 50000
    available_h = FOOTER_TOP - board_top - 200000

    # Top metrics bar
    metrics_h = 350000
    add_rect(slide, m, board_top, total_w, metrics_h,
             fill_color=tint(t.primary, 0.92 if not s.dark_mode else 0.12),
             corner_radius=30000)

    total_cards = sum(len(col.get("cards", [])) for col in columns[:n])
    metric_w = total_w // (n + 1)

    # Overall total metric
    add_textbox(slide, "TOTAL", m + 40000, board_top + 40000,
                metric_w - 80000, 140000,
                font_size=8, bold=True, color=t.secondary, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, str(total_cards), m + 40000, board_top + 160000,
                metric_w - 80000, 160000,
                font_size=18, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)

    # Per-column metrics in top bar
    for ci, col_data in enumerate(columns[:n]):
        accent = p[ci % len(p)]
        metric_left = m + (ci + 1) * metric_w

        add_textbox(slide, col_data["title"].upper(), metric_left, board_top + 40000,
                    metric_w - 40000, 140000,
                    font_size=8, bold=True, color=t.secondary, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, str(len(col_data.get("cards", []))),
                    metric_left, board_top + 160000,
                    metric_w - 40000, 160000,
                    font_size=18, bold=True, color=accent, alignment=PP_ALIGN.CENTER)

    # Kanban columns below metrics bar
    col_top = board_top + metrics_h + 100000
    col_available_h = available_h - metrics_h - 100000

    for ci, col_data in enumerate(columns[:n]):
        accent = p[ci % len(p)]
        left = m + ci * (col_w + gap)

        # Dense column container with thin border
        add_rect(slide, left, col_top, col_w, col_available_h,
                 fill_color=tint(t.primary, 0.97 if not s.dark_mode else 0.08),
                 line_color=tint(t.secondary, 0.7), line_width=0.5,
                 corner_radius=30000)

        # Compact color header bar
        add_rect(slide, left, col_top, col_w, 50000, fill_color=accent,
                 corner_radius=30000)

        # Column title - compact
        add_textbox(slide, col_data["title"], left + 50000, col_top + 70000,
                    col_w - 100000, 240000,
                    font_size=10, bold=True, color=t.primary,
                    alignment=PP_ALIGN.LEFT)

        # Dense compact cards
        card_start_y = col_top + 340000
        card_h = 240000
        card_gap_v = 60000

        for j, card_text in enumerate(col_data.get("cards", [])[:7]):
            card_top = card_start_y + j * (card_h + card_gap_v)
            if card_top + card_h > col_top + col_available_h - 60000:
                break

            # Compact tile
            tile_fill = tint(accent, 0.93 if not s.dark_mode else 0.12)
            add_rect(slide, left + 30000, card_top, col_w - 60000, card_h,
                     fill_color=tile_fill, corner_radius=20000)

            # Status dot
            status_colors = [t.accent, p[1 % len(p)], p[2 % len(p)]]
            dot_color = status_colors[j % len(status_colors)]
            add_circle(slide, left + 70000, card_top + card_h // 2, 18000,
                       fill_color=dot_color)

            add_textbox(slide, card_text, left + 110000, card_top + 20000,
                        col_w - 180000, card_h - 40000,
                        font_size=9, color=t.primary,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Compact progress indicator at bottom
        count = len(col_data.get("cards", []))
        pct = count / max(total_cards, 1)
        bar_y = col_top + col_available_h - 80000
        bar_w = col_w - 80000
        add_rect(slide, left + 40000, bar_y, bar_w, 30000,
                 fill_color=tint(accent, 0.85 if not s.dark_mode else 0.15),
                 corner_radius=20000)
        fill_w = int(bar_w * pct)
        if fill_w > 0:
            add_rect(slide, left + 40000, bar_y, fill_w, 30000,
                     fill_color=accent, corner_radius=20000)
