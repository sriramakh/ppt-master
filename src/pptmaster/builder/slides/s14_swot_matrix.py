"""Slide 14 — SWOT Matrix: 2x2 colored quadrants."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from pptmaster.builder.design_system import SLIDE_W, SLIDE_H, MARGIN, FONT_CARD_HEADING, FONT_BODY, FONT_CAPTION, CONTENT_TOP
from pptmaster.builder.helpers import add_textbox, add_rect, add_line, add_gold_accent_line, add_bullet_list, add_styled_card, add_slide_title, add_dark_bg, add_background
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    p = t.palette
    s = t.ux_style

    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("swot_title", "SWOT Analysis"), theme=t)

    bullet_text_color = "#FFFFFF" if s.dark_mode else t.primary

    swot = c.get("swot", {
        "strengths": ["Brand recognition", "Strong team", "Financial position"],
        "weaknesses": ["Limited reach", "Legacy tech", "Turnover"],
        "opportunities": ["Market expansion", "Acquisitions", "Digital demand"],
        "threats": ["Competition", "Regulation", "Uncertainty"],
    })

    # In dark mode, use deeper tints so quadrant backgrounds still work
    tint_factor = 0.25 if s.dark_mode else 0.85
    quadrants = [
        ("Strengths", tint(p[1], tint_factor), p[1], swot.get("strengths", [])),
        ("Weaknesses", tint(p[2], tint_factor), p[2], swot.get("weaknesses", [])),
        ("Opportunities", tint(p[0], tint_factor), p[0], swot.get("opportunities", [])),
        ("Threats", tint(t.accent, tint_factor), t.accent, swot.get("threats", [])),
    ]

    m = int(MARGIN * s.margin_factor)
    gap = int(150000 * s.gap_factor)
    total_w = SLIDE_W - 2 * m
    total_h = 4300000
    quad_w = (total_w - gap) // 2
    quad_h = (total_h - gap) // 2
    start_y = content_y + 100000

    for i, (title, bg_color, accent, items) in enumerate(quadrants):
        row, col = divmod(i, 2)
        left = m + col * (quad_w + gap)
        top = start_y + row * (quad_h + gap)

        add_rect(slide, left, top, quad_w, quad_h, fill_color=bg_color, corner_radius=60000)
        add_rect(slide, left, top, quad_w, 60000, fill_color=accent)
        add_textbox(slide, title, left + 200000, top + 150000, quad_w - 400000, 400000,
                    font_size=FONT_CARD_HEADING, bold=True, color=accent)
        add_bullet_list(slide, items, left + 200000, top + 550000, quad_w - 400000, quad_h - 700000,
                        font_size=13, color=bullet_text_color, bullet_color=accent, spacing=6)


# ── Variant: scholarly ────────────────────────────────────────────────

def _scholarly(slide, t):
    """White bg, thin-bordered 2x2 grid with 'Figure 1:' caption."""
    c = t.content
    m = MARGIN

    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, c.get("swot_title", "SWOT Analysis"), theme=t)

    # Figure caption
    add_textbox(slide, "Figure 1: SWOT Analysis", m, content_y + 20000,
                SLIDE_W - 2 * m, 300000,
                font_size=FONT_CAPTION, italic=True, color=t.secondary)

    swot = c.get("swot", {
        "strengths": ["Brand recognition", "Strong team", "Financial position"],
        "weaknesses": ["Limited reach", "Legacy tech", "Turnover"],
        "opportunities": ["Market expansion", "Acquisitions", "Digital demand"],
        "threats": ["Competition", "Regulation", "Uncertainty"],
    })

    labels = ["Strengths", "Weaknesses", "Opportunities", "Threats"]
    keys = ["strengths", "weaknesses", "opportunities", "threats"]
    rule_color = tint(t.secondary, 0.5)

    total_w = SLIDE_W - 2 * m
    gap = 150000
    quad_w = (total_w - gap) // 2
    quad_h = 1900000
    start_y = content_y + 380000

    for i, (label, key) in enumerate(zip(labels, keys)):
        row, col = divmod(i, 2)
        left = m + col * (quad_w + gap)
        top = start_y + row * (quad_h + gap)
        items = swot.get(key, [])

        # Thin border cell (no fill)
        add_rect(slide, left, top, quad_w, quad_h,
                 line_color=rule_color, line_width=0.75)
        # Bold header
        add_textbox(slide, label, left + 150000, top + 80000, quad_w - 300000, 300000,
                    font_size=FONT_CARD_HEADING, bold=True, color=t.primary)
        # Thin rule below header
        add_line(slide, left + 150000, top + 380000, left + quad_w - 150000, top + 380000,
                 color=rule_color, width=0.5)
        # Items as simple text with thin rule separators
        item_y = top + 440000
        for j, item in enumerate(items):
            add_textbox(slide, item, left + 150000, item_y, quad_w - 300000, 250000,
                        font_size=11, color=t.primary)
            item_y += 300000
            if j < len(items) - 1:
                add_line(slide, left + 150000, item_y - 50000,
                         left + quad_w - 150000, item_y - 50000,
                         color=tint(t.secondary, 0.75), width=0.25)


# ── Variant: laboratory ──────────────────────────────────────────────

def _laboratory(slide, t):
    """Dark bg, SWOT quadrants as dark cards with colored left borders."""
    c = t.content
    p = t.palette
    m = MARGIN

    add_background(slide, t.primary)
    content_y = add_slide_title(slide, c.get("swot_title", "SWOT Analysis"), theme=t)

    swot = c.get("swot", {
        "strengths": ["Brand recognition", "Strong team", "Financial position"],
        "weaknesses": ["Limited reach", "Legacy tech", "Turnover"],
        "opportunities": ["Market expansion", "Acquisitions", "Digital demand"],
        "threats": ["Competition", "Regulation", "Uncertainty"],
    })

    border_colors = [p[0] if len(p) > 0 else t.accent,
                     p[1] if len(p) > 1 else t.accent,
                     p[2] if len(p) > 2 else t.accent,
                     p[3] if len(p) > 3 else t.accent]
    labels = ["Strengths", "Weaknesses", "Opportunities", "Threats"]
    keys = ["strengths", "weaknesses", "opportunities", "threats"]

    total_w = SLIDE_W - 2 * m
    gap = 150000
    quad_w = (total_w - gap) // 2
    quad_h = 1900000
    start_y = content_y + 100000
    card_fill = tint(t.primary, 0.15)

    for i, (label, key, bdr_color) in enumerate(zip(labels, keys, border_colors)):
        row, col = divmod(i, 2)
        left = m + col * (quad_w + gap)
        top = start_y + row * (quad_h + gap)
        items = swot.get(key, [])

        # Dark card background
        add_rect(slide, left, top, quad_w, quad_h, fill_color=card_fill, corner_radius=30000)
        # Colored left border
        add_rect(slide, left, top, 50000, quad_h, fill_color=bdr_color)
        # Header in accent color
        add_textbox(slide, label, left + 150000, top + 80000, quad_w - 300000, 300000,
                    font_size=FONT_CARD_HEADING, bold=True, color=bdr_color)
        # Bullet items in white
        add_bullet_list(slide, items, left + 150000, top + 420000,
                        quad_w - 300000, quad_h - 520000,
                        font_size=12, color="#FFFFFF", bullet_color=bdr_color, spacing=6)


# ── Variant: dashboard ───────────────────────────────────────────────

def _dashboard(slide, t):
    """White bg with accent header band, compact 2x2 SWOT grid."""
    c = t.content
    p = t.palette
    m = int(MARGIN * 0.85)

    add_background(slide, "#FFFFFF")

    # Accent header band
    add_rect(slide, 0, 0, SLIDE_W, 750000, fill_color=t.primary)
    add_textbox(slide, c.get("swot_title", "SWOT Analysis"), m, 150000,
                SLIDE_W - 2 * m, 450000,
                font_size=22, bold=True, color="#FFFFFF")
    add_gold_accent_line(slide, m, 700000, 600000, color=t.accent)

    swot = c.get("swot", {
        "strengths": ["Brand recognition", "Strong team", "Financial position"],
        "weaknesses": ["Limited reach", "Legacy tech", "Turnover"],
        "opportunities": ["Market expansion", "Acquisitions", "Digital demand"],
        "threats": ["Competition", "Regulation", "Uncertainty"],
    })

    labels = ["Strengths", "Weaknesses", "Opportunities", "Threats"]
    keys = ["strengths", "weaknesses", "opportunities", "threats"]
    quad_colors = [p[0] if len(p) > 0 else t.accent,
                   p[1] if len(p) > 1 else t.accent,
                   p[2] if len(p) > 2 else t.accent,
                   p[3] if len(p) > 3 else t.accent]

    total_w = SLIDE_W - 2 * m
    gap = 100000
    quad_w = (total_w - gap) // 2
    quad_h = 2600000
    start_y = 850000

    for i, (label, key, qcolor) in enumerate(zip(labels, keys, quad_colors)):
        row, col = divmod(i, 2)
        left = m + col * (quad_w + gap)
        top = start_y + row * (quad_h + gap)
        items = swot.get(key, [])

        # Light tinted background card
        add_rect(slide, left, top, quad_w, quad_h,
                 fill_color=tint(qcolor, 0.90), corner_radius=50000)
        # Colored top accent bar
        add_rect(slide, left, top, quad_w, 45000, fill_color=qcolor)
        # Compact header
        add_textbox(slide, label, left + 120000, top + 70000, quad_w - 240000, 280000,
                    font_size=14, bold=True, color=qcolor)
        # Dense bullet list
        add_bullet_list(slide, items, left + 120000, top + 370000,
                        quad_w - 240000, quad_h - 450000,
                        font_size=11, color=t.primary, bullet_color=qcolor, spacing=4)
