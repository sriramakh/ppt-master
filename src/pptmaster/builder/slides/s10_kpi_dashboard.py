"""Slide 10 — KPI Dashboard: 11 unique visual variants dispatched by UX style."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP,
    card_positions, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_card, add_styled_card, add_circle, add_line,
    add_gold_accent_line, add_slide_title, add_dark_bg, add_background,
    add_diamond, add_hexagon,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    s = t.ux_style
    c = t.content

    dispatch = {
        "cards-row": _cards_row,
        "stat-line": _stat_line,
        "large-numbers": _large_numbers,
        "elevated-cards": _elevated_cards,
        "dark-cards": _dark_cards,
        "split-kpi": _split_kpi,
        "angular-cards": _angular_cards,
        "editorial-stats": _editorial_stats,
        "gradient-cards": _gradient_cards,
        "retro-badges": _retro_badges,
        "magazine-blocks": _magazine_blocks,
        "scholarly-stats": _scholarly_stats,
        "laboratory-metrics": _laboratory_metrics,
        "dashboard-tiles": _dashboard_tiles,
    }
    builder = dispatch.get(s.kpi, _cards_row)
    builder(slide, t, c)


def _get_kpis(c):
    """Extract KPI data from content dict with sensible defaults."""
    return c.get("kpis", [
        ("Revenue", "$850M", "+23%", 0.85, "\u2191"),
        ("Profit Margin", "18.5%", "+2.1%", 0.72, "\u2191"),
        ("Customer Churn", "2.1%", "-0.8%", 0.21, "\u2193"),
        ("NPS Score", "72", "+5pts", 0.72, "\u2191"),
    ])


def _trend_color(t, arrow):
    """Return green for up-arrow, red for down-arrow."""
    p = t.palette
    return p[1] if "\u2191" in arrow else p[2]


def _progress_bar(slide, left, top, width, height, progress, accent, bg_color):
    """Draw a filled progress bar."""
    add_rect(slide, left, top, width, height, fill_color=bg_color, corner_radius=50000)
    fill_w = int(width * min(progress, 1.0))
    if fill_w > 0:
        add_rect(slide, left, top, fill_w, height, fill_color=accent, corner_radius=50000)


# ── Variant 1: CLASSIC — 4 cards in a row, accent bar top, progress bars ──

def _cards_row(slide, t, c):
    content_top = add_slide_title(slide, c.get("kpi_title", "KPI Dashboard"), theme=t)
    kpis = _get_kpis(c)
    p = t.palette

    positions = card_positions(4, 1, top=content_top + 100000, gap=200000, card_height=4200000)

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:4]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        add_styled_card(slide, left, top, width, height, theme=t, accent_color=accent)
        add_textbox(slide, label, left + 100000, top + 250000, width - 200000, 350000,
                    font_size=13, bold=True, color=t.secondary, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, value, left + 50000, top + 700000, width - 100000, 700000,
                    font_size=36, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)

        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", left + 50000, top + 1400000,
                    width - 100000, 350000,
                    font_size=14, bold=True, color=tc, alignment=PP_ALIGN.CENTER)

        bar_left = left + 150000
        bar_top = top + 1900000
        bar_w = width - 300000
        _progress_bar(slide, bar_left, bar_top, bar_w, 100000, progress, accent,
                      tint(accent, 0.8))

        add_textbox(slide, f"{int(progress * 100)}% of target", left + 100000, bar_top + 150000,
                    width - 200000, 250000, font_size=10, color=t.secondary,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, "vs. previous quarter", left + 100000, top + height - 500000,
                    width - 200000, 250000, font_size=9, color=t.secondary,
                    alignment=PP_ALIGN.CENTER)


# ── Variant 2: MINIMAL — clean numbers on a line, thin separators ────────

def _stat_line(slide, t, c):
    content_top = add_slide_title(slide, c.get("kpi_title", "KPI Dashboard"), theme=t)
    kpis = _get_kpis(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = min(len(kpis), 4)
    total_w = SLIDE_W - 2 * m
    col_w = total_w // n
    stat_top = content_top + 600000

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:n]):
        left = m + i * col_w
        center = left + col_w // 2

        # Thin vertical separator between items
        if i > 0:
            sep_x = left
            add_line(slide, sep_x, stat_top, sep_x, stat_top + 2400000,
                     color=tint(t.secondary, 0.75), width=0.5)

        # Label above
        add_textbox(slide, label.upper(), left + 80000, stat_top, col_w - 160000, 300000,
                    font_size=10, bold=True, color=t.secondary, alignment=PP_ALIGN.CENTER)

        # Large value
        add_textbox(slide, value, left + 40000, stat_top + 400000, col_w - 80000, 800000,
                    font_size=40, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)

        # Trend
        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", left + 80000, stat_top + 1300000,
                    col_w - 160000, 350000,
                    font_size=12, color=tc, alignment=PP_ALIGN.CENTER)

        # Thin horizontal rule below trend
        rule_left = left + col_w // 4
        rule_w = col_w // 2
        add_line(slide, rule_left, stat_top + 1800000, rule_left + rule_w, stat_top + 1800000,
                 color=tint(t.secondary, 0.8), width=0.5)

        # Progress as text
        add_textbox(slide, f"{int(progress * 100)}%", left + 80000, stat_top + 1950000,
                    col_w - 160000, 300000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.CENTER)


# ── Variant 3: BOLD — oversized numbers centered, thick accent bars ──────

def _large_numbers(slide, t, c):
    content_top = add_slide_title(slide, c.get("kpi_title", "KPI Dashboard").upper(), theme=t)
    kpis = _get_kpis(c)
    p = t.palette

    n = min(len(kpis), 4)
    total_w = SLIDE_W - 2 * MARGIN
    col_w = total_w // n
    gap = 120000
    block_w = col_w - gap
    num_top = content_top + 200000

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:n]):
        left = MARGIN + i * col_w + gap // 2
        accent = p[i % len(p)]

        # Thick accent bar on left
        add_rect(slide, left, num_top, 80000, 3200000, fill_color=accent)

        # Label
        add_textbox(slide, label.upper(), left + 180000, num_top + 100000,
                    block_w - 200000, 350000,
                    font_size=12, bold=True, color=t.secondary, alignment=PP_ALIGN.CENTER)

        # Oversized value
        add_textbox(slide, value, left + 100000, num_top + 500000,
                    block_w - 120000, 1200000,
                    font_size=48, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)

        # Trend
        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", left + 100000, num_top + 1700000,
                    block_w - 120000, 400000,
                    font_size=16, bold=True, color=tc, alignment=PP_ALIGN.CENTER)

        # Thick accent underline
        add_rect(slide, left + 180000, num_top + 2300000, block_w - 360000, 60000,
                 fill_color=accent)

        # Progress text large
        add_textbox(slide, f"{int(progress * 100)}%", left + 100000, num_top + 2500000,
                    block_w - 120000, 400000,
                    font_size=28, bold=True, color=tint(t.primary, 0.3),
                    alignment=PP_ALIGN.CENTER)


# ── Variant 4: ELEVATED — very rounded cards, heavy shadow, progress bars ─

def _elevated_cards(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("kpi_title", "KPI Dashboard"), theme=t)
    kpis = _get_kpis(c)
    p = t.palette

    positions = card_positions(4, 1, top=content_top + 200000, gap=250000, card_height=4000000)

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:4]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        add_styled_card(slide, left, top, width, height, theme=t, accent_color=accent)

        text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
        sub_color = tint(t.primary, 0.6) if t.ux_style.dark_mode else t.secondary

        # Icon circle at top
        circle_x = left + width // 2
        circle_y = top + 400000
        add_circle(slide, circle_x, circle_y, 200000, fill_color=tint(accent, 0.2))
        add_circle(slide, circle_x, circle_y, 120000, fill_color=accent)

        # Label
        add_textbox(slide, label, left + 80000, top + 700000, width - 160000, 300000,
                    font_size=12, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER)

        # Value
        add_textbox(slide, value, left + 40000, top + 1050000, width - 80000, 700000,
                    font_size=34, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)

        # Trend
        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", left + 80000, top + 1750000,
                    width - 160000, 300000,
                    font_size=13, bold=True, color=tc, alignment=PP_ALIGN.CENTER)

        # Progress bar
        bar_left = left + 120000
        bar_top = top + 2200000
        bar_w = width - 240000
        _progress_bar(slide, bar_left, bar_top, bar_w, 120000, progress, accent,
                      tint(accent, 0.75))

        # Progress label
        add_textbox(slide, f"{int(progress * 100)}% of target", left + 80000, bar_top + 200000,
                    width - 160000, 250000,
                    font_size=10, color=sub_color, alignment=PP_ALIGN.CENTER)


# ── Variant 5: DARK — dark bg cards, thin border, neon-bright values ─────

def _dark_cards(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("kpi_title", "KPI Dashboard").upper(), theme=t)
    kpis = _get_kpis(c)
    p = t.palette

    positions = card_positions(4, 1, top=content_top + 200000, gap=200000, card_height=4000000)

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:4]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]
        card_fill = tint(t.primary, 0.08)

        # Card with thin border
        add_rect(slide, left, top, width, height, fill_color=card_fill,
                 line_color=tint(t.primary, 0.2), line_width=0.75, corner_radius=60000)

        # Accent bar top (thin, neon)
        add_rect(slide, left + 1, top, width - 2, 50000, fill_color=accent)

        # Label
        add_textbox(slide, label.upper(), left + 80000, top + 300000, width - 160000, 300000,
                    font_size=11, bold=True, color=tint(t.primary, 0.5),
                    alignment=PP_ALIGN.CENTER)

        # Neon value
        add_textbox(slide, value, left + 40000, top + 700000, width - 80000, 800000,
                    font_size=38, bold=True, color=accent, alignment=PP_ALIGN.CENTER)

        # Trend with glow effect (lighter tint behind)
        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", left + 80000, top + 1550000,
                    width - 160000, 350000,
                    font_size=14, bold=True, color=tc, alignment=PP_ALIGN.CENTER)

        # Progress bar on dark bg
        bar_left = left + 120000
        bar_top = top + 2100000
        bar_w = width - 240000
        _progress_bar(slide, bar_left, bar_top, bar_w, 80000, progress, accent,
                      tint(t.primary, 0.15))

        # Progress text
        add_textbox(slide, f"{int(progress * 100)}%", left + 80000, bar_top + 150000,
                    width - 160000, 250000,
                    font_size=10, color=tint(t.primary, 0.4), alignment=PP_ALIGN.CENTER)

        # Separator line near bottom
        add_line(slide, left + 150000, top + height - 600000,
                 left + width - 150000, top + height - 600000,
                 color=tint(t.primary, 0.15), width=0.5)

        # Context text
        add_textbox(slide, "vs. previous quarter", left + 80000, top + height - 500000,
                    width - 160000, 250000,
                    font_size=9, color=tint(t.primary, 0.35), alignment=PP_ALIGN.CENTER)


# ── Variant 6: SPLIT — left/right split, 2 KPIs per side ────────────────

def _split_kpi(slide, t, c):
    content_top = add_slide_title(slide, c.get("kpi_title", "KPI Dashboard"), theme=t)
    kpis = _get_kpis(c)
    p = t.palette

    mid = SLIDE_W // 2
    divider_x = mid - 30000

    # Center divider line
    add_line(slide, divider_x, content_top + 200000, divider_x, FOOTER_TOP - 300000,
             color=t.accent, width=2)

    card_h = 1800000
    gap = 250000
    left_margin = MARGIN
    right_margin = mid + 60000
    card_w_left = mid - MARGIN - 100000
    card_w_right = SLIDE_W - mid - MARGIN - 60000

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:4]):
        accent = p[i % len(p)]
        is_left = i < 2
        row = i % 2

        if is_left:
            card_left = left_margin
            card_w = card_w_left
        else:
            card_left = right_margin
            card_w = card_w_right

        card_top = content_top + 300000 + row * (card_h + gap)

        add_styled_card(slide, card_left, card_top, card_w, card_h,
                        theme=t, accent_color=accent)

        # Label
        add_textbox(slide, label, card_left + 120000, card_top + 150000,
                    card_w - 240000, 280000,
                    font_size=12, bold=True, color=t.secondary, alignment=PP_ALIGN.LEFT)

        # Value (right-aligned large)
        add_textbox(slide, value, card_left + 120000, card_top + 500000,
                    card_w - 240000, 600000,
                    font_size=34, bold=True, color=t.primary, alignment=PP_ALIGN.RIGHT)

        # Trend
        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", card_left + 120000, card_top + 1100000,
                    card_w // 2, 300000,
                    font_size=13, bold=True, color=tc, alignment=PP_ALIGN.LEFT)

        # Progress bar
        bar_left = card_left + card_w // 2 + 60000
        bar_w = card_w // 2 - 180000
        _progress_bar(slide, bar_left, card_top + 1200000, bar_w, 80000,
                      progress, accent, tint(accent, 0.8))


# ── Variant 7: GEO — sharp cards, bottom accent bar, diamond markers ─────

def _angular_cards(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("kpi_title", "KPI Dashboard"), theme=t)
    kpis = _get_kpis(c)
    p = t.palette

    positions = card_positions(4, 1, top=content_top + 200000, gap=180000, card_height=4000000)

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:4]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        # Sharp card with bottom accent
        add_styled_card(slide, left, top, width, height, theme=t, accent_color=accent)

        # Diamond marker at top center
        add_diamond(slide, left + width // 2, top + 350000, 100000, fill_color=accent)

        text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
        sub_color = tint(t.primary, 0.5) if t.ux_style.dark_mode else t.secondary

        # Label
        add_textbox(slide, label.upper(), left + 60000, top + 600000,
                    width - 120000, 300000,
                    font_size=11, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER)

        # Value
        add_textbox(slide, value, left + 40000, top + 950000, width - 80000, 800000,
                    font_size=36, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)

        # Trend
        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", left + 60000, top + 1750000,
                    width - 120000, 350000,
                    font_size=13, bold=True, color=tc, alignment=PP_ALIGN.CENTER)

        # Angular progress bar (sharp corners)
        bar_left = left + 100000
        bar_top = top + 2250000
        bar_w = width - 200000
        add_rect(slide, bar_left, bar_top, bar_w, 80000,
                 fill_color=tint(accent, 0.8 if not t.ux_style.dark_mode else 0.15))
        fill_w = int(bar_w * min(progress, 1.0))
        if fill_w > 0:
            add_rect(slide, bar_left, bar_top, fill_w, 80000, fill_color=accent)

        # Progress text
        add_textbox(slide, f"{int(progress * 100)}% of target", left + 60000, bar_top + 150000,
                    width - 120000, 250000,
                    font_size=10, color=sub_color, alignment=PP_ALIGN.CENTER)

        # Decorative bottom line
        add_line(slide, left + 100000, top + height - 400000,
                 left + width - 100000, top + height - 400000,
                 color=accent, width=0.5)


# ── Variant 8: EDITORIAL — clean left-aligned stats with thin rules ──────

def _editorial_stats(slide, t, c):
    content_top = add_slide_title(slide, c.get("kpi_title", "KPI Dashboard"), theme=t)
    kpis = _get_kpis(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    # Two-column layout, 2 KPIs per column
    col_w = (SLIDE_W - 2 * m - 400000) // 2
    stat_h = 1600000
    gap_v = 300000

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:4]):
        col = i // 2
        row = i % 2
        left = m + col * (col_w + 400000)
        top = content_top + 300000 + row * (stat_h + gap_v)

        # Thin top rule
        add_line(slide, left, top, left + col_w, top,
                 color=tint(t.secondary, 0.7), width=0.5)

        # Short accent mark
        add_line(slide, left, top, left + 500000, top, color=t.accent, width=1.5)

        # Label
        add_textbox(slide, label, left, top + 120000, col_w, 300000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.LEFT)

        # Value — left-aligned, large
        add_textbox(slide, value, left, top + 450000, col_w, 600000,
                    font_size=36, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT)

        # Trend — left-aligned beside value
        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", left, top + 1050000,
                    col_w // 2, 300000,
                    font_size=12, bold=True, color=tc, alignment=PP_ALIGN.LEFT)

        # Progress as text on the right side
        add_textbox(slide, f"{int(progress * 100)}% of target", left + col_w // 2, top + 1050000,
                    col_w // 2, 300000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.RIGHT)

    # Bottom rule
    add_line(slide, m, FOOTER_TOP - 400000, SLIDE_W - m, FOOTER_TOP - 400000,
             color=tint(t.secondary, 0.8), width=0.5)


# ── Variant 9: GRADIENT — rounded cards, no accent bar, soft shadow ──────

def _gradient_cards(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("kpi_title", "KPI Dashboard"), theme=t)
    kpis = _get_kpis(c)
    p = t.palette

    positions = card_positions(4, 1, top=content_top + 200000, gap=250000, card_height=4000000)

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:4]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        # Card with no accent bar, soft shadow, very rounded
        add_styled_card(slide, left, top, width, height, theme=t)

        text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
        sub_color = tint(t.primary, 0.55) if t.ux_style.dark_mode else t.secondary

        # Subtle color circle behind value
        circle_x = left + width // 2
        circle_y = top + 1100000
        add_circle(slide, circle_x, circle_y, 500000, fill_color=tint(accent, 0.88))

        # Label
        add_textbox(slide, label, left + 80000, top + 300000, width - 160000, 300000,
                    font_size=12, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER)

        # Value centered
        add_textbox(slide, value, left + 40000, top + 750000, width - 80000, 700000,
                    font_size=36, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)

        # Trend
        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", left + 80000, top + 1550000,
                    width - 160000, 300000,
                    font_size=13, bold=True, color=tc, alignment=PP_ALIGN.CENTER)

        # Progress circle-text (no bar, display as large percentage)
        add_textbox(slide, f"{int(progress * 100)}%", left + 80000, top + 2000000,
                    width - 160000, 500000,
                    font_size=24, bold=True, color=tint(accent, 0.3),
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, "of target", left + 80000, top + 2400000,
                    width - 160000, 250000,
                    font_size=10, color=sub_color, alignment=PP_ALIGN.CENTER)


# ── Variant 10: RETRO — rounded bordered cards with decorative elements ──

def _retro_badges(slide, t, c):
    add_background(slide, t.light_bg)
    content_top = add_slide_title(slide, c.get("kpi_title", "KPI Dashboard"), theme=t)
    kpis = _get_kpis(c)
    p = t.palette

    positions = card_positions(4, 1, top=content_top + 200000, gap=220000, card_height=4000000)

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:4]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        # Double-bordered card
        add_rect(slide, left, top, width, height, fill_color="#FFFFFF",
                 line_color=accent, line_width=2.0, corner_radius=100000)
        add_rect(slide, left + 40000, top + 40000, width - 80000, height - 80000,
                 line_color=accent, line_width=0.75, corner_radius=80000)

        # Decorative circles in corners
        add_circle(slide, left + 120000, top + 120000, 40000, fill_color=accent)
        add_circle(slide, left + width - 120000, top + 120000, 40000, fill_color=accent)

        # Label with decorative dashes
        add_textbox(slide, f"\u2014 {label} \u2014", left + 60000, top + 350000,
                    width - 120000, 350000,
                    font_size=12, bold=True, color=accent, alignment=PP_ALIGN.CENTER)

        # Value
        add_textbox(slide, value, left + 40000, top + 800000, width - 80000, 800000,
                    font_size=36, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)

        # Decorative rule
        rule_left = left + width // 4
        add_line(slide, rule_left, top + 1650000, rule_left + width // 2, top + 1650000,
                 color=accent, width=1.5)

        # Trend
        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", left + 80000, top + 1800000,
                    width - 160000, 350000,
                    font_size=14, bold=True, color=tc, alignment=PP_ALIGN.CENTER)

        # Progress bar
        bar_left = left + 120000
        bar_top = top + 2300000
        bar_w = width - 240000
        _progress_bar(slide, bar_left, bar_top, bar_w, 100000, progress, accent,
                      tint(accent, 0.7))

        # Progress text
        add_textbox(slide, f"{int(progress * 100)}%", left + 80000, bar_top + 180000,
                    width - 160000, 250000,
                    font_size=10, color=t.secondary, alignment=PP_ALIGN.CENTER)


# ── Variant 11: MAGAZINE — full-width color blocks stacked ──────────────

def _magazine_blocks(slide, t, c):
    content_top = add_slide_title(slide, c.get("kpi_title", "KPI Dashboard"), theme=t)
    kpis = _get_kpis(c)
    p = t.palette

    n = min(len(kpis), 4)
    available_h = FOOTER_TOP - content_top - 200000
    block_h = available_h // n
    gap = 50000
    actual_h = block_h - gap

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:n]):
        accent = p[i % len(p)]
        block_top = content_top + 100000 + i * block_h
        block_left = MARGIN

        # Full-width color block
        add_rect(slide, block_left, block_top, SLIDE_W - 2 * MARGIN, actual_h,
                 fill_color=accent)

        # Label on left
        add_textbox(slide, label.upper(), block_left + 200000, block_top + 80000,
                    2500000, actual_h - 160000,
                    font_size=13, bold=True, color="#FFFFFF", alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Value centered
        val_left = block_left + 2800000
        add_textbox(slide, value, val_left, block_top + 80000,
                    3000000, actual_h - 160000,
                    font_size=36, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Trend on right
        tc_bg = shade(accent, 0.15)
        trend_left = val_left + 3200000
        trend_w = 1800000
        add_rect(slide, trend_left, block_top + (actual_h - 350000) // 2,
                 trend_w, 350000, fill_color=tc_bg, corner_radius=50000)
        add_textbox(slide, f"{arrow} {trend_val}", trend_left, block_top + 80000,
                    trend_w, actual_h - 160000,
                    font_size=14, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Progress on far right
        pct_left = trend_left + trend_w + 200000
        pct_w = SLIDE_W - MARGIN - pct_left - 100000
        if pct_w > 400000:
            add_textbox(slide, f"{int(progress * 100)}%", pct_left, block_top + 80000,
                        pct_w, actual_h - 160000,
                        font_size=20, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant 12: SCHOLARLY — white bg, centered text, thin rules ──────────

def _scholarly_stats(slide, t, c):
    add_background(slide, "#FFFFFF")
    content_top = add_slide_title(slide, c.get("kpi_title", "Key Metrics"), theme=t)
    kpis = _get_kpis(c)

    n = min(len(kpis), 4)
    total_w = SLIDE_W - 2 * MARGIN
    col_w = total_w // n
    stat_top = content_top + 400000

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:n]):
        left = MARGIN + i * col_w
        cx = left + col_w // 2

        # Thin vertical separator between columns
        if i > 0:
            add_line(slide, left, stat_top, left, stat_top + 3000000,
                     color=tint(t.secondary, 0.8), width=0.5)

        # Figure caption "Figure N:" style label
        add_textbox(slide, f"Figure {i + 1}:", left + 60000, stat_top, col_w - 120000, 250000,
                    font_size=9, bold=True, color=t.secondary, alignment=PP_ALIGN.CENTER)

        # Thin rule under caption
        rule_left = left + col_w // 4
        rule_right = left + col_w - col_w // 4
        add_line(slide, rule_left, stat_top + 280000, rule_right, stat_top + 280000,
                 color=tint(t.secondary, 0.7), width=0.5)

        # Label
        add_textbox(slide, label, left + 40000, stat_top + 400000, col_w - 80000, 300000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.CENTER)

        # Large centered value
        add_textbox(slide, value, left + 20000, stat_top + 800000, col_w - 40000, 700000,
                    font_size=38, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)

        # Thin rule below value
        add_line(slide, rule_left, stat_top + 1600000, rule_right, stat_top + 1600000,
                 color=tint(t.secondary, 0.7), width=0.5)

        # Trend
        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", left + 60000, stat_top + 1700000,
                    col_w - 120000, 300000,
                    font_size=12, color=tc, alignment=PP_ALIGN.CENTER)

        # Progress as simple text
        add_textbox(slide, f"({int(progress * 100)}% of target)", left + 60000, stat_top + 2050000,
                    col_w - 120000, 250000,
                    font_size=9, color=t.secondary, alignment=PP_ALIGN.CENTER)

    # Bottom rule
    add_line(slide, MARGIN, FOOTER_TOP - 400000, SLIDE_W - MARGIN, FOOTER_TOP - 400000,
             color=tint(t.secondary, 0.7), width=0.5)


# ── Variant 13: LABORATORY — dark bg, glowing accent top borders ─────────

def _laboratory_metrics(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("kpi_title", "Key Metrics").upper(), theme=t)
    kpis = _get_kpis(c)
    p = t.palette

    positions = card_positions(4, 1, top=content_top + 200000, gap=200000, card_height=4000000)

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:4]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]
        card_fill = tint(t.primary, 0.08)

        # Dark card with thin border
        add_rect(slide, left, top, width, height, fill_color=card_fill,
                 line_color=tint(t.primary, 0.2), line_width=0.5, corner_radius=50000)

        # Glowing accent top border (double layer for glow effect)
        add_rect(slide, left + 1, top, width - 2, 60000,
                 fill_color=tint(accent, 0.3))
        add_rect(slide, left + 1, top, width - 2, 35000,
                 fill_color=accent)

        # Label
        add_textbox(slide, label.upper(), left + 80000, top + 300000, width - 160000, 300000,
                    font_size=10, bold=True, color=tint(t.primary, 0.45),
                    alignment=PP_ALIGN.CENTER)

        # Large value in accent color
        add_textbox(slide, value, left + 40000, top + 700000, width - 80000, 800000,
                    font_size=40, bold=True, color=accent, alignment=PP_ALIGN.CENTER)

        # Trend
        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", left + 80000, top + 1550000,
                    width - 160000, 350000,
                    font_size=14, bold=True, color=tc, alignment=PP_ALIGN.CENTER)

        # Progress bar on dark bg
        bar_left = left + 120000
        bar_top = top + 2100000
        bar_w = width - 240000
        _progress_bar(slide, bar_left, bar_top, bar_w, 70000, progress, accent,
                      tint(t.primary, 0.15))

        # Progress text
        add_textbox(slide, f"{int(progress * 100)}%", left + 80000, bar_top + 130000,
                    width - 160000, 250000,
                    font_size=10, color=tint(t.primary, 0.4), alignment=PP_ALIGN.CENTER)

        # Bottom separator
        add_line(slide, left + 120000, top + height - 500000,
                 left + width - 120000, top + height - 500000,
                 color=tint(t.primary, 0.15), width=0.5)

        # Context label
        add_textbox(slide, "vs. previous quarter", left + 80000, top + height - 400000,
                    width - 160000, 250000,
                    font_size=9, color=tint(t.primary, 0.35), alignment=PP_ALIGN.CENTER)


# ── Variant 14: DASHBOARD — dense 2×N grid, small tiles, header band ─────

def _dashboard_tiles(slide, t, c):
    add_background(slide, "#FFFFFF")
    content_top = add_slide_title(slide, c.get("kpi_title", "KPI Dashboard"), theme=t)
    kpis = _get_kpis(c)
    p = t.palette

    # Header accent band
    band_h = 45000
    add_rect(slide, 0, content_top, SLIDE_W, band_h, fill_color=t.accent)

    n = min(len(kpis), 8)
    cols = min(n, 4)
    rows = (n + cols - 1) // cols
    m = int(MARGIN * 0.8)
    gap = 120000
    total_w = SLIDE_W - 2 * m
    tile_w = (total_w - (cols - 1) * gap) // cols
    available_h = FOOTER_TOP - content_top - band_h - 400000
    tile_h = (available_h - (rows - 1) * gap) // rows
    grid_top = content_top + band_h + 150000

    for i, (label, value, trend_val, progress, arrow) in enumerate(kpis[:n]):
        accent = p[i % len(p)]
        col = i % cols
        row = i // cols
        tile_left = m + col * (tile_w + gap)
        tile_top = grid_top + row * (tile_h + gap)

        # Compact tile with accent left bar
        add_rect(slide, tile_left, tile_top, tile_w, tile_h,
                 fill_color=tint(accent, 0.95), corner_radius=40000)
        add_rect(slide, tile_left, tile_top + 20000, 30000, tile_h - 40000,
                 fill_color=accent)

        # Label — small
        add_textbox(slide, label, tile_left + 60000, tile_top + 60000,
                    tile_w - 120000, 220000,
                    font_size=9, bold=True, color=t.secondary, alignment=PP_ALIGN.LEFT)

        # Value — prominent
        add_textbox(slide, value, tile_left + 60000, tile_top + 300000,
                    tile_w - 120000, 500000,
                    font_size=28, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT)

        # Trend
        tc = _trend_color(t, arrow)
        add_textbox(slide, f"{arrow} {trend_val}", tile_left + 60000, tile_top + 800000,
                    (tile_w - 120000) // 2, 250000,
                    font_size=11, bold=True, color=tc, alignment=PP_ALIGN.LEFT)

        # Progress bar
        bar_left = tile_left + tile_w // 2
        bar_w = tile_w // 2 - 80000
        bar_top_pos = tile_top + 870000
        _progress_bar(slide, bar_left, bar_top_pos, bar_w, 60000, progress, accent,
                      tint(accent, 0.8))

        # Progress text
        add_textbox(slide, f"{int(progress * 100)}%", bar_left, bar_top_pos + 80000,
                    bar_w, 200000,
                    font_size=8, color=t.secondary, alignment=PP_ALIGN.RIGHT)
