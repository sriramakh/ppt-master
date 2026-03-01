"""AI-powered presentation builder — LLM picks slides, generates content, builds PPTX."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation as PptxPresentation
from pptx.util import Emu

from pptmaster.builder.design_system import SLIDE_W, SLIDE_H
from pptmaster.builder.helpers import add_themed_footer, add_themed_slide_number
from pptmaster.builder.themes import TemplateTheme, get_theme, THEME_MAP

# ── Slide-type → builder import mapping ──────────────────────────────

_BUILDER_MAP: dict[str, str] = {
    "company_overview": "pptmaster.builder.slides.s04_company_overview",
    "our_values": "pptmaster.builder.slides.s05_our_values",
    "team_leadership": "pptmaster.builder.slides.s06_team_leadership",
    "key_facts": "pptmaster.builder.slides.s07_key_facts",
    "sources": "pptmaster.builder.slides.s08_sources",
    "executive_summary": "pptmaster.builder.slides.s09_executive_summary",
    "kpi_dashboard": "pptmaster.builder.slides.s10_kpi_dashboard",
    "process_linear": "pptmaster.builder.slides.s11_process_linear",
    "process_circular": "pptmaster.builder.slides.s12_process_circular",
    "roadmap_timeline": "pptmaster.builder.slides.s13_roadmap_timeline",
    "swot_matrix": "pptmaster.builder.slides.s14_swot_matrix",
    "bar_chart": "pptmaster.builder.slides.s16_bar_chart",
    "line_chart": "pptmaster.builder.slides.s17_line_chart",
    "pie_chart": "pptmaster.builder.slides.s18_pie_chart",
    "comparison": "pptmaster.builder.slides.s19_comparison",
    "data_table": "pptmaster.builder.slides.s20_data_table",
    "two_column": "pptmaster.builder.slides.s23_two_column",
    "three_column": "pptmaster.builder.slides.s24_three_column",
    "highlight_quote": "pptmaster.builder.slides.s25_highlight_quote",
    "infographic_dashboard": "pptmaster.builder.slides.s26_infographic_dashboard",
    "next_steps": "pptmaster.builder.slides.s28_next_steps",
    "call_to_action": "pptmaster.builder.slides.s29_call_to_action",
    # New slide types (s31-s40)
    "funnel_diagram": "pptmaster.builder.slides.s31_funnel_diagram",
    "pyramid_hierarchy": "pptmaster.builder.slides.s32_pyramid_hierarchy",
    "venn_diagram": "pptmaster.builder.slides.s33_venn_diagram",
    "hub_spoke": "pptmaster.builder.slides.s34_hub_spoke",
    "milestone_roadmap": "pptmaster.builder.slides.s35_milestone_roadmap",
    "kanban_board": "pptmaster.builder.slides.s36_kanban_board",
    "matrix_quadrant": "pptmaster.builder.slides.s37_matrix_quadrant",
    "gauge_dashboard": "pptmaster.builder.slides.s38_gauge_dashboard",
    "icon_grid": "pptmaster.builder.slides.s39_icon_grid",
    "risk_matrix": "pptmaster.builder.slides.s40_risk_matrix",
}

# Slides that take company_name as kwarg
_COMPANY_NAME_SLIDES = {"company_overview", "call_to_action"}


def _get_builder(slide_type: str):
    """Lazily import and return the build() function for a slide type."""
    import importlib
    module_path = _BUILDER_MAP[slide_type]
    mod = importlib.import_module(module_path)
    return mod.build


def _build_selective_pptx(
    gen_result: dict[str, Any],
    theme: TemplateTheme,
    company_name: str,
    output_path: Path,
) -> Path:
    """Build a PPTX with only the AI-selected slides.

    Slide order:
      1. Cover
      2. TOC (auto-generated from sections)
      For each section:
        3. Section divider
        4..N. Content slides in the section
      N+1. Thank You
    """
    from pptmaster.builder.slides.s01_cover import build as build_cover
    from pptmaster.builder.slides.s02_toc import build as build_toc
    from pptmaster.builder.slides.s03_section_divider import build as build_divider
    from pptmaster.builder.slides.s30_thank_you import build as build_thankyou

    sections = gen_result["sections"]
    content = gen_result["content"]

    # Inject content into theme
    theme.content = content
    theme.company_name = company_name
    if content.get("cover_subtitle"):
        theme.tagline = content["cover_subtitle"]

    # Build TOC sections list for the TOC slide
    toc_sections = []
    for i, sec in enumerate(sections):
        toc_sections.append([f"{i+1:02d}", sec["title"], sec["subtitle"]])
    theme.content["toc_sections"] = toc_sections

    # Create presentation
    prs = PptxPresentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    blank_layout = prs.slide_layouts[6]
    style = theme.ux_style

    slide_num = 0

    def _add_slide(builder_fn, is_dark: bool, **kwargs):
        nonlocal slide_num
        slide_num += 1
        slide = prs.slides.add_slide(blank_layout)
        builder_fn(slide, theme=theme, **kwargs)
        effective_dark = is_dark or style.dark_mode
        if slide_num > 1:
            add_themed_footer(slide, company_name=company_name,
                              dark_bg=effective_dark, theme=theme)
            add_themed_slide_number(slide, slide_num, dark_bg=effective_dark,
                                   theme=theme)

    # 1. Cover
    _add_slide(build_cover, True, company_name=company_name)

    # 2. TOC
    _add_slide(build_toc, False, company_name=company_name)

    # 3. Sections with dividers
    for i, sec in enumerate(sections):
        # Section divider
        _add_slide(
            build_divider, True,
            section_number=f"{i+1:02d}",
            section_title=sec["title"],
            subtitle=sec["subtitle"],
        )
        # Content slides in this section
        for slide_type in sec["slides"]:
            if slide_type not in _BUILDER_MAP:
                continue
            builder_fn = _get_builder(slide_type)
            kwargs = {}
            if slide_type in _COMPANY_NAME_SLIDES:
                kwargs["company_name"] = company_name
            _add_slide(builder_fn, False, **kwargs)

    # 4. Thank You
    _add_slide(build_thankyou, False, company_name=company_name)

    output_path = Path(output_path)
    prs.save(str(output_path))
    return output_path


# ── Public API ────────────────────────────────────────────────────────

def ai_build_presentation(
    topic: str,
    company_name: str = "Acme Corp",
    industry: str = "",
    audience: str = "",
    theme_key: str = "corporate",
    output_path: str | Path = "output.pptx",
    additional_context: str = "",
    provider: str = "minimax",
    api_key: str | None = None,
) -> Path:
    """Generate a presentation with AI-selected slides and content.

    The LLM decides which of the 22 available slide types to include,
    organizes them into sections, and generates all content.  Cover,
    TOC, section dividers, and Thank You are always included.

    Returns the output Path.
    """
    from pptmaster.content.builder_content_gen import generate_builder_content

    # 1. Load base theme (colors, UX style, font)
    theme = get_theme(theme_key)

    # 2. Generate: LLM picks slides + writes content
    gen_result = generate_builder_content(
        topic=topic,
        company_name=company_name,
        industry=industry or theme.industry,
        audience=audience,
        additional_context=additional_context,
        provider=provider,
        api_key=api_key,
    )

    selected = gen_result["selected_slides"]
    n_sections = len(gen_result["sections"])
    # total = cover + toc + n_sections dividers + len(selected) content + thank_you
    total = 2 + n_sections + len(selected) + 1
    print(f"  AI selected {len(selected)} content slides in {n_sections} sections ({total} total)")

    # 3. Build PPTX
    return _build_selective_pptx(gen_result, theme, company_name, Path(output_path))


def build_from_content(
    content_dict: dict[str, Any],
    company_name: str = "Acme Corp",
    theme_key: str = "corporate",
    output_path: str | Path = "output.pptx",
) -> Path:
    """Build a PPTX from a pre-made content dict — no LLM call needed.

    The content_dict must follow the same structure that the AI content
    generator produces:

        {
            "selected_slides": ["company_overview", "key_facts", ...],
            "sections": [
                {"title": "...", "subtitle": "...", "slides": ["company_overview", ...]},
                ...
            ],
            "content": {
                "cover_title": "...",
                "overview_mission": "...",
                ...all slide content keys...
            }
        }

    If ``sections`` is omitted, a single default section wrapping all
    selected slides is created automatically.

    Returns the output Path.
    """
    from pptmaster.content.builder_content_gen import _validate_content

    theme = get_theme(theme_key)

    # Normalise structure
    selected = content_dict.get("selected_slides", [])
    sections = content_dict.get("sections", [])
    content = content_dict.get("content", content_dict)

    # If no explicit sections, auto-wrap all slides in one section
    if not sections:
        sections = [{
            "title": content.get("cover_title", "Overview"),
            "subtitle": "",
            "slides": selected,
        }]

    # Validate and clamp content
    content = _validate_content({"content": content})

    gen_result = {
        "selected_slides": selected,
        "sections": sections,
        "content": content,
    }

    n_sections = len(sections)
    total = 2 + n_sections + len(selected) + 1
    print(f"  Building {len(selected)} content slides in {n_sections} sections ({total} total)")

    return _build_selective_pptx(gen_result, theme, company_name, Path(output_path))


def list_available_themes() -> list[dict[str, Any]]:
    """Return metadata for each of the 11 available themes."""
    results = []
    for key, theme in THEME_MAP.items():
        results.append({
            "key": theme.key,
            "industry": theme.industry,
            "company_name": theme.company_name,
            "tagline": theme.tagline,
            "primary_color": theme.primary,
            "accent_color": theme.accent,
            "font": theme.font,
            "ux_style": theme.ux_style.name,
        })
    return results
