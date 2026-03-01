"""Catalog and classify all slide layouts from a template.

Uses a 16-rule classifier based on placeholder types and layout name patterns
to assign each layout a content_category for matching during composition.
"""

from __future__ import annotations

import re
from typing import Any

from pptx import Presentation
from pptx.util import Emu

from pptmaster.models import LayoutInfo, PlaceholderInfo


# ── Placeholder type mapping ───────────────────────────────────────────

_PH_TYPE_MAP = {
    0: "title",       # PP_PLACEHOLDER.TITLE
    1: "generic",     # PP_PLACEHOLDER.CENTER_TITLE / BODY / OBJECT
    2: "body",        # PP_PLACEHOLDER.BODY
    3: "body",        # PP_PLACEHOLDER.SUBTITLE (treat as body)
    5: "body",        # PP_PLACEHOLDER.FOOTER (body subtype)
    6: "sldNum",      # PP_PLACEHOLDER.SLIDE_NUMBER
    7: "body",        # PP_PLACEHOLDER.BODY (variant)
    10: "dt",         # PP_PLACEHOLDER.DATE
    11: "ftr",        # PP_PLACEHOLDER.FOOTER
    12: "sldNum",     # PP_PLACEHOLDER.SLIDE_NUMBER
    13: "pic",        # PP_PLACEHOLDER.PICTURE
    14: "body",       # PP_PLACEHOLDER.OBJECT
    15: "title",      # PP_PLACEHOLDER.TITLE
    18: "pic",        # PP_PLACEHOLDER.PICTURE (variant)
}


def _extract_placeholder(ph: Any) -> PlaceholderInfo:
    """Extract placeholder info from a python-pptx placeholder shape."""
    ph_type_raw = ph.placeholder_format.type
    # Map the enumeration to our simpler type string
    if ph_type_raw is not None:
        ph_type = _PH_TYPE_MAP.get(int(ph_type_raw), "generic")
    else:
        ph_type = "generic"

    # Override based on XML inspection for more accuracy
    sp_elem = ph._element
    nvPr = sp_elem.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr")
    if nvPr is not None:
        ph_elem = nvPr.find("{http://schemas.openxmlformats.org/presentationml/2006/main}ph")
        if ph_elem is not None:
            xml_type = ph_elem.get("type", "")
            if xml_type == "title" or xml_type == "ctrTitle":
                ph_type = "title"
            elif xml_type == "body" or xml_type == "subTitle":
                ph_type = "body"
            elif xml_type == "pic":
                ph_type = "pic"
            elif xml_type == "ftr":
                ph_type = "ftr"
            elif xml_type == "dt":
                ph_type = "dt"
            elif xml_type == "sldNum":
                ph_type = "sldNum"

    is_vertical = False
    txBody = sp_elem.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr")
    if txBody is not None and txBody.get("vert") in ("vert", "vert270", "wordArtVert"):
        is_vertical = True

    return PlaceholderInfo(
        idx=ph.placeholder_format.idx,
        ph_type=ph_type,
        name=ph.name or "",
        left=ph.left or 0,
        top=ph.top or 0,
        width=ph.width or 0,
        height=ph.height or 0,
        is_vertical=is_vertical,
    )


# ── Layout classification rules ───────────────────────────────────────

_CATEGORY_PATTERNS: list[tuple[str, re.Pattern[str] | None, Any]] = [
    # (category, name_pattern, placeholder_check_func)
]


def _classify_layout(name: str, placeholders: list[PlaceholderInfo]) -> str:
    """Classify a layout into a content category using 16 rules.

    Categories: title, content_text, content_image_right, content_image_left,
    multi_column_2, multi_column_3, multi_column_4, device_mockup,
    thank_you, divider, profile_grid, hexagon, full_image,
    mission_vision, blank_canvas, utility
    """
    name_lower = name.lower()

    # Count placeholder types (excluding footer/date/slidenum)
    content_phs = [p for p in placeholders if p.ph_type not in ("ftr", "dt", "sldNum")]
    title_phs = [p for p in content_phs if p.ph_type == "title"]
    body_phs = [p for p in content_phs if p.ph_type == "body"]
    pic_phs = [p for p in content_phs if p.ph_type == "pic"]
    generic_phs = [p for p in content_phs if p.ph_type == "generic"]

    # Rule 1: Title slides
    if re.search(r"\btitle\s*[123]?\b", name_lower) and not re.search(r"(only|solo|imagen|text|vertical)", name_lower):
        if len(content_phs) <= 3 and len(body_phs) <= 1:
            return "title"

    # Rule 2: Thank you slides
    if re.search(r"thank\s*you", name_lower):
        return "thank_you"

    # Rule 3: Divider slides
    if re.search(r"divider|encabezado\s*de\s*secci", name_lower):
        return "divider"

    # Rule 4: Mission-vision
    if re.search(r"mission.*vision", name_lower):
        return "mission_vision"

    # Rule 5: Hexagon layout
    if "hexagon" in name_lower:
        return "hexagon"

    # Rule 6: Full image
    if re.search(r"full\s*image", name_lower):
        return "full_image"

    # Rule 7: Profile grid
    if re.search(r"(profile|multiple.*profile)", name_lower):
        return "profile_grid"

    # Rule 8: Device mockup (laptop)
    if re.search(r"laptop|device|phone|tablet", name_lower):
        return "device_mockup"

    # Rule 9: Blank canvas
    if re.search(r"\bblank\b|en\s*blanco", name_lower):
        return "blank_canvas"

    # Rule 10: Title only / Solo el titulo
    if re.search(r"(title|titulo)\s*(only|solo)", name_lower) or re.search(r"solo\s*el\s*titulo", name_lower):
        return "blank_canvas"

    # Rule 11: 4 columns
    if re.search(r"4\s*column", name_lower) or len(pic_phs) >= 4:
        return "multi_column_4"

    # Rule 12: 3 columns
    if re.search(r"3\s*column", name_lower) or (len(pic_phs) == 3 and len(body_phs) >= 3):
        return "multi_column_3"

    # Rule 13: Half and 2 columns / 2-object layouts
    if re.search(r"(half.*column|dos\s*objetos|comparacion|two\s*level)", name_lower):
        return "multi_column_2"

    # Rule 14: Image right
    if re.search(r"image\s*right|imagen.*right", name_lower):
        return "content_image_right"

    # Rule 15: Image left / half image
    if re.search(r"image\s*left|half\s*image|imagen.*left", name_lower):
        return "content_image_left"

    # Rule 16: Text-only / content with generic placeholders
    if re.search(r"text\s*only|contenido|picture.*caption|titulo.*texto", name_lower):
        return "content_text"

    # Fallback: classify by placeholder composition
    if len(pic_phs) >= 1 and len(body_phs) >= 1:
        return "content_image_right"
    if len(body_phs) >= 2 or len(generic_phs) >= 2:
        return "multi_column_2"
    if len(body_phs) == 1 or len(generic_phs) == 1:
        return "content_text"

    return "utility"


def extract_layouts(prs: Presentation) -> list[LayoutInfo]:
    """Extract and classify all layouts from a presentation."""
    layouts: list[LayoutInfo] = []

    for master in prs.slide_masters:
        for i, layout in enumerate(master.slide_layouts):
            phs: list[PlaceholderInfo] = []
            for ph in layout.placeholders:
                phs.append(_extract_placeholder(ph))

            name = layout.name or f"Layout {i}"
            category = _classify_layout(name, phs)

            # Check for background fill
            has_bg = False
            bg_color = ""
            bg_elem = layout._element.find(
                ".//{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
            )
            if bg_elem is not None:
                solid = bg_elem.find(
                    ".//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
                )
                if solid is not None:
                    has_bg = True
                    bg_color = f"#{solid.get('val', '')}"

            layouts.append(LayoutInfo(
                index=i,
                name=name,
                content_category=category,
                placeholders=phs,
                has_background_fill=has_bg,
                background_color=bg_color,
            ))

    return layouts
