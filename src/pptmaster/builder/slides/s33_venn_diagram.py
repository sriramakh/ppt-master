"""Slide 33 — Venn Diagram: 14 unique visual variants dispatched by UX style."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_circle, add_line, add_hexagon,
    add_semi_transparent_circle, add_gold_accent_line, add_slide_title,
    add_dark_bg, add_background, add_styled_card,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    s = t.ux_style
    c = t.content
    p = t.palette

    dispatch = {
        "overlapping": _overlapping,
        "labeled-rings": _labeled_rings,
        "bold-overlap": _bold_overlap,
        "soft-overlap": _soft_overlap,
        "glow-overlap": _glow_overlap,
        "split-venn": _split_venn,
        "hex-overlap": _hex_overlap,
        "editorial-venn": _editorial_venn,
        "gradient-venn": _gradient_venn,
        "retro-venn": _retro_venn,
        "creative-venn": _creative_venn,
        "scholarly-venn": _scholarly_venn,
        "laboratory-venn": _laboratory_venn,
        "dashboard-venn": _dashboard_venn,
    }
    builder = dispatch.get(s.venn, _overlapping)
    builder(slide, t, c, p)


# ── Content helpers ───────────────────────────────────────────────────

def _get_sets(c):
    """Extract venn sets from content dict with sensible defaults."""
    return c.get("venn_sets", [
        ("Innovation", "Cutting-edge technology and R&D"),
        ("Experience", "Deep industry expertise and talent"),
        ("Trust", "Proven track record with clients"),
    ])


def _get_overlap(c):
    """Extract overlap label."""
    return c.get("venn_overlap", "Our Competitive Advantage")


def _circle_centers_2(cx, cy, radius):
    """Return two circle centers for 2-set Venn, horizontally overlapping."""
    offset = int(radius * 0.7)
    return [(cx - offset, cy), (cx + offset, cy)]


def _circle_centers_3(cx, cy, radius):
    """Return three circle centers for 3-set Venn in triangle formation."""
    offset = int(radius * 0.6)
    top_y = cy - int(offset * 0.55)
    bot_y = cy + int(offset * 0.55)
    return [
        (cx, top_y),                          # top center
        (cx - int(offset * 0.85), bot_y),     # bottom left
        (cx + int(offset * 0.85), bot_y),     # bottom right
    ]


def _venn_colors(t, p, n):
    """Return n distinct colors from palette or accent."""
    base = [t.accent] + list(p[:5])
    colors = []
    for i in range(n):
        colors.append(base[i % len(base)])
    return colors


# ── Variant 1: CLASSIC — overlapping semi-transparent circles ────────

def _overlapping(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)

    # Venn area
    area_h = FOOTER_TOP - content_top - 200000
    cx = SLIDE_W // 2
    cy = content_top + area_h // 2
    radius = min(area_h // 2 - 100000, 1400000)

    centers = _circle_centers_3(cx, cy, radius) if n >= 3 else _circle_centers_2(cx, cy, radius)

    for i, (label, desc) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        add_semi_transparent_circle(slide, x, y, radius,
                                    fill_color=colors[i], alpha=55000)
        # Label outside the circle
        if n >= 3:
            # Push labels outward from center
            dx = x - cx
            dy = y - cy
            lx = x + int(dx * 0.9) - 700000
            ly = y + int(dy * 0.9) - 200000
        else:
            lx = x - 700000
            ly = y - radius - 400000
        text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
        add_textbox(slide, label, lx, ly, 1600000, 300000,
                    font_size=12, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, lx - 100000, ly + 300000, 1800000, 350000,
                    font_size=10, color=tint(t.primary, 0.6) if not t.ux_style.dark_mode
                    else tint(t.primary, 0.5),
                    alignment=PP_ALIGN.CENTER)

    # Overlap label in center
    add_textbox(slide, overlap, cx - 1100000, cy - 180000, 2200000, 360000,
                font_size=11, bold=True,
                color="#FFFFFF" if t.ux_style.dark_mode else t.primary,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant 2: MINIMAL — thin circle outlines, clean labels ─────────

def _labeled_rings(slide, t, c, p):
    add_dark_bg(slide, t)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)

    area_h = FOOTER_TOP - content_top - 200000
    cx = SLIDE_W // 2
    cy = content_top + area_h // 2
    radius = min(area_h // 2 - 100000, 1300000)

    centers = _circle_centers_3(cx, cy, radius) if n >= 3 else _circle_centers_2(cx, cy, radius)

    for i, (label, desc) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        # Thin ring outline, no fill
        add_circle(slide, x, y, radius,
                   fill_color=t.light_bg if not s.dark_mode else tint(t.primary, 0.08),
                   line_color=colors[i], line_width=1.5)
        # Clean label at circle center offset
        if n >= 3:
            dx = x - cx
            dy = y - cy
            lx = x + int(dx * 0.5) - 600000
            ly = y + int(dy * 0.5) - 120000
        else:
            lx = x - 600000
            ly = y - 120000
        text_color = t.primary if not s.dark_mode else "#FFFFFF"
        add_textbox(slide, label, lx, ly, 1200000, 240000,
                    font_size=12, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

    # Minimal overlap text
    sub_color = t.secondary if not s.dark_mode else tint(t.primary, 0.6)
    add_textbox(slide, overlap, cx - 900000, cy - 140000, 1800000, 280000,
                font_size=10, color=sub_color, alignment=PP_ALIGN.CENTER)


# ── Variant 3: BOLD — large bold overlapping colored circles ─────────

def _bold_overlap(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)

    area_h = FOOTER_TOP - content_top - 200000
    cx = SLIDE_W // 2
    cy = content_top + area_h // 2
    radius = min(area_h // 2 - 50000, 1500000)

    centers = _circle_centers_3(cx, cy, radius) if n >= 3 else _circle_centers_2(cx, cy, radius)

    for i, (label, desc) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        add_semi_transparent_circle(slide, x, y, radius,
                                    fill_color=colors[i], alpha=45000)
        # Bold label inside circle, pushed outward
        if n >= 3:
            dx = x - cx
            dy = y - cy
            lx = x + int(dx * 0.6) - 700000
            ly = y + int(dy * 0.6) - 200000
        else:
            lx = x - 700000
            ly = y - 200000
        add_textbox(slide, label.upper(), lx, ly, 1400000, 300000,
                    font_size=16, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Bold overlap banner
    banner_w = 2400000
    banner_h = 400000
    add_rect(slide, cx - banner_w // 2, cy - banner_h // 2,
             banner_w, banner_h, fill_color=t.accent)
    add_textbox(slide, overlap.upper(),
                cx - banner_w // 2 + 50000, cy - banner_h // 2 + 20000,
                banner_w - 100000, banner_h - 40000,
                font_size=13, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant 4: ELEVATED — soft shadowed circles, floating labels ─────

def _soft_overlap(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)

    area_h = FOOTER_TOP - content_top - 200000
    cx = SLIDE_W // 2
    cy = content_top + area_h // 2
    radius = min(area_h // 2 - 100000, 1350000)

    centers = _circle_centers_3(cx, cy, radius) if n >= 3 else _circle_centers_2(cx, cy, radius)

    # Draw shadow circles first
    for i in range(len(sets[:len(centers)])):
        x, y = centers[i]
        add_circle(slide, x + 40000, y + 40000, radius,
                   fill_color=shade(t.light_bg, 0.12))

    # Then semi-transparent circles
    for i, (label, desc) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        add_semi_transparent_circle(slide, x, y, radius,
                                    fill_color=tint(colors[i], 0.3), alpha=40000)

        # Floating label card outside
        if n >= 3:
            dx = x - cx
            dy = y - cy
            card_x = x + int(dx * 1.1) - 600000
            card_y = y + int(dy * 1.1) - 180000
        else:
            card_x = x - 600000
            card_y = y - radius - 500000
        add_styled_card(slide, card_x, card_y, 1200000, 360000,
                        theme=t, accent_color=colors[i])
        add_textbox(slide, label, card_x + 40000, card_y + 30000,
                    1120000, 160000,
                    font_size=12, bold=True, color=t.primary,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, card_x + 40000, card_y + 180000,
                    1200000, 320000,
                    font_size=9, color=t.secondary,
                    alignment=PP_ALIGN.CENTER)

    # Overlap label
    add_textbox(slide, overlap, cx - 900000, cy - 140000, 1800000, 280000,
                font_size=10, bold=True, color=t.accent,
                alignment=PP_ALIGN.CENTER)


# ── Variant 5: DARK — glowing semi-transparent circles ───────────────

def _glow_overlap(slide, t, c, p):
    add_background(slide, t.primary)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)

    area_h = FOOTER_TOP - content_top - 200000
    cx = SLIDE_W // 2
    cy = content_top + area_h // 2
    radius = min(area_h // 2 - 100000, 1400000)

    centers = _circle_centers_3(cx, cy, radius) if n >= 3 else _circle_centers_2(cx, cy, radius)

    # Glow halos (larger, more transparent)
    for i in range(len(sets[:len(centers)])):
        x, y = centers[i]
        add_semi_transparent_circle(slide, x, y, radius + 100000,
                                    fill_color=colors[i], alpha=80000)

    # Main glowing circles
    for i, (label, desc) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        add_semi_transparent_circle(slide, x, y, radius,
                                    fill_color=colors[i], alpha=55000)

        # Label in neon accent
        if n >= 3:
            dx = x - cx
            dy = y - cy
            lx = x + int(dx * 0.8) - 700000
            ly = y + int(dy * 0.8) - 200000
        else:
            lx = x - 700000
            ly = y - radius - 450000
        add_textbox(slide, label, lx, ly, 1400000, 250000,
                    font_size=13, bold=True, color=colors[i],
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, lx - 100000, ly + 250000, 1600000, 380000,
                    font_size=9, color=tint(t.primary, 0.5),
                    alignment=PP_ALIGN.CENTER)

    # Center overlap with accent glow
    add_rect(slide, cx - 1000000, cy - 150000, 2000000, 300000,
             fill_color=shade(t.primary, 0.15), corner_radius=80000)
    add_textbox(slide, overlap, cx - 950000, cy - 130000, 1900000, 260000,
                font_size=11, bold=True, color=t.accent,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Corner brackets
    bsz = 180000
    bc = tint(t.accent, 0.3)
    add_line(slide, MARGIN, MARGIN, MARGIN + bsz, MARGIN, color=bc, width=0.5)
    add_line(slide, MARGIN, MARGIN, MARGIN, MARGIN + bsz, color=bc, width=0.5)
    add_line(slide, SLIDE_W - MARGIN - bsz, FOOTER_TOP - 200000,
             SLIDE_W - MARGIN, FOOTER_TOP - 200000, color=bc, width=0.5)
    add_line(slide, SLIDE_W - MARGIN, FOOTER_TOP - 200000 - bsz,
             SLIDE_W - MARGIN, FOOTER_TOP - 200000, color=bc, width=0.5)


# ── Variant 6: SPLIT — Venn diagram left, descriptions right ─────────

def _split_venn(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)
    s = t.ux_style

    # Left side: Venn diagram
    split_x = int(SLIDE_W * 0.52)
    area_h = FOOTER_TOP - content_top - 200000
    cx = split_x // 2 + MARGIN // 2
    cy = content_top + area_h // 2
    radius = min(area_h // 2 - 100000, 1100000)

    centers = _circle_centers_3(cx, cy, radius) if n >= 3 else _circle_centers_2(cx, cy, radius)

    for i, (label, _) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        add_semi_transparent_circle(slide, x, y, radius,
                                    fill_color=colors[i], alpha=50000)
        # Small label inside
        if n >= 3:
            dx = x - cx
            dy = y - cy
            lx = x + int(dx * 0.5) - 400000
            ly = y + int(dy * 0.5) - 100000
        else:
            lx = x - 400000
            ly = y - 100000
        tc = "#FFFFFF" if s.dark_mode else t.primary
        add_textbox(slide, label, lx, ly, 800000, 200000,
                    font_size=11, bold=True, color=tc,
                    alignment=PP_ALIGN.CENTER)

    # Right side: descriptions in cards
    right_x = split_x + 200000
    desc_w = SLIDE_W - right_x - MARGIN
    card_h = min(600000, (area_h - 200000 * n) // n)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if s.dark_mode else t.secondary

    for i, (label, desc) in enumerate(sets):
        card_y = content_top + 100000 + i * (card_h + 200000)
        # Color dot
        add_circle(slide, right_x + 80000, card_y + card_h // 2, 50000,
                   fill_color=colors[i])
        add_textbox(slide, label, right_x + 200000, card_y + 20000,
                    desc_w - 200000, 250000,
                    font_size=13, bold=True, color=text_color)
        add_textbox(slide, desc, right_x + 200000, card_y + 260000,
                    desc_w - 200000, card_h - 300000,
                    font_size=10, color=sub_color)

    # Overlap label
    add_gold_accent_line(slide, right_x, content_top + n * (card_h + 200000) + 200000,
                         desc_w, color=t.accent)
    add_textbox(slide, overlap, right_x,
                content_top + n * (card_h + 200000) + 350000,
                desc_w, 250000,
                font_size=12, bold=True, color=t.accent)


# ── Variant 7: GEO — overlapping hexagons instead of circles ─────────

def _hex_overlap(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)
    s = t.ux_style

    area_h = FOOTER_TOP - content_top - 200000
    cx = SLIDE_W // 2
    cy = content_top + area_h // 2
    hex_size = min(area_h // 2 - 100000, 1200000)

    # Hexagon positions (same formation as circles)
    if n >= 3:
        offset = int(hex_size * 0.55)
        positions = [
            (cx, cy - int(offset * 0.55)),
            (cx - int(offset * 0.85), cy + int(offset * 0.55)),
            (cx + int(offset * 0.85), cy + int(offset * 0.55)),
        ]
    else:
        offset = int(hex_size * 0.65)
        positions = [(cx - offset, cy), (cx + offset, cy)]

    # Draw hexagons with semi-transparent fill
    for i, (label, desc) in enumerate(sets[:len(positions)]):
        hx, hy = positions[i]
        add_hexagon(slide, hx, hy, hex_size, fill_color=tint(colors[i], 0.4))
        # Smaller inner hexagon
        add_hexagon(slide, hx, hy, int(hex_size * 0.7),
                    fill_color=tint(colors[i], 0.2))

        # Label
        if n >= 3:
            dx = hx - cx
            dy = hy - cy
            lx = hx + int(dx * 0.8) - 700000
            ly = hy + int(dy * 0.8) - 200000
        else:
            lx = hx - 700000
            ly = hy - hex_size - 400000
        text_color = "#FFFFFF" if s.dark_mode else t.primary
        add_textbox(slide, label, lx, ly, 1400000, 250000,
                    font_size=13, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, lx - 100000, ly + 250000, 1600000, 380000,
                    font_size=9,
                    color=tint(t.primary, 0.5) if s.dark_mode else t.secondary,
                    alignment=PP_ALIGN.CENTER)

    # Overlap: diamond in center
    from pptmaster.builder.helpers import add_diamond
    add_diamond(slide, cx, cy, 120000, fill_color=t.accent)
    add_textbox(slide, overlap, cx - 1000000, cy + 180000, 2000000, 300000,
                font_size=11, bold=True,
                color="#FFFFFF" if s.dark_mode else t.primary,
                alignment=PP_ALIGN.CENTER)


# ── Variant 8: EDITORIAL — elegant thin circles, editorial typography ─

def _editorial_venn(slide, t, c, p):
    add_dark_bg(slide, t)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)

    area_h = FOOTER_TOP - content_top - 400000
    cx = SLIDE_W // 2
    cy = content_top + area_h // 2 + 100000
    radius = min(area_h // 2 - 100000, 1250000)

    centers = _circle_centers_3(cx, cy, radius) if n >= 3 else _circle_centers_2(cx, cy, radius)

    for i, (label, desc) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        # Elegant thin circle — just outline
        add_circle(slide, x, y, radius,
                   fill_color="#FFFFFF" if not s.dark_mode else tint(t.primary, 0.05),
                   line_color=tint(colors[i], 0.3), line_width=0.75)
        # Italic label with rule
        if n >= 3:
            dx = x - cx
            dy = y - cy
            lx = x + int(dx * 0.55) - 600000
            ly = y + int(dy * 0.55) - 150000
        else:
            lx = x - 600000
            ly = y - 150000
        text_color = t.primary if not s.dark_mode else "#FFFFFF"
        add_textbox(slide, label, lx, ly, 1200000, 200000,
                    font_size=14, bold=False, italic=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)
        # Short rule under label
        add_line(slide, lx + 400000, ly + 210000, lx + 800000, ly + 210000,
                 color=colors[i], width=0.5)
        add_textbox(slide, desc, lx - 100000, ly + 240000, 1500000, 350000,
                    font_size=9, color=t.secondary if not s.dark_mode
                    else tint(t.primary, 0.55),
                    alignment=PP_ALIGN.CENTER)

    # Overlap: small italic text
    add_textbox(slide, overlap, cx - 900000, cy - 80000, 1800000, 240000,
                font_size=10, italic=True, color=t.accent,
                alignment=PP_ALIGN.CENTER)


# ── Variant 9: GRADIENT — gradient-tinted overlapping circles ────────

def _gradient_venn(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)
    s = t.ux_style

    area_h = FOOTER_TOP - content_top - 200000
    cx = SLIDE_W // 2
    cy = content_top + area_h // 2
    radius = min(area_h // 2 - 80000, 1400000)

    centers = _circle_centers_3(cx, cy, radius) if n >= 3 else _circle_centers_2(cx, cy, radius)

    # Simulated gradient: concentric rings per circle (outer light, inner deeper)
    for i, (label, desc) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        for ring in range(4, -1, -1):
            r = radius - ring * (radius // 8)
            tint_factor = 0.6 - ring * 0.1
            c_ring = tint(colors[i], max(tint_factor, 0.1))
            add_semi_transparent_circle(slide, x, y, r,
                                        fill_color=c_ring, alpha=50000 + ring * 5000)

        # Label outside
        if n >= 3:
            dx = x - cx
            dy = y - cy
            lx = x + int(dx * 0.9) - 700000
            ly = y + int(dy * 0.9) - 200000
        else:
            lx = x - 700000
            ly = y - radius - 400000
        text_color = t.primary if not s.dark_mode else "#FFFFFF"
        add_textbox(slide, label, lx, ly, 1400000, 250000,
                    font_size=13, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, lx - 100000, ly + 250000, 1600000, 380000,
                    font_size=9,
                    color=t.secondary if not s.dark_mode else tint(t.primary, 0.5),
                    alignment=PP_ALIGN.CENTER)

    # Overlap: rounded pill label
    pill_w = 2200000
    pill_h = 320000
    add_rect(slide, cx - pill_w // 2, cy - pill_h // 2,
             pill_w, pill_h, fill_color=t.accent, corner_radius=160000)
    add_textbox(slide, overlap, cx - pill_w // 2 + 40000, cy - pill_h // 2 + 20000,
                pill_w - 80000, pill_h - 40000,
                font_size=11, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant 10: RETRO — vintage styled with decorative borders ───────

def _retro_venn(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)

    area_h = FOOTER_TOP - content_top - 200000
    cx = SLIDE_W // 2
    cy = content_top + area_h // 2
    radius = min(area_h // 2 - 100000, 1300000)

    centers = _circle_centers_3(cx, cy, radius) if n >= 3 else _circle_centers_2(cx, cy, radius)

    # Decorative border frame
    frame_pad = 150000
    frame_left = MARGIN + frame_pad
    frame_top = content_top + frame_pad
    frame_w = SLIDE_W - 2 * MARGIN - 2 * frame_pad
    frame_h = FOOTER_TOP - content_top - 2 * frame_pad - 200000
    add_rect(slide, frame_left, frame_top, frame_w, frame_h,
             line_color=t.accent, line_width=2)
    add_rect(slide, frame_left + 40000, frame_top + 40000,
             frame_w - 80000, frame_h - 80000,
             line_color=t.accent, line_width=0.75)

    # Vintage circles: thick outlines, warm fill
    for i, (label, desc) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        add_circle(slide, x, y, radius,
                   fill_color=tint(colors[i], 0.7),
                   line_color=colors[i], line_width=2.5)
        # Inner decorative ring
        add_circle(slide, x, y, radius - 60000,
                   fill_color=tint(colors[i], 0.7),
                   line_color=colors[i], line_width=0.5)

        # Label with vintage style
        if n >= 3:
            dx = x - cx
            dy = y - cy
            lx = x + int(dx * 0.6) - 600000
            ly = y + int(dy * 0.6) - 160000
        else:
            lx = x - 600000
            ly = y - 160000
        add_textbox(slide, label.upper(), lx, ly, 1200000, 200000,
                    font_size=12, bold=True, color=t.primary,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, lx - 100000, ly + 220000, 1500000, 350000,
                    font_size=9, color=t.secondary,
                    alignment=PP_ALIGN.CENTER)

    # Overlap: decorative badge
    badge_r = 300000
    add_circle(slide, cx, cy, badge_r,
               fill_color=t.accent, line_color=shade(t.accent, 0.15), line_width=2)
    add_circle(slide, cx, cy, badge_r - 30000,
               fill_color=t.accent, line_color="#FFFFFF", line_width=0.5)
    add_textbox(slide, overlap, cx - 800000, cy + badge_r + 80000, 1600000, 300000,
                font_size=10, bold=True, color=t.accent,
                alignment=PP_ALIGN.CENTER)

    # Small dot decorations in corners of frame
    dot_r = 25000
    for dx, dy in [(frame_left + 80000, frame_top + 80000),
                   (frame_left + frame_w - 80000, frame_top + 80000),
                   (frame_left + 80000, frame_top + frame_h - 80000),
                   (frame_left + frame_w - 80000, frame_top + frame_h - 80000)]:
        add_circle(slide, dx, dy, dot_r, fill_color=t.accent)


# ── Variant 11: MAGAZINE — creative asymmetric circles, oversized ─────

def _creative_venn(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)
    s = t.ux_style

    area_h = FOOTER_TOP - content_top - 200000
    cx = SLIDE_W // 2
    cy = content_top + area_h // 2

    # Asymmetric: vary radius per circle for creative feel
    base_r = min(area_h // 2 - 80000, 1400000)
    radii = [base_r, int(base_r * 0.9), int(base_r * 1.05)][:n]

    if n >= 3:
        offset = int(base_r * 0.55)
        centers = [
            (cx - int(offset * 0.3), cy - int(offset * 0.6)),
            (cx - int(offset * 0.9), cy + int(offset * 0.5)),
            (cx + int(offset * 0.9), cy + int(offset * 0.2)),
        ]
    else:
        offset = int(base_r * 0.7)
        centers = [(cx - offset, cy), (cx + offset, cy - 200000)]

    for i, (label, desc) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        r = radii[i]
        add_semi_transparent_circle(slide, x, y, r,
                                    fill_color=colors[i], alpha=50000)

        # Oversized label
        if n >= 3:
            dx = x - cx
            dy = y - cy
            lx = x + int(dx * 0.7) - 800000
            ly = y + int(dy * 0.7) - 200000
        else:
            lx = x - 800000
            ly = y - r - 450000
        text_color = "#FFFFFF" if s.dark_mode else t.primary
        add_textbox(slide, label, lx, ly, 1600000, 350000,
                    font_size=18, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

    # Overlap: oversized accent text
    add_textbox(slide, overlap, cx - 1200000, cy - 180000, 2400000, 360000,
                font_size=14, bold=True, color=t.accent,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant 12: SCHOLARLY — "Figure N:" caption, thin outlines ───────

def _scholarly_venn(slide, t, c, p):
    add_dark_bg(slide, t)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)

    # Figure caption
    text_color = t.primary if not s.dark_mode else "#FFFFFF"
    sub_color = t.secondary if not s.dark_mode else tint(t.primary, 0.55)
    add_textbox(slide, "Figure 1: Conceptual Overlap Model",
                m, content_top, SLIDE_W - 2 * m, 280000,
                font_size=10, italic=True, color=sub_color)

    fig_top = content_top + 350000
    area_h = FOOTER_TOP - fig_top - 500000
    cx = SLIDE_W // 2
    cy = fig_top + area_h // 2
    radius = min(area_h // 2 - 80000, 1250000)

    centers = _circle_centers_3(cx, cy, radius) if n >= 3 else _circle_centers_2(cx, cy, radius)

    # Thin scholarly outlines (single weight)
    for i, (label, desc) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        bg = "#FFFFFF" if not s.dark_mode else tint(t.primary, 0.06)
        add_circle(slide, x, y, radius,
                   fill_color=bg, line_color=t.secondary, line_width=0.75)

        # Numbered label: (a), (b), (c) ...
        letter = chr(ord('a') + i)
        if n >= 3:
            dx = x - cx
            dy = y - cy
            lx = x + int(dx * 0.55) - 550000
            ly = y + int(dy * 0.55) - 180000
        else:
            lx = x - 550000
            ly = y - 180000
        add_textbox(slide, f"({letter}) {label}", lx, ly, 1400000, 280000,
                    font_size=11, bold=False, color=text_color,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, lx - 200000, ly + 290000, 1800000, 300000,
                    font_size=9, italic=True, color=sub_color,
                    alignment=PP_ALIGN.CENTER)

    # Overlap: thin dashed box label
    add_line(slide, cx - 900000, cy - 140000, cx + 900000, cy - 140000,
             color=t.secondary, width=0.5, dash="dash")
    add_textbox(slide, overlap, cx - 900000, cy - 120000, 1800000, 280000,
                font_size=10, color=text_color, alignment=PP_ALIGN.CENTER)
    add_line(slide, cx - 900000, cy + 160000, cx + 900000, cy + 160000,
             color=t.secondary, width=0.5, dash="dash")

    # Source note at bottom
    add_textbox(slide, "Note: Overlap region represents combined strategic value.",
                m, FOOTER_TOP - 350000, SLIDE_W - 2 * m, 200000,
                font_size=8, italic=True, color=sub_color)


# ── Variant 13: LABORATORY — dark bg, data-coded circles, grid ───────

def _laboratory_venn(slide, t, c, p):
    add_background(slide, t.primary)

    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)

    # Subtle grid overlay
    grid_spacing = 500000
    grid_color = tint(t.primary, 0.08)
    for x in range(0, SLIDE_W + 1, grid_spacing):
        add_line(slide, x, content_top, x, FOOTER_TOP - 200000,
                 color=grid_color, width=0.25)
    for y in range(content_top, FOOTER_TOP - 200000 + 1, grid_spacing):
        add_line(slide, MARGIN, y, SLIDE_W - MARGIN, y,
                 color=grid_color, width=0.25)

    area_h = FOOTER_TOP - content_top - 300000
    cx = SLIDE_W // 2
    cy = content_top + area_h // 2
    radius = min(area_h // 2 - 80000, 1350000)

    centers = _circle_centers_3(cx, cy, radius) if n >= 3 else _circle_centers_2(cx, cy, radius)

    # Data-coded circles with dashed outlines
    for i, (label, desc) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        add_semi_transparent_circle(slide, x, y, radius,
                                    fill_color=colors[i], alpha=65000)
        # Dashed outline ring
        add_circle(slide, x, y, radius,
                   fill_color=tint(t.primary, 0.02),
                   line_color=colors[i], line_width=1)

        # Data-style label: monospace feel, left-aligned
        if n >= 3:
            dx = x - cx
            dy = y - cy
            lx = x + int(dx * 0.8) - 700000
            ly = y + int(dy * 0.8) - 220000
        else:
            lx = x - 700000
            ly = y - radius - 450000

        # Color-coded index tag
        tag_w = 60000
        add_rect(slide, lx + 350000, ly - 10000, tag_w, 380000,
                 fill_color=colors[i])
        add_textbox(slide, f"SET-{chr(ord('A') + i)}: {label}",
                    lx + 350000 + tag_w + 40000, ly, 1800000, 380000,
                    font_size=9, bold=True, color=colors[i])
        add_textbox(slide, desc, lx + 350000 + tag_w + 40000, ly + 390000,
                    1800000, 380000,
                    font_size=8, color=tint(t.primary, 0.45))

    # Overlap: evidence-style box
    box_w = 3200000
    box_h = 500000
    add_rect(slide, cx - box_w // 2, cy - box_h // 2,
             box_w, box_h,
             fill_color=shade(t.primary, 0.15),
             line_color=t.accent, line_width=0.75)
    add_textbox(slide, f"OVERLAP: {overlap}",
                cx - box_w // 2 + 40000, cy - box_h // 2 + 20000,
                box_w - 80000, box_h - 40000,
                font_size=10, bold=True, color=t.accent,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Coordinate markers
    add_textbox(slide, f"n={n} sets", SLIDE_W - MARGIN - 800000,
                FOOTER_TOP - 350000, 800000, 200000,
                font_size=8, color=tint(t.primary, 0.35),
                alignment=PP_ALIGN.RIGHT)


# ── Variant 14: DASHBOARD — compact venn with metric annotations ─────

def _dashboard_venn(slide, t, c, p):
    add_dark_bg(slide, t)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)
    content_top = add_slide_title(
        slide, c.get("venn_title", "Synergy Analysis"), theme=t,
    )
    sets = _get_sets(c)
    overlap = _get_overlap(c)
    n = len(sets)
    colors = _venn_colors(t, p, n)
    text_color = t.primary if not s.dark_mode else "#FFFFFF"
    sub_color = t.secondary if not s.dark_mode else tint(t.primary, 0.5)

    # Dashboard header bar
    hdr_h = 60000
    add_rect(slide, m, content_top, SLIDE_W - 2 * m, hdr_h,
             fill_color=t.accent)
    add_textbox(slide, "OVERLAP ANALYSIS",
                m + 80000, content_top + 5000,
                2000000, hdr_h - 10000,
                font_size=8, bold=True, color="#FFFFFF")

    # Compact Venn in left panel
    panel_top = content_top + hdr_h + 60000
    panel_w = int((SLIDE_W - 2 * m) * 0.6)
    area_h = FOOTER_TOP - panel_top - 300000
    vcx = m + panel_w // 2
    vcy = panel_top + area_h // 2
    radius = min(area_h // 2 - 60000, 1100000)

    # Panel background
    add_rect(slide, m, panel_top, panel_w, area_h,
             fill_color=tint(t.primary, 0.05) if s.dark_mode
             else tint(t.light_bg, 0.3),
             line_color=tint(t.secondary, 0.7) if not s.dark_mode
             else tint(t.primary, 0.2),
             line_width=0.5)

    centers = _circle_centers_3(vcx, vcy, radius) if n >= 3 else _circle_centers_2(vcx, vcy, radius)

    for i, (label, _) in enumerate(sets[:len(centers)]):
        x, y = centers[i]
        add_semi_transparent_circle(slide, x, y, radius,
                                    fill_color=colors[i], alpha=55000)
        # Small label inside
        if n >= 3:
            dx = x - vcx
            dy = y - vcy
            lx = x + int(dx * 0.45) - 450000
            ly = y + int(dy * 0.45) - 100000
        else:
            lx = x - 450000
            ly = y - 100000
        add_textbox(slide, label, lx, ly, 900000, 200000,
                    font_size=10, bold=True,
                    color="#FFFFFF" if s.dark_mode else t.primary,
                    alignment=PP_ALIGN.CENTER)

    # Overlap label in center of venn
    add_textbox(slide, overlap, vcx - 800000, vcy - 120000, 1600000, 240000,
                font_size=9, bold=True, color=t.accent,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Right panel: metric tiles
    right_x = m + panel_w + 100000
    tile_w = SLIDE_W - m - right_x
    tile_h = (area_h - 100000 * (n - 1)) // n

    for i, (label, desc) in enumerate(sets):
        tile_y = panel_top + i * (tile_h + 100000)
        # Tile bg
        tile_bg = tint(t.primary, 0.08) if s.dark_mode else "#FFFFFF"
        add_rect(slide, right_x, tile_y, tile_w, tile_h,
                 fill_color=tile_bg,
                 line_color=tint(t.secondary, 0.7) if not s.dark_mode
                 else tint(t.primary, 0.2),
                 line_width=0.5)
        # Color indicator strip
        add_rect(slide, right_x, tile_y, 50000, tile_h,
                 fill_color=colors[i])
        # Metric-style content
        add_textbox(slide, label, right_x + 120000, tile_y + 30000,
                    tile_w - 180000, 200000,
                    font_size=11, bold=True, color=text_color)
        add_textbox(slide, desc, right_x + 120000, tile_y + 230000,
                    tile_w - 180000, tile_h - 280000,
                    font_size=9, color=sub_color)
