"""PPT Master MCP Server — 10 tools for AI agents to build presentations.

Start with:
    python -m pptmaster.mcp_server
or via Agent settings.json:
    {
      "mcpServers": {
        "pptmaster": {
          "command": "/path/to/.venv/bin/python",
          "args": ["-m", "pptmaster.mcp_server"],
          "cwd": "/Volumes/Seagate Expansion Drive/Seagate AI/PPT Master"
        }
      }
    }
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

# Ensure src/ is on the path when invoked directly
_ROOT = Path(__file__).resolve().parent.parent.parent
_SRC = _ROOT / "src"
if str(_SRC) not in sys.path:
    sys.path.insert(0, str(_SRC))

from mcp.server.fastmcp import FastMCP

mcp = FastMCP(
    "PPT Master",
    instructions=(
        "AI-powered executive presentation generator. "
        "14 visual themes × 40 slide types + 144 icons across 14 categories.\n\n"
        "Typical workflow:\n"
        "  1. list_themes() — choose a visual theme\n"
        "  2. list_slide_types() — explore available slide layouts\n"
        "  3. build_presentation() — AI generates + builds full PPTX in one call\n"
        "     OR generate_content() → edit JSON → build_from_content()\n"
        "  4. Optionally: build_single_slide() to prototype one slide\n"
        "     list_icons() / get_icon_path() to look up icon assets\n\n"
        "All output_path args accept absolute or relative paths. "
        "Relative paths are resolved from the project root."
    ),
)

_PPT_ROOT = Path(__file__).resolve().parent.parent.parent


def _resolve(path: str) -> Path:
    p = Path(path)
    if not p.is_absolute():
        p = _PPT_ROOT / p
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


# ── Tool 1: list_themes ──────────────────────────────────────────────────────

@mcp.tool()
def list_themes() -> str:
    """List all 14 available visual themes with industry, style, and colors.

    Returns a formatted table. Use the 'key' value in other tools.
    Themes: corporate, healthcare, technology, finance, education,
    sustainability, luxury, startup, government, realestate, creative,
    academic, research, report.
    """
    from pptmaster.builder.themes import THEME_MAP, get_theme

    rows = []
    for key in sorted(THEME_MAP.keys()):
        t = get_theme(key)
        rows.append({
            "key": key,
            "industry": t.industry,
            "ux_style": t.ux_style.name,
            "dark_mode": t.ux_style.dark_mode,
            "primary": t.primary,
            "accent": t.accent,
            "font": t.font,
        })

    lines = [f"{'KEY':<16} {'INDUSTRY':<22} {'STYLE':<12} {'DARK':<6} {'PRIMARY':<10} {'ACCENT':<10} FONT"]
    lines.append("-" * 90)
    for r in rows:
        lines.append(
            f"{r['key']:<16} {r['industry']:<22} {r['ux_style']:<12} "
            f"{'yes' if r['dark_mode'] else 'no':<6} {r['primary']:<10} {r['accent']:<10} {r['font']}"
        )
    lines.append(f"\n{len(rows)} themes available.")
    return "\n".join(lines)


# ── Tool 2: list_slide_types ─────────────────────────────────────────────────

@mcp.tool()
def list_slide_types() -> str:
    """List all 32 selectable slide types with descriptions and required content keys.

    The 'type' field is used in build_single_slide() and in content JSON passed to
    build_from_content(). Always-included slides (cover, toc, section dividers,
    thank_you) are not listed here — they are added automatically.
    """
    from pptmaster.content.builder_content_gen import SLIDE_CATALOG

    lines = [f"{'TYPE':<24} DESCRIPTION"]
    lines.append("-" * 80)
    for slide_type, meta in SLIDE_CATALOG.items():
        desc = meta.get("description", "")
        lines.append(f"{slide_type:<24} {desc}")

    lines.append(f"\n{len(SLIDE_CATALOG)} slide types. Always-included: cover, toc, section_divider×7, thank_you.")
    lines.append("Use list_slide_types(verbose=True)... or call this tool and request content_keys for a type.")
    return "\n".join(lines)


# ── Tool 3: get_slide_type_details ───────────────────────────────────────────

@mcp.tool()
def get_slide_type_details(slide_type: str) -> str:
    """Get full details for a slide type: description, content keys, and defaults.

    Args:
        slide_type: One of the 32 slide type keys from list_slide_types().

    Returns:
        JSON with description, required content_keys, and example values.
    """
    from pptmaster.content.builder_content_gen import SLIDE_CATALOG

    if slide_type not in SLIDE_CATALOG:
        available = ", ".join(sorted(SLIDE_CATALOG.keys()))
        return f"Unknown slide type: '{slide_type}'. Available: {available}"

    meta = SLIDE_CATALOG[slide_type]
    result = {
        "type": slide_type,
        "description": meta.get("description", ""),
        "content_keys": meta.get("content_keys", []),
        "note": (
            "All content keys are optional — defaults are used for missing keys. "
            "Pass content as JSON to build_single_slide() or include in build_from_content()."
        ),
    }
    return json.dumps(result, indent=2)


# ── Tool 4: get_theme_details ────────────────────────────────────────────────

@mcp.tool()
def get_theme_details(theme_key: str) -> str:
    """Get full details for a specific theme: palette, UX style, fonts, and sample content.

    Args:
        theme_key: One of the 14 theme keys from list_themes().

    Returns:
        JSON with color palette, UX style variants, font, and content preview.
    """
    from pptmaster.builder.themes import get_theme, THEME_MAP

    if theme_key not in THEME_MAP:
        available = ", ".join(sorted(THEME_MAP.keys()))
        return f"Unknown theme: '{theme_key}'. Available: {available}"

    t = get_theme(theme_key)
    s = t.ux_style

    result = {
        "key": theme_key,
        "industry": t.industry,
        "company_name": t.company_name,
        "tagline": t.tagline,
        "font": t.font,
        "colors": {
            "primary": t.primary,
            "secondary": t.secondary,
            "accent": t.accent,
            "light_bg": t.light_bg,
            "palette": list(t.palette),
        },
        "ux_style": {
            "name": s.name,
            "dark_mode": s.dark_mode,
            "margin_factor": s.margin_factor,
            "cover": s.cover,
            "kpi": s.kpi,
            "chart": s.chart,
            "venn": s.venn,
            "hub_spoke": s.hub_spoke,
            "funnel": s.funnel,
            "pyramid": s.pyramid,
            "milestone": s.milestone,
            "kanban": s.kanban,
            "matrix": s.matrix,
            "gauge": s.gauge,
            "risk": s.risk,
        },
        "sample_content_keys": list(t.content.keys())[:20],
    }
    return json.dumps(result, indent=2)


# ── Tool 5: build_template ───────────────────────────────────────────────────

@mcp.tool()
def build_template(
    theme_key: str = "corporate",
    output_path: str = "output.pptx",
    company_name: str = "Acme Corp",
) -> str:
    """Build a full 40-slide themed template PPTX. No LLM call needed.

    Produces all 40 slides with default content for the chosen theme.
    Useful as a starting point, reference deck, or blank slate to fill in.

    Args:
        theme_key: One of the 14 theme keys (e.g. "corporate", "healthcare").
        output_path: Where to save the .pptx file.
        company_name: Company name used in cover and branding.

    Returns:
        Absolute path to the saved PPTX file.
    """
    from pptmaster.builder.template_builder import build_template as _build
    from pptmaster.builder.themes import get_theme, THEME_MAP

    if theme_key not in THEME_MAP:
        available = ", ".join(sorted(THEME_MAP.keys()))
        return f"Unknown theme: '{theme_key}'. Available: {available}"

    theme = get_theme(theme_key)
    out = _resolve(output_path)
    result = _build(out, company_name, theme=theme)
    return (
        f"Template built successfully.\n"
        f"Path: {result}\n"
        f"Theme: {theme_key} ({theme.industry}) | Slides: 40 | Company: {company_name}"
    )


# ── Tool 6: build_presentation ───────────────────────────────────────────────

@mcp.tool()
def build_presentation(
    topic: str,
    company_name: str = "Acme Corp",
    theme_key: str = "corporate",
    output_path: str = "output.pptx",
    audience: str = "",
    industry: str = "",
    additional_context: str = "",
) -> str:
    """AI-generate a presentation: LLM picks slides + writes content, then builds PPTX.

    The AI selects 8-30 of the 32 slide types based on the topic, generates all
    content, and assembles a professionally styled PPTX.

    Args:
        topic: Main subject of the presentation (e.g. "Q3 Financial Results").
        company_name: Company name for cover and branding.
        theme_key: Visual theme (use list_themes() to see options).
        output_path: Where to save the .pptx file.
        audience: Target audience (e.g. "Board of Directors", "Sales Team").
        industry: Industry context — auto-inferred from theme if empty.
        additional_context: Extra instructions for the AI (e.g. "focus on cost reduction").

    Returns:
        Path to the PPTX file and a summary of what was generated.
    """
    from pptmaster.builder.ai_builder import ai_build_presentation
    from pptmaster.builder.themes import THEME_MAP

    if theme_key not in THEME_MAP:
        available = ", ".join(sorted(THEME_MAP.keys()))
        return f"Unknown theme: '{theme_key}'. Available: {available}"

    out = _resolve(output_path)
    result = ai_build_presentation(
        topic=topic,
        company_name=company_name,
        industry=industry,
        audience=audience,
        theme_key=theme_key,
        output_path=str(out),
        additional_context=additional_context,
    )
    return (
        f"Presentation generated successfully.\n"
        f"Path: {result}\n"
        f"Topic: {topic} | Theme: {theme_key} | Company: {company_name}"
    )


# ── Tool 7: generate_content ─────────────────────────────────────────────────

@mcp.tool()
def generate_content(
    topic: str,
    company_name: str = "Acme Corp",
    industry: str = "",
    audience: str = "",
    additional_context: str = "",
) -> str:
    """Generate slide content as JSON without building a PPTX.

    Returns a JSON object with:
      - selected_slides: list of slide type keys the AI chose
      - sections: section divider titles
      - content: dict of content for each slide (titles, bullets, metrics, etc.)

    Edit the JSON, then pass it to build_from_content() to produce the PPTX.

    Args:
        topic: Main subject of the presentation.
        company_name: Company name for branding.
        industry: Industry context (helps AI tailor terminology).
        audience: Target audience.
        additional_context: Extra instructions for the AI.

    Returns:
        JSON string with selected_slides, sections, and content.
    """
    from pptmaster.content.builder_content_gen import generate_builder_content

    result = generate_builder_content(
        topic=topic,
        company_name=company_name,
        industry=industry,
        audience=audience,
        additional_context=additional_context,
    )
    return json.dumps(result, indent=2, ensure_ascii=False)


# ── Tool 8: build_from_content ───────────────────────────────────────────────

@mcp.tool()
def build_from_content(
    content_json: str,
    theme_key: str = "corporate",
    company_name: str = "Acme Corp",
    output_path: str = "output.pptx",
) -> str:
    """Build a PPTX from a pre-generated (or hand-crafted) content JSON dict.

    Use this after generate_content() if you want to review/edit the content
    before building. All content keys are optional — defaults fill in gaps.

    The content_json must be either:
      (a) Full spec with selected_slides + sections + content (from generate_content())
      (b) Flat content dict with any subset of content keys

    Args:
        content_json: JSON string with slide content.
        theme_key: Visual theme to apply.
        company_name: Company name for branding.
        output_path: Where to save the .pptx file.

    Returns:
        Path to the generated PPTX file.
    """
    from pptmaster.builder.themes import get_theme, THEME_MAP
    from pptmaster.builder.ai_builder import build_from_content as _build

    if theme_key not in THEME_MAP:
        available = ", ".join(sorted(THEME_MAP.keys()))
        return f"Unknown theme: '{theme_key}'. Available: {available}"

    try:
        data = json.loads(content_json)
    except json.JSONDecodeError as e:
        return f"Invalid JSON: {e}"

    out = _resolve(output_path)
    result = _build(
        content_dict=data,
        company_name=company_name,
        theme_key=theme_key,
        output_path=str(out),
    )
    return (
        f"Presentation built from content.\n"
        f"Path: {result}\n"
        f"Theme: {theme_key} | Company: {company_name}"
    )


# ── Tool 9: build_single_slide ───────────────────────────────────────────────

@mcp.tool()
def build_single_slide(
    slide_type: str,
    theme_key: str = "corporate",
    content_json: str = "{}",
    output_path: str = "slide_preview.pptx",
) -> str:
    """Build a single slide as a standalone PPTX for rapid prototyping.

    Great for testing slide designs or generating a specific slide type
    without building a full presentation.

    Args:
        slide_type: One of the 32 slide types from list_slide_types().
        theme_key: Visual theme to apply.
        content_json: Optional JSON dict with content keys for this slide.
            Use get_slide_type_details() to see what keys are available.
        output_path: Where to save the preview PPTX.

    Returns:
        Path to the single-slide PPTX.

    Example:
        build_single_slide(
            slide_type="hub_spoke",
            theme_key="technology",
            content_json='{"hub_title": "Our Platform", "hub_center": "AI Core",
                           "hub_spokes": [["Analytics", "Real-time insights"],
                                          ["Security", "Zero-trust model"]]}'
        )
    """
    from pptmaster.content.builder_content_gen import SLIDE_CATALOG
    from pptmaster.builder.themes import get_theme, THEME_MAP
    from pptmaster.builder.ai_builder import _BUILDER_MAP
    from pptx import Presentation
    from pptx.util import Inches

    if slide_type not in SLIDE_CATALOG:
        available = ", ".join(sorted(SLIDE_CATALOG.keys()))
        return f"Unknown slide type: '{slide_type}'. Available:\n{available}"

    if theme_key not in THEME_MAP:
        available = ", ".join(sorted(THEME_MAP.keys()))
        return f"Unknown theme: '{theme_key}'. Available: {available}"

    try:
        content_data = json.loads(content_json)
    except json.JSONDecodeError as e:
        return f"Invalid content_json: {e}"

    theme = get_theme(theme_key)
    # Merge provided content into theme content
    theme.content.update(content_data)

    # Build single-slide PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)

    builder_fn = _BUILDER_MAP.get(slide_type)
    if builder_fn is None:
        return f"No builder registered for slide type: '{slide_type}'"

    builder_fn(slide, theme=theme)

    out = _resolve(output_path)
    prs.save(str(out))
    return (
        f"Single slide built successfully.\n"
        f"Path: {out}\n"
        f"Slide type: {slide_type} | Theme: {theme_key}"
    )


# ── Tool 10: list_icons ──────────────────────────────────────────────────────

@mcp.tool()
def list_icons(category: str = "") -> str:
    """List available icons, optionally filtered by category.

    144 flat PNG icons across 14 categories. Use get_icon_path() to get the
    file path for embedding in a presentation.

    Args:
        category: Optional category filter. Leave empty to see all categories.
            Categories: analytics, business, communication, creative, education,
            finance, healthcare, logistics, people, real-estate, security,
            sustainability, technology.

    Returns:
        List of icon names, organized by category.
    """
    from pptmaster.assets.raster_icon_manager import RasterIconManager

    mgr = RasterIconManager()

    if category:
        # Filter to category
        icons = [i for i in mgr.icons.values() if i.category.lower() == category.lower()]
        if not icons:
            cats = sorted(set(i.category for i in mgr.icons.values()))
            return f"No icons found for category '{category}'. Available: {', '.join(cats)}"
        lines = [f"Icons in '{category}' ({len(icons)} total):"]
        for icon in sorted(icons, key=lambda i: i.name):
            lines.append(f"  {icon.name}")
        return "\n".join(lines)

    # All categories
    cats: dict[str, list[str]] = {}
    for icon in mgr.icons.values():
        cats.setdefault(icon.category, []).append(icon.name)

    lines = [f"Available icons ({len(mgr.icons)} total across {len(cats)} categories):\n"]
    for cat in sorted(cats.keys()):
        names = sorted(cats[cat])
        lines.append(f"  {cat} ({len(names)}): {', '.join(names)}")
    return "\n".join(lines)


# ── Tool 11: get_icon_path ───────────────────────────────────────────────────

@mcp.tool()
def get_icon_path(query: str) -> str:
    """Find an icon by name or keyword and return its absolute file path.

    Searches icon names for the query string (case-insensitive, partial match).

    Args:
        query: Icon name or keyword (e.g. "briefcase", "chart", "heart").

    Returns:
        Absolute path to the PNG file, or a list of matches if multiple found.
    """
    from pptmaster.assets.raster_icon_manager import RasterIconManager

    mgr = RasterIconManager()
    matches = mgr.search(query)

    if not matches:
        return f"No icons found matching '{query}'. Use list_icons() to browse all icons."

    if len(matches) == 1:
        icon = matches[0]
        return (
            f"Icon found: {icon.name}\n"
            f"Category: {icon.category}\n"
            f"Path: {icon.path}"
        )

    lines = [f"Multiple icons match '{query}' ({len(matches)} results):"]
    for icon in matches[:10]:
        lines.append(f"  [{icon.category}] {icon.name} — {icon.path}")
    if len(matches) > 10:
        lines.append(f"  ... and {len(matches) - 10} more. Refine your query.")
    return "\n".join(lines)


if __name__ == "__main__":
    mcp.run()
