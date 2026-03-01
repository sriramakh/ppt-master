"""Slide 40 — Risk Matrix (Heatmap): 14 unique visual variants dispatched by UX style."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_circle, add_line, add_octagon,
    add_color_cell, add_gold_accent_line, add_slide_title,
    add_dark_bg, add_background, add_styled_card,
)
from pptmaster.assets.color_utils import tint, shade


# ── Risk severity color mapping ──────────────────────────────────────
_RISK_COLORS = {
    "low": "#10B981",       # green
    "medium": "#F59E0B",    # amber
    "high": "#EF4444",      # red
    "critical": "#991B1B",  # dark red
}

_SEVERITY_ORDER = ["low", "medium", "high", "critical"]

# 3x3 grid cell colors — rows = impact (low→high top→bottom),
# cols = likelihood (low→high left→right).
# Cell[row][col] gives composite risk color.
_GRID_3x3 = [
    # low likelihood   med likelihood   high likelihood
    ["#10B981",        "#F59E0B",        "#EF4444"],       # high impact
    ["#A7F3D0",        "#FCD34D",        "#F59E0B"],       # medium impact
    ["#D1FAE5",        "#A7F3D0",        "#FCD34D"],       # low impact
]

# Map severity string → (row, col) position on the 3x3 grid
_SEVERITY_GRID_POS = {
    "low":      (2, 0),   # low impact, low likelihood
    "medium":   (1, 1),   # medium impact, medium likelihood
    "high":     (0, 2),   # high impact, high likelihood
    "critical": (0, 2),   # top-right danger zone (shares with high, offset)
}


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    p = t.palette
    s = t.ux_style

    dispatch = {
        "heatmap":            _heatmap,
        "dot-matrix":         _dot_matrix,
        "bold-heatmap":       _bold_heatmap,
        "floating-heatmap":   _floating_heatmap,
        "dark-heatmap":       _dark_heatmap,
        "split-risk":         _split_risk,
        "angular-heatmap":    _angular_heatmap,
        "editorial-heatmap":  _editorial_heatmap,
        "gradient-heatmap":   _gradient_heatmap,
        "retro-heatmap":      _retro_heatmap,
        "creative-heatmap":   _creative_heatmap,
        "scholarly-heatmap":  _scholarly_heatmap,
        "laboratory-heatmap": _laboratory_heatmap,
        "dashboard-heatmap":  _dashboard_heatmap,
    }
    builder = dispatch.get(s.risk, _heatmap)
    builder(slide, t, c)


# ── Helpers ──────────────────────────────────────────────────────────

def _get_items(c):
    """Extract risk items from content dict with sensible defaults."""
    return c.get("risk_items", [
        ("Data Breach", "critical", "Unauthorized access to sensitive data"),
        ("Supply Chain", "high", "Key vendor disruption risk"),
        ("Compliance", "medium", "Regulatory non-compliance"),
        ("Talent", "medium", "Key personnel retention"),
        ("Market Shift", "high", "Competitive landscape change"),
        ("Technology", "low", "Legacy system failure"),
    ])


def _risk_color(severity: str) -> str:
    """Return color for a severity level."""
    return _RISK_COLORS.get(severity, _RISK_COLORS["medium"])


def _severity_label(severity: str) -> str:
    """Return a readable label for severity."""
    return severity.upper()


def _assign_grid_positions(items):
    """Assign each risk item to a 3x3 grid position, spreading duplicates.

    Returns list of (name, severity, description, row, col).
    """
    # Track occupied cells to offset duplicates
    used = {}
    result = []
    for name, severity, desc in items:
        base_row, base_col = _SEVERITY_GRID_POS.get(severity, (1, 1))
        key = (base_row, base_col)
        count = used.get(key, 0)
        # Nudge duplicates to adjacent cells
        row, col = base_row, base_col
        if count == 1:
            col = min(base_col + 1, 2)
            if (row, col) == (base_row, base_col):
                row = min(base_row + 1, 2)
        elif count == 2:
            row = min(base_row + 1, 2)
        elif count >= 3:
            row = max(base_row - 1, 0)
            col = max(base_col - 1, 0)
        used[key] = count + 1
        result.append((name, severity, desc, row, col))
    return result


def _draw_axis_labels(slide, grid_left, grid_top, grid_w, grid_h,
                      x_label, y_label, text_color, accent_color):
    """Draw X and Y axis labels around the grid."""
    # X-axis label (bottom center)
    add_textbox(slide, x_label,
                grid_left, grid_top + grid_h + 80000,
                grid_w, 300000,
                font_size=12, bold=True, color=text_color,
                alignment=PP_ALIGN.CENTER)
    # X-axis arrow
    add_line(slide, grid_left + 200000, grid_top + grid_h + 60000,
             grid_left + grid_w - 200000, grid_top + grid_h + 60000,
             color=accent_color, width=1.0)

    # Y-axis label (left, rotated via vertical text)
    add_textbox(slide, y_label,
                grid_left - 500000, grid_top + grid_h // 2 - 150000,
                450000, 300000,
                font_size=12, bold=True, color=text_color,
                alignment=PP_ALIGN.CENTER)
    # Y-axis arrow
    add_line(slide, grid_left - 60000, grid_top + grid_h - 200000,
             grid_left - 60000, grid_top + 200000,
             color=accent_color, width=1.0)


def _draw_grid_labels(slide, grid_left, grid_top, cell_w, cell_h,
                      text_color):
    """Draw 'Low / Med / High' labels on axes of a 3x3 grid."""
    col_labels = ["Low", "Med", "High"]
    row_labels = ["High", "Med", "Low"]
    for i, lbl in enumerate(col_labels):
        add_textbox(slide, lbl,
                    grid_left + i * cell_w, grid_top - 250000,
                    cell_w, 220000,
                    font_size=9, color=text_color, alignment=PP_ALIGN.CENTER)
    for i, lbl in enumerate(row_labels):
        add_textbox(slide, lbl,
                    grid_left - 450000, grid_top + i * cell_h,
                    400000, cell_h,
                    font_size=9, color=text_color, alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)


# ═════════════════════════════════════════════════════════════════════
# Variant 1: CLASSIC — 3x3 color-coded grid with risk item markers
# ═════════════════════════════════════════════════════════════════════

def _heatmap(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("risk_title", "Risk Assessment Matrix"), theme=t)

    s = t.ux_style
    items = _get_items(c)
    x_label = c.get("risk_x_label", "Likelihood")
    y_label = c.get("risk_y_label", "Impact")
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    # Grid dimensions
    m = int(MARGIN * s.margin_factor)
    grid_left = m + 550000
    grid_top = content_y + 250000
    grid_w = 6000000
    grid_h = 4200000
    cell_w = grid_w // 3
    cell_h = grid_h // 3

    # Draw 3x3 colored grid
    for row in range(3):
        for col in range(3):
            cx = grid_left + col * cell_w
            cy = grid_top + row * cell_h
            color = _GRID_3x3[row][col]
            cell_fill = tint(color, 0.15) if s.dark_mode else tint(color, 0.5)
            add_color_cell(slide, cx, cy, cell_w, cell_h,
                           fill_color=cell_fill,
                           line_color=tint(text_color, 0.8), line_width=0.5)

    # Place risk items
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w + cell_w // 2
        cy = grid_top + row * cell_h + cell_h // 2
        marker_color = _risk_color(severity)
        add_circle(slide, cx, cy, 180000, fill_color=marker_color)
        add_textbox(slide, name,
                    cx - cell_w // 2 + 30000, cy + 200000,
                    cell_w - 60000, 250000,
                    font_size=8, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

    # Axis labels
    _draw_axis_labels(slide, grid_left, grid_top, grid_w, grid_h,
                      x_label, y_label, text_color, t.accent)
    _draw_grid_labels(slide, grid_left, grid_top, cell_w, cell_h, sub_color)

    # Legend on the right
    legend_left = grid_left + grid_w + 400000
    legend_top = grid_top + 200000
    add_textbox(slide, "Severity", legend_left, legend_top - 300000,
                2500000, 280000,
                font_size=12, bold=True, color=text_color)
    for i, sev in enumerate(_SEVERITY_ORDER):
        y = legend_top + i * 380000
        add_rect(slide, legend_left, y, 200000, 200000,
                 fill_color=_risk_color(sev), corner_radius=30000)
        add_textbox(slide, _severity_label(sev),
                    legend_left + 280000, y, 1500000, 200000,
                    font_size=10, color=text_color,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)


# ═════════════════════════════════════════════════════════════════════
# Variant 2: MINIMAL — Clean grid with dot indicators
# ═════════════════════════════════════════════════════════════════════

def _dot_matrix(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("risk_title", "Risk Assessment Matrix"), theme=t)

    s = t.ux_style
    items = _get_items(c)
    x_label = c.get("risk_x_label", "Likelihood")
    y_label = c.get("risk_y_label", "Impact")
    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    # Minimal grid — thin lines only, no cell fill
    grid_left = m + 600000
    grid_top = content_y + 300000
    grid_w = 5400000
    grid_h = 3600000
    cell_w = grid_w // 3
    cell_h = grid_h // 3
    line_color = tint(t.secondary, 0.75) if not s.dark_mode else tint(t.primary, 0.25)

    # Horizontal rules
    for row in range(4):
        y = grid_top + row * cell_h
        add_line(slide, grid_left, y, grid_left + grid_w, y,
                 color=line_color, width=0.5)
    # Vertical rules
    for col in range(4):
        x = grid_left + col * cell_w
        add_line(slide, x, grid_top, x, grid_top + grid_h,
                 color=line_color, width=0.5)

    # Place risk items as small dots with labels
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w + cell_w // 2
        cy = grid_top + row * cell_h + cell_h // 2
        dot_color = _risk_color(severity)
        add_circle(slide, cx, cy, 100000, fill_color=dot_color)
        add_textbox(slide, name,
                    cx - cell_w // 2 + 20000, cy + 120000,
                    cell_w - 40000, 200000,
                    font_size=8, color=text_color, alignment=PP_ALIGN.CENTER)

    # Axis labels
    _draw_axis_labels(slide, grid_left, grid_top, grid_w, grid_h,
                      x_label, y_label, sub_color, t.accent)
    _draw_grid_labels(slide, grid_left, grid_top, cell_w, cell_h, sub_color)

    # Right side — ordered item list
    list_left = grid_left + grid_w + 500000
    list_top = grid_top
    add_textbox(slide, "Risk Register", list_left, list_top - 50000,
                3000000, 280000,
                font_size=11, bold=True, color=text_color)
    add_line(slide, list_left, list_top + 250000,
             list_left + 2000000, list_top + 250000,
             color=line_color, width=0.5)

    for i, (name, severity, desc) in enumerate(items[:6]):
        y = list_top + 350000 + i * 550000
        add_circle(slide, list_left + 80000, y + 80000, 50000,
                   fill_color=_risk_color(severity))
        add_textbox(slide, name, list_left + 200000, y,
                    2800000, 200000,
                    font_size=10, bold=True, color=text_color)
        add_textbox(slide, desc, list_left + 200000, y + 200000,
                    2800000, 200000,
                    font_size=8, color=sub_color)


# ═════════════════════════════════════════════════════════════════════
# Variant 3: BOLD — Bold colored blocks with large labels
# ═════════════════════════════════════════════════════════════════════

def _bold_heatmap(slide, t, c):
    add_dark_bg(slide, t)
    title = c.get("risk_title", "Risk Assessment Matrix")
    content_y = add_slide_title(slide, title.upper(), theme=t)

    s = t.ux_style
    items = _get_items(c)
    text_color = "#FFFFFF" if s.dark_mode else t.primary

    # Full-width bold grid, 3x3 with thick gaps
    m = int(MARGIN * s.margin_factor)
    grid_left = m + 500000
    grid_top = content_y + 200000
    grid_w = SLIDE_W - 2 * m - 1000000
    grid_h = 4400000
    cell_w = grid_w // 3
    cell_h = grid_h // 3
    gap = 40000

    # Draw 3x3 bold colored blocks
    for row in range(3):
        for col in range(3):
            cx = grid_left + col * cell_w + gap // 2
            cy = grid_top + row * cell_h + gap // 2
            w = cell_w - gap
            h = cell_h - gap
            color = _GRID_3x3[row][col]
            cell_fill = tint(color, 0.1) if s.dark_mode else color
            add_rect(slide, cx, cy, w, h, fill_color=cell_fill)

    # Bold severity labels in cells
    severity_map = {
        (0, 0): "MEDIUM", (0, 1): "HIGH", (0, 2): "CRITICAL",
        (1, 0): "LOW", (1, 1): "MEDIUM", (1, 2): "HIGH",
        (2, 0): "MINIMAL", (2, 1): "LOW", (2, 2): "MEDIUM",
    }
    for (row, col), label in severity_map.items():
        cx = grid_left + col * cell_w + gap // 2
        cy = grid_top + row * cell_h + gap // 2
        w = cell_w - gap
        h = cell_h - gap
        add_textbox(slide, label,
                    cx, cy + h - 350000, w, 300000,
                    font_size=11, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

    # Place risk items as thick accent markers
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w + cell_w // 2
        cy = grid_top + row * cell_h + cell_h // 3
        marker_color = _risk_color(severity)
        # Bold circle marker
        add_circle(slide, cx, cy, 200000, fill_color=marker_color,
                   line_color="#FFFFFF", line_width=3)
        add_textbox(slide, name,
                    cx - cell_w // 2 + 50000, cy - 120000,
                    cell_w - 100000, 240000,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Thick accent line under title (bold style)
    add_rect(slide, grid_left, content_y + 100000, grid_w, 60000,
             fill_color=t.accent)


# ═════════════════════════════════════════════════════════════════════
# Variant 4: ELEVATED — Floating card grid with shadow
# ═════════════════════════════════════════════════════════════════════

def _floating_heatmap(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("risk_title", "Risk Assessment Matrix"), theme=t)

    s = t.ux_style
    items = _get_items(c)
    x_label = c.get("risk_x_label", "Likelihood")
    y_label = c.get("risk_y_label", "Impact")
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    m = int(MARGIN * s.margin_factor)
    grid_left = m + 550000
    grid_top = content_y + 300000
    grid_w = 5600000
    grid_h = 4000000
    cell_w = grid_w // 3
    cell_h = grid_h // 3
    gap = 60000

    # Floating card for each cell
    for row in range(3):
        for col in range(3):
            cx = grid_left + col * cell_w + gap // 2
            cy = grid_top + row * cell_h + gap // 2
            w = cell_w - gap
            h = cell_h - gap
            color = _GRID_3x3[row][col]
            cell_fill = tint(color, 0.2) if s.dark_mode else tint(color, 0.55)
            add_styled_card(slide, cx, cy, w, h, theme=t,
                            fill_color=cell_fill, accent_color=color)

    # Place risk items
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w + cell_w // 2
        cy = grid_top + row * cell_h + cell_h // 2
        marker_color = _risk_color(severity)
        add_circle(slide, cx, cy - 50000, 150000, fill_color=marker_color)
        add_textbox(slide, name,
                    cx - cell_w // 2 + 80000, cy + 120000,
                    cell_w - 160000, 220000,
                    font_size=8, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

    # Axis labels
    _draw_axis_labels(slide, grid_left, grid_top, grid_w, grid_h,
                      x_label, y_label, text_color, t.accent)
    _draw_grid_labels(slide, grid_left, grid_top, cell_w, cell_h, sub_color)

    # Legend card on the right
    legend_left = grid_left + grid_w + 350000
    legend_top = grid_top + 100000
    legend_w = SLIDE_W - legend_left - m
    legend_h = 2800000
    add_styled_card(slide, legend_left, legend_top, legend_w, legend_h,
                    theme=t, accent_color=t.accent)
    add_textbox(slide, "Risk Levels", legend_left + 120000, legend_top + 150000,
                legend_w - 240000, 280000,
                font_size=11, bold=True, color=text_color)
    for i, sev in enumerate(_SEVERITY_ORDER):
        y = legend_top + 550000 + i * 480000
        add_circle(slide, legend_left + 220000, y + 80000, 80000,
                   fill_color=_risk_color(sev))
        add_textbox(slide, sev.capitalize(),
                    legend_left + 380000, y, legend_w - 500000, 200000,
                    font_size=10, color=text_color,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)


# ═════════════════════════════════════════════════════════════════════
# Variant 5: DARK — Dark bg, glowing colored cells
# ═════════════════════════════════════════════════════════════════════

def _dark_heatmap(slide, t, c):
    add_background(slide, t.primary)
    title = c.get("risk_title", "Risk Assessment Matrix")
    content_y = add_slide_title(slide, title.upper(), theme=t)

    items = _get_items(c)
    x_label = c.get("risk_x_label", "Likelihood")
    y_label = c.get("risk_y_label", "Impact")
    text_color = "#FFFFFF"
    sub_color = tint(t.primary, 0.5)

    m = MARGIN
    grid_left = m + 550000
    grid_top = content_y + 300000
    grid_w = 5600000
    grid_h = 4000000
    cell_w = grid_w // 3
    cell_h = grid_h // 3

    # Dark bg grid with neon-glow cells
    for row in range(3):
        for col in range(3):
            cx = grid_left + col * cell_w
            cy = grid_top + row * cell_h
            color = _GRID_3x3[row][col]
            # Deep tint for dark mode glow
            cell_fill = tint(color, 0.08)
            add_rect(slide, cx + 15000, cy + 15000,
                     cell_w - 30000, cell_h - 30000,
                     fill_color=cell_fill,
                     line_color=tint(color, 0.2), line_width=0.75,
                     corner_radius=60000)

    # Thin neon accent line at top of grid
    add_rect(slide, grid_left, grid_top - 30000, grid_w, 25000,
             fill_color=t.accent)

    # Place risk items as glowing dots
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w + cell_w // 2
        cy = grid_top + row * cell_h + cell_h // 2
        marker_color = _risk_color(severity)
        # Outer glow ring
        add_circle(slide, cx, cy, 180000, fill_color=tint(marker_color, 0.12))
        # Inner dot
        add_circle(slide, cx, cy, 110000, fill_color=marker_color)
        add_textbox(slide, name,
                    cx - cell_w // 2 + 20000, cy + 200000,
                    cell_w - 40000, 200000,
                    font_size=8, bold=True, color=marker_color,
                    alignment=PP_ALIGN.CENTER)

    # Axis labels
    _draw_axis_labels(slide, grid_left, grid_top, grid_w, grid_h,
                      x_label, y_label, sub_color, t.accent)
    _draw_grid_labels(slide, grid_left, grid_top, cell_w, cell_h, sub_color)

    # Right side legend
    legend_left = grid_left + grid_w + 400000
    legend_top = grid_top + 200000
    add_textbox(slide, "SEVERITY", legend_left, legend_top - 300000,
                2500000, 280000,
                font_size=10, bold=True, color=tint(t.primary, 0.5))
    for i, sev in enumerate(_SEVERITY_ORDER):
        y = legend_top + i * 380000
        sev_color = _risk_color(sev)
        add_circle(slide, legend_left + 100000, y + 80000, 60000,
                   fill_color=sev_color)
        add_textbox(slide, sev.upper(),
                    legend_left + 240000, y, 1800000, 200000,
                    font_size=9, color=sev_color,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)


# ═════════════════════════════════════════════════════════════════════
# Variant 6: SPLIT — Heatmap left, risk list right
# ═════════════════════════════════════════════════════════════════════

def _split_risk(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("risk_title", "Risk Assessment Matrix"), theme=t)

    s = t.ux_style
    items = _get_items(c)
    x_label = c.get("risk_x_label", "Likelihood")
    y_label = c.get("risk_y_label", "Impact")
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    mid = SLIDE_W // 2 - 100000

    # Left side: compact 3x3 heatmap
    grid_left = MARGIN + 500000
    grid_top = content_y + 300000
    grid_w = mid - MARGIN - 700000
    grid_h = 3800000
    cell_w = grid_w // 3
    cell_h = grid_h // 3

    for row in range(3):
        for col in range(3):
            cx = grid_left + col * cell_w
            cy = grid_top + row * cell_h
            color = _GRID_3x3[row][col]
            cell_fill = tint(color, 0.15) if s.dark_mode else tint(color, 0.5)
            add_color_cell(slide, cx, cy, cell_w, cell_h,
                           fill_color=cell_fill,
                           line_color=tint(text_color, 0.85), line_width=0.5)

    # Place items on grid
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w + cell_w // 2
        cy = grid_top + row * cell_h + cell_h // 2
        add_circle(slide, cx, cy, 120000, fill_color=_risk_color(severity))

    _draw_axis_labels(slide, grid_left, grid_top, grid_w, grid_h,
                      x_label, y_label, sub_color, t.accent)

    # Center divider
    divider_x = mid + 50000
    add_line(slide, divider_x, content_y + 200000, divider_x, FOOTER_TOP - 300000,
             color=t.accent, width=2)

    # Right side: risk item list with severity badges
    list_left = mid + 200000
    list_top = content_y + 200000
    list_w = SLIDE_W - list_left - MARGIN

    add_textbox(slide, "Risk Register", list_left, list_top,
                list_w, 300000,
                font_size=14, bold=True, color=text_color)
    add_gold_accent_line(slide, list_left, list_top + 320000,
                         1500000, color=t.accent)

    for i, (name, severity, desc) in enumerate(items[:6]):
        y = list_top + 500000 + i * 650000
        sev_color = _risk_color(severity)

        # Severity badge
        badge_w = 700000
        add_rect(slide, list_left, y, badge_w, 260000,
                 fill_color=sev_color, corner_radius=40000)
        add_textbox(slide, severity.upper(),
                    list_left, y, badge_w, 260000,
                    font_size=8, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Name + description
        add_textbox(slide, name, list_left + badge_w + 150000, y,
                    list_w - badge_w - 200000, 260000,
                    font_size=11, bold=True, color=text_color,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, desc, list_left + badge_w + 150000, y + 280000,
                    list_w - badge_w - 200000, 250000,
                    font_size=9, color=sub_color)


# ═════════════════════════════════════════════════════════════════════
# Variant 7: GEO — Angular geometric grid with octagon markers
# ═════════════════════════════════════════════════════════════════════

def _angular_heatmap(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("risk_title", "Risk Assessment Matrix"), theme=t)

    s = t.ux_style
    items = _get_items(c)
    x_label = c.get("risk_x_label", "Likelihood")
    y_label = c.get("risk_y_label", "Impact")
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if s.dark_mode else t.secondary

    m = int(MARGIN * s.margin_factor)
    grid_left = m + 550000
    grid_top = content_y + 300000
    grid_w = 5800000
    grid_h = 4000000
    cell_w = grid_w // 3
    cell_h = grid_h // 3

    # Sharp-cornered grid cells with thick borders
    for row in range(3):
        for col in range(3):
            cx = grid_left + col * cell_w
            cy = grid_top + row * cell_h
            color = _GRID_3x3[row][col]
            cell_fill = tint(color, 0.12) if s.dark_mode else tint(color, 0.6)
            add_rect(slide, cx, cy, cell_w, cell_h,
                     fill_color=cell_fill,
                     line_color=tint(t.primary, 0.3) if s.dark_mode else t.secondary,
                     line_width=1.5)

    # Bottom accent bar (geometric style)
    add_rect(slide, grid_left, grid_top + grid_h, grid_w, 50000,
             fill_color=t.accent)

    # Place risk items as octagon markers
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w + cell_w // 2
        cy = grid_top + row * cell_h + cell_h // 2
        marker_color = _risk_color(severity)
        add_octagon(slide, cx, cy - 30000, 160000,
                    fill_color=marker_color, text=name[:3].upper(),
                    text_color="#FFFFFF", font_size=8)
        add_textbox(slide, name,
                    cx - cell_w // 2 + 20000, cy + 160000,
                    cell_w - 40000, 200000,
                    font_size=8, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

    # Axis labels
    _draw_axis_labels(slide, grid_left, grid_top, grid_w, grid_h,
                      x_label, y_label, text_color, t.accent)
    _draw_grid_labels(slide, grid_left, grid_top, cell_w, cell_h, sub_color)

    # Legend with octagon markers
    legend_left = grid_left + grid_w + 400000
    legend_top = grid_top + 200000
    add_textbox(slide, "Severity", legend_left, legend_top - 300000,
                2000000, 280000,
                font_size=11, bold=True, color=text_color)
    for i, sev in enumerate(_SEVERITY_ORDER):
        y = legend_top + i * 420000
        add_octagon(slide, legend_left + 100000, y + 100000, 80000,
                    fill_color=_risk_color(sev))
        add_textbox(slide, sev.capitalize(),
                    legend_left + 280000, y, 1500000, 220000,
                    font_size=10, color=text_color,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)


# ═════════════════════════════════════════════════════════════════════
# Variant 8: EDITORIAL — Elegant thin-ruled grid
# ═════════════════════════════════════════════════════════════════════

def _editorial_heatmap(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("risk_title", "Risk Assessment Matrix"), theme=t)

    s = t.ux_style
    items = _get_items(c)
    x_label = c.get("risk_x_label", "Likelihood")
    y_label = c.get("risk_y_label", "Impact")
    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    rule_color = tint(t.secondary, 0.7) if not s.dark_mode else tint(t.primary, 0.25)

    # Elegant minimal grid — thin rules, no fill, colored text labels
    grid_left = m + 550000
    grid_top = content_y + 350000
    grid_w = 5400000
    grid_h = 3600000
    cell_w = grid_w // 3
    cell_h = grid_h // 3

    # Short accent mark at top
    add_line(slide, grid_left, content_y + 200000,
             grid_left + 500000, content_y + 200000,
             color=t.accent, width=1.5)

    # Thin horizontal rules
    for row in range(4):
        y = grid_top + row * cell_h
        w = 1.0 if row == 0 or row == 3 else 0.5
        add_line(slide, grid_left, y, grid_left + grid_w, y,
                 color=rule_color, width=w)
    # Thin vertical rules
    for col in range(4):
        x = grid_left + col * cell_w
        w = 1.0 if col == 0 or col == 3 else 0.5
        add_line(slide, x, grid_top, x, grid_top + grid_h,
                 color=rule_color, width=w)

    # Place risk items as colored text (no shapes, editorial purity)
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w
        cy = grid_top + row * cell_h
        sev_color = _risk_color(severity)
        # Small colored underline accent
        add_line(slide, cx + 60000, cy + cell_h - 80000,
                 cx + 350000, cy + cell_h - 80000,
                 color=sev_color, width=1.0)
        add_textbox(slide, name,
                    cx + 60000, cy + 60000,
                    cell_w - 120000, 280000,
                    font_size=10, bold=True, color=text_color)
        add_textbox(slide, severity.capitalize(),
                    cx + 60000, cy + 340000,
                    cell_w - 120000, 200000,
                    font_size=8, italic=True, color=sev_color)

    # Axis labels — editorial style (italic)
    add_textbox(slide, x_label,
                grid_left + grid_w // 2 - 500000, grid_top + grid_h + 100000,
                1000000, 280000,
                font_size=11, italic=True, color=sub_color,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, y_label,
                grid_left - 500000, grid_top + grid_h // 2 - 140000,
                450000, 280000,
                font_size=11, italic=True, color=sub_color,
                alignment=PP_ALIGN.CENTER)

    # Bottom editorial rule
    add_line(slide, m, FOOTER_TOP - 400000, SLIDE_W - m, FOOTER_TOP - 400000,
             color=rule_color, width=0.5)


# ═════════════════════════════════════════════════════════════════════
# Variant 9: GRADIENT — Gradient-tinted grid cells
# ═════════════════════════════════════════════════════════════════════

def _gradient_heatmap(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("risk_title", "Risk Assessment Matrix"), theme=t)

    s = t.ux_style
    items = _get_items(c)
    x_label = c.get("risk_x_label", "Likelihood")
    y_label = c.get("risk_y_label", "Impact")
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    m = int(MARGIN * s.margin_factor)
    grid_left = m + 550000
    grid_top = content_y + 350000
    grid_w = 5800000
    grid_h = 4000000
    cell_w = grid_w // 3
    cell_h = grid_h // 3
    gap = 30000

    # Highly-rounded gradient-tinted cells
    for row in range(3):
        for col in range(3):
            cx = grid_left + col * cell_w + gap // 2
            cy = grid_top + row * cell_h + gap // 2
            w = cell_w - gap
            h = cell_h - gap
            color = _GRID_3x3[row][col]
            # Softer gradient tint
            factor = 0.25 if s.dark_mode else 0.65
            cell_fill = tint(color, factor)
            add_rect(slide, cx, cy, w, h,
                     fill_color=cell_fill, corner_radius=120000)

    # Place risk items as soft circles
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w + cell_w // 2
        cy = grid_top + row * cell_h + cell_h // 2
        marker_color = _risk_color(severity)
        # Soft outer ring
        add_circle(slide, cx, cy - 20000, 190000, fill_color=tint(marker_color, 0.5))
        # Inner fill
        add_circle(slide, cx, cy - 20000, 130000, fill_color=marker_color)
        add_textbox(slide, name,
                    cx - cell_w // 2 + 30000, cy + 180000,
                    cell_w - 60000, 200000,
                    font_size=8, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

    # Axis labels
    _draw_axis_labels(slide, grid_left, grid_top, grid_w, grid_h,
                      x_label, y_label, sub_color, t.accent)
    _draw_grid_labels(slide, grid_left, grid_top, cell_w, cell_h, sub_color)

    # Legend as soft badges on right
    legend_left = grid_left + grid_w + 350000
    legend_top = grid_top + 400000
    for i, sev in enumerate(_SEVERITY_ORDER):
        y = legend_top + i * 480000
        sev_color = _risk_color(sev)
        add_rect(slide, legend_left, y, 1800000, 350000,
                 fill_color=tint(sev_color, 0.5 if not s.dark_mode else 0.15),
                 corner_radius=100000)
        add_circle(slide, legend_left + 200000, y + 175000, 80000,
                   fill_color=sev_color)
        add_textbox(slide, sev.capitalize(),
                    legend_left + 380000, y, 1300000, 350000,
                    font_size=10, color=text_color,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)


# ═════════════════════════════════════════════════════════════════════
# Variant 10: RETRO — Vintage styled grid with borders
# ═════════════════════════════════════════════════════════════════════

def _retro_heatmap(slide, t, c):
    add_background(slide, t.light_bg)
    content_y = add_slide_title(slide, c.get("risk_title", "Risk Assessment Matrix"), theme=t)

    items = _get_items(c)
    x_label = c.get("risk_x_label", "Likelihood")
    y_label = c.get("risk_y_label", "Impact")
    text_color = t.primary
    sub_color = t.secondary

    m = MARGIN
    grid_left = m + 550000
    grid_top = content_y + 350000
    grid_w = 5600000
    grid_h = 3800000
    cell_w = grid_w // 3
    cell_h = grid_h // 3

    # Outer double border
    add_rect(slide, grid_left - 30000, grid_top - 30000,
             grid_w + 60000, grid_h + 60000,
             line_color=t.accent, line_width=2.0)
    add_rect(slide, grid_left - 60000, grid_top - 60000,
             grid_w + 120000, grid_h + 120000,
             line_color=t.accent, line_width=0.75)

    # Grid cells with warm tints
    for row in range(3):
        for col in range(3):
            cx = grid_left + col * cell_w
            cy = grid_top + row * cell_h
            color = _GRID_3x3[row][col]
            cell_fill = tint(color, 0.55)
            add_rect(slide, cx, cy, cell_w, cell_h,
                     fill_color=cell_fill,
                     line_color=shade(t.accent, 0.1), line_width=1.0)

    # Decorative corner circles
    for dx, dy in [(0, 0), (grid_w, 0), (0, grid_h), (grid_w, grid_h)]:
        add_circle(slide, grid_left + dx, grid_top + dy, 50000,
                   fill_color=t.accent)

    # Place risk items with vintage markers
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w + cell_w // 2
        cy = grid_top + row * cell_h + cell_h // 2
        marker_color = _risk_color(severity)
        # Double circle marker
        add_circle(slide, cx, cy, 170000, fill_color=marker_color,
                   line_color="#FFFFFF", line_width=2)
        add_circle(slide, cx, cy, 120000,
                   fill_color="#FFFFFF", line_color=marker_color, line_width=1.5)
        add_textbox(slide, name,
                    cx - cell_w // 2 + 30000, cy + 190000,
                    cell_w - 60000, 200000,
                    font_size=8, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

    # Axis labels with dashes
    add_textbox(slide, f"\u2014 {x_label} \u2014",
                grid_left, grid_top + grid_h + 100000,
                grid_w, 300000,
                font_size=11, bold=True, color=t.accent,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, y_label,
                grid_left - 550000, grid_top + grid_h // 2 - 140000,
                500000, 280000,
                font_size=11, bold=True, color=t.accent,
                alignment=PP_ALIGN.CENTER)

    # Legend with decorative rule
    legend_left = grid_left + grid_w + 500000
    legend_top = grid_top + 200000
    add_textbox(slide, "\u2014 Legend \u2014", legend_left, legend_top - 300000,
                2000000, 280000,
                font_size=11, bold=True, color=t.accent,
                alignment=PP_ALIGN.CENTER)
    for i, sev in enumerate(_SEVERITY_ORDER):
        y = legend_top + i * 400000
        sev_color = _risk_color(sev)
        add_circle(slide, legend_left + 100000, y + 100000, 70000,
                   fill_color=sev_color, line_color="#FFFFFF", line_width=1)
        add_textbox(slide, sev.capitalize(),
                    legend_left + 260000, y, 1500000, 220000,
                    font_size=10, color=text_color,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)


# ═════════════════════════════════════════════════════════════════════
# Variant 11: MAGAZINE — Creative asymmetric risk display
# ═════════════════════════════════════════════════════════════════════

def _creative_heatmap(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("risk_title", "Risk Assessment Matrix"), theme=t)

    s = t.ux_style
    items = _get_items(c)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    m = int(MARGIN * s.margin_factor)

    # Asymmetric layout: large critical zone top-right, smaller elsewhere
    # Full-width risk bands — stacked horizontal bars per severity
    band_left = m
    band_top = content_y + 200000
    band_w = SLIDE_W - 2 * m
    available_h = FOOTER_TOP - band_top - 400000

    # Group items by severity
    severity_groups = {}
    for name, severity, desc in items:
        severity_groups.setdefault(severity, []).append((name, desc))

    # Draw severity bands in order from critical to low
    ordered = [s for s in reversed(_SEVERITY_ORDER) if s in severity_groups]
    n_bands = max(len(ordered), 1)
    band_h = available_h // n_bands
    gap = 40000

    for i, sev in enumerate(ordered):
        y = band_top + i * band_h
        sev_color = _risk_color(sev)
        actual_h = band_h - gap

        # Full-width color band
        add_rect(slide, band_left, y, band_w, actual_h,
                 fill_color=sev_color)

        # Severity label (large, left)
        add_textbox(slide, sev.upper(),
                    band_left + 200000, y + 50000,
                    1800000, actual_h - 100000,
                    font_size=24, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Risk items listed across the band
        group = severity_groups[sev]
        item_left = band_left + 2200000
        item_w = (band_w - 2400000) // max(len(group), 1)
        for j, (name, desc) in enumerate(group):
            x = item_left + j * item_w
            # Semi-transparent item block
            add_rect(slide, x + 30000, y + 30000,
                     item_w - 60000, actual_h - 60000,
                     fill_color=shade(sev_color, 0.15),
                     corner_radius=0)
            add_textbox(slide, name,
                        x + 60000, y + 50000,
                        item_w - 120000, 280000,
                        font_size=12, bold=True, color="#FFFFFF")
            add_textbox(slide, desc,
                        x + 60000, y + 340000,
                        item_w - 120000, actual_h - 420000,
                        font_size=9, color=tint("#FFFFFF", 0.3))


# ═════════════════════════════════════════════════════════════════════
# Variant 12: SCHOLARLY — Figure-captioned, thin rules, labeled
# ═════════════════════════════════════════════════════════════════════

def _scholarly_heatmap(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("risk_title", "Risk Assessment Matrix"), theme=t)

    s = t.ux_style
    items = _get_items(c)
    x_label = c.get("risk_x_label", "Likelihood")
    y_label = c.get("risk_y_label", "Impact")
    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    rule_color = tint(t.secondary, 0.7) if not s.dark_mode else tint(t.primary, 0.2)

    # Figure number + caption (scholarly convention)
    add_textbox(slide, "Figure 1.", m, content_y + 100000,
                700000, 250000,
                font_size=10, bold=True, color=text_color)
    add_textbox(slide, "Risk probability-impact assessment matrix",
                m + 700000, content_y + 100000,
                5000000, 250000,
                font_size=10, italic=True, color=sub_color)

    # Thin accent mark (scholarly: very short)
    add_line(slide, m, content_y + 380000, m + 400000, content_y + 380000,
             color=t.accent, width=1.0)

    # Grid with thin rules — no cell fills, scholarly purity
    grid_left = m + 600000
    grid_top = content_y + 500000
    grid_w = 5000000
    grid_h = 3400000
    cell_w = grid_w // 3
    cell_h = grid_h // 3

    # Horizontal rules (outer thicker)
    for row in range(4):
        y = grid_top + row * cell_h
        w = 1.0 if row == 0 or row == 3 else 0.5
        add_line(slide, grid_left, y, grid_left + grid_w, y,
                 color=rule_color, width=w)
    # Vertical rules
    for col in range(4):
        x = grid_left + col * cell_w
        w = 1.0 if col == 0 or col == 3 else 0.5
        add_line(slide, x, grid_top, x, grid_top + grid_h,
                 color=rule_color, width=w)

    # Very subtle fill tints (scholarly: barely visible)
    for row in range(3):
        for col in range(3):
            cx = grid_left + col * cell_w + 5000
            cy = grid_top + row * cell_h + 5000
            color = _GRID_3x3[row][col]
            factor = 0.08 if s.dark_mode else 0.85
            add_rect(slide, cx, cy, cell_w - 10000, cell_h - 10000,
                     fill_color=tint(color, factor))

    # Place items as numbered annotations
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w + 80000
        cy = grid_top + row * cell_h + 80000
        sev_color = _risk_color(severity)
        # Numbered annotation circle
        add_circle(slide, cx + 120000, cy + 120000, 100000,
                   fill_color=sev_color)
        add_textbox(slide, str(i + 1),
                    cx + 20000, cy + 20000, 200000, 200000,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, name,
                    cx + 250000, cy + 60000,
                    cell_w - 380000, 200000,
                    font_size=9, color=text_color)

    # Axis labels (italic, scholarly style)
    add_textbox(slide, x_label,
                grid_left + grid_w // 2 - 600000, grid_top + grid_h + 80000,
                1200000, 250000,
                font_size=10, italic=True, color=sub_color,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, y_label,
                grid_left - 520000, grid_top + grid_h // 2 - 125000,
                470000, 250000,
                font_size=10, italic=True, color=sub_color,
                alignment=PP_ALIGN.CENTER)

    # Right side: numbered legend (footnote style)
    legend_left = grid_left + grid_w + 500000
    legend_top = grid_top + 50000
    legend_w = SLIDE_W - legend_left - m
    add_textbox(slide, "Risk Items", legend_left, legend_top,
                legend_w, 250000,
                font_size=10, bold=True, color=text_color)
    add_line(slide, legend_left, legend_top + 260000,
             legend_left + legend_w, legend_top + 260000,
             color=rule_color, width=0.5)

    for i, (name, severity, desc) in enumerate(items[:6]):
        y = legend_top + 350000 + i * 520000
        sev_color = _risk_color(severity)
        add_textbox(slide, f"{i + 1}.",
                    legend_left, y, 200000, 200000,
                    font_size=9, bold=True, color=sev_color)
        add_textbox(slide, name, legend_left + 200000, y,
                    legend_w - 200000, 200000,
                    font_size=9, bold=True, color=text_color)
        add_textbox(slide, desc, legend_left + 200000, y + 200000,
                    legend_w - 200000, 200000,
                    font_size=8, color=sub_color)


# ═════════════════════════════════════════════════════════════════════
# Variant 13: LABORATORY — Dark bg, data-coded cells
# ═════════════════════════════════════════════════════════════════════

def _laboratory_heatmap(slide, t, c):
    add_background(slide, t.primary)
    content_y = add_slide_title(slide, c.get("risk_title", "Risk Assessment Matrix"), theme=t)

    items = _get_items(c)
    x_label = c.get("risk_x_label", "Likelihood")
    y_label = c.get("risk_y_label", "Impact")
    text_color = "#FFFFFF"
    sub_color = tint(t.primary, 0.45)
    border_color = tint(t.primary, 0.2)

    m = MARGIN
    grid_left = m + 600000
    grid_top = content_y + 350000
    grid_w = 5400000
    grid_h = 3800000
    cell_w = grid_w // 3
    cell_h = grid_h // 3

    # Lab-style: color-coded left border accent on each row
    for row in range(3):
        row_color = _GRID_3x3[row][1]  # middle col color for row accent
        add_rect(slide, grid_left - 40000, grid_top + row * cell_h,
                 30000, cell_h, fill_color=row_color)

    # Grid cells with data-coded fills
    for row in range(3):
        for col in range(3):
            cx = grid_left + col * cell_w
            cy = grid_top + row * cell_h
            color = _GRID_3x3[row][col]
            cell_fill = tint(color, 0.06)
            add_rect(slide, cx, cy, cell_w, cell_h,
                     fill_color=cell_fill,
                     line_color=border_color, line_width=0.75,
                     corner_radius=30000)

    # Data-coded risk score labels in corners of cells
    score_map = {
        (0, 0): "6", (0, 1): "8", (0, 2): "9",
        (1, 0): "3", (1, 1): "5", (1, 2): "7",
        (2, 0): "1", (2, 1): "2", (2, 2): "4",
    }
    for (row, col), score in score_map.items():
        cx = grid_left + col * cell_w
        cy = grid_top + row * cell_h
        add_textbox(slide, score,
                    cx + cell_w - 280000, cy + 30000,
                    250000, 200000,
                    font_size=14, bold=True,
                    color=tint(_GRID_3x3[row][col], 0.3),
                    alignment=PP_ALIGN.RIGHT)

    # Place risk items as data points
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w + cell_w // 2
        cy = grid_top + row * cell_h + cell_h // 2
        marker_color = _risk_color(severity)
        # Outer data ring
        add_circle(slide, cx, cy - 20000, 160000,
                   fill_color=tint(marker_color, 0.1),
                   line_color=marker_color, line_width=1.5)
        # Inner point
        add_circle(slide, cx, cy - 20000, 70000, fill_color=marker_color)
        add_textbox(slide, name,
                    cx - cell_w // 2 + 30000, cy + 160000,
                    cell_w - 60000, 200000,
                    font_size=8, bold=True, color=marker_color,
                    alignment=PP_ALIGN.CENTER)

    # Accent line at top
    add_rect(slide, grid_left, grid_top - 50000, grid_w, 30000,
             fill_color=t.accent)

    # Axis labels
    add_textbox(slide, x_label.upper(),
                grid_left, grid_top + grid_h + 100000,
                grid_w, 280000,
                font_size=9, bold=True, color=sub_color,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, y_label.upper(),
                grid_left - 550000, grid_top + grid_h // 2 - 140000,
                500000, 280000,
                font_size=9, bold=True, color=sub_color,
                alignment=PP_ALIGN.CENTER)

    # Right panel: data summary
    panel_left = grid_left + grid_w + 400000
    panel_top = grid_top + 50000
    panel_w = SLIDE_W - panel_left - m
    panel_h = grid_h - 100000

    add_rect(slide, panel_left, panel_top, panel_w, panel_h,
             fill_color=tint(t.primary, 0.06),
             line_color=border_color, line_width=0.75,
             corner_radius=30000)
    # Left accent stripe
    add_rect(slide, panel_left, panel_top, 30000, panel_h,
             fill_color=t.accent)

    add_textbox(slide, "Risk Summary", panel_left + 100000, panel_top + 100000,
                panel_w - 200000, 280000,
                font_size=11, bold=True, color=text_color)
    add_line(slide, panel_left + 100000, panel_top + 400000,
             panel_left + panel_w - 100000, panel_top + 400000,
             color=border_color, width=0.5)

    for i, (name, severity, desc) in enumerate(items[:6]):
        y = panel_top + 500000 + i * 500000
        sev_color = _risk_color(severity)
        # Severity code
        add_rect(slide, panel_left + 100000, y, 50000, 200000,
                 fill_color=sev_color)
        add_textbox(slide, name, panel_left + 220000, y,
                    panel_w - 340000, 200000,
                    font_size=9, bold=True, color=text_color)
        add_textbox(slide, severity.upper(), panel_left + 220000, y + 200000,
                    panel_w - 340000, 180000,
                    font_size=7, color=sev_color)


# ═════════════════════════════════════════════════════════════════════
# Variant 14: DASHBOARD — Dense compact risk matrix with metrics
# ═════════════════════════════════════════════════════════════════════

def _dashboard_heatmap(slide, t, c):
    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("risk_title", "Risk Assessment Matrix"), theme=t)

    s = t.ux_style
    items = _get_items(c)
    x_label = c.get("risk_x_label", "Likelihood")
    y_label = c.get("risk_y_label", "Impact")
    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary
    border_color = tint(t.secondary, 0.75) if not s.dark_mode else tint(t.primary, 0.2)

    # Top: severity count metrics strip
    strip_top = content_y + 50000
    strip_h = 600000
    counts = {"low": 0, "medium": 0, "high": 0, "critical": 0}
    for _, sev, _ in items:
        if sev in counts:
            counts[sev] += 1

    n_sev = len(_SEVERITY_ORDER)
    metric_w = (SLIDE_W - 2 * m) // n_sev
    for i, sev in enumerate(_SEVERITY_ORDER):
        x = m + i * metric_w
        sev_color = _risk_color(sev)
        # Tile background
        tile_fill = tint(sev_color, 0.15) if s.dark_mode else tint(sev_color, 0.85)
        add_rect(slide, x + 20000, strip_top, metric_w - 40000, strip_h,
                 fill_color=tile_fill,
                 line_color=border_color, line_width=0.5,
                 corner_radius=50000)
        # Count number
        add_textbox(slide, str(counts[sev]),
                    x + 20000, strip_top + 50000,
                    metric_w - 40000, 300000,
                    font_size=24, bold=True, color=sev_color,
                    alignment=PP_ALIGN.CENTER)
        # Label
        add_textbox(slide, sev.capitalize(),
                    x + 20000, strip_top + 350000,
                    metric_w - 40000, 200000,
                    font_size=9, color=sub_color,
                    alignment=PP_ALIGN.CENTER)

    # Compact 3x3 grid below metrics
    grid_left = m + 550000
    grid_top = strip_top + strip_h + 250000
    grid_w = 4800000
    grid_h = 3200000
    cell_w = grid_w // 3
    cell_h = grid_h // 3

    for row in range(3):
        for col in range(3):
            cx = grid_left + col * cell_w
            cy = grid_top + row * cell_h
            color = _GRID_3x3[row][col]
            cell_fill = tint(color, 0.12) if s.dark_mode else tint(color, 0.55)
            add_rect(slide, cx, cy, cell_w, cell_h,
                     fill_color=cell_fill,
                     line_color=border_color, line_width=0.5,
                     corner_radius=50000)

    # Place items as compact dots
    positioned = _assign_grid_positions(items)
    for i, (name, severity, desc, row, col) in enumerate(positioned[:6]):
        cx = grid_left + col * cell_w + cell_w // 2
        cy = grid_top + row * cell_h + cell_h // 2
        marker_color = _risk_color(severity)
        add_circle(slide, cx, cy, 120000, fill_color=marker_color)
        add_textbox(slide, name[:10],
                    cx - cell_w // 2 + 15000, cy + 140000,
                    cell_w - 30000, 180000,
                    font_size=7, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

    # Axis labels (compact)
    add_textbox(slide, x_label,
                grid_left, grid_top + grid_h + 50000,
                grid_w, 250000,
                font_size=9, bold=True, color=sub_color,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, y_label,
                grid_left - 480000, grid_top + grid_h // 2 - 125000,
                430000, 250000,
                font_size=9, bold=True, color=sub_color,
                alignment=PP_ALIGN.CENTER)
    _draw_grid_labels(slide, grid_left, grid_top, cell_w, cell_h, sub_color)

    # Right panel: compact risk list
    panel_left = grid_left + grid_w + 350000
    panel_top = grid_top
    panel_w = SLIDE_W - panel_left - m

    add_textbox(slide, "Risk Details", panel_left, panel_top - 50000,
                panel_w, 250000,
                font_size=10, bold=True, color=text_color)
    add_line(slide, panel_left, panel_top + 220000,
             panel_left + panel_w, panel_top + 220000,
             color=border_color, width=0.5)

    row_h = (grid_h - 300000) // max(len(items[:6]), 1)
    for i, (name, severity, desc) in enumerate(items[:6]):
        y = panel_top + 300000 + i * row_h
        sev_color = _risk_color(severity)

        # Compact severity indicator
        add_rect(slide, panel_left, y + 20000, 60000, 160000,
                 fill_color=sev_color, corner_radius=20000)
        add_textbox(slide, name, panel_left + 120000, y,
                    panel_w - 130000, 200000,
                    font_size=9, bold=True, color=text_color)
        add_textbox(slide, desc, panel_left + 120000, y + 200000,
                    panel_w - 130000, 180000,
                    font_size=7, color=sub_color)

        # Separator
        if i < len(items[:6]) - 1:
            sep_y = y + row_h - 30000
            add_line(slide, panel_left, sep_y,
                     panel_left + panel_w, sep_y,
                     color=border_color, width=0.25)
