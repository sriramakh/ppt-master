"""Slide 12 — Process Circular: 4 nodes in diamond pattern + center circle."""

from __future__ import annotations

import math

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP
from pptmaster.builder.helpers import add_textbox, add_rect, add_circle, add_line, add_gold_accent_line, add_styled_card, add_slide_title, add_dark_bg, add_background
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    p = t.palette
    s = t.ux_style

    _dispatch = {
        "scholarly-numbered": _scholarly_numbered,
        "laboratory-flow": _laboratory_flow,
        "dashboard-steps": _dashboard_steps,
    }
    _fn = _dispatch.get(getattr(s, "process", None))
    if _fn:
        return _fn(slide, t, c)

    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("cycle_title", "Continuous Improvement Cycle"), theme=t)

    connector_color = tint(t.primary, 0.4) if s.dark_mode else t.secondary

    cx = SLIDE_W // 2
    cy = content_y + 2200000
    spread = 2000000
    node_radius = 380000

    phases = c.get("cycle_phases", ["Plan", "Execute", "Review", "Improve"])
    colors = [p[i % len(p)] for i in range(4)]
    positions = [(cx, cy - spread), (cx + spread, cy), (cx, cy + spread), (cx - spread, cy)]

    # Connectors
    for i in range(4):
        ni = (i + 1) % 4
        add_line(slide, positions[i][0], positions[i][1], positions[ni][0], positions[ni][1],
                 color=connector_color, width=1.5)

    # Center circle
    center_radius = 300000
    add_circle(slide, cx, cy, center_radius, fill_color=t.accent)
    add_textbox(slide, "Core\nProcess", cx - center_radius, cy - center_radius,
                center_radius * 2, center_radius * 2,
                font_size=13, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Nodes
    for i, (label, color, (nx, ny)) in enumerate(zip(phases, colors, positions)):
        add_circle(slide, nx, ny, node_radius, fill_color=color)
        add_textbox(slide, label, nx - node_radius, ny - node_radius, node_radius * 2, node_radius * 2,
                    font_size=16, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant: SCHOLARLY — white bg, circular text layout, thin lines ──────

def _scholarly_numbered(slide, t, c):
    p = t.palette

    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, c.get("cycle_title", "Research Cycle"), theme=t)

    phases = c.get("cycle_phases", ["Plan", "Execute", "Review", "Improve"])
    n = len(phases)

    cx = SLIDE_W // 2
    cy = content_y + 2200000
    spread = 1600000

    # Place labels in a circle with thin connecting lines
    positions = []
    for i in range(n):
        angle = -math.pi / 2 + i * (2 * math.pi / n)
        nx = cx + int(spread * math.cos(angle))
        ny = cy + int(spread * math.sin(angle))
        positions.append((nx, ny))

    # Thin connecting lines between adjacent nodes
    for i in range(n):
        ni = (i + 1) % n
        add_line(slide, positions[i][0], positions[i][1],
                 positions[ni][0], positions[ni][1],
                 color=tint(t.secondary, 0.6), width=0.75)

    # Center label
    add_textbox(slide, "Cycle", cx - 400000, cy - 150000, 800000, 300000,
                font_size=11, bold=True, color=t.secondary, alignment=PP_ALIGN.CENTER)

    # Small circle at center
    add_circle(slide, cx, cy, 60000, fill_color=tint(t.secondary, 0.8))

    # Node labels as numbered text with small dot marker
    for i, (label, (nx, ny)) in enumerate(zip(phases, positions)):
        # Small dot marker
        add_circle(slide, nx, ny, 50000,
                   fill_color="#FFFFFF", line_color=t.primary, line_width=1.0)

        # Number inside dot
        add_textbox(slide, str(i + 1), nx - 50000, ny - 50000, 100000, 100000,
                    font_size=8, bold=True, color=t.primary,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Label positioned outside the circle
        offset_x = nx - cx
        offset_y = ny - cy
        label_x = nx + int(offset_x * 0.35) - 500000
        label_y = ny + int(offset_y * 0.35) - 150000

        add_textbox(slide, f"{i + 1}. {label}", label_x, label_y, 1000000, 300000,
                    font_size=12, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)


# ── Variant: LABORATORY — dark bg, accent-colored connectors ─────────────

def _laboratory_flow(slide, t, c):
    p = t.palette

    add_background(slide, t.primary)
    content_y = add_slide_title(slide, c.get("cycle_title", "Iterative Process").upper(), theme=t)

    phases = c.get("cycle_phases", ["Plan", "Execute", "Review", "Improve"])
    n = len(phases)
    colors = [p[i % len(p)] for i in range(n)]

    cx = SLIDE_W // 2
    cy = content_y + 2200000
    spread = 1800000
    node_radius = 350000

    positions = []
    for i in range(n):
        angle = -math.pi / 2 + i * (2 * math.pi / n)
        nx = cx + int(spread * math.cos(angle))
        ny = cy + int(spread * math.sin(angle))
        positions.append((nx, ny))

    # Accent-colored connector lines
    for i in range(n):
        ni = (i + 1) % n
        accent = colors[i]
        add_line(slide, positions[i][0], positions[i][1],
                 positions[ni][0], positions[ni][1],
                 color=tint(accent, 0.3), width=2)

    # Center circle
    center_r = 280000
    add_circle(slide, cx, cy, center_r, fill_color=tint(t.primary, 0.15))
    add_circle(slide, cx, cy, center_r - 40000, fill_color=t.accent)
    add_textbox(slide, "Core\nProcess", cx - center_r, cy - center_r,
                center_r * 2, center_r * 2,
                font_size=12, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Dark cards at node positions
    for i, (label, color, (nx, ny)) in enumerate(zip(phases, colors, positions)):
        card_fill = tint(t.primary, 0.1)

        add_rect(slide, nx - node_radius, ny - node_radius,
                 node_radius * 2, node_radius * 2,
                 fill_color=card_fill, line_color=tint(t.primary, 0.2),
                 line_width=0.5, corner_radius=50000)

        # Accent top bar
        add_rect(slide, nx - node_radius + 1, ny - node_radius,
                 node_radius * 2 - 2, 35000,
                 fill_color=color)

        # Number
        add_textbox(slide, f"0{i + 1}", nx - node_radius + 40000, ny - node_radius + 80000,
                    node_radius * 2 - 80000, 250000,
                    font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)

        # Label
        add_textbox(slide, label, nx - node_radius + 40000, ny - 50000,
                    node_radius * 2 - 80000, 300000,
                    font_size=13, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER)


# ── Variant: DASHBOARD — compact diagram, header band, tighter spacing ───

def _dashboard_steps(slide, t, c):
    p = t.palette

    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, c.get("cycle_title", "Continuous Improvement Cycle"), theme=t)

    # Header accent band
    band_h = 45000
    add_rect(slide, 0, content_y, SLIDE_W, band_h, fill_color=t.accent)

    phases = c.get("cycle_phases", ["Plan", "Execute", "Review", "Improve"])
    n = len(phases)
    colors = [p[i % len(p)] for i in range(n)]

    cx = SLIDE_W // 2
    cy = content_y + band_h + 1900000
    spread = 1500000
    node_radius = 300000

    positions = []
    for i in range(n):
        angle = -math.pi / 2 + i * (2 * math.pi / n)
        nx = cx + int(spread * math.cos(angle))
        ny = cy + int(spread * math.sin(angle))
        positions.append((nx, ny))

    # Thin connectors
    for i in range(n):
        ni = (i + 1) % n
        add_line(slide, positions[i][0], positions[i][1],
                 positions[ni][0], positions[ni][1],
                 color=tint(t.secondary, 0.5), width=1)

    # Small center dot
    add_circle(slide, cx, cy, 100000, fill_color=t.accent)
    add_textbox(slide, "Core", cx - 100000, cy - 100000, 200000, 200000,
                font_size=8, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Compact colored circles
    for i, (label, color, (nx, ny)) in enumerate(zip(phases, colors, positions)):
        add_circle(slide, nx, ny, node_radius, fill_color=color)

        # Step number
        add_textbox(slide, str(i + 1), nx - node_radius, ny - node_radius,
                    node_radius * 2, int(node_radius * 0.9),
                    font_size=18, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.BOTTOM)

        # Label
        add_textbox(slide, label, nx - node_radius, ny - int(node_radius * 0.1),
                    node_radius * 2, int(node_radius * 1.1),
                    font_size=12, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.TOP)
