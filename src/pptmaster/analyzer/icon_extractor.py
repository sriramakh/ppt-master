"""Extract vector icons from toolkit PPTX files.

Icons are stored as grouped shapes in the toolkit slides.
We serialize each group's XML for later cloning into target presentations.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation

from pptmaster.models import IconInfo

_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def extract_icons(toolkit_path: str | Path) -> list[IconInfo]:
    """Extract all vector icon shapes from a toolkit PPTX.

    Each slide is treated as a category. Group shapes and individual
    shapes (excluding placeholders) are captured as icons.
    """
    toolkit_path = Path(toolkit_path)
    if not toolkit_path.exists():
        return []

    prs = Presentation(str(toolkit_path))
    icons: list[IconInfo] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        category = _guess_category(slide, slide_num)

        # Process all shapes on the slide
        sp_tree = slide._element.find(f".//{{{_NS_P}}}cSld/{{{_NS_P}}}spTree")
        if sp_tree is None:
            continue

        shape_idx = 0
        for child in sp_tree:
            tag = etree.QName(child.tag).localname

            # Skip non-shape elements (nvGrpSpPr, grpSpPr)
            if tag not in ("sp", "grpSp", "cxnSp"):
                continue

            # Skip placeholder shapes
            nv_pr = child.find(f".//{{{_NS_P}}}nvPr/{{{_NS_P}}}ph")
            if nv_pr is not None:
                continue

            name = _get_shape_name(child)
            xml_snippet = etree.tostring(child, encoding="unicode")

            # Get dimensions
            width, height = _get_shape_size(child)

            keywords = _generate_keywords(name, category)

            icons.append(IconInfo(
                name=name or f"icon_{slide_num}_{shape_idx}",
                source_slide=slide_num,
                shape_index=shape_idx,
                category=category,
                keywords=keywords,
                xml_snippet=xml_snippet,
                width=width,
                height=height,
            ))
            shape_idx += 1

    return icons


def _guess_category(slide: Any, slide_num: int) -> str:
    """Try to determine the icon category from slide title or content."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and len(text) < 50:
                return text.lower()
    return f"slide_{slide_num}"


def _get_shape_name(elem: etree._Element) -> str:
    """Extract shape name from nvSpPr/cNvPr or nvGrpSpPr/cNvPr."""
    for ns in [_NS_P, _NS_A]:
        cnv_pr = elem.find(f".//{{{ns}}}cNvPr")
        if cnv_pr is not None:
            return cnv_pr.get("name", "")
    return ""


def _get_shape_size(elem: etree._Element) -> tuple[int, int]:
    """Get width and height from xfrm/ext."""
    xfrm = elem.find(f".//{{{_NS_A}}}xfrm")
    if xfrm is not None:
        ext = xfrm.find(f"{{{_NS_A}}}ext")
        if ext is not None:
            return int(ext.get("cx", 0)), int(ext.get("cy", 0))
    return 0, 0


def _generate_keywords(name: str, category: str) -> list[str]:
    """Generate search keywords from shape name and category."""
    words: set[str] = set()

    # Split name by non-alpha characters
    for part in re.split(r"[^a-zA-Z]+", name):
        if len(part) > 2:
            words.add(part.lower())

    # Add category words
    for part in re.split(r"[^a-zA-Z]+", category):
        if len(part) > 2:
            words.add(part.lower())

    return sorted(words)
