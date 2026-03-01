"""Slide 38 -- Gauge Dashboard: 14 unique visual variants dispatched by UX style.

Each variant renders 3-4 donut-style gauge indicators using layered circles
(track -> fill -> inner cutout) for a reliable cross-platform gauge visual.
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_circle, add_line, add_block_arc,
    add_gold_accent_line, add_slide_title, add_dark_bg, add_background,
    add_styled_card, add_color_cell,
)
from pptmaster.assets.color_utils import tint, shade


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    s = t.ux_style
    c = t.content
    p = t.palette

    dispatch = {
        "arc-row": _arc_row,
        "thin-arc": _thin_arc,
        "bold-arc": _bold_arc,
        "floating-arc": _floating_arc,
        "glow-arc": _glow_arc,
        "split-gauge": _split_gauge,
        "angular-arc": _angular_arc,
        "editorial-arc": _editorial_arc,
        "gradient-arc": _gradient_arc,
        "retro-arc": _retro_arc,
        "creative-arc": _creative_arc,
        "scholarly-arc": _scholarly_arc,
        "laboratory-arc": _laboratory_arc,
        "dashboard-arc": _dashboard_arc,
    }
    builder = dispatch.get(s.gauge, _arc_row)
    builder(slide, t, c)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _get_gauges(c):
    """Extract gauge data with sensible defaults."""
    return c.get("gauges", [
        ("Revenue Target", "$8.2M / $10M", 0.82),
        ("Customer Satisfaction", "94%", 0.94),
        ("Sprint Velocity", "42 / 50 pts", 0.84),
        ("Uptime SLA", "99.95%", 0.999),
    ])


def _draw_donut_gauge(
    slide,
    cx: int,
    cy: int,
    outer_r: int,
    inner_r: int,
    progress: float,
    *,
    track_color: str,
    fill_color: str,
    center_color: str,
    value_text: str = "",
    value_color: str = "#FFFFFF",
    value_size: int = 18,
    label_text: str = "",
    label_color: str = "#AAAAAA",
    label_size: int = 10,
    label_below: bool = True,
    pct_text: str = "",
    pct_color: str = "",
    pct_size: int = 11,
) -> None:
    """Draw a donut-style gauge using block arcs to accurately show progress.

    Layers (bottom to top):
      1. Block arc (360°) in track_color — the empty background ring
      2. Block arc (progress×360°) in fill_color — the filled portion
      3. Inner circle in center_color — creates the donut hole
      4. Value text centered in the donut hole
      5. Optional label below (or above) the gauge
    """
    progress = max(0.0, min(1.0, progress))
    thickness = (outer_r - inner_r) / outer_r

    # 1. Full ring background (track)
    add_block_arc(slide, cx, cy, outer_r,
                  fill_color=track_color,
                  start_angle=0.0, sweep_angle=360.0,
                  thickness=thickness)

    # 2. Partial arc fill — shows actual progress clockwise from top
    sweep = progress * 360.0
    if sweep > 1.0:
        add_block_arc(slide, cx, cy, outer_r,
                      fill_color=fill_color,
                      start_angle=270.0, sweep_angle=sweep,
                      thickness=thickness)

    # 3. Inner circle cutout — creates the donut center
    add_circle(slide, cx, cy, inner_r, fill_color=center_color)

    # 4. Value text centered in the hole
    if value_text:
        tw = inner_r * 2
        th = int(inner_r * 0.7)
        add_textbox(
            slide, value_text,
            cx - inner_r, cy - th // 2,
            tw, th,
            font_size=value_size, bold=True, color=value_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

    # 4b. Optional percentage text below value in the donut hole
    if pct_text:
        pw = inner_r * 2
        ph = int(inner_r * 0.45)
        add_textbox(
            slide, pct_text,
            cx - inner_r, cy + int(inner_r * 0.2),
            pw, ph,
            font_size=pct_size, bold=False, color=pct_color or value_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

    # 5. Label below or above the gauge
    if label_text:
        lw = outer_r * 3
        lh = 350000
        if label_below:
            ly = cy + outer_r + 60000
        else:
            ly = cy - outer_r - lh - 30000
        add_textbox(
            slide, label_text,
            cx - lw // 2, ly, lw, lh,
            font_size=label_size, bold=True, color=label_color,
            alignment=PP_ALIGN.CENTER,
        )


def _draw_progress_ring(
    slide,
    cx: int,
    cy: int,
    outer_r: int,
    thickness: int,
    progress: float,
    *,
    track_color: str,
    fill_color: str,
    center_color: str,
) -> None:
    """Draw a ring gauge using block arcs to accurately show the metric level."""
    progress = max(0.0, min(1.0, progress))
    inner_r = outer_r - thickness
    arc_thickness = thickness / outer_r

    # Track (full ring)
    add_block_arc(slide, cx, cy, outer_r,
                  fill_color=track_color,
                  start_angle=0.0, sweep_angle=360.0,
                  thickness=arc_thickness)

    # Fill arc (shows progress)
    sweep = progress * 360.0
    if sweep > 1.0:
        add_block_arc(slide, cx, cy, outer_r,
                      fill_color=fill_color,
                      start_angle=270.0, sweep_angle=sweep,
                      thickness=arc_thickness)

    # Center cutout
    add_circle(slide, cx, cy, inner_r, fill_color=center_color)


# ---------------------------------------------------------------------------
# Variant 1: CLASSIC (arc-row) -- 3-4 donut arc gauges in a row
# ---------------------------------------------------------------------------

def _arc_row(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("gauge_title", "Performance Gauges"), theme=t)
    gauges = _get_gauges(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    center_bg = t.primary if s.dark_mode else "#FFFFFF"

    n = min(len(gauges), 4)
    total_w = SLIDE_W - 2 * MARGIN
    col_w = total_w // n
    gauge_cy = content_top + 1800000
    outer_r = min(col_w // 2 - 100000, 750000)
    inner_r = int(outer_r * 0.65)

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        cx = MARGIN + i * col_w + col_w // 2
        track_c = tint(accent, 0.8) if not s.dark_mode else tint(t.primary, 0.2)

        _draw_donut_gauge(
            slide, cx, gauge_cy, outer_r, inner_r, progress,
            track_color=track_c, fill_color=accent, center_color=center_bg,
            value_text=value_display, value_color=text_color, value_size=16,
            label_text=label, label_color=sub_color, label_size=11,
            pct_text=f"{int(progress * 100)}%", pct_color=accent, pct_size=10,
        )


# ---------------------------------------------------------------------------
# Variant 2: MINIMAL (thin-arc) -- Thin arcs with clean typography
# ---------------------------------------------------------------------------

def _thin_arc(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("gauge_title", "Performance Gauges"), theme=t)
    gauges = _get_gauges(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    center_bg = t.primary if s.dark_mode else t.light_bg

    n = min(len(gauges), 4)
    total_w = SLIDE_W - 2 * m
    col_w = total_w // n
    gauge_cy = content_top + 1600000
    outer_r = min(col_w // 2 - 150000, 650000)
    inner_r = int(outer_r * 0.82)  # thin ring

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        cx = m + i * col_w + col_w // 2
        track_c = tint(t.secondary, 0.85) if not s.dark_mode else tint(t.primary, 0.15)

        # Thin ring
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=track_c)
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=accent)
        add_circle(slide, cx, gauge_cy, inner_r, fill_color=center_bg)

        # Value in center -- minimal, clean
        tw = inner_r * 2
        add_textbox(
            slide, value_display,
            cx - inner_r, gauge_cy - int(inner_r * 0.25),
            tw, int(inner_r * 0.5),
            font_size=14, bold=True, color=text_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Label below gauge
        add_textbox(
            slide, label.upper(),
            cx - col_w // 2 + 40000, gauge_cy + outer_r + 80000,
            col_w - 80000, 250000,
            font_size=9, bold=True, color=sub_color,
            alignment=PP_ALIGN.CENTER,
        )

        # Thin horizontal rule below label
        rule_w = 400000
        rule_y = gauge_cy + outer_r + 380000
        add_line(slide, cx - rule_w // 2, rule_y, cx + rule_w // 2, rule_y,
                 color=tint(t.secondary, 0.7), width=0.5)

        # Percentage below rule
        add_textbox(
            slide, f"{int(progress * 100)}%",
            cx - 300000, rule_y + 60000,
            600000, 250000,
            font_size=11, color=accent, alignment=PP_ALIGN.CENTER,
        )


# ---------------------------------------------------------------------------
# Variant 3: BOLD (bold-arc) -- Large bold arcs with oversized values
# ---------------------------------------------------------------------------

def _bold_arc(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(
        slide, c.get("gauge_title", "Performance Gauges").upper(), theme=t
    )
    gauges = _get_gauges(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    center_bg = t.primary if s.dark_mode else "#FFFFFF"

    n = min(len(gauges), 4)
    total_w = SLIDE_W - 2 * MARGIN
    col_w = total_w // n
    gauge_cy = content_top + 1700000
    outer_r = min(col_w // 2 - 60000, 800000)
    inner_r = int(outer_r * 0.60)

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        cx = MARGIN + i * col_w + col_w // 2
        track_c = tint(accent, 0.75) if not s.dark_mode else tint(t.primary, 0.15)

        # Bold thick ring
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=track_c)
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=accent)
        add_circle(slide, cx, gauge_cy, inner_r, fill_color=center_bg)

        # Thick accent bar on left of column
        bar_left = MARGIN + i * col_w + 20000
        add_rect(slide, bar_left, content_top + 100000, 70000, 4400000, fill_color=accent)

        # Oversized percentage in center
        add_textbox(
            slide, f"{int(progress * 100)}%",
            cx - inner_r, gauge_cy - int(inner_r * 0.35),
            inner_r * 2, int(inner_r * 0.7),
            font_size=28, bold=True, color=text_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Label UPPERCASE below
        add_textbox(
            slide, label.upper(),
            cx - col_w // 2 + 60000, gauge_cy + outer_r + 100000,
            col_w - 120000, 300000,
            font_size=12, bold=True, color=accent,
            alignment=PP_ALIGN.CENTER,
        )

        # Value display below label
        add_textbox(
            slide, value_display,
            cx - col_w // 2 + 60000, gauge_cy + outer_r + 420000,
            col_w - 120000, 300000,
            font_size=11, color=text_color, alignment=PP_ALIGN.CENTER,
        )

        # Bold accent underline
        ul_w = col_w - 300000
        add_rect(
            slide, cx - ul_w // 2, gauge_cy + outer_r + 770000,
            ul_w, 50000, fill_color=accent,
        )


# ---------------------------------------------------------------------------
# Variant 4: ELEVATED (floating-arc) -- Arcs on floating cards with shadow
# ---------------------------------------------------------------------------

def _floating_arc(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("gauge_title", "Performance Gauges"), theme=t)
    gauges = _get_gauges(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    card_bg = tint(t.primary, 0.12) if s.dark_mode else "#FFFFFF"
    center_bg = card_bg

    n = min(len(gauges), 4)
    gap = 250000
    card_w = (SLIDE_W - 2 * MARGIN - (n - 1) * gap) // n
    card_h = FOOTER_TOP - content_top - 400000
    card_top = content_top + 150000

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        left = MARGIN + i * (card_w + gap)

        # Floating card with shadow
        add_styled_card(slide, left, card_top, card_w, card_h, theme=t, accent_color=accent)

        # Gauge centered in card
        cx = left + card_w // 2
        gauge_cy = card_top + card_h // 2 - 250000
        outer_r = min(card_w // 2 - 120000, 550000)
        inner_r = int(outer_r * 0.65)
        track_c = tint(accent, 0.8) if not s.dark_mode else tint(t.primary, 0.2)

        add_circle(slide, cx, gauge_cy, outer_r, fill_color=track_c)
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=accent)
        add_circle(slide, cx, gauge_cy, inner_r, fill_color=center_bg)

        # Value in center
        add_textbox(
            slide, f"{int(progress * 100)}%",
            cx - inner_r, gauge_cy - int(inner_r * 0.3),
            inner_r * 2, int(inner_r * 0.6),
            font_size=20, bold=True, color=text_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Label below gauge
        label_y = gauge_cy + outer_r + 80000
        add_textbox(
            slide, label,
            left + 60000, label_y, card_w - 120000, 300000,
            font_size=11, bold=True, color=sub_color,
            alignment=PP_ALIGN.CENTER,
        )

        # Value display at bottom of card
        add_textbox(
            slide, value_display,
            left + 60000, label_y + 320000, card_w - 120000, 280000,
            font_size=10, color=text_color, alignment=PP_ALIGN.CENTER,
        )


# ---------------------------------------------------------------------------
# Variant 5: DARK (glow-arc) -- Dark bg, glowing neon arcs
# ---------------------------------------------------------------------------

def _glow_arc(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(
        slide, c.get("gauge_title", "Performance Gauges").upper(), theme=t
    )
    gauges = _get_gauges(c)
    p = t.palette

    n = min(len(gauges), 4)
    total_w = SLIDE_W - 2 * MARGIN
    col_w = total_w // n
    gauge_cy = content_top + 1700000
    outer_r = min(col_w // 2 - 100000, 700000)
    inner_r = int(outer_r * 0.68)
    center_bg = t.primary

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        cx = MARGIN + i * col_w + col_w // 2

        # Glow halo (larger faint circle behind)
        glow_r = outer_r + 60000
        add_circle(slide, cx, gauge_cy, glow_r, fill_color=tint(accent, 0.08))

        # Dark track
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=tint(t.primary, 0.12))

        # Neon fill ring
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=accent)

        # Center cutout
        add_circle(slide, cx, gauge_cy, inner_r, fill_color=center_bg)

        # Neon value
        add_textbox(
            slide, f"{int(progress * 100)}%",
            cx - inner_r, gauge_cy - int(inner_r * 0.35),
            inner_r * 2, int(inner_r * 0.5),
            font_size=22, bold=True, color=accent,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Value display in faint white
        add_textbox(
            slide, value_display,
            cx - inner_r, gauge_cy + int(inner_r * 0.1),
            inner_r * 2, int(inner_r * 0.35),
            font_size=9, color=tint(t.primary, 0.45),
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Label below -- uppercase, neon accent color
        add_textbox(
            slide, label.upper(),
            cx - col_w // 2 + 40000, gauge_cy + outer_r + 120000,
            col_w - 80000, 300000,
            font_size=10, bold=True, color=accent,
            alignment=PP_ALIGN.CENTER,
        )

        # Thin neon line under label
        line_w = 500000
        line_y = gauge_cy + outer_r + 460000
        add_line(slide, cx - line_w // 2, line_y, cx + line_w // 2, line_y,
                 color=accent, width=1.0)


# ---------------------------------------------------------------------------
# Variant 6: SPLIT (split-gauge) -- Gauges left, summary stats right
# ---------------------------------------------------------------------------

def _split_gauge(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("gauge_title", "Performance Gauges"), theme=t)
    gauges = _get_gauges(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    center_bg = t.primary if s.dark_mode else "#FFFFFF"

    mid_x = SLIDE_W // 2 - 60000

    # Vertical divider
    add_line(slide, mid_x, content_top + 100000, mid_x, FOOTER_TOP - 200000,
             color=t.accent, width=2)

    # Left side: 2x2 grid of gauges
    n = min(len(gauges), 4)
    left_w = mid_x - MARGIN - 100000
    gauge_cols = 2
    gauge_rows = (n + 1) // 2
    cell_w = left_w // gauge_cols
    avail_h = FOOTER_TOP - content_top - 500000
    cell_h = avail_h // gauge_rows
    outer_r = min(cell_w // 2 - 80000, cell_h // 2 - 120000, 500000)
    inner_r = int(outer_r * 0.62)

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        col = i % gauge_cols
        row = i // gauge_cols
        cx = MARGIN + col * cell_w + cell_w // 2
        cy = content_top + 200000 + row * cell_h + cell_h // 2 - 60000
        track_c = tint(accent, 0.8) if not s.dark_mode else tint(t.primary, 0.18)

        add_circle(slide, cx, cy, outer_r, fill_color=track_c)
        add_circle(slide, cx, cy, outer_r, fill_color=accent)
        add_circle(slide, cx, cy, inner_r, fill_color=center_bg)

        # Percentage in center
        add_textbox(
            slide, f"{int(progress * 100)}%",
            cx - inner_r, cy - int(inner_r * 0.3),
            inner_r * 2, int(inner_r * 0.6),
            font_size=16, bold=True, color=text_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Label below
        add_textbox(
            slide, label,
            cx - cell_w // 2 + 20000, cy + outer_r + 40000,
            cell_w - 40000, 250000,
            font_size=9, bold=True, color=sub_color,
            alignment=PP_ALIGN.CENTER,
        )

    # Right side: summary stat cards
    right_left = mid_x + 120000
    right_w = SLIDE_W - MARGIN - right_left
    stat_h = 800000
    stat_gap = 200000
    stat_top = content_top + 200000

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        sy = stat_top + i * (stat_h + stat_gap)

        # Accent left bar
        add_rect(slide, right_left, sy, 60000, stat_h, fill_color=accent)

        # Label
        add_textbox(
            slide, label,
            right_left + 120000, sy + 60000, right_w - 180000, 250000,
            font_size=11, bold=True, color=sub_color,
            alignment=PP_ALIGN.LEFT,
        )

        # Value
        add_textbox(
            slide, value_display,
            right_left + 120000, sy + 320000, right_w - 180000, 350000,
            font_size=22, bold=True, color=text_color,
            alignment=PP_ALIGN.LEFT,
        )

        # Progress bar
        bar_left = right_left + 120000
        bar_top = sy + stat_h - 120000
        bar_w = right_w - 180000
        add_rect(slide, bar_left, bar_top, bar_w, 50000,
                 fill_color=tint(accent, 0.8) if not s.dark_mode else tint(t.primary, 0.15),
                 corner_radius=30000)
        fill_w = int(bar_w * min(progress, 1.0))
        if fill_w > 0:
            add_rect(slide, bar_left, bar_top, fill_w, 50000,
                     fill_color=accent, corner_radius=30000)


# ---------------------------------------------------------------------------
# Variant 7: GEO (angular-arc) -- Angular/geometric styled gauges
# ---------------------------------------------------------------------------

def _angular_arc(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("gauge_title", "Performance Gauges"), theme=t)
    gauges = _get_gauges(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if s.dark_mode else t.secondary
    center_bg = t.primary if s.dark_mode else "#FFFFFF"

    n = min(len(gauges), 4)
    total_w = SLIDE_W - 2 * MARGIN
    col_w = total_w // n
    gauge_cy = content_top + 1700000
    outer_r = min(col_w // 2 - 100000, 700000)
    inner_r = int(outer_r * 0.65)

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        cx = MARGIN + i * col_w + col_w // 2
        track_c = tint(accent, 0.75) if not s.dark_mode else tint(t.primary, 0.15)

        # Geometric border square behind gauge
        sq_size = outer_r + 60000
        add_rect(
            slide, cx - sq_size, gauge_cy - sq_size,
            sq_size * 2, sq_size * 2,
            line_color=accent, line_width=1.5,
        )

        # Gauge circles
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=track_c)
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=accent)
        add_circle(slide, cx, gauge_cy, inner_r, fill_color=center_bg)

        # Corner accents on the square
        corner_sz = 80000
        for dx, dy in [(-1, -1), (1, -1), (-1, 1), (1, 1)]:
            corner_x = cx + dx * sq_size
            corner_y = gauge_cy + dy * sq_size
            add_rect(
                slide, corner_x - corner_sz // 2, corner_y - corner_sz // 2,
                corner_sz, corner_sz, fill_color=accent,
            )

        # Percentage in center
        add_textbox(
            slide, f"{int(progress * 100)}%",
            cx - inner_r, gauge_cy - int(inner_r * 0.3),
            inner_r * 2, int(inner_r * 0.6),
            font_size=20, bold=True, color=text_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Label UPPERCASE below square
        add_textbox(
            slide, label.upper(),
            cx - col_w // 2 + 40000, gauge_cy + sq_size + 80000,
            col_w - 80000, 280000,
            font_size=10, bold=True, color=sub_color,
            alignment=PP_ALIGN.CENTER,
        )

        # Value below label
        add_textbox(
            slide, value_display,
            cx - col_w // 2 + 40000, gauge_cy + sq_size + 380000,
            col_w - 80000, 250000,
            font_size=10, color=accent, alignment=PP_ALIGN.CENTER,
        )


# ---------------------------------------------------------------------------
# Variant 8: EDITORIAL (editorial-arc) -- Elegant thin arcs, serif feel
# ---------------------------------------------------------------------------

def _editorial_arc(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("gauge_title", "Performance Gauges"), theme=t)
    gauges = _get_gauges(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    center_bg = t.primary if s.dark_mode else t.light_bg

    # Two-column layout, 2 gauges per column
    n = min(len(gauges), 4)
    col_w = (SLIDE_W - 2 * m - 400000) // 2
    row_h = (FOOTER_TOP - content_top - 400000) // 2
    outer_r = min(row_h // 2 - 200000, col_w // 4 - 50000, 450000)
    inner_r = int(outer_r * 0.78)  # thin ring

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        col = i % 2
        row = i // 2
        cell_left = m + col * (col_w + 400000)
        cell_top = content_top + 150000 + row * row_h

        # Thin top rule
        add_line(slide, cell_left, cell_top, cell_left + col_w, cell_top,
                 color=tint(t.secondary, 0.7), width=0.5)
        # Short accent mark
        add_line(slide, cell_left, cell_top, cell_left + 400000, cell_top,
                 color=t.accent, width=1.5)

        # Gauge on left of cell
        cx = cell_left + outer_r + 80000
        cy = cell_top + row_h // 2
        track_c = tint(t.secondary, 0.85) if not s.dark_mode else tint(t.primary, 0.12)

        add_circle(slide, cx, cy, outer_r, fill_color=track_c)
        add_circle(slide, cx, cy, outer_r, fill_color=accent)
        add_circle(slide, cx, cy, inner_r, fill_color=center_bg)

        # Percentage in center
        add_textbox(
            slide, f"{int(progress * 100)}%",
            cx - inner_r, cy - int(inner_r * 0.25),
            inner_r * 2, int(inner_r * 0.5),
            font_size=15, bold=True, color=text_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Label + value to the right of gauge
        text_left = cx + outer_r + 120000
        text_w = cell_left + col_w - text_left
        add_textbox(
            slide, label,
            text_left, cy - 300000, text_w, 250000,
            font_size=11, color=sub_color, alignment=PP_ALIGN.LEFT,
        )
        add_textbox(
            slide, value_display,
            text_left, cy - 50000, text_w, 350000,
            font_size=22, bold=True, color=text_color, alignment=PP_ALIGN.LEFT,
        )
        # Thin rule under value
        add_line(slide, text_left, cy + 300000, text_left + text_w, cy + 300000,
                 color=tint(t.secondary, 0.75), width=0.5)


# ---------------------------------------------------------------------------
# Variant 9: GRADIENT (gradient-arc) -- Gradient-colored arcs
# ---------------------------------------------------------------------------

def _gradient_arc(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("gauge_title", "Performance Gauges"), theme=t)
    gauges = _get_gauges(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary
    center_bg = t.primary if s.dark_mode else "#FFFFFF"

    n = min(len(gauges), 4)
    total_w = SLIDE_W - 2 * MARGIN
    col_w = total_w // n
    gauge_cy = content_top + 1700000
    outer_r = min(col_w // 2 - 100000, 720000)
    inner_r = int(outer_r * 0.62)

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        accent2 = p[(i + 1) % len(p)]  # secondary color for gradient feel
        cx = MARGIN + i * col_w + col_w // 2
        track_c = tint(accent, 0.88) if not s.dark_mode else tint(t.primary, 0.1)

        # Soft halo
        halo_r = outer_r + 40000
        add_circle(slide, cx, gauge_cy, halo_r, fill_color=tint(accent, 0.92))

        # Track
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=track_c)

        # Main ring (primary accent)
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=accent)

        # Inner secondary ring for gradient hint
        mid_r = int(outer_r * 0.82)
        add_circle(slide, cx, gauge_cy, mid_r, fill_color=accent2)

        # Center cutout
        add_circle(slide, cx, gauge_cy, inner_r, fill_color=center_bg)

        # Value text
        add_textbox(
            slide, f"{int(progress * 100)}%",
            cx - inner_r, gauge_cy - int(inner_r * 0.3),
            inner_r * 2, int(inner_r * 0.45),
            font_size=20, bold=True, color=text_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Small value display text
        add_textbox(
            slide, value_display,
            cx - inner_r, gauge_cy + int(inner_r * 0.1),
            inner_r * 2, int(inner_r * 0.3),
            font_size=9, color=sub_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Label below
        add_textbox(
            slide, label,
            cx - col_w // 2 + 40000, gauge_cy + outer_r + 120000,
            col_w - 80000, 300000,
            font_size=11, bold=True, color=sub_color,
            alignment=PP_ALIGN.CENTER,
        )


# ---------------------------------------------------------------------------
# Variant 10: RETRO (retro-arc) -- Vintage styled gauges with borders
# ---------------------------------------------------------------------------

def _retro_arc(slide, t, c):
    add_background(slide, t.light_bg)
    content_top = add_slide_title(slide, c.get("gauge_title", "Performance Gauges"), theme=t)
    gauges = _get_gauges(c)
    p = t.palette

    n = min(len(gauges), 4)
    total_w = SLIDE_W - 2 * MARGIN
    col_w = total_w // n
    gauge_cy = content_top + 1800000
    outer_r = min(col_w // 2 - 120000, 680000)
    inner_r = int(outer_r * 0.60)

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        cx = MARGIN + i * col_w + col_w // 2

        # Double-ring border (retro badge look)
        border_r = outer_r + 50000
        add_circle(slide, cx, gauge_cy, border_r,
                   fill_color=t.light_bg, line_color=accent, line_width=2.0)
        add_circle(slide, cx, gauge_cy, outer_r,
                   fill_color=t.light_bg, line_color=accent, line_width=0.75)

        # Fill ring
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=accent)

        # Center cutout with border
        add_circle(slide, cx, gauge_cy, inner_r,
                   fill_color="#FFFFFF", line_color=accent, line_width=0.75)

        # Decorative dots at top and bottom
        add_circle(slide, cx, gauge_cy - border_r - 30000, 25000, fill_color=accent)
        add_circle(slide, cx, gauge_cy + border_r + 30000, 25000, fill_color=accent)

        # Percentage in center with dashes
        add_textbox(
            slide, f"\u2014 {int(progress * 100)}% \u2014",
            cx - inner_r, gauge_cy - int(inner_r * 0.3),
            inner_r * 2, int(inner_r * 0.5),
            font_size=16, bold=True, color=t.primary,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Value below percentage
        add_textbox(
            slide, value_display,
            cx - inner_r, gauge_cy + int(inner_r * 0.15),
            inner_r * 2, int(inner_r * 0.35),
            font_size=9, color=t.secondary,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Label below gauge with retro decorative dashes
        add_textbox(
            slide, f"\u2022 {label} \u2022",
            cx - col_w // 2 + 40000, gauge_cy + border_r + 80000,
            col_w - 80000, 300000,
            font_size=11, bold=True, color=accent,
            alignment=PP_ALIGN.CENTER,
        )


# ---------------------------------------------------------------------------
# Variant 11: MAGAZINE (creative-arc) -- Creative oversized, varying sizes
# ---------------------------------------------------------------------------

def _creative_arc(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("gauge_title", "Performance Gauges"), theme=t)
    gauges = _get_gauges(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    center_bg = t.primary if s.dark_mode else "#FFFFFF"

    n = min(len(gauges), 4)

    # Creative layout: varying sizes, staggered positions
    # First gauge is hero-sized, rest are smaller
    avail_w = SLIDE_W - 2 * MARGIN
    avail_h = FOOTER_TOP - content_top - 200000

    # Hero gauge (first) -- large, left-center
    if n >= 1:
        label, value_display, progress = gauges[0]
        accent = p[0 % len(p)]
        hero_r = min(avail_h // 2 - 150000, int(avail_w * 0.22))
        hero_inner = int(hero_r * 0.58)
        hx = MARGIN + int(avail_w * 0.25)
        hy = content_top + avail_h // 2
        track_c = tint(accent, 0.8) if not s.dark_mode else tint(t.primary, 0.15)

        add_circle(slide, hx, hy, hero_r, fill_color=track_c)
        add_circle(slide, hx, hy, hero_r, fill_color=accent)
        add_circle(slide, hx, hy, hero_inner, fill_color=center_bg)

        add_textbox(
            slide, f"{int(progress * 100)}%",
            hx - hero_inner, hy - int(hero_inner * 0.35),
            hero_inner * 2, int(hero_inner * 0.45),
            font_size=28, bold=True, color=text_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide, value_display,
            hx - hero_inner, hy + int(hero_inner * 0.05),
            hero_inner * 2, int(hero_inner * 0.35),
            font_size=11, color=sub_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide, label,
            hx - hero_r, hy + hero_r + 60000,
            hero_r * 2, 300000,
            font_size=13, bold=True, color=accent,
            alignment=PP_ALIGN.CENTER,
        )

    # Remaining gauges: smaller, stacked on right
    right_left = MARGIN + int(avail_w * 0.55)
    right_w = SLIDE_W - MARGIN - right_left
    small_r = min(right_w // 4, avail_h // (max(n - 1, 1) * 2) - 80000, 380000)
    small_inner = int(small_r * 0.60)

    if n > 1:
        small_n = n - 1
        slot_h = avail_h // small_n

        for j in range(small_n):
            idx = j + 1
            label, value_display, progress = gauges[idx]
            accent = p[idx % len(p)]
            cx = right_left + right_w // 2
            cy = content_top + 100000 + j * slot_h + slot_h // 2
            track_c = tint(accent, 0.8) if not s.dark_mode else tint(t.primary, 0.15)

            add_circle(slide, cx, cy, small_r, fill_color=track_c)
            add_circle(slide, cx, cy, small_r, fill_color=accent)
            add_circle(slide, cx, cy, small_inner, fill_color=center_bg)

            add_textbox(
                slide, f"{int(progress * 100)}%",
                cx - small_inner, cy - int(small_inner * 0.3),
                small_inner * 2, int(small_inner * 0.6),
                font_size=16, bold=True, color=text_color,
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
            )

            # Label to the left of gauge
            add_textbox(
                slide, label,
                right_left - 100000, cy - 150000,
                cx - right_left - small_r - 20000, 300000,
                font_size=10, bold=True, color=sub_color,
                alignment=PP_ALIGN.RIGHT, vertical_anchor=MSO_ANCHOR.MIDDLE,
            )


# ---------------------------------------------------------------------------
# Variant 12: SCHOLARLY (scholarly-arc) -- Figure-captioned, thin clean arcs
# ---------------------------------------------------------------------------

def _scholarly_arc(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("gauge_title", "Performance Gauges"), theme=t)
    gauges = _get_gauges(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    center_bg = t.primary if s.dark_mode else t.light_bg

    n = min(len(gauges), 4)
    total_w = SLIDE_W - 2 * m
    col_w = total_w // n
    gauge_cy = content_top + 1500000
    outer_r = min(col_w // 2 - 160000, 550000)
    inner_r = int(outer_r * 0.80)  # very thin ring

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        cx = m + i * col_w + col_w // 2
        track_c = tint(t.secondary, 0.88) if not s.dark_mode else tint(t.primary, 0.1)

        # Thin scholarly ring
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=track_c)
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=accent)
        add_circle(slide, cx, gauge_cy, inner_r, fill_color=center_bg)

        # Percentage in center
        add_textbox(
            slide, f"{int(progress * 100)}%",
            cx - inner_r, gauge_cy - int(inner_r * 0.2),
            inner_r * 2, int(inner_r * 0.4),
            font_size=14, bold=True, color=text_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Figure caption style: "Fig. N" + label + value
        fig_y = gauge_cy + outer_r + 100000
        fig_label = f"Fig. {i + 1}. {label}"
        add_textbox(
            slide, fig_label,
            cx - col_w // 2 + 60000, fig_y,
            col_w - 120000, 250000,
            font_size=9, italic=True, color=sub_color,
            alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, value_display,
            cx - col_w // 2 + 60000, fig_y + 250000,
            col_w - 120000, 250000,
            font_size=11, bold=True, color=text_color,
            alignment=PP_ALIGN.CENTER,
        )

        # Thin rule below caption
        rule_y = fig_y + 540000
        rule_w = col_w - 250000
        add_line(slide, cx - rule_w // 2, rule_y, cx + rule_w // 2, rule_y,
                 color=tint(t.secondary, 0.8), width=0.5)

    # Bottom attribution line
    add_textbox(
        slide, "Source: Internal performance metrics",
        m, FOOTER_TOP - 350000, total_w, 200000,
        font_size=8, italic=True, color=sub_color, alignment=PP_ALIGN.LEFT,
    )


# ---------------------------------------------------------------------------
# Variant 13: LABORATORY (laboratory-arc) -- Dark bg, data-coded, grid overlay
# ---------------------------------------------------------------------------

def _laboratory_arc(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("gauge_title", "Performance Gauges"), theme=t)
    gauges = _get_gauges(c)
    p = t.palette

    n = min(len(gauges), 4)
    total_w = SLIDE_W - 2 * MARGIN
    col_w = total_w // n
    gauge_cy = content_top + 1700000
    outer_r = min(col_w // 2 - 100000, 680000)
    inner_r = int(outer_r * 0.68)
    center_bg = t.primary

    # Grid overlay lines (laboratory feel)
    grid_color = tint(t.primary, 0.08)
    grid_spacing = 400000
    for gx in range(MARGIN, SLIDE_W - MARGIN, grid_spacing):
        add_line(slide, gx, content_top, gx, FOOTER_TOP - 100000,
                 color=grid_color, width=0.25)
    for gy in range(content_top, FOOTER_TOP, grid_spacing):
        add_line(slide, MARGIN, gy, SLIDE_W - MARGIN, gy,
                 color=grid_color, width=0.25)

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        cx = MARGIN + i * col_w + col_w // 2
        track_c = tint(t.primary, 0.12)

        # Data ID label top-left of gauge area
        id_label = f"METRIC-{i + 1:02d}"
        add_textbox(
            slide, id_label,
            cx - col_w // 2 + 40000, content_top + 80000,
            col_w - 80000, 200000,
            font_size=7, bold=True, color=tint(t.primary, 0.3),
            alignment=PP_ALIGN.LEFT,
        )

        # Ring gauge
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=track_c)
        add_circle(slide, cx, gauge_cy, outer_r, fill_color=accent)
        add_circle(slide, cx, gauge_cy, inner_r, fill_color=center_bg)

        # Left-border accent on each column
        col_left = MARGIN + i * col_w
        add_rect(slide, col_left, content_top + 60000, 40000,
                 FOOTER_TOP - content_top - 200000, fill_color=accent)

        # Percentage (data readout style)
        add_textbox(
            slide, f"{int(progress * 100)}%",
            cx - inner_r, gauge_cy - int(inner_r * 0.35),
            inner_r * 2, int(inner_r * 0.5),
            font_size=20, bold=True, color=accent,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Value display (smaller readout)
        add_textbox(
            slide, value_display,
            cx - inner_r, gauge_cy + int(inner_r * 0.1),
            inner_r * 2, int(inner_r * 0.3),
            font_size=9, color=tint(t.primary, 0.4),
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Label below gauge (data label style)
        add_textbox(
            slide, label.upper(),
            cx - col_w // 2 + 40000, gauge_cy + outer_r + 100000,
            col_w - 80000, 250000,
            font_size=9, bold=True, color=tint(t.primary, 0.45),
            alignment=PP_ALIGN.CENTER,
        )

        # Status indicator bar
        status_y = gauge_cy + outer_r + 400000
        status_w = col_w - 200000
        add_rect(slide, cx - status_w // 2, status_y, status_w, 35000,
                 fill_color=tint(t.primary, 0.1))
        fill_w = int(status_w * min(progress, 1.0))
        if fill_w > 0:
            add_rect(slide, cx - status_w // 2, status_y, fill_w, 35000,
                     fill_color=accent)


# ---------------------------------------------------------------------------
# Variant 14: DASHBOARD (dashboard-arc) -- Dense compact gauges with tiles
# ---------------------------------------------------------------------------

def _dashboard_arc(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("gauge_title", "Performance Gauges"), theme=t)
    gauges = _get_gauges(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if s.dark_mode else t.secondary
    tile_bg = tint(t.primary, 0.1) if s.dark_mode else "#FFFFFF"
    center_bg = tile_bg

    n = min(len(gauges), 4)
    gap = int(150000 * s.gap_factor)
    tile_w = (SLIDE_W - 2 * m - (n - 1) * gap) // n
    tile_h = FOOTER_TOP - content_top - 300000
    tile_top = content_top + 100000

    for i, (label, value_display, progress) in enumerate(gauges[:n]):
        accent = p[i % len(p)]
        left = m + i * (tile_w + gap)

        # Compact tile
        border_c = tint(t.primary, 0.2) if s.dark_mode else tint(t.secondary, 0.8)
        add_rect(slide, left, tile_top, tile_w, tile_h,
                 fill_color=tile_bg, line_color=border_c, line_width=0.5,
                 corner_radius=s.card_radius)

        # Thin accent header bar at top of tile
        add_rect(slide, left, tile_top, tile_w, 45000, fill_color=accent)

        # Label at top
        add_textbox(
            slide, label,
            left + 50000, tile_top + 80000, tile_w - 100000, 280000,
            font_size=10, bold=True, color=sub_color,
            alignment=PP_ALIGN.CENTER,
        )

        # Compact gauge in center of tile
        cx = left + tile_w // 2
        gauge_area_top = tile_top + 400000
        gauge_area_h = tile_h - 1200000
        cy = gauge_area_top + gauge_area_h // 2
        outer_r = min(tile_w // 2 - 80000, gauge_area_h // 2 - 40000, 500000)
        inner_r = int(outer_r * 0.62)
        track_c = tint(accent, 0.8) if not s.dark_mode else tint(t.primary, 0.15)

        add_circle(slide, cx, cy, outer_r, fill_color=track_c)
        add_circle(slide, cx, cy, outer_r, fill_color=accent)
        add_circle(slide, cx, cy, inner_r, fill_color=center_bg)

        # Percentage in center
        add_textbox(
            slide, f"{int(progress * 100)}%",
            cx - inner_r, cy - int(inner_r * 0.3),
            inner_r * 2, int(inner_r * 0.6),
            font_size=16, bold=True, color=text_color,
            alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        # Value display below gauge
        val_y = cy + outer_r + 60000
        add_textbox(
            slide, value_display,
            left + 40000, val_y, tile_w - 80000, 280000,
            font_size=12, bold=True, color=text_color,
            alignment=PP_ALIGN.CENTER,
        )

        # Mini progress bar at bottom of tile
        bar_left = left + 60000
        bar_w = tile_w - 120000
        bar_y = tile_top + tile_h - 200000
        add_rect(slide, bar_left, bar_y, bar_w, 40000,
                 fill_color=track_c, corner_radius=25000)
        fill_w = int(bar_w * min(progress, 1.0))
        if fill_w > 0:
            add_rect(slide, bar_left, bar_y, fill_w, 40000,
                     fill_color=accent, corner_radius=25000)

        # Status text at very bottom
        add_textbox(
            slide, f"{int(progress * 100)}% complete",
            left + 40000, bar_y + 60000, tile_w - 80000, 200000,
            font_size=7, color=sub_color, alignment=PP_ALIGN.CENTER,
        )
