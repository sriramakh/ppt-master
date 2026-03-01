"""Image placement in picture placeholders."""

from __future__ import annotations

import io
from pathlib import Path

from pptx.util import Emu

from pptmaster.assets.image_handler import emu_to_pixels, resize_to_fit
from pptmaster.models import DesignProfile


def place_image(
    slide,
    image_path: str | Path,
    placeholder_idx: int | None = None,
    profile: DesignProfile | None = None,
) -> None:
    """Place an image into a slide's picture placeholder.

    If placeholder_idx is given, uses that specific placeholder.
    Otherwise, uses the first available picture placeholder.
    """
    target_ph = None

    if placeholder_idx is not None:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == placeholder_idx:
                target_ph = ph
                break

    if target_ph is None:
        # Find first picture placeholder
        for ph in slide.placeholders:
            if _is_picture_placeholder(ph):
                target_ph = ph
                break

    if target_ph is None:
        # No picture placeholder — add as free-standing image
        _add_freestanding_image(slide, image_path, profile)
        return

    # Insert image into placeholder
    target_ph.insert_picture(str(image_path))


def _is_picture_placeholder(ph) -> bool:
    """Check if a placeholder is a picture placeholder."""
    # Check via XML for accuracy
    from lxml import etree
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    nvPr = ph._element.find(f".//{{{NS_P}}}nvPr")
    if nvPr is not None:
        ph_elem = nvPr.find(f"{{{NS_P}}}ph")
        if ph_elem is not None and ph_elem.get("type") == "pic":
            return True

    # Also check python-pptx type
    try:
        ph_type = ph.placeholder_format.type
        if ph_type is not None and int(ph_type) == 18:  # PP_PLACEHOLDER.PICTURE
            return True
    except Exception:
        pass

    return False


def _add_freestanding_image(
    slide,
    image_path: str | Path,
    profile: DesignProfile | None,
) -> None:
    """Add an image as a free-standing shape (not in a placeholder)."""
    if profile:
        # Place in the right half of the slide
        left = Emu(int(profile.slide_size.width * 0.55))
        top = Emu(int(profile.slide_size.height * 0.2))
        width = Emu(int(profile.slide_size.width * 0.4))
        height = Emu(int(profile.slide_size.height * 0.6))
    else:
        left = Emu(5000000)
        top = Emu(2000000)
        width = Emu(8000000)
        height = Emu(6000000)

    slide.shapes.add_picture(str(image_path), left, top, width, height)
