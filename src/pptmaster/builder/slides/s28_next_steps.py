"""Slide 28 — Next Steps: Vertical timeline, 4 action items with owners."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, FONT_CARD_HEADING, CONTENT_TOP, FOOTER_TOP,
)
from pptmaster.builder.helpers import (
    add_textbox, add_circle, add_line, add_card, add_gold_accent_line,
    add_styled_card, add_slide_title, add_dark_bg, add_background, add_rect,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    p = t.palette
    s = t.ux_style

    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("next_steps_title", "Next Steps & Action Items"), theme=t)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    actions = c.get("next_steps", [
        ("Action 1", "Description", "Owner", "Date"),
        ("Action 2", "Description", "Owner", "Date"),
        ("Action 3", "Description", "Owner", "Date"),
        ("Action 4", "Description", "Owner", "Date"),
    ])

    m = int(MARGIN * s.margin_factor)
    timeline_x = m + 300000
    content_left = timeline_x + 500000
    content_w = 8500000
    start_y = content_y + 200000
    item_h = int(1050000 * s.gap_factor)
    marker_radius = 100000

    add_line(slide, timeline_x, start_y, timeline_x, start_y + (len(actions) - 1) * item_h,
             color=t.accent, width=2.5)

    for i, (title, desc, owner, date) in enumerate(actions[:4]):
        y = start_y + i * item_h
        accent = p[i % len(p)]

        add_circle(slide, timeline_x, y, marker_radius, fill_color=t.accent)
        add_styled_card(slide, content_left, y - 200000, content_w, 850000,
                        theme=t, accent_color=accent)
        add_textbox(slide, title, content_left + 200000, y - 160000, 5000000, 350000,
                    font_size=FONT_CARD_HEADING, bold=True, color=text_color)
        add_textbox(slide, desc, content_left + 200000, y + 150000, 5000000, 350000,
                    font_size=12, color=sub_color)
        add_textbox(slide, owner, content_left + content_w - 2500000, y - 160000, 2300000, 300000,
                    font_size=11, color=accent, bold=True, alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, date, content_left + content_w - 2500000, y + 150000, 2300000, 300000,
                    font_size=10, color=sub_color, alignment=PP_ALIGN.RIGHT)


# ── Scholarly variant: white bg, numbered recommendations, thin rules ─────

def _scholarly(slide, t):
    c = t.content
    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, c.get("next_steps_title", "Recommendations"), theme=t)

    actions = c.get("next_steps", [
        ("Action 1", "Description", "Owner", "Date"),
        ("Action 2", "Description", "Owner", "Date"),
        ("Action 3", "Description", "Owner", "Date"),
        ("Action 4", "Description", "Owner", "Date"),
    ])

    total_w = SLIDE_W - 2 * MARGIN
    start_y = content_y + 200000
    item_h = 950000

    for i, (title, desc, owner, date) in enumerate(actions[:4]):
        y = start_y + i * item_h

        # Thin rule above each item (skip first if at top)
        add_line(slide, MARGIN, y, MARGIN + total_w, y,
                 color=tint(t.secondary, 0.7), width=0.5)

        # Numbered label
        add_textbox(slide, f"{i + 1}.", MARGIN, y + 100000, 300000, 350000,
                    font_size=FONT_CARD_HEADING, bold=True, color=t.primary)

        # Title and description
        add_textbox(slide, title, MARGIN + 350000, y + 100000, total_w - 3000000, 350000,
                    font_size=FONT_CARD_HEADING, bold=True, color=t.primary)
        add_textbox(slide, desc, MARGIN + 350000, y + 420000, total_w - 3000000, 350000,
                    font_size=12, color=t.secondary)

        # Owner and date right-aligned
        add_textbox(slide, owner, MARGIN + total_w - 2500000, y + 100000, 2500000, 300000,
                    font_size=11, color=t.primary, alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, date, MARGIN + total_w - 2500000, y + 420000, 2500000, 300000,
                    font_size=10, color=t.secondary, alignment=PP_ALIGN.RIGHT)

    # Final thin rule
    final_y = start_y + len(actions[:4]) * item_h
    add_line(slide, MARGIN, final_y, MARGIN + total_w, final_y,
             color=tint(t.secondary, 0.7), width=0.5)


# ── Laboratory variant: dark bg, numbered dark cards with accent borders ──

def _laboratory(slide, t):
    c = t.content
    p = t.palette
    add_background(slide, t.primary)
    content_y = add_slide_title(slide, c.get("next_steps_title", "Action Items"), theme=t)

    actions = c.get("next_steps", [
        ("Action 1", "Description", "Owner", "Date"),
        ("Action 2", "Description", "Owner", "Date"),
        ("Action 3", "Description", "Owner", "Date"),
        ("Action 4", "Description", "Owner", "Date"),
    ])

    total_w = SLIDE_W - 2 * MARGIN
    start_y = content_y + 200000
    item_h = 1000000
    card_h = 850000

    for i, (title, desc, owner, date) in enumerate(actions[:4]):
        y = start_y + i * item_h
        accent = p[i % len(p)]

        # Dark card
        add_rect(slide, MARGIN, y, total_w, card_h,
                 fill_color=tint(t.primary, 0.08), corner_radius=50000)
        # Colored left accent border
        add_rect(slide, MARGIN, y, 40000, card_h, fill_color=accent)

        # Number badge
        add_textbox(slide, f"{i + 1}", MARGIN + 120000, y + 100000, 300000, 350000,
                    font_size=FONT_CARD_HEADING, bold=True, color=accent)

        # Title and description
        add_textbox(slide, title, MARGIN + 450000, y + 100000, total_w - 3200000, 350000,
                    font_size=FONT_CARD_HEADING, bold=True, color="#FFFFFF")
        add_textbox(slide, desc, MARGIN + 450000, y + 420000, total_w - 3200000, 350000,
                    font_size=12, color=tint(t.primary, 0.5))

        # Owner and date
        add_textbox(slide, owner, MARGIN + total_w - 2500000, y + 100000, 2300000, 300000,
                    font_size=11, color=accent, bold=True, alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, date, MARGIN + total_w - 2500000, y + 420000, 2300000, 300000,
                    font_size=10, color=tint(t.primary, 0.4), alignment=PP_ALIGN.RIGHT)


# ── Dashboard variant: header band, compact checklist cards ───────────────

def _dashboard(slide, t):
    c = t.content
    p = t.palette
    add_background(slide, "#FFFFFF")

    # Accent header band
    add_rect(slide, 0, 0, SLIDE_W, 80000, fill_color=t.accent)

    content_y = add_slide_title(slide, c.get("next_steps_title", "Next Steps & Action Items"), theme=t)

    actions = c.get("next_steps", [
        ("Action 1", "Description", "Owner", "Date"),
        ("Action 2", "Description", "Owner", "Date"),
        ("Action 3", "Description", "Owner", "Date"),
        ("Action 4", "Description", "Owner", "Date"),
    ])

    total_w = SLIDE_W - 2 * MARGIN
    start_y = content_y + 100000
    # Compact: fit more items
    available_h = FOOTER_TOP - start_y - 100000
    item_h = min(900000, available_h // max(len(actions[:5]), 1))
    card_h = item_h - 80000

    for i, (title, desc, owner, date) in enumerate(actions[:5]):
        y = start_y + i * item_h
        accent = p[i % len(p)]

        # Shadow
        add_rect(slide, MARGIN + 15000, y + 15000, total_w, card_h,
                 fill_color=tint(t.secondary, 0.88), corner_radius=35000)
        add_styled_card(slide, MARGIN, y, total_w, card_h, theme=t, accent_color=accent)

        # Checkbox-style marker
        add_textbox(slide, "\u2610", MARGIN + 100000, y + 80000, 300000, 300000,
                    font_size=18, color=accent)

        # Title
        add_textbox(slide, title, MARGIN + 400000, y + 80000, total_w - 3500000, 300000,
                    font_size=15, bold=True, color=t.primary)
        # Description
        add_textbox(slide, desc, MARGIN + 400000, y + 370000, total_w - 3500000, 300000,
                    font_size=11, color=t.secondary)

        # Owner and date compact
        add_textbox(slide, f"{owner} | {date}", MARGIN + total_w - 2800000, y + 80000,
                    2600000, 300000,
                    font_size=10, color=t.secondary, alignment=PP_ALIGN.RIGHT)
