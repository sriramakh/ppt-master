"""Chart generation — bar, line, pie, donut, area charts styled with template colors."""

from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Pt, Emu

from pptmaster.assets.color_utils import hex_to_rgb
from pptmaster.models import ChartSpec, DesignProfile


_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "donut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
}


def add_chart(
    slide,
    chart_spec: ChartSpec,
    profile: DesignProfile,
    left: int | None = None,
    top: int | None = None,
    width: int | None = None,
    height: int | None = None,
) -> None:
    """Add a chart to a slide with template-consistent styling.

    Args:
        slide: The pptx slide object.
        chart_spec: Chart specification with data.
        profile: Design profile for styling.
        left/top/width/height: Optional position in EMU. Defaults to centered.
    """
    # Default position: centered, taking ~70% of slide
    sw = profile.slide_size.width
    sh = profile.slide_size.height

    if left is None:
        left = int(sw * 0.1)
    if top is None:
        top = int(sh * 0.25)
    if width is None:
        width = int(sw * 0.8)
    if height is None:
        height = int(sh * 0.6)

    chart_type = _CHART_TYPE_MAP.get(chart_spec.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # Build chart data
    chart_data = CategoryChartData()
    chart_data.categories = chart_spec.categories

    for series in chart_spec.series:
        chart_data.add_series(series.get("name", ""), series.get("values", []))

    # Add chart to slide
    chart_frame = slide.shapes.add_chart(
        chart_type,
        Emu(left), Emu(top), Emu(width), Emu(height),
        chart_data,
    )
    chart = chart_frame.chart

    # Set chart title if provided
    if chart_spec.title:
        chart.has_title = True
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = chart_spec.title
        for para in chart.chart_title.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = profile.fonts.minor
                run.font.size = Pt(12)
                r, g, b = hex_to_rgb(profile.colors.dk2)
                run.font.color.rgb = RGBColor(r, g, b)
    else:
        chart.has_title = False

    # Style the chart
    _apply_chart_style(chart, chart_spec, profile)


def _apply_chart_style(chart, spec: ChartSpec, profile: DesignProfile) -> None:
    """Apply template-consistent styling to a chart."""
    accent_colors = profile.colors.accent_colors

    # Color the series / points with accent colors
    plot = chart.plots[0]
    is_pie = spec.chart_type in ("pie", "donut")

    if is_pie:
        # Pie/donut: one series, color each POINT (slice) individually
        series = plot.series[0]
        for idx, point in enumerate(series.points):
            color_hex = accent_colors[idx % len(accent_colors)]
            r, g, b = hex_to_rgb(color_hex)
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor(r, g, b)
    else:
        # Bar/line/etc: color each series
        for i, series in enumerate(plot.series):
            color_hex = accent_colors[i % len(accent_colors)]
            r, g, b = hex_to_rgb(color_hex)
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor(r, g, b)

            # For line charts, also set line color
            if hasattr(series.format, "line"):
                series.format.line.color.rgb = RGBColor(r, g, b)
                series.format.line.width = Pt(2.5)

    # Remove gridlines for clean look
    try:
        value_axis = chart.value_axis
        if value_axis.major_gridlines:
            value_axis.major_gridlines.format.line.fill.background()
        value_axis.has_minor_gridlines = False
        _style_axis_font(value_axis, profile)
    except (ValueError, AttributeError):
        pass  # Pie/donut charts have no value axis

    try:
        cat_axis = chart.category_axis
        _style_axis_font(cat_axis, profile)
    except (ValueError, AttributeError):
        pass  # Pie/donut charts have no category axis

    # Legend — show for multi-series charts and pie/donut (where legend = category names)
    show_legend = spec.show_legend and (len(spec.series) > 1 or is_pie)
    if show_legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = profile.fonts.minor
        chart.legend.font.size = Pt(10)
    else:
        chart.has_legend = False

    # Data labels — always show for pie/donut, otherwise respect spec
    show_labels = spec.show_data_labels or is_pie
    if show_labels:
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.name = profile.fonts.minor
        data_labels.font.size = Pt(10)
        r, g, b = hex_to_rgb(profile.colors.dk2)
        data_labels.font.color.rgb = RGBColor(r, g, b)
        if is_pie:
            data_labels.show_percentage = True
            data_labels.show_category_name = False


def _style_axis_font(axis, profile: DesignProfile) -> None:
    """Apply template font to axis labels."""
    try:
        tf = axis.tick_labels
        tf.font.name = profile.fonts.minor
        tf.font.size = Pt(10)
        r, g, b = hex_to_rgb(profile.colors.dk2)
        tf.font.color.rgb = RGBColor(r, g, b)
    except Exception:
        pass
