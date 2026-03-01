"""Orchestrator — builds a 40-slide corporate template from scratch."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from pptmaster.builder.design_system import SLIDE_W, SLIDE_H
from pptmaster.builder.helpers import add_themed_footer, add_themed_slide_number


def build_template(
    output_path: str | Path = "corporate_template_40.pptx",
    company_name: str = "[Company Name]",
    theme=None,
) -> Path:
    """Build a 40-slide Fortune 500 corporate template.

    Args:
        output_path: Where to save the PPTX file.
        company_name: Company name placeholder for footers and title slides.
        theme: Optional TemplateTheme for colors and content. Uses default if None.

    Returns:
        Path to the saved presentation.
    """
    if theme is None:
        from pptmaster.builder.themes import DEFAULT_THEME
        theme = DEFAULT_THEME

    # Use theme's company name if not explicitly overridden
    if company_name == "[Company Name]":
        company_name = theme.company_name

    output_path = Path(output_path)

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    blank_layout = prs.slide_layouts[6]

    # Import all slide builders
    from pptmaster.builder.slides.s01_cover import build as build_s01
    from pptmaster.builder.slides.s02_toc import build as build_s02
    from pptmaster.builder.slides.s03_section_divider import build as build_s03
    from pptmaster.builder.slides.s04_company_overview import build as build_s04
    from pptmaster.builder.slides.s05_our_values import build as build_s05
    from pptmaster.builder.slides.s06_team_leadership import build as build_s06
    from pptmaster.builder.slides.s07_key_facts import build as build_s07
    from pptmaster.builder.slides.s09_executive_summary import build as build_s09
    from pptmaster.builder.slides.s10_kpi_dashboard import build as build_s10
    from pptmaster.builder.slides.s11_process_linear import build as build_s11
    from pptmaster.builder.slides.s12_process_circular import build as build_s12
    from pptmaster.builder.slides.s13_roadmap_timeline import build as build_s13
    from pptmaster.builder.slides.s14_swot_matrix import build as build_s14
    from pptmaster.builder.slides.s16_bar_chart import build as build_s16
    from pptmaster.builder.slides.s17_line_chart import build as build_s17
    from pptmaster.builder.slides.s18_pie_chart import build as build_s18
    from pptmaster.builder.slides.s19_comparison import build as build_s19
    from pptmaster.builder.slides.s20_data_table import build as build_s20
    from pptmaster.builder.slides.s22_text_image import build as build_s22
    from pptmaster.builder.slides.s23_two_column import build as build_s23
    from pptmaster.builder.slides.s24_three_column import build as build_s24
    from pptmaster.builder.slides.s25_highlight_quote import build as build_s25
    from pptmaster.builder.slides.s26_infographic_dashboard import build as build_s26
    from pptmaster.builder.slides.s28_next_steps import build as build_s28
    from pptmaster.builder.slides.s29_call_to_action import build as build_s29
    from pptmaster.builder.slides.s30_thank_you import build as build_s30
    # New slide builders (s31-s40)
    from pptmaster.builder.slides.s31_funnel_diagram import build as build_s31
    from pptmaster.builder.slides.s32_pyramid_hierarchy import build as build_s32
    from pptmaster.builder.slides.s33_venn_diagram import build as build_s33
    from pptmaster.builder.slides.s34_hub_spoke import build as build_s34
    from pptmaster.builder.slides.s35_milestone_roadmap import build as build_s35
    from pptmaster.builder.slides.s36_kanban_board import build as build_s36
    from pptmaster.builder.slides.s37_matrix_quadrant import build as build_s37
    from pptmaster.builder.slides.s38_gauge_dashboard import build as build_s38
    from pptmaster.builder.slides.s39_icon_grid import build as build_s39
    from pptmaster.builder.slides.s40_risk_matrix import build as build_s40

    # Get section divider content from theme (now 7 sections)
    dividers = theme.content.get("section_dividers", [
        ("01", "About Us", "Who we are and what we stand for"),
        ("02", "Strategy", "Our path to sustainable growth"),
        ("03", "Process & Planning", "How we execute and deliver"),
        ("04", "Data & Insights", "Metrics that drive decisions"),
        ("05", "Planning", "Roadmap, milestones, and risk management"),
        ("06", "Deliverables", "Tangible outcomes and milestones"),
        ("07", "Next Steps", "Actions and accountability"),
    ])

    # Ensure we have at least 7 dividers (pad if theme only has 5)
    while len(dividers) < 7:
        idx = len(dividers) + 1
        dividers.append((f"{idx:02d}", f"Section {idx}", ""))

    # 40-slide sequence: (builder_fn, kwargs, is_dark_bg)
    slide_builders = [
        # 1. Cover
        (build_s01, {"company_name": company_name}, True),
        # 2. TOC
        (build_s02, {"company_name": company_name}, False),
        # 3. Section: About Us
        (build_s03, {"section_number": dividers[0][0], "section_title": dividers[0][1], "subtitle": dividers[0][2]}, True),
        # 4. Company Overview
        (build_s04, {"company_name": company_name}, False),
        # 5. Our Values
        (build_s05, {}, False),
        # 6. Team Leadership
        (build_s06, {}, False),
        # 7. Key Facts
        (build_s07, {}, False),
        # 8. Section: Strategy
        (build_s03, {"section_number": dividers[1][0], "section_title": dividers[1][1], "subtitle": dividers[1][2]}, True),
        # 9. Executive Summary
        (build_s09, {}, False),
        # 10. KPI Dashboard
        (build_s10, {}, False),
        # 11. SWOT Matrix
        (build_s14, {}, False),
        # 12. Matrix Quadrant (NEW)
        (build_s37, {}, False),
        # 13. Venn Diagram (NEW)
        (build_s33, {}, False),
        # 14. Section: Process & Planning (NEW SECTION)
        (build_s03, {"section_number": dividers[2][0], "section_title": dividers[2][1], "subtitle": dividers[2][2]}, True),
        # 15. Process Linear
        (build_s11, {}, False),
        # 16. Process Circular
        (build_s12, {}, False),
        # 17. Roadmap Timeline
        (build_s13, {}, False),
        # 18. Funnel Diagram (NEW)
        (build_s31, {}, False),
        # 19. Pyramid Hierarchy (NEW)
        (build_s32, {}, False),
        # 20. Hub & Spoke (NEW)
        (build_s34, {}, False),
        # 21. Section: Data & Insights
        (build_s03, {"section_number": dividers[3][0], "section_title": dividers[3][1], "subtitle": dividers[3][2]}, True),
        # 22. Bar Chart
        (build_s16, {}, False),
        # 23. Line Chart
        (build_s17, {}, False),
        # 24. Pie Chart
        (build_s18, {}, False),
        # 25. Comparison
        (build_s19, {}, False),
        # 26. Data Table
        (build_s20, {}, False),
        # 27. Gauge Dashboard (NEW)
        (build_s38, {}, False),
        # 28. Section: Planning (NEW SECTION)
        (build_s03, {"section_number": dividers[4][0], "section_title": dividers[4][1], "subtitle": dividers[4][2]}, True),
        # 29. Milestone Roadmap (NEW)
        (build_s35, {}, False),
        # 30. Kanban Board (NEW)
        (build_s36, {}, False),
        # 31. Risk Matrix (NEW)
        (build_s40, {}, False),
        # 32. Section: Deliverables
        (build_s03, {"section_number": dividers[5][0], "section_title": dividers[5][1], "subtitle": dividers[5][2]}, True),
        # 33. Two Column
        (build_s23, {}, False),
        # 34. Three Column
        (build_s24, {}, False),
        # 35. Highlight Quote
        (build_s25, {}, False),
        # 36. Infographic Dashboard
        (build_s26, {}, False),
        # 37. Icon Grid (NEW)
        (build_s39, {}, False),
        # 38. Next Steps
        (build_s28, {}, False),
        # 39. Call to Action
        (build_s29, {"company_name": company_name}, True),
        # 40. Thank You
        (build_s30, {"company_name": company_name}, False),
    ]

    # Check if this theme uses all-dark mode
    style = theme.ux_style

    for i, (builder_fn, kwargs, is_dark) in enumerate(slide_builders, start=1):
        slide = prs.slides.add_slide(blank_layout)
        builder_fn(slide, theme=theme, **kwargs)

        # In dark_mode themes, all slides get dark footers
        effective_dark = is_dark or style.dark_mode

        if i > 1:
            add_themed_footer(slide, company_name=company_name,
                              dark_bg=effective_dark, theme=theme)
            add_themed_slide_number(slide, i, dark_bg=effective_dark, theme=theme)

    prs.save(str(output_path))
    return output_path


def build_all_templates(output_dir: str | Path = ".") -> list[Path]:
    """Build all 14 templates (default + 13 industry themes).

    Returns list of paths to generated PPTX files.
    """
    from pptmaster.builder.themes import DEFAULT_THEME, ALL_THEMES

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    results = []

    all_themes = [DEFAULT_THEME] + ALL_THEMES

    for theme in all_themes:
        filename = f"template_{theme.key}.pptx"
        path = output_dir / filename
        build_template(path, theme.company_name, theme=theme)
        results.append(path)

    return results
