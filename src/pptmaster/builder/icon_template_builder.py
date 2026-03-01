"""Build a standalone icon reference PPTX — one slide per category."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from pptmaster.builder.design_system import SLIDE_W, SLIDE_H, MARGIN
from pptmaster.builder.helpers import add_background, add_textbox, add_rect
from pptmaster.assets.icon_generator import ICON_CATALOG, CATEGORY_COLORS
from pptmaster.assets.raster_icon_manager import get_icon_manager


def build_icon_template(
    output_path: str | Path = "icon_toolkit.pptx",
    icon_dir: str | Path | None = None,
) -> Path:
    """Build a reference PPTX showing all generated icons.

    One slide per category, icons in a 4x3 grid.

    Args:
        output_path: Where to save the PPTX.
        icon_dir: Directory with generated icons. Uses default if None.

    Returns:
        Path to saved PPTX.
    """
    output_path = Path(output_path)
    manager = get_icon_manager(icon_dir) if icon_dir else get_icon_manager()

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    blank = prs.slide_layouts[6]

    # Title slide
    title_slide = prs.slides.add_slide(blank)
    add_background(title_slide, "#1B2A4A")
    add_textbox(title_slide, "Icon Toolkit",
                MARGIN, SLIDE_H // 2 - 600000, SLIDE_W - 2 * MARGIN, 800000,
                font_size=44, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER)
    add_textbox(title_slide, f"{sum(len(v) for v in ICON_CATALOG.values())} Professional Icons  |  {len(ICON_CATALOG)} Categories",
                MARGIN, SLIDE_H // 2 + 300000, SLIDE_W - 2 * MARGIN, 400000,
                font_size=18, color="#C8A951", alignment=PP_ALIGN.CENTER)

    # One slide per category
    for category, icon_names in ICON_CATALOG.items():
        slide = prs.slides.add_slide(blank)
        add_background(slide, "#FFFFFF")

        cat_color = CATEGORY_COLORS.get(category, "#1B2A4A")

        # Category header bar
        add_rect(slide, 0, 0, SLIDE_W, 800000, fill_color=cat_color)
        add_textbox(slide, category.replace("-", " ").title(),
                    MARGIN, 150000, SLIDE_W - 2 * MARGIN, 500000,
                    font_size=28, bold=True, color="#FFFFFF")

        # Grid: 4 columns x 3 rows
        cols = 4
        rows = 3
        grid_top = 1000000
        grid_left = MARGIN
        grid_w = SLIDE_W - 2 * MARGIN
        grid_h = SLIDE_H - grid_top - 400000

        cell_w = grid_w // cols
        cell_h = grid_h // rows
        icon_size = min(cell_w, cell_h) * 55 // 100  # 55% of cell

        for idx, name in enumerate(icon_names[:cols * rows]):
            row, col = divmod(idx, cols)
            cx = grid_left + col * cell_w
            cy = grid_top + row * cell_h

            # Try to place actual icon image
            icon_path = manager.get_icon_path(name)
            if icon_path and icon_path.exists():
                ix = cx + (cell_w - icon_size) // 2
                iy = cy + 50000
                try:
                    slide.shapes.add_picture(
                        str(icon_path),
                        Emu(ix), Emu(iy),
                        Emu(icon_size), Emu(icon_size),
                    )
                except Exception:
                    # Fallback: gray placeholder rect
                    add_rect(slide, ix, iy, icon_size, icon_size,
                             fill_color="#E0E0E0")
            else:
                # Gray placeholder
                ix = cx + (cell_w - icon_size) // 2
                iy = cy + 50000
                add_rect(slide, ix, iy, icon_size, icon_size,
                         fill_color="#E0E0E0")

            # Label below icon
            label_y = cy + icon_size + 100000
            add_textbox(slide, name.replace("-", " "),
                        cx, label_y, cell_w, 300000,
                        font_size=9, color="#333333", alignment=PP_ALIGN.CENTER)

    prs.save(str(output_path))
    print(f"Icon template saved: {output_path} ({len(ICON_CATALOG) + 1} slides)")
    return output_path
