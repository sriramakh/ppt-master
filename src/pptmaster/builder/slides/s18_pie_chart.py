"""Slide 18 — Pie/Donut Chart: Donut chart left + manual legend right."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from pptmaster.builder.design_system import SLIDE_W, SLIDE_H, MARGIN, FONT_CAPTION, CONTENT_TOP
from pptmaster.builder.helpers import add_textbox, add_rect, add_line, add_gold_accent_line, add_styled_card, add_slide_title, add_dark_bg, add_background
from pptmaster.builder.slides._chart_helpers import profile_from_theme
from pptmaster.models import ChartSpec
from pptmaster.composer.chart_renderer import add_chart
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    p = t.palette

    s = t.ux_style

    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)

    content_top = add_slide_title(slide, c.get("pie_title", "Revenue Distribution"), theme=t)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    spec = ChartSpec(
        chart_type="donut", title="",
        categories=c.get("pie_categories", ["Enterprise", "Mid-Market", "SMB", "Gov", "Partners"]),
        series=[{"name": "Share", "values": c.get("pie_values", [42, 28, 15, 10, 5])}],
        show_legend=False, show_data_labels=True,
    )

    chart_w = int((SLIDE_W - 2 * MARGIN) * 0.60)
    add_chart(slide, spec, profile_from_theme(t),
              left=MARGIN, top=content_top + 100000, width=chart_w, height=4500000)

    # Manual legend
    legend_left = MARGIN + chart_w + 300000
    legend_w = SLIDE_W - legend_left - MARGIN
    legend_top = content_top + 500000

    legend_colors = p
    legend_labels = c.get("pie_legend", ["Cat 1", "Cat 2", "Cat 3", "Cat 4", "Cat 5"])

    for i, (label, color) in enumerate(zip(legend_labels, legend_colors)):
        y = legend_top + i * 500000
        add_rect(slide, legend_left, y, 200000, 200000, fill_color=color, corner_radius=30000)
        add_textbox(slide, label, legend_left + 300000, y - 20000, legend_w - 400000, 250000,
                    font_size=13, color=text_color)


# ── Variant: scholarly ────────────────────────────────────────────────

def _scholarly(slide, t):
    """White bg, 'Figure 4:' caption, chart centered with rules above/below."""
    c = t.content
    p = t.palette
    m = MARGIN
    title = c.get("pie_title", "Revenue Distribution")

    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, title, theme=t)

    # Figure caption
    caption = f"Figure 4: {title}"
    add_textbox(slide, caption, m, content_y + 20000,
                SLIDE_W - 2 * m, 300000,
                font_size=FONT_CAPTION, italic=True, color=t.secondary)

    chart_top = content_y + 350000
    chart_h = 4200000
    rule_color = tint(t.secondary, 0.5)

    # Thin rule above chart
    add_line(slide, m, chart_top - 30000, SLIDE_W - m, chart_top - 30000,
             color=rule_color, width=0.5)

    spec = ChartSpec(
        chart_type="donut", title="",
        categories=c.get("pie_categories", ["Enterprise", "Mid-Market", "SMB", "Gov", "Partners"]),
        series=[{"name": "Share", "values": c.get("pie_values", [42, 28, 15, 10, 5])}],
        show_legend=True, show_data_labels=True,
    )

    # Center the chart
    chart_w = int((SLIDE_W - 2 * m) * 0.65)
    chart_left = m + (SLIDE_W - 2 * m - chart_w) // 2
    add_chart(slide, spec, profile_from_theme(t),
              left=chart_left, top=chart_top, width=chart_w, height=chart_h)

    # Thin rule below chart
    add_line(slide, m, chart_top + chart_h + 30000, SLIDE_W - m, chart_top + chart_h + 30000,
             color=rule_color, width=0.5)

    # Source caption
    add_textbox(slide, "Source: Internal data", m, chart_top + chart_h + 60000,
                SLIDE_W - 2 * m, 250000,
                font_size=9, italic=True, color=tint(t.secondary, 0.6))


# ── Variant: laboratory ──────────────────────────────────────────────

def _laboratory(slide, t):
    """Dark bg, accent line below title, donut chart + legend."""
    c = t.content
    p = t.palette
    m = MARGIN
    title = c.get("pie_title", "Revenue Distribution")

    add_background(slide, t.primary)
    content_y = add_slide_title(slide, title, theme=t)

    # Colored accent line below title
    add_line(slide, m, content_y + 30000, SLIDE_W - m, content_y + 30000,
             color=t.accent, width=2.0)

    spec = ChartSpec(
        chart_type="donut", title="",
        categories=c.get("pie_categories", ["Enterprise", "Mid-Market", "SMB", "Gov", "Partners"]),
        series=[{"name": "Share", "values": c.get("pie_values", [42, 28, 15, 10, 5])}],
        show_legend=False, show_data_labels=True,
    )

    chart_w = int((SLIDE_W - 2 * m) * 0.60)
    add_chart(slide, spec, profile_from_theme(t),
              left=m, top=content_y + 150000, width=chart_w, height=4500000)

    # Manual legend (white text on dark bg)
    legend_left = m + chart_w + 300000
    legend_w = SLIDE_W - legend_left - m
    legend_top = content_y + 600000

    legend_colors = p
    legend_labels = c.get("pie_legend", ["Cat 1", "Cat 2", "Cat 3", "Cat 4", "Cat 5"])

    for i, (label, color) in enumerate(zip(legend_labels, legend_colors)):
        y = legend_top + i * 500000
        add_rect(slide, legend_left, y, 200000, 200000, fill_color=color, corner_radius=30000)
        add_textbox(slide, label, legend_left + 300000, y - 20000, legend_w - 400000, 250000,
                    font_size=13, color="#FFFFFF")


# ── Variant: dashboard ───────────────────────────────────────────────

def _dashboard(slide, t):
    """White bg with accent header band, chart rendered larger."""
    c = t.content
    p = t.palette
    m = int(MARGIN * 0.85)
    title = c.get("pie_title", "Revenue Distribution")

    add_background(slide, "#FFFFFF")

    # Accent header band
    add_rect(slide, 0, 0, SLIDE_W, 750000, fill_color=t.primary)
    add_textbox(slide, title, m, 150000, SLIDE_W - 2 * m, 450000,
                font_size=22, bold=True, color="#FFFFFF")
    add_gold_accent_line(slide, m, 700000, 600000, color=t.accent)

    spec = ChartSpec(
        chart_type="donut", title="",
        categories=c.get("pie_categories", ["Enterprise", "Mid-Market", "SMB", "Gov", "Partners"]),
        series=[{"name": "Share", "values": c.get("pie_values", [42, 28, 15, 10, 5])}],
        show_legend=False, show_data_labels=True,
    )

    chart_w = int((SLIDE_W - 2 * m) * 0.60)
    add_chart(slide, spec, profile_from_theme(t),
              left=m, top=850000, width=chart_w, height=5500000)

    # Compact legend
    legend_left = m + chart_w + 200000
    legend_w = SLIDE_W - legend_left - m
    legend_top = 1300000

    legend_colors = p
    legend_labels = c.get("pie_legend", ["Cat 1", "Cat 2", "Cat 3", "Cat 4", "Cat 5"])

    for i, (label, color) in enumerate(zip(legend_labels, legend_colors)):
        y = legend_top + i * 450000
        add_rect(slide, legend_left, y, 180000, 180000, fill_color=color, corner_radius=25000)
        add_textbox(slide, label, legend_left + 260000, y - 15000, legend_w - 350000, 220000,
                    font_size=12, color=t.primary)
