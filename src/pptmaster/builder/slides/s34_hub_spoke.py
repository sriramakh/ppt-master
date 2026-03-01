"""Slide 34 — Hub & Spoke Diagram: 14 unique visual variants dispatched by UX style."""

from __future__ import annotations

import math

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_circle, add_line, add_hexagon, add_pentagon,
    add_connector_arrow, add_gold_accent_line, add_slide_title,
    add_dark_bg, add_background, add_styled_card,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    s = t.ux_style
    c = t.content
    p = t.palette

    dispatch = {
        "radial": _radial,
        "clean-lines": _clean_lines,
        "pentagon-hub": _pentagon_hub,
        "floating-hub": _floating_hub,
        "dark-radial": _dark_radial,
        "split-hub": _split_hub,
        "hex-hub": _hex_hub,
        "editorial-hub": _editorial_hub,
        "gradient-hub": _gradient_hub,
        "retro-hub": _retro_hub,
        "creative-hub": _creative_hub,
        "scholarly-hub": _scholarly_hub,
        "laboratory-hub": _laboratory_hub,
        "dashboard-hub": _dashboard_hub,
    }
    builder = dispatch.get(s.hub_spoke, _radial)
    builder(slide, t, c, p)


# ── Data helpers ────────────────────────────────────────────────────────

def _get_spokes(c):
    """Extract hub-spoke data from content dict with sensible defaults."""
    return c.get("hub_spokes", [
        ("Analytics", "Real-time data insights"),
        ("Security", "Enterprise-grade protection"),
        ("Integration", "Seamless API connectivity"),
        ("Automation", "Workflow optimization"),
        ("Support", "24/7 expert assistance"),
        ("Scale", "Global infrastructure"),
    ])


def _get_center(c):
    """Extract center hub label."""
    return c.get("hub_center", "Our Platform")


def _spoke_positions(n, center_x, center_y, orbit_radius):
    """Return list of (x, y) positions for n spokes evenly spaced in a circle."""
    positions = []
    for i in range(n):
        angle = 2 * math.pi * i / n - math.pi / 2  # start from top
        sx = center_x + int(orbit_radius * math.cos(angle))
        sy = center_y + int(orbit_radius * math.sin(angle))
        positions.append((sx, sy))
    return positions


# ═══════════════════════════════════════════════════════════════════════
# Variant 1: CLASSIC (radial) — Central filled circle + radial satellites
# ═══════════════════════════════════════════════════════════════════════

def _radial(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    s = t.ux_style
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    # Layout geometry
    area_top = content_top + 100000
    area_h = FOOTER_TOP - area_top - 200000
    cx = SLIDE_W // 2
    cy = area_top + area_h // 2
    hub_r = 450000
    spoke_r = 280000
    orbit_r = 1700000

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Connectors
    for sx, sy in positions:
        add_connector_arrow(slide, cx, cy, sx, sy,
                            color=tint(t.accent, 0.5), width=1.5,
                            head_end=False, tail_end=False)

    # Hub
    add_circle(slide, cx, cy, hub_r, fill_color=t.accent)
    add_textbox(slide, center_label,
                cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2,
                font_size=16, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Spokes
    for i, ((sx, sy), (label, desc)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        add_circle(slide, sx, sy, spoke_r, fill_color=accent)
        add_textbox(slide, label,
                    sx - spoke_r, sy - spoke_r, spoke_r * 2, spoke_r * 2,
                    font_size=11, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        # Description below spoke
        add_textbox(slide, desc,
                    sx - 700000, sy + spoke_r + 40000, 1400000, 380000,
                    font_size=9, color=sub_color, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# Variant 2: MINIMAL (clean-lines) — Thin connectors, no-fill circles
# ═══════════════════════════════════════════════════════════════════════

def _clean_lines(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    s = t.ux_style
    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    line_color = tint(t.secondary, 0.6) if not s.dark_mode else tint(t.primary, 0.3)
    bg_color = t.primary if s.dark_mode else t.light_bg

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    area_top = content_top + 100000
    area_h = FOOTER_TOP - area_top - 200000
    cx = SLIDE_W // 2
    cy = area_top + area_h // 2
    hub_r = 400000
    spoke_r = 240000
    orbit_r = 1650000

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Thin connectors
    for sx, sy in positions:
        add_line(slide, cx, cy, sx, sy, color=line_color, width=0.5)

    # Hub — outline only (bg fill to simulate transparency)
    add_circle(slide, cx, cy, hub_r, fill_color=bg_color,
               line_color=t.accent, line_width=1.5)
    add_textbox(slide, center_label,
                cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2,
                font_size=14, bold=True, color=text_color,
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Spoke circles — outline only (bg fill to simulate transparency)
    for i, ((sx, sy), (label, desc)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        add_circle(slide, sx, sy, spoke_r, fill_color=bg_color,
                   line_color=accent, line_width=1)
        add_textbox(slide, label,
                    sx - spoke_r, sy - spoke_r, spoke_r * 2, spoke_r * 2,
                    font_size=10, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, desc,
                    sx - 700000, sy + spoke_r + 30000, 1400000, 380000,
                    font_size=8, color=sub_color, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# Variant 3: BOLD (pentagon-hub) — Pentagon center + bold colored blocks
# ═══════════════════════════════════════════════════════════════════════

def _pentagon_hub(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    s = t.ux_style
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    area_top = content_top + 100000
    area_h = FOOTER_TOP - area_top - 200000
    cx = SLIDE_W // 2
    cy = area_top + area_h // 2
    hub_size = 500000
    spoke_block_w = 700000
    spoke_block_h = 500000
    orbit_r = 1800000

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Bold connectors
    for sx, sy in positions:
        add_connector_arrow(slide, cx, cy, sx, sy,
                            color=t.accent, width=3,
                            head_end=True, tail_end=False)

    # Pentagon hub
    add_pentagon(slide, cx, cy, hub_size, fill_color=t.accent,
                 text=center_label.upper(), text_color="#FFFFFF", font_size=14)

    # Bold rectangular spoke blocks
    for i, ((sx, sy), (label, desc)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        bx = sx - spoke_block_w // 2
        by = sy - spoke_block_h // 2
        add_rect(slide, bx, by, spoke_block_w, spoke_block_h, fill_color=accent)
        add_textbox(slide, label.upper(),
                    bx + 40000, by + 50000, spoke_block_w - 80000, 200000,
                    font_size=12, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc,
                    bx + 30000, by + 240000, spoke_block_w - 60000, 260000,
                    font_size=7, color=tint("#FFFFFF", 0.3),
                    alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# Variant 4: ELEVATED (floating-hub) — Floating shadowed circles
# ═══════════════════════════════════════════════════════════════════════

def _floating_hub(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    s = t.ux_style
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    area_top = content_top + 100000
    area_h = FOOTER_TOP - area_top - 200000
    cx = SLIDE_W // 2
    cy = area_top + area_h // 2
    hub_r = 480000
    spoke_r = 360000
    orbit_r = 1800000
    shadow_offset = 45000

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Connectors (dashed, elegant)
    for sx, sy in positions:
        add_line(slide, cx, cy, sx, sy,
                 color=tint(t.accent, 0.4), width=1, dash="dash")

    # Hub shadow + circle
    shadow_c = shade(t.light_bg, 0.15) if not s.dark_mode else shade(t.primary, 0.3)
    add_circle(slide, cx + shadow_offset, cy + shadow_offset, hub_r,
               fill_color=shadow_c)
    add_circle(slide, cx, cy, hub_r, fill_color="#FFFFFF",
               line_color=t.accent, line_width=2)
    add_textbox(slide, center_label,
                cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2,
                font_size=15, bold=True, color=t.primary,
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Spoke shadows + circles
    for i, ((sx, sy), (label, desc)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        add_circle(slide, sx + shadow_offset, sy + shadow_offset, spoke_r,
                   fill_color=shadow_c)
        add_circle(slide, sx, sy, spoke_r, fill_color="#FFFFFF",
                   line_color=accent, line_width=1.5)
        add_textbox(slide, label,
                    sx - spoke_r, sy - spoke_r + 30000, spoke_r * 2, spoke_r,
                    font_size=11, bold=True, color=t.primary,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, desc,
                    sx - spoke_r + 20000, sy + 20000, spoke_r * 2 - 40000, spoke_r - 30000,
                    font_size=7, color=t.secondary,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.TOP)


# ═══════════════════════════════════════════════════════════════════════
# Variant 5: DARK (dark-radial) — Dark bg, glowing hub, neon connectors
# ═══════════════════════════════════════════════════════════════════════

def _dark_radial(slide, t, c, p):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    area_top = content_top + 100000
    area_h = FOOTER_TOP - area_top - 200000
    cx = SLIDE_W // 2
    cy = area_top + area_h // 2
    hub_r = 460000
    spoke_r = 280000
    orbit_r = 1700000

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Neon glow ring behind hub
    add_circle(slide, cx, cy, hub_r + 80000,
               fill_color=tint(t.accent, 0.15))

    # Neon connectors (double-line glow effect)
    for sx, sy in positions:
        add_line(slide, cx, cy, sx, sy,
                 color=tint(t.accent, 0.2), width=4)
        add_line(slide, cx, cy, sx, sy,
                 color=t.accent, width=1.5)

    # Hub
    add_circle(slide, cx, cy, hub_r, fill_color=shade(t.primary, 0.2),
               line_color=t.accent, line_width=2)
    add_textbox(slide, center_label,
                cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2,
                font_size=15, bold=True, color=t.accent,
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Spokes with glow
    for i, ((sx, sy), (label, desc)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        add_circle(slide, sx, sy, spoke_r + 50000,
                   fill_color=tint(accent, 0.1))
        add_circle(slide, sx, sy, spoke_r,
                   fill_color=shade(t.primary, 0.15),
                   line_color=accent, line_width=1.5)
        add_textbox(slide, label,
                    sx - spoke_r, sy - spoke_r, spoke_r * 2, spoke_r * 2,
                    font_size=10, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, desc,
                    sx - 700000, sy + spoke_r + 50000, 1400000, 380000,
                    font_size=8, color=tint(t.primary, 0.5),
                    alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# Variant 6: SPLIT (split-hub) — Hub on left, details list on right
# ═══════════════════════════════════════════════════════════════════════

def _split_hub(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    s = t.ux_style
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    # Left half: hub diagram
    area_top = content_top + 100000
    area_h = FOOTER_TOP - area_top - 200000
    left_w = int(SLIDE_W * 0.50)
    cx = left_w // 2 + MARGIN // 2
    cy = area_top + area_h // 2
    hub_r = 350000
    spoke_r = 200000
    orbit_r = 1100000

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Connectors
    for sx, sy in positions:
        add_connector_arrow(slide, cx, cy, sx, sy,
                            color=tint(t.accent, 0.4), width=1.5,
                            head_end=False, tail_end=False)

    # Hub
    add_circle(slide, cx, cy, hub_r, fill_color=t.accent)
    add_textbox(slide, center_label,
                cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2,
                font_size=13, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Spoke circles
    for i, ((sx, sy), (label, _)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        add_circle(slide, sx, sy, spoke_r, fill_color=accent)
        add_textbox(slide, label,
                    sx - spoke_r, sy - spoke_r, spoke_r * 2, spoke_r * 2,
                    font_size=9, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Vertical divider
    div_x = left_w + MARGIN // 2
    add_line(slide, div_x, area_top + 100000, div_x, FOOTER_TOP - 300000,
             color=tint(t.accent, 0.3), width=1)

    # Right half: detail list
    detail_left = div_x + 300000
    detail_w = SLIDE_W - detail_left - MARGIN
    item_h = (FOOTER_TOP - area_top - 400000) // max(n, 1)

    for i, (label, desc) in enumerate(spokes):
        accent = p[i % len(p)]
        iy = area_top + 200000 + i * item_h

        # Accent dot
        add_circle(slide, detail_left + 60000, iy + 80000, 40000, fill_color=accent)

        add_textbox(slide, label,
                    detail_left + 180000, iy, detail_w - 180000, 250000,
                    font_size=13, bold=True, color=text_color)
        add_textbox(slide, desc,
                    detail_left + 180000, iy + 260000, detail_w - 180000, 250000,
                    font_size=10, color=sub_color)


# ═══════════════════════════════════════════════════════════════════════
# Variant 7: GEO (hex-hub) — Hexagonal center with hex satellites
# ═══════════════════════════════════════════════════════════════════════

def _hex_hub(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    s = t.ux_style
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    area_top = content_top + 100000
    area_h = FOOTER_TOP - area_top - 200000
    cx = SLIDE_W // 2
    cy = area_top + area_h // 2
    hub_size = 450000
    spoke_size = 300000
    orbit_r = 1700000

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Angular connectors
    for sx, sy in positions:
        add_line(slide, cx, cy, sx, sy, color=tint(t.accent, 0.4), width=1.5)

    # Center hexagon
    add_hexagon(slide, cx, cy, hub_size, fill_color=t.accent)
    add_textbox(slide, center_label,
                cx - hub_size, cy - hub_size, hub_size * 2, hub_size * 2,
                font_size=14, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Spoke hexagons
    for i, ((sx, sy), (label, desc)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        add_hexagon(slide, sx, sy, spoke_size, fill_color=accent)
        add_textbox(slide, label,
                    sx - spoke_size, sy - spoke_size, spoke_size * 2, spoke_size * 2,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, desc,
                    sx - 700000, sy + spoke_size + 40000, 1400000, 380000,
                    font_size=8, color=sub_color, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# Variant 8: EDITORIAL (editorial-hub) — Elegant thin circles, editorial type
# ═══════════════════════════════════════════════════════════════════════

def _editorial_hub(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    s = t.ux_style
    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    bg_color = t.primary if s.dark_mode else "#FFFFFF"

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    area_top = content_top + 100000
    area_h = FOOTER_TOP - area_top - 200000
    cx = SLIDE_W // 2
    cy = area_top + area_h // 2
    hub_r = 420000
    spoke_r = 320000
    orbit_r = 1750000

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Thin elegant lines
    for sx, sy in positions:
        add_line(slide, cx, cy, sx, sy,
                 color=tint(t.secondary, 0.7), width=0.5)

    # Hub — thin double ring (bg fill to simulate outline-only)
    add_circle(slide, cx, cy, hub_r + 40000, fill_color=bg_color,
               line_color=t.accent, line_width=0.5)
    add_circle(slide, cx, cy, hub_r, fill_color=bg_color,
               line_color=t.accent, line_width=1.5)
    add_textbox(slide, center_label,
                cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2,
                font_size=14, bold=False, italic=True, color=text_color,
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Spoke — thin single ring with editorial-style labels (bg fill)
    for i, ((sx, sy), (label, desc)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        add_circle(slide, sx, sy, spoke_r, fill_color=bg_color,
                   line_color=accent, line_width=1)
        add_textbox(slide, label,
                    sx - spoke_r, sy - spoke_r + 20000, spoke_r * 2, spoke_r,
                    font_size=10, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, desc,
                    sx - spoke_r + 20000, sy + 20000, spoke_r * 2 - 40000, spoke_r - 30000,
                    font_size=7, italic=True, color=sub_color,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.TOP)

    # Editorial thin rule under title
    add_line(slide, m, content_top - 20000, m + 500000, content_top - 20000,
             color=t.accent, width=0.5)


# ═══════════════════════════════════════════════════════════════════════
# Variant 9: GRADIENT (gradient-hub) — Gradient-tinted hub + spoke circles
# ═══════════════════════════════════════════════════════════════════════

def _gradient_hub(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    s = t.ux_style
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    area_top = content_top + 100000
    area_h = FOOTER_TOP - area_top - 200000
    cx = SLIDE_W // 2
    cy = area_top + area_h // 2
    hub_r = 480000
    spoke_r = 350000
    orbit_r = 1800000

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Large subtle tinted background circle
    add_circle(slide, cx, cy, orbit_r + spoke_r + 100000,
               fill_color=tint(t.accent, 0.05) if not s.dark_mode else tint(t.primary, 0.08))

    # Gradient-band connectors (two-tone)
    for i, (sx, sy) in enumerate(positions):
        accent = p[i % len(p)]
        add_line(slide, cx, cy, sx, sy, color=tint(accent, 0.3), width=3)
        add_line(slide, cx, cy, sx, sy, color=tint(accent, 0.6), width=1)

    # Hub — gradient feel via layered circles
    add_circle(slide, cx, cy, hub_r + 30000, fill_color=tint(t.accent, 0.3))
    add_circle(slide, cx, cy, hub_r, fill_color=t.accent)
    add_textbox(slide, center_label,
                cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2,
                font_size=15, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Spoke — layered gradient tint
    for i, ((sx, sy), (label, desc)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        add_circle(slide, sx, sy, spoke_r + 25000, fill_color=tint(accent, 0.3))
        add_circle(slide, sx, sy, spoke_r, fill_color=accent)
        add_textbox(slide, label,
                    sx - spoke_r, sy - spoke_r, spoke_r * 2, spoke_r,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.BOTTOM)
        add_textbox(slide, desc,
                    sx - spoke_r + 20000, sy, spoke_r * 2 - 40000, spoke_r,
                    font_size=7, color=tint("#FFFFFF", 0.3),
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.TOP)


# ═══════════════════════════════════════════════════════════════════════
# Variant 10: RETRO (retro-hub) — Vintage styled hub, decorative connectors
# ═══════════════════════════════════════════════════════════════════════

def _retro_hub(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    s = t.ux_style
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    bg_color = t.primary if s.dark_mode else t.light_bg

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    area_top = content_top + 100000
    area_h = FOOTER_TOP - area_top - 200000
    cx = SLIDE_W // 2
    cy = area_top + area_h // 2
    hub_r = 450000
    spoke_r = 280000
    orbit_r = 1700000

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Decorative dashed connectors
    for sx, sy in positions:
        add_line(slide, cx, cy, sx, sy,
                 color=tint(t.accent, 0.5), width=2, dash="lgDash")

    # Hub — retro double ring with thick border (bg fill for outer ring)
    add_circle(slide, cx, cy, hub_r + 50000, fill_color=bg_color,
               line_color=t.accent, line_width=3)
    add_circle(slide, cx, cy, hub_r, fill_color=t.accent,
               line_color=shade(t.accent, 0.2), line_width=2)
    add_textbox(slide, center_label,
                cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2,
                font_size=14, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Decorative dots along the orbit path
    dot_count = n * 3
    for i in range(dot_count):
        angle = 2 * math.pi * i / dot_count - math.pi / 2
        dx = cx + int(orbit_r * 0.65 * math.cos(angle))
        dy = cy + int(orbit_r * 0.65 * math.sin(angle))
        add_circle(slide, dx, dy, 15000, fill_color=tint(t.accent, 0.5))

    # Spoke — retro bordered circles
    for i, ((sx, sy), (label, desc)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        add_circle(slide, sx, sy, spoke_r, fill_color=accent,
                   line_color=shade(accent, 0.2), line_width=2)
        add_textbox(slide, label,
                    sx - spoke_r, sy - spoke_r, spoke_r * 2, spoke_r * 2,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, desc,
                    sx - 700000, sy + spoke_r + 40000, 1400000, 380000,
                    font_size=8, color=sub_color, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# Variant 11: MAGAZINE (creative-hub) — Asymmetric layout, varying sizes
# ═══════════════════════════════════════════════════════════════════════

def _creative_hub(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    s = t.ux_style
    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    area_top = content_top + 100000
    area_h = FOOTER_TOP - area_top - 200000

    # Offset hub to left-of-center for asymmetry
    cx = int(SLIDE_W * 0.38)
    cy = area_top + area_h // 2 + 50000
    hub_r = 500000
    orbit_r = 1600000

    # Varying spoke sizes for creative feel (min 320000 to fit descriptions)
    base_sizes = [380000, 340000, 360000, 320000, 350000, 330000, 370000, 340000]

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Connectors — varied widths
    for i, (sx, sy) in enumerate(positions):
        w = 1.0 + (i % 3) * 0.5
        add_line(slide, cx, cy, sx, sy,
                 color=tint(t.accent, 0.35), width=w)

    # Hub — oversized for magazine emphasis
    add_circle(slide, cx, cy, hub_r, fill_color=t.accent)
    add_textbox(slide, center_label,
                cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2,
                font_size=18, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Spokes — varying sizes
    for i, ((sx, sy), (label, desc)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        sr = base_sizes[i % len(base_sizes)]
        add_circle(slide, sx, sy, sr, fill_color=accent)
        add_textbox(slide, label,
                    sx - sr, sy - sr, sr * 2, sr,
                    font_size=10 if sr < 300000 else 11, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.BOTTOM)
        add_textbox(slide, desc,
                    sx - sr + 20000, sy, sr * 2 - 40000, sr - 10000,
                    font_size=7, color=tint("#FFFFFF", 0.3),
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.TOP)


# ═══════════════════════════════════════════════════════════════════════
# Variant 12: SCHOLARLY (scholarly-hub) — Figure caption, numbered spokes
# ═══════════════════════════════════════════════════════════════════════

def _scholarly_hub(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    s = t.ux_style
    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    line_color = tint(t.secondary, 0.6) if not s.dark_mode else tint(t.primary, 0.3)
    bg_color = t.primary if s.dark_mode else "#FFFFFF"

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    area_top = content_top + 200000
    area_h = FOOTER_TOP - area_top - 600000  # room for caption
    cx = SLIDE_W // 2
    cy = area_top + area_h // 2
    hub_r = 380000
    spoke_r = 310000
    orbit_r = 1650000

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Thin scholarly lines
    for sx, sy in positions:
        add_line(slide, cx, cy, sx, sy, color=line_color, width=0.75)

    # Hub — thin outline, scholarly (bg fill to simulate outline-only)
    add_circle(slide, cx, cy, hub_r, fill_color=bg_color,
               line_color=t.primary if not s.dark_mode else tint(t.primary, 0.5),
               line_width=1.5)
    add_textbox(slide, center_label,
                cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2,
                font_size=13, bold=False, color=text_color,
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Numbered spokes — thin outlines (bg fill to simulate outline-only)
    for i, ((sx, sy), (label, desc)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        add_circle(slide, sx, sy, spoke_r, fill_color=bg_color,
                   line_color=accent, line_width=1)
        # Number in circle
        num_r = 70000
        add_circle(slide, sx, sy - spoke_r + num_r + 20000, num_r,
                   fill_color=accent)
        add_textbox(slide, str(i + 1),
                    sx - num_r, sy - spoke_r + 20000, num_r * 2, num_r * 2,
                    font_size=9, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, label,
                    sx - spoke_r, sy - spoke_r + num_r * 2 + 40000,
                    spoke_r * 2, spoke_r,
                    font_size=9, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, desc,
                    sx - spoke_r + 15000, sy + 30000,
                    spoke_r * 2 - 30000, spoke_r - 40000,
                    font_size=7, color=sub_color,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.TOP)

    # Figure caption at bottom
    caption = f"Fig. 1 \u2014 {center_label}: Interconnected capability model (n={n})"
    add_textbox(slide, caption,
                m, FOOTER_TOP - 380000, SLIDE_W - 2 * m, 300000,
                font_size=9, italic=True, color=sub_color,
                alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# Variant 13: LABORATORY (laboratory-hub) — Dark bg, data-coded nodes, grid
# ═══════════════════════════════════════════════════════════════════════

def _laboratory_hub(slide, t, c, p):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    area_top = content_top + 100000
    area_h = FOOTER_TOP - area_top - 200000
    cx = SLIDE_W // 2
    cy = area_top + area_h // 2
    hub_r = 420000
    spoke_r = 320000
    orbit_r = 1750000

    grid_color = tint(t.primary, 0.08)

    # Background grid
    grid_step = 400000
    for gx in range(0, SLIDE_W + 1, grid_step):
        add_line(slide, gx, area_top, gx, FOOTER_TOP - 100000,
                 color=grid_color, width=0.25)
    for gy in range(area_top, FOOTER_TOP, grid_step):
        add_line(slide, MARGIN // 2, gy, SLIDE_W - MARGIN // 2, gy,
                 color=grid_color, width=0.25)

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Data connectors — solid thin
    for sx, sy in positions:
        add_line(slide, cx, cy, sx, sy,
                 color=tint(t.accent, 0.3), width=1)

    # Hub — data-coded (bordered, dark fill)
    add_circle(slide, cx, cy, hub_r,
               fill_color=shade(t.primary, 0.3),
               line_color=t.accent, line_width=1.5)
    add_textbox(slide, center_label,
                cx - hub_r, cy - hub_r + 50000, hub_r * 2, hub_r - 50000,
                font_size=13, bold=True, color=t.accent,
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
    # Simulated data label
    add_textbox(slide, f"n={n} nodes",
                cx - hub_r, cy + 50000, hub_r * 2, hub_r - 50000,
                font_size=8, color=tint(t.primary, 0.4),
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Spoke — color-coded data nodes
    for i, ((sx, sy), (label, desc)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        # Data-code color bar on left side of node
        add_circle(slide, sx, sy, spoke_r,
                   fill_color=shade(t.primary, 0.25),
                   line_color=accent, line_width=1.5)
        add_rect(slide, sx - spoke_r, sy - spoke_r,
                 30000, spoke_r * 2, fill_color=accent)
        add_textbox(slide, f"[{i + 1}] {label}",
                    sx - spoke_r + 40000, sy - spoke_r + 30000,
                    spoke_r * 2 - 50000, spoke_r,
                    font_size=9, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, desc,
                    sx - spoke_r + 40000, sy + 10000,
                    spoke_r * 2 - 50000, spoke_r - 30000,
                    font_size=7, color=tint(t.primary, 0.5),
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.TOP)

    # Corner brackets (lab style)
    bsz = 150000
    bt = area_top + 50000
    bb = FOOTER_TOP - 150000
    bl = MARGIN // 2
    br = SLIDE_W - MARGIN // 2
    bcolor = t.accent
    add_rect(slide, bl, bt, bsz, 3, fill_color=bcolor)
    add_rect(slide, bl, bt, 3, bsz, fill_color=bcolor)
    add_rect(slide, br - bsz, bb, bsz, 3, fill_color=bcolor)
    add_rect(slide, br - 3, bb - bsz, 3, bsz, fill_color=bcolor)


# ═══════════════════════════════════════════════════════════════════════
# Variant 14: DASHBOARD (dashboard-hub) — Compact hub with metric annotations
# ═══════════════════════════════════════════════════════════════════════

def _dashboard_hub(slide, t, c, p):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("hub_title", "Core Capabilities"), theme=t)

    s = t.ux_style
    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    card_fill = tint(t.primary, 0.12) if s.dark_mode else "#FFFFFF"

    spokes = _get_spokes(c)
    center_label = _get_center(c)
    n = len(spokes)

    # Compact layout — hub occupies left 55%, metric cards on right
    area_top = content_top + 50000
    area_h = FOOTER_TOP - area_top - 150000
    hub_area_w = int(SLIDE_W * 0.55)
    cx = hub_area_w // 2 + m // 2
    cy = area_top + area_h // 2
    hub_r = 320000
    spoke_r = 200000
    orbit_r = 1050000

    positions = _spoke_positions(n, cx, cy, orbit_r)

    # Connectors
    for sx, sy in positions:
        add_line(slide, cx, cy, sx, sy,
                 color=tint(t.accent, 0.3), width=1)

    # Compact hub
    add_circle(slide, cx, cy, hub_r, fill_color=t.accent)
    add_textbox(slide, center_label,
                cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2,
                font_size=11, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Compact spoke nodes
    for i, ((sx, sy), (label, _)) in enumerate(zip(positions, spokes)):
        accent = p[i % len(p)]
        add_circle(slide, sx, sy, spoke_r, fill_color=accent)
        add_textbox(slide, label,
                    sx - spoke_r, sy - spoke_r, spoke_r * 2, spoke_r * 2,
                    font_size=8, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Right panel — metric annotation cards
    panel_left = hub_area_w + 100000
    panel_w = SLIDE_W - panel_left - m
    card_h = (area_h - (n - 1) * 60000) // max(n, 1)
    card_h = min(card_h, 650000)

    for i, (label, desc) in enumerate(spokes):
        accent = p[i % len(p)]
        cy_card = area_top + i * (card_h + 60000)

        add_styled_card(slide, panel_left, cy_card, panel_w, card_h,
                        theme=t, accent_color=accent)

        # Metric-style number
        add_textbox(slide, f"0{i + 1}",
                    panel_left + 60000, cy_card + 40000, 250000, card_h - 80000,
                    font_size=16, bold=True, color=accent,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Label + description
        txt_left = panel_left + 330000
        txt_w = panel_w - 390000
        add_textbox(slide, label,
                    txt_left, cy_card + 40000, txt_w, int(card_h * 0.45),
                    font_size=10, bold=True, color=text_color)
        add_textbox(slide, desc,
                    txt_left, cy_card + int(card_h * 0.45) + 20000,
                    txt_w, int(card_h * 0.45),
                    font_size=8, color=sub_color)
