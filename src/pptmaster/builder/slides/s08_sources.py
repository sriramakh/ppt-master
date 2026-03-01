"""Slide 8 — Sources & References: Numbered bibliography list."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from pptmaster.builder.design_system import SLIDE_W, SLIDE_H, MARGIN, FONT_BODY, FONT_CARD_HEADING, FONT_CAPTION, CONTENT_TOP
from pptmaster.builder.helpers import (
    add_textbox, add_styled_card, add_slide_title, add_dark_bg,
    add_background, add_rect, add_line,
)
from pptmaster.assets.color_utils import tint, shade


# ── Scholarly variant ─────────────────────────────────────────────────
def _scholarly(slide, t) -> None:
    """References — numbered bibliography list with hanging indent, thin top rule."""
    c = t.content
    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, c.get("sources_title", "References"), theme=t)

    text_color = t.primary
    sub_color = t.secondary
    m = MARGIN

    # Thin top rule
    rule_y = content_y + 50000
    add_line(slide, m, rule_y, SLIDE_W - m, rule_y, color=t.secondary, width=0.5)

    sources = c.get("sources_list", [
        "World Economic Forum (2025). Global Risks Report.",
        "McKinsey & Company (2025). The State of AI.",
        "IPCC (2024). Sixth Assessment Report.",
        "Harvard Business Review (2025). Leadership in Transition.",
    ])

    y = rule_y + 200000
    content_w = SLIDE_W - 2 * m
    num_w = 250000
    row_h = min(500000, 4500000 // max(len(sources), 1))

    for i, ref in enumerate(sources[:8]):
        # Number — flush left
        add_textbox(slide, f"[{i + 1}]", m, y, num_w, row_h,
                    font_size=13, bold=True, color=text_color)
        # Reference text — hanging indent
        add_textbox(slide, str(ref), m + num_w + 50000, y,
                    content_w - num_w - 100000, row_h,
                    font_size=13, color=sub_color)
        y += row_h


# ── Laboratory variant ────────────────────────────────────────────────
def _laboratory(slide, t) -> None:
    """References & Citations — numbered list on dark bg, accent colored numbers."""
    c = t.content
    add_background(slide, t.primary)
    content_y = add_slide_title(
        slide, c.get("sources_title", "References & Citations"), theme=t,
    )

    text_color = "#FFFFFF"
    sub_color = tint(t.primary, 0.6)
    m = MARGIN

    sources = c.get("sources_list", [
        "World Economic Forum (2025). Global Risks Report.",
        "McKinsey & Company (2025). The State of AI.",
        "IPCC (2024). Sixth Assessment Report.",
        "Harvard Business Review (2025). Leadership in Transition.",
    ])

    # Thin accent rule under title area
    rule_y = content_y + 50000
    add_line(slide, m, rule_y, SLIDE_W - m, rule_y, color=t.accent, width=1.0)

    y = rule_y + 200000
    content_w = SLIDE_W - 2 * m
    num_w = 300000
    row_h = min(500000, 4500000 // max(len(sources), 1))

    for i, ref in enumerate(sources[:8]):
        # Accent-colored number
        add_textbox(slide, f"{i + 1}.", m + 100000, y, num_w, row_h,
                    font_size=14, bold=True, color=t.accent)
        # Reference text
        add_textbox(slide, str(ref), m + num_w + 150000, y,
                    content_w - num_w - 250000, row_h,
                    font_size=13, color=sub_color)
        y += row_h


# ── Dashboard variant ─────────────────────────────────────────────────
def _dashboard(slide, t) -> None:
    """Header band + compact two-column numbered list to fit more sources."""
    c = t.content
    s = t.ux_style
    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("sources_title", "Sources & References"), theme=t,
    )

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    # Accent header band
    add_rect(slide, 0, content_y - 200000, SLIDE_W, 400000, fill_color=t.accent)

    sources = c.get("sources_list", [
        "World Economic Forum (2025). Global Risks Report.",
        "McKinsey & Company (2025). The State of AI.",
        "IPCC (2024). Sixth Assessment Report.",
        "Harvard Business Review (2025). Leadership in Transition.",
    ])

    m = int(MARGIN * 0.7)
    list_y = content_y + 350000
    avail_h = SLIDE_H - list_y - 400000
    col_gap = 300000
    col_w = (SLIDE_W - 2 * m - col_gap) // 2
    num_w = 220000

    # Split sources into two columns
    n = len(sources[:8])
    mid = (n + 1) // 2
    row_h = min(450000, avail_h // max(mid, 1))

    for i, ref in enumerate(sources[:8]):
        if i < mid:
            cx = m
            ry = list_y + i * row_h
        else:
            cx = m + col_w + col_gap
            ry = list_y + (i - mid) * row_h

        add_textbox(slide, f"{i + 1}.", cx, ry, num_w, row_h,
                    font_size=12, bold=True, color=t.accent)
        add_textbox(slide, str(ref), cx + num_w, ry, col_w - num_w, row_h,
                    font_size=11, color=text_color)


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    s = t.ux_style

    # ── Style dispatch ────────────────────────────────────────────────
    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)
    content_y = add_slide_title(
        slide, c.get("sources_title", "Sources & References"), theme=t,
    )

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    sources = c.get("sources_list", [
        "World Economic Forum (2025). Global Risks Report.",
        "McKinsey & Company (2025). The State of AI.",
        "IPCC (2024). Sixth Assessment Report.",
        "Harvard Business Review (2025). Leadership in Transition.",
    ])

    m = int(MARGIN * s.margin_factor)
    card_left = m
    card_w = SLIDE_W - 2 * m
    card_top = content_y + 100000
    card_h = 4800000

    add_styled_card(slide, card_left, card_top, card_w, card_h,
                    theme=t, accent_color=t.accent)

    # Render each reference as a numbered entry
    n = len(sources)
    row_h = min(500000, card_h // max(n, 1))
    y = card_top + 200000

    for i, ref in enumerate(sources[:8]):
        # Number badge
        num_text = f"{i + 1}."
        add_textbox(slide, num_text,
                    card_left + 200000, y, 300000, row_h,
                    font_size=14, bold=True, color=t.accent,
                    alignment=PP_ALIGN.RIGHT)
        # Reference text
        add_textbox(slide, str(ref),
                    card_left + 550000, y, card_w - 750000, row_h,
                    font_size=13, color=text_color,
                    alignment=PP_ALIGN.LEFT)
        y += row_h
