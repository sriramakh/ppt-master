"""Shared drawing primitives for the corporate template builder."""

from __future__ import annotations

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Emu

from pptmaster.assets.color_utils import hex_to_rgb, tint, shade
from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, NAVY, GOLD, SLATE, WHITE,
    LIGHT_GRAY, FONT_FAMILY, FONT_FALLBACK, FONT_SLIDE_NUMBER,
    FONT_CAPTION, FOOTER_TOP, FOOTER_HEIGHT,
)

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


# ── Font helper ────────────────────────────────────────────────────────

def _font_name() -> str:
    """Return the preferred font (Inter), with note that Calibri is fallback."""
    return FONT_FAMILY


# ── Text primitives ───────────────────────────────────────────────────

def add_textbox(
    slide,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    font_size: int = 16,
    bold: bool = False,
    italic: bool = False,
    color: str = NAVY,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = "",
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> object:
    """Add a text box with styled text. Returns the shape."""
    txbox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txbox.text_frame
    tf.word_wrap = True

    # Vertical alignment
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass

    # Set vertical anchor on the text frame body properties
    try:
        txBody = txbox._element.find(f".//{{{_NS_A}}}bodyPr")
        if txBody is not None:
            anchor_map = {
                MSO_ANCHOR.TOP: "t",
                MSO_ANCHOR.MIDDLE: "ctr",
                MSO_ANCHOR.BOTTOM: "b",
            }
            txBody.set("anchor", anchor_map.get(vertical_anchor, "t"))
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name or _font_name()
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    r, g, b = hex_to_rgb(color)
    run.font.color.rgb = RGBColor(r, g, b)

    return txbox


def add_multiline_textbox(
    slide,
    lines: list[tuple[str, int, bool, str]],
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    line_spacing: float = 1.2,
) -> object:
    """Add a textbox with multiple lines, each with its own style.

    Args:
        lines: List of (text, font_size_pt, bold, color_hex) tuples.
    """
    txbox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, (text, font_size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_before = Pt(int(font_size * 0.3))
        run = p.add_run()
        run.text = text
        run.font.name = _font_name()
        run.font.size = Pt(font_size)
        run.font.bold = bold
        r, g, b = hex_to_rgb(color)
        run.font.color.rgb = RGBColor(r, g, b)

    return txbox


def add_bullet_list(
    slide,
    bullets: list[str],
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    font_size: int = 16,
    color: str = NAVY,
    bullet_color: str = GOLD,
    bullet_char: str = "\u2022",
    spacing: int = 6,
) -> object:
    """Add a bulleted list text box."""
    txbox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        p.space_after = Pt(spacing // 2)
        p.level = 0

        # Set bullet character via XML
        pPr = p._pPr
        if pPr is None:
            pPr = etree.SubElement(p._p, f"{{{_NS_A}}}pPr")
        buFont = etree.SubElement(pPr, f"{{{_NS_A}}}buFont")
        buFont.set("typeface", _font_name())
        buChar = etree.SubElement(pPr, f"{{{_NS_A}}}buChar")
        buChar.set("char", bullet_char)

        # Bullet color
        br, bg, bb = hex_to_rgb(bullet_color)
        buClr = etree.SubElement(pPr, f"{{{_NS_A}}}buClr")
        srgbClr = etree.SubElement(buClr, f"{{{_NS_A}}}srgbClr")
        srgbClr.set("val", f"{br:02X}{bg:02X}{bb:02X}")

        run = p.add_run()
        run.text = bullet_text
        run.font.name = _font_name()
        run.font.size = Pt(font_size)
        r, g, b = hex_to_rgb(color)
        run.font.color.rgb = RGBColor(r, g, b)

    return txbox


# ── Shape primitives ──────────────────────────────────────────────────

def add_rect(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill_color: str = "",
    line_color: str = "",
    line_width: float = 0,
    corner_radius: int = 0,
) -> object:
    """Add a rectangle (optionally rounded). Returns the shape."""
    if corner_radius > 0:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(left), Emu(top), Emu(width), Emu(height),
        )
        # Adjust corner radius (0-50000 range in shape adjustments)
        try:
            shape.adjustments[0] = min(corner_radius / max(width, height), 0.1)
        except Exception:
            pass
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(left), Emu(top), Emu(width), Emu(height),
        )

    if fill_color:
        r, g, b = hex_to_rgb(fill_color)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
    else:
        shape.fill.background()

    if line_color and line_width > 0:
        r, g, b = hex_to_rgb(line_color)
        shape.line.color.rgb = RGBColor(r, g, b)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()

    return shape


def add_circle(
    slide,
    center_x: int,
    center_y: int,
    radius: int,
    *,
    fill_color: str = GOLD,
    line_color: str = "",
    line_width: float = 0,
) -> object:
    """Add a circle centered at (center_x, center_y)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(center_x - radius), Emu(center_y - radius),
        Emu(radius * 2), Emu(radius * 2),
    )

    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)

    if line_color and line_width > 0:
        r, g, b = hex_to_rgb(line_color)
        shape.line.color.rgb = RGBColor(r, g, b)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()

    return shape


def add_line(
    slide,
    start_x: int,
    start_y: int,
    end_x: int,
    end_y: int,
    *,
    color: str = GOLD,
    width: float = 2.0,
    dash: str = "",
) -> object:
    """Add a straight line connector."""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Emu(start_x), Emu(start_y),
        Emu(end_x), Emu(end_y),
    )
    r, g, b = hex_to_rgb(color)
    connector.line.color.rgb = RGBColor(r, g, b)
    connector.line.width = Pt(width)

    if dash:
        ln = connector._element.find(f".//{{{_NS_A}}}ln")
        if ln is not None:
            prstDash = etree.SubElement(ln, f"{{{_NS_A}}}prstDash")
            prstDash.set("val", dash)

    return connector


# ── Compound primitives ───────────────────────────────────────────────

def add_card(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill_color: str = WHITE,
    accent_color: str = "",
    accent_side: str = "top",
    accent_thickness: int = 4,
    shadow: bool = True,
    corner_radius: int = 80000,
) -> object:
    """Add a card with optional accent bar and shadow.

    Args:
        accent_side: "top", "left", or "" for no accent.
        accent_thickness: Accent bar thickness in pt.
    """
    # Shadow: offset darker rect behind
    if shadow:
        shadow_offset = 27432  # ~0.03"
        add_rect(
            slide,
            left + shadow_offset, top + shadow_offset,
            width, height,
            fill_color=shade(LIGHT_GRAY, 0.15),
            corner_radius=corner_radius,
        )

    # Main card
    card = add_rect(
        slide, left, top, width, height,
        fill_color=fill_color,
        corner_radius=corner_radius,
    )

    # Accent bar
    if accent_color:
        if accent_side == "top":
            add_rect(
                slide, left, top, width, Pt(accent_thickness).emu,
                fill_color=accent_color,
            )
        elif accent_side == "left":
            add_rect(
                slide, left, top, Pt(accent_thickness).emu, height,
                fill_color=accent_color,
            )

    return card


def add_chevron(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill_color: str = NAVY,
    text: str = "",
    text_color: str = WHITE,
    font_size: int = 14,
) -> object:
    """Add a chevron shape with optional centered text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Emu(left), Emu(top), Emu(width), Emu(height),
    )
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = _font_name()
        run.font.size = Pt(font_size)
        run.font.bold = True
        r, g, b = hex_to_rgb(text_color)
        run.font.color.rgb = RGBColor(r, g, b)

        # Vertical center
        try:
            bodyPr = shape._element.find(f".//{{{_NS_A}}}bodyPr")
            if bodyPr is not None:
                bodyPr.set("anchor", "ctr")
        except Exception:
            pass


    return shape


def add_image_placeholder(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    label: str = "Image Placeholder",
    bg_color: str = LIGHT_GRAY,
) -> object:
    """Add a placeholder rectangle where an image would go."""
    shape = add_rect(
        slide, left, top, width, height,
        fill_color=bg_color,
        line_color=SLATE,
        line_width=1,
    )

    # Dashed border
    ln = shape._element.find(f".//{{{_NS_A}}}ln")
    if ln is not None:
        prstDash = etree.SubElement(ln, f"{{{_NS_A}}}prstDash")
        prstDash.set("val", "dash")

    # Center label
    add_textbox(
        slide, label,
        left + width // 4, top + height // 2 - 100000,
        width // 2, 200000,
        font_size=12, color=SLATE, alignment=PP_ALIGN.CENTER,
    )

    return shape


# ── Decorative elements ───────────────────────────────────────────────

def add_gold_accent_line(
    slide,
    left: int,
    top: int,
    width: int,
    *,
    thickness: float = 3.0,
    color: str = "",
) -> object:
    """Add a horizontal accent line (defaults to gold)."""
    line_color = color or GOLD
    return add_line(slide, left, top, left + width, top, color=line_color, width=thickness)


def add_footer_bar(
    slide,
    *,
    company_name: str = "",
    dark_bg: bool = False,
) -> None:
    """Add a footer bar with company name at the bottom of a slide."""
    bar_color = shade(NAVY, 0.2) if dark_bg else LIGHT_GRAY
    text_color = tint(NAVY, 0.6) if dark_bg else SLATE

    add_rect(
        slide, 0, FOOTER_TOP, SLIDE_W, FOOTER_HEIGHT,
        fill_color=bar_color,
    )

    if company_name:
        add_textbox(
            slide, company_name,
            MARGIN, FOOTER_TOP + 100000,
            3000000, 260000,
            font_size=FONT_CAPTION, color=text_color,
            alignment=PP_ALIGN.LEFT,
        )


def add_slide_number(slide, number: int, *, dark_bg: bool = False) -> None:
    """Add a slide number in the bottom-right corner."""
    text_color = tint(NAVY, 0.6) if dark_bg else SLATE
    add_textbox(
        slide, str(number),
        SLIDE_W - MARGIN - 400000, FOOTER_TOP + 100000,
        400000, 260000,
        font_size=FONT_SLIDE_NUMBER, color=text_color,
        alignment=PP_ALIGN.RIGHT,
    )


# ── Z-order helper ────────────────────────────────────────────────────

def send_to_back(slide, shape) -> None:
    """Move a shape to the back of the z-order."""
    sp_tree = slide._element.find(f".//{{{_NS_P}}}cSld/{{{_NS_P}}}spTree")
    if sp_tree is not None:
        sp = shape._element
        sp_tree.remove(sp)
        sp_tree.insert(2, sp)  # After nvGrpSpPr and grpSpPr


def add_background(slide, color: str = NAVY) -> None:
    """Set the full-slide background to a solid color."""
    shape = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=color)
    send_to_back(slide, shape)


# ═══════════════════════════════════════════════════════════════════════
# Style-aware primitives (read UXStyle from theme)
# ═══════════════════════════════════════════════════════════════════════

def _get_style(theme):
    """Extract UXStyle from theme, with fallback."""
    if theme is None:
        from pptmaster.builder.ux_styles import CLASSIC
        return CLASSIC
    return getattr(theme, "ux_style", None) or __import__(
        "pptmaster.builder.ux_styles", fromlist=["CLASSIC"]
    ).CLASSIC


def add_styled_card(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    theme=None,
    accent_color: str = "",
    fill_color: str = "",
) -> object:
    """Add a card whose appearance is determined by the theme's UX style.

    Automatically applies: corner radius, shadow, accent position/style,
    border, and fill color based on the theme's UXStyle.
    """
    s = _get_style(theme)

    # Determine fill
    if not fill_color:
        if s.dark_mode:
            fill_color = tint(theme.primary, 0.15) if theme else "#2A2A3E"
        else:
            fill_color = WHITE

    # Shadow
    if s.card_shadow:
        shadow_c = shade(fill_color, 0.2) if s.dark_mode else shade(LIGHT_GRAY, 0.15)
        add_rect(slide, left + s.card_shadow_offset, top + s.card_shadow_offset,
                 width, height, fill_color=shadow_c, corner_radius=s.card_radius)

    # Main card
    if s.card_border:
        border_c = tint(theme.secondary, 0.6) if theme else SLATE
        if s.dark_mode and theme:
            border_c = tint(theme.primary, 0.3)
        card = add_rect(slide, left, top, width, height, fill_color=fill_color,
                        line_color=border_c, line_width=s.card_border_width,
                        corner_radius=s.card_radius)
    else:
        card = add_rect(slide, left, top, width, height, fill_color=fill_color,
                        corner_radius=s.card_radius)

    # Accent bar
    if accent_color and s.card_accent != "none":
        thick = Pt(4).emu
        if s.card_accent == "top":
            add_rect(slide, left, top, width, thick, fill_color=accent_color)
        elif s.card_accent == "left":
            add_rect(slide, left, top, thick, height, fill_color=accent_color)
        elif s.card_accent == "bottom":
            add_rect(slide, left, top + height - thick, width, thick, fill_color=accent_color)
        elif s.card_accent == "full":
            add_rect(slide, left, top, width, height,
                     line_color=accent_color, line_width=2, corner_radius=s.card_radius)

    return card


def add_slide_title(
    slide,
    title: str,
    *,
    theme=None,
    subtitle: str = "",
) -> int:
    """Add a styled slide title + accent line based on UX style.

    Returns the y-position where content should start below the title.
    """
    from pptmaster.builder.design_system import MARGIN, FONT_SLIDE_TITLE, CONTENT_TOP

    s = _get_style(theme)
    t = theme
    m = int(MARGIN * s.margin_factor)
    font_size = int(FONT_SLIDE_TITLE * s.title_scale)
    title_text = title.upper() if s.title_caps else title

    text_color = t.primary if t else NAVY
    accent_c = t.accent if t else GOLD
    sub_color = t.secondary if t else SLATE

    if s.dark_mode and t:
        text_color = "#FFFFFF"
        sub_color = tint(t.primary, 0.6)

    align = PP_ALIGN.CENTER if s.title_align == "center" else PP_ALIGN.LEFT
    title_w = SLIDE_W - 2 * m

    add_textbox(slide, title_text, m, 274320, title_w, 600000,
                font_size=font_size, bold=True, color=text_color, alignment=align)

    accent_y = 880000
    if s.title_accent:
        if s.title_align == "center":
            line_left = SLIDE_W // 2 - s.title_accent_width // 2
        else:
            line_left = m
        add_gold_accent_line(slide, line_left, accent_y, s.title_accent_width, color=accent_c)

    if subtitle:
        add_textbox(slide, subtitle, m, accent_y + 100000, title_w, 350000,
                    font_size=13, color=sub_color, alignment=align)
        return accent_y + 500000

    return CONTENT_TOP


def add_dark_bg(slide, theme) -> None:
    """Add dark background if theme uses dark mode, otherwise light tint."""
    s = _get_style(theme)
    if s.dark_mode:
        add_background(slide, theme.primary if theme else NAVY)
    elif theme and theme.light_bg != "#FFFFFF":
        add_background(slide, theme.light_bg)


def add_themed_footer(
    slide,
    *,
    company_name: str = "",
    dark_bg: bool = False,
    theme=None,
) -> None:
    """Footer bar that respects theme colors and dark mode."""
    s = _get_style(theme)
    t = theme

    if s.dark_mode or dark_bg:
        bar_color = shade(t.primary, 0.2) if t else shade(NAVY, 0.2)
        text_color = tint(t.primary, 0.5) if t else tint(NAVY, 0.6)
    else:
        bar_color = LIGHT_GRAY
        text_color = t.secondary if t else SLATE

    add_rect(slide, 0, FOOTER_TOP, SLIDE_W, FOOTER_HEIGHT, fill_color=bar_color)

    if company_name:
        add_textbox(slide, company_name, int(MARGIN * s.margin_factor),
                    FOOTER_TOP + 100000, 3000000, 260000,
                    font_size=FONT_CAPTION, color=text_color, alignment=PP_ALIGN.LEFT)


def add_themed_slide_number(slide, number: int, *, dark_bg: bool = False, theme=None) -> None:
    """Slide number that respects theme colors and dark mode."""
    s = _get_style(theme)
    t = theme
    if s.dark_mode or dark_bg:
        text_color = tint(t.primary, 0.5) if t else tint(NAVY, 0.6)
    else:
        text_color = t.secondary if t else SLATE
    add_textbox(slide, str(number), SLIDE_W - MARGIN - 400000, FOOTER_TOP + 100000,
                400000, 260000, font_size=FONT_SLIDE_NUMBER, color=text_color,
                alignment=PP_ALIGN.RIGHT)


def add_triangle(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill_color: str = GOLD,
    rotation: float = 0,
) -> object:
    """Add a triangle shape (for geometric styles)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height),
    )
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    if rotation:
        shape.rotation = rotation
    return shape


def add_hexagon(
    slide,
    center_x: int,
    center_y: int,
    size: int,
    *,
    fill_color: str = GOLD,
) -> object:
    """Add a hexagon shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Emu(center_x - size), Emu(center_y - size),
        Emu(size * 2), Emu(size * 2),
    )
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    return shape


def add_diamond(
    slide,
    center_x: int,
    center_y: int,
    size: int,
    *,
    fill_color: str = GOLD,
) -> object:
    """Add a diamond shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        Emu(center_x - size), Emu(center_y - size),
        Emu(size * 2), Emu(size * 2),
    )
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    return shape


# ── Additional shape primitives ──────────────────────────────────────

def add_pentagon(
    slide,
    center_x: int,
    center_y: int,
    size: int,
    *,
    fill_color: str = GOLD,
    text: str = "",
    text_color: str = WHITE,
    font_size: int = 12,
) -> object:
    """Add a pentagon shape centered at (center_x, center_y)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.PENTAGON,
        Emu(center_x - size), Emu(center_y - size),
        Emu(size * 2), Emu(size * 2),
    )
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = _font_name()
        run.font.size = Pt(font_size)
        run.font.bold = True
        tr, tg, tb = hex_to_rgb(text_color)
        run.font.color.rgb = RGBColor(tr, tg, tb)
        bodyPr = shape._element.find(f".//{{{_NS_A}}}bodyPr")
        if bodyPr is not None:
            bodyPr.set("anchor", "ctr")

    return shape


def add_octagon(
    slide,
    center_x: int,
    center_y: int,
    size: int,
    *,
    fill_color: str = GOLD,
    text: str = "",
    text_color: str = WHITE,
    font_size: int = 12,
) -> object:
    """Add an octagon shape centered at (center_x, center_y)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OCTAGON,
        Emu(center_x - size), Emu(center_y - size),
        Emu(size * 2), Emu(size * 2),
    )
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = _font_name()
        run.font.size = Pt(font_size)
        run.font.bold = True
        tr, tg, tb = hex_to_rgb(text_color)
        run.font.color.rgb = RGBColor(tr, tg, tb)
        bodyPr = shape._element.find(f".//{{{_NS_A}}}bodyPr")
        if bodyPr is not None:
            bodyPr.set("anchor", "ctr")

    return shape


def add_block_arc(
    slide,
    center_x: int,
    center_y: int,
    radius: int,
    *,
    fill_color: str = GOLD,
    start_angle: float = 0.0,
    sweep_angle: float = 270.0,
    thickness: float = 0.25,
) -> object:
    """Add a block arc (donut arc) shape for gauge/donut visuals.

    Args:
        start_angle: Start angle in degrees (0 = top, clockwise).
        sweep_angle: Arc sweep in degrees.
        thickness: Arc thickness as fraction of radius (0.0-0.5).
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.BLOCK_ARC,
        Emu(center_x - radius), Emu(center_y - radius),
        Emu(radius * 2), Emu(radius * 2),
    )
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    # Adjust arc: adj1 = start angle, adj2 = sweep, adj3 = thickness
    try:
        shape.adjustments[0] = start_angle / 360.0
        shape.adjustments[1] = sweep_angle / 360.0
        shape.adjustments[2] = thickness
    except (IndexError, Exception):
        pass
    return shape


def add_funnel_tier(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill_color: str = GOLD,
    text: str = "",
    text_color: str = WHITE,
    font_size: int = 14,
) -> object:
    """Add a trapezoid shape representing a funnel tier."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        Emu(left), Emu(top), Emu(width), Emu(height),
    )
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = _font_name()
        run.font.size = Pt(font_size)
        run.font.bold = True
        tr, tg, tb = hex_to_rgb(text_color)
        run.font.color.rgb = RGBColor(tr, tg, tb)
        bodyPr = shape._element.find(f".//{{{_NS_A}}}bodyPr")
        if bodyPr is not None:
            bodyPr.set("anchor", "ctr")
    return shape


def add_notched_arrow(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill_color: str = GOLD,
    text: str = "",
    text_color: str = WHITE,
    font_size: int = 12,
) -> object:
    """Add a notched right arrow shape (for milestone roadmaps)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.NOTCHED_RIGHT_ARROW,
        Emu(left), Emu(top), Emu(width), Emu(height),
    )
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = _font_name()
        run.font.size = Pt(font_size)
        run.font.bold = True
        tr, tg, tb = hex_to_rgb(text_color)
        run.font.color.rgb = RGBColor(tr, tg, tb)
        bodyPr = shape._element.find(f".//{{{_NS_A}}}bodyPr")
        if bodyPr is not None:
            bodyPr.set("anchor", "ctr")
    return shape


def add_tear(
    slide,
    center_x: int,
    center_y: int,
    size: int,
    *,
    fill_color: str = GOLD,
    rotation: float = 0,
) -> object:
    """Add a tear/droplet shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.TEAR,
        Emu(center_x - size), Emu(center_y - size),
        Emu(size * 2), Emu(size * 2),
    )
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    if rotation:
        shape.rotation = rotation
    return shape


def add_snip_rect(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill_color: str = WHITE,
    line_color: str = "",
    line_width: float = 0,
) -> object:
    """Add a snipped-corner rectangle (for kanban cards)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.SNIP_2_SAME_RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height),
    )
    if fill_color:
        r, g, b = hex_to_rgb(fill_color)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
    else:
        shape.fill.background()
    if line_color and line_width > 0:
        r, g, b = hex_to_rgb(line_color)
        shape.line.color.rgb = RGBColor(r, g, b)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_round_rect_2(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill_color: str = GOLD,
    text: str = "",
    text_color: str = WHITE,
    font_size: int = 11,
) -> object:
    """Add a round-2-same rectangle (for milestone markers)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height),
    )
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = _font_name()
        run.font.size = Pt(font_size)
        tr, tg, tb = hex_to_rgb(text_color)
        run.font.color.rgb = RGBColor(tr, tg, tb)
        bodyPr = shape._element.find(f".//{{{_NS_A}}}bodyPr")
        if bodyPr is not None:
            bodyPr.set("anchor", "ctr")
    return shape


# ── Compound helpers ─────────────────────────────────────────────────

def add_icon_image(
    slide,
    image_path: str,
    left: int,
    top: int,
    width: int,
    height: int,
) -> object:
    """Place a PNG icon image on the slide."""
    from pathlib import Path
    p = Path(image_path)
    if not p.exists():
        # Fallback: draw a placeholder circle instead
        return add_circle(slide, left + width // 2, top + height // 2,
                         min(width, height) // 2, fill_color=LIGHT_GRAY)
    return slide.shapes.add_picture(str(p), Emu(left), Emu(top), Emu(width), Emu(height))


# ── Icon keyword → PNG resolver ─────────────────────────────────────

# Map common short keywords to icon names in the generated icon set.
_ICON_KEYWORD_MAP: dict[str, str] = {
    # business
    "briefcase": "briefcase", "handshake": "handshake", "building": "building",
    "presentation": "presentation", "strategy": "strategy", "target": "target",
    "growth": "growth", "partnership": "partnership", "meeting": "meeting",
    "contract": "contract", "leadership": "leadership", "innovation": "innovation",
    # technology
    "computer": "computer", "cloud": "cloud", "database": "database",
    "network": "network", "security-shield": "security-shield", "mobile": "mobile",
    "code": "code", "ai": "ai-brain", "brain": "ai-brain", "server": "server",
    "iot": "iot-device", "circuit": "circuit", "robot": "robot",
    # finance
    "dollar": "dollar", "money": "dollar", "wallet": "wallet", "bank": "bank",
    "calculator": "calculator", "invoice": "invoice", "coins": "coins",
    "budget": "budget", "audit": "audit", "stock": "stock-market",
    # healthcare
    "heart": "heart", "stethoscope": "stethoscope", "hospital": "hospital",
    "pill": "pill", "microscope": "microscope", "dna": "dna",
    "ambulance": "ambulance", "vaccine": "vaccine", "health": "brain-health",
    # education
    "book": "book", "graduation": "graduation-cap", "pencil": "pencil",
    "globe": "globe", "lightbulb": "lightbulb", "bulb": "lightbulb",
    "chemistry": "chemistry", "atom": "atom", "telescope": "telescope",
    "library": "library", "certificate": "certificate",
    # sustainability
    "leaf": "leaf", "solar": "solar-panel", "wind": "wind-turbine",
    "recycle": "recycle", "earth": "earth", "water": "water-drop",
    "tree": "tree", "eco": "eco-house", "green": "green-energy",
    # communication
    "email": "email", "chat": "chat-bubble", "phone": "phone",
    "megaphone": "megaphone", "video": "video-call", "satellite": "satellite",
    "podcast": "podcast", "newsletter": "newsletter", "social": "social-media",
    # people
    "person": "person", "team": "team", "community": "community",
    "diversity": "diversity", "family": "family", "mentor": "mentor",
    "users": "team", "customer": "customer", "employee": "employee",
    "speaker": "speaker", "scientist": "scientist",
    # analytics
    "chart": "bar-chart", "bar-chart": "bar-chart", "pie": "pie-chart",
    "line-graph": "line-graph", "dashboard": "dashboard", "funnel": "funnel",
    "gauge": "gauge", "magnifying": "magnifying-glass", "search": "magnifying-glass",
    "data": "data-flow", "statistics": "statistics", "report": "report",
    "analytics": "bar-chart", "heatmap": "heatmap",
    # logistics
    "truck": "truck", "warehouse": "warehouse", "shipping": "shipping",
    "airplane": "airplane", "train": "train", "package": "package",
    "delivery": "delivery", "forklift": "forklift", "container": "container",
    # security
    "lock": "lock", "key": "key", "firewall": "firewall", "eye": "eye",
    "fingerprint": "fingerprint", "shield": "shield-check", "camera": "camera",
    "alarm": "alarm", "vpn": "vpn", "encryption": "encryption",
    "security": "shield-check",
    # creative
    "palette": "palette", "music": "music-note", "film": "film",
    "brush": "brush", "design": "design-grid", "sparkle": "sparkle",
    "typography": "typography", "color": "color-wheel",
    # real-estate
    "house": "house", "apartment": "apartment", "crane": "crane",
    "blueprint": "blueprint", "location": "location-pin", "pin": "location-pin",
    "construction": "construction", "property": "property-value",
    # common aliases
    "trophy": "certificate", "award": "certificate", "star": "sparkle",
    "lightning": "green-energy", "bolt": "green-energy", "speed": "gauge",
    "performance": "gauge", "world": "globe", "global": "globe",
    "idea": "lightbulb", "protect": "shield-check", "safe": "lock",
    "connect": "network", "link": "network", "group": "team",
    "scale": "growth", "expand": "growth", "profit": "chart-up",
    "revenue": "chart-up", "savings": "piggy-bank", "invest": "stock-market",
    "compliance": "shield-check", "risk": "alarm", "quality": "certificate",
    "support": "chat-bubble", "service": "headset",
}

_icon_mgr_cache = None


def _get_icon_mgr():
    """Lazy-load the icon manager singleton."""
    global _icon_mgr_cache
    if _icon_mgr_cache is None:
        try:
            from pptmaster.assets.raster_icon_manager import get_icon_manager
            _icon_mgr_cache = get_icon_manager()
        except Exception:
            _icon_mgr_cache = False  # sentinel: unavailable
    return _icon_mgr_cache if _icon_mgr_cache is not False else None


def resolve_icon_path(keyword: str) -> str | None:
    """Resolve an icon keyword to a PNG file path, or None if unavailable."""
    mgr = _get_icon_mgr()
    if not mgr or not mgr.available:
        return None

    # Direct name match
    mapped = _ICON_KEYWORD_MAP.get(keyword.lower())
    if mapped:
        p = mgr.get_icon_path(mapped)
        if p:
            return str(p)

    # Try exact name
    p = mgr.get_icon_path(keyword.lower())
    if p:
        return str(p)

    # Try search
    results = mgr.search(keyword.lower())
    if results:
        return results[0]["path"]

    return None


def draw_icon_or_placeholder(
    slide, keyword: str, cx: int, cy: int, radius: int,
    accent_color: str, *, letter: str = ""
) -> object:
    """Place a real PNG icon if available, else draw colored circle + letter.

    The icon is sized to fit within the circle area (80% of diameter)
    and centered at (cx, cy).
    """
    icon_path = resolve_icon_path(keyword)
    if icon_path:
        icon_size = int(radius * 1.6)  # 80% of diameter
        left = cx - icon_size // 2
        top = cy - icon_size // 2
        return add_icon_image(slide, icon_path, left, top, icon_size, icon_size)

    # Fallback: colored circle with letter
    char = (letter or keyword[:1]).upper()
    add_circle(slide, cx, cy, radius, fill_color=accent_color)
    return add_textbox(slide, char, cx - radius, cy - radius,
                       radius * 2, radius * 2,
                       font_size=max(10, radius // 14000), bold=True, color="#FFFFFF",
                       alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)


def add_semi_transparent_circle(
    slide,
    center_x: int,
    center_y: int,
    radius: int,
    *,
    fill_color: str = GOLD,
    alpha: int = 50000,
) -> object:
    """Add a circle with semi-transparent fill (for Venn diagrams).

    Args:
        alpha: Transparency in EMU-style (0=opaque, 100000=fully transparent).
              50000 = 50% transparent.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(center_x - radius), Emu(center_y - radius),
        Emu(radius * 2), Emu(radius * 2),
    )
    r, g, b = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()

    # Set alpha transparency via lxml
    solidFill = shape._element.find(f".//{{{_NS_A}}}solidFill")
    if solidFill is not None:
        srgb = solidFill.find(f"{{{_NS_A}}}srgbClr")
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, f"{{{_NS_A}}}alpha")
            alpha_elem.set("val", str(alpha))
    return shape


def add_connector_arrow(
    slide,
    start_x: int,
    start_y: int,
    end_x: int,
    end_y: int,
    *,
    color: str = GOLD,
    width: float = 1.5,
    head_end: bool = True,
    tail_end: bool = False,
) -> object:
    """Add a line connector with arrowhead (for hub-spoke connectors)."""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Emu(start_x), Emu(start_y),
        Emu(end_x), Emu(end_y),
    )
    r, g, b = hex_to_rgb(color)
    connector.line.color.rgb = RGBColor(r, g, b)
    connector.line.width = Pt(width)

    # Add arrowheads via lxml
    ln = connector._element.find(f".//{{{_NS_A}}}ln")
    if ln is not None:
        if head_end:
            headEnd = etree.SubElement(ln, f"{{{_NS_A}}}headEnd")
            headEnd.set("type", "triangle")
            headEnd.set("w", "med")
            headEnd.set("len", "med")
        if tail_end:
            tailEnd = etree.SubElement(ln, f"{{{_NS_A}}}tailEnd")
            tailEnd.set("type", "triangle")
            tailEnd.set("w", "med")
            tailEnd.set("len", "med")
    return connector


def add_color_cell(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill_color: str = GOLD,
    text: str = "",
    text_color: str = WHITE,
    font_size: int = 11,
    bold: bool = False,
    line_color: str = "",
    line_width: float = 0,
) -> object:
    """Add a filled rectangle with centered text (for matrices, heatmaps).

    Convenience wrapper combining add_rect + add_textbox in one call.
    """
    shape = add_rect(slide, left, top, width, height,
                     fill_color=fill_color, line_color=line_color,
                     line_width=line_width)
    if text:
        add_textbox(slide, text,
                    left + 30000, top + 10000,
                    width - 60000, height - 20000,
                    font_size=font_size, bold=bold, color=text_color,
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)
