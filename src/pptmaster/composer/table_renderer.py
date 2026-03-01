"""Table generation — themed tables with accent header row."""

from __future__ import annotations

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.util import Pt, Emu, Inches

from pptmaster.assets.color_utils import hex_to_rgb, tint
from pptmaster.models import DesignProfile, TableSpec

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def add_table(
    slide,
    table_spec: TableSpec,
    profile: DesignProfile,
    left: int | None = None,
    top: int | None = None,
    width: int | None = None,
    height: int | None = None,
) -> None:
    """Add a styled table to a slide.

    Args:
        slide: The pptx slide object.
        table_spec: Table specification with headers and rows.
        profile: Design profile for styling.
        left/top/width/height: Optional position in EMU.
    """
    if not table_spec.headers and not table_spec.rows:
        return

    sw = profile.slide_size.width
    sh = profile.slide_size.height

    num_cols = len(table_spec.headers) if table_spec.headers else (
        len(table_spec.rows[0]) if table_spec.rows else 0
    )
    num_rows = len(table_spec.rows) + (1 if table_spec.headers else 0)

    if num_cols == 0 or num_rows == 0:
        return

    if left is None:
        left = int(sw * 0.08)
    if top is None:
        top = int(sh * 0.25)
    if width is None:
        width = int(sw * 0.84)
    if height is None:
        height = min(int(sh * 0.65), num_rows * 600000)  # ~0.65" per row

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Emu(left), Emu(top), Emu(width), Emu(height),
    )
    table = table_shape.table

    # Set column widths
    if table_spec.column_widths:
        total = sum(table_spec.column_widths)
        for i, cw in enumerate(table_spec.column_widths):
            if i < num_cols:
                table.columns[i].width = Emu(int(width * cw / total))
    else:
        col_width = width // num_cols
        for i in range(num_cols):
            table.columns[i].width = Emu(col_width)

    # Fill header row
    row_idx = 0
    if table_spec.headers:
        for col_idx, header in enumerate(table_spec.headers):
            if col_idx < num_cols:
                cell = table.cell(0, col_idx)
                cell.text = header
                _style_header_cell(cell, profile)
        row_idx = 1

    # Remove default table styling (banding) for clean custom look
    try:
        tbl_elem = table._tbl
        tblPr = tbl_elem.find(f"{{{_NS_A}}}tblPr")
        if tblPr is not None:
            tblPr.set("bandRow", "0")
            tblPr.set("firstRow", "0")
            tblPr.set("lastRow", "0")
    except Exception:
        pass

    # Fill data rows
    for data_row in table_spec.rows:
        if row_idx >= num_rows:
            break
        for col_idx, value in enumerate(data_row):
            if col_idx < num_cols:
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                _style_data_cell(cell, profile, row_idx % 2 == 0)
        row_idx += 1


def _style_header_cell(cell, profile: DesignProfile) -> None:
    """Style a header cell with accent background and white text."""
    r, g, b = hex_to_rgb(profile.colors.accent1)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(r, g, b)

    # Vertical centering
    from pptx.enum.text import MSO_ANCHOR
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Add cell margins for readability
    _set_cell_margins(cell, left=Emu(120000), right=Emu(80000),
                      top=Emu(50000), bottom=Emu(50000))

    for paragraph in cell.text_frame.paragraphs:
        paragraph.space_before = Pt(4)
        paragraph.space_after = Pt(4)
        for run in paragraph.runs:
            run.font.name = profile.fonts.minor
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)


def _style_data_cell(cell, profile: DesignProfile, is_even: bool) -> None:
    """Style a data cell with alternating row colors."""
    if is_even:
        # Light tint of accent color
        tinted = tint(profile.colors.accent1, 0.90)
        r, g, b = hex_to_rgb(tinted)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(r, g, b)
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Vertical centering
    from pptx.enum.text import MSO_ANCHOR
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Add cell margins for readability
    _set_cell_margins(cell, left=Emu(120000), right=Emu(80000),
                      top=Emu(40000), bottom=Emu(40000))

    text_r, text_g, text_b = hex_to_rgb(profile.colors.dk2)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.space_before = Pt(3)
        paragraph.space_after = Pt(3)
        for run in paragraph.runs:
            run.font.name = profile.fonts.minor
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(text_r, text_g, text_b)


def _set_cell_margins(cell, left=None, right=None, top=None, bottom=None) -> None:
    """Set cell margins via XML (python-pptx doesn't expose this directly)."""
    try:
        tc = cell._tc
        tcPr = tc.find(f"{{{_NS_A}}}tcPr")
        if tcPr is None:
            tcPr = etree.SubElement(tc, f"{{{_NS_A}}}tcPr")
        if left is not None:
            tcPr.set("marL", str(int(left)))
        if right is not None:
            tcPr.set("marR", str(int(right)))
        if top is not None:
            tcPr.set("marT", str(int(top)))
        if bottom is not None:
            tcPr.set("marB", str(int(bottom)))
    except Exception:
        pass
