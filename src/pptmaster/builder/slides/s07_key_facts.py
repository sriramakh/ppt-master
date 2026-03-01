"""Slide 7 — Key Facts: 6 metric cards in 3x2 grid."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from pptmaster.builder.design_system import SLIDE_W, SLIDE_H, MARGIN, FONT_SLIDE_TITLE, FONT_BODY, FONT_CARD_HEADING, FONT_CAPTION, CONTENT_TOP, card_positions, col_span
from pptmaster.builder.helpers import add_textbox, add_card, add_gold_accent_line, add_styled_card, add_slide_title, add_dark_bg, add_background, add_rect, add_line
from pptmaster.assets.color_utils import tint, shade


# ── Scholarly variant ─────────────────────────────────────────────────
def _scholarly(slide, t) -> None:
    """Key Data — table-style layout with thin rules, no colored stat cards."""
    c = t.content
    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, c.get("key_facts_title", "Key Data"), theme=t)

    text_color = t.primary
    sub_color = t.secondary
    m = MARGIN

    # Thin top rule
    rule_y = content_y + 50000
    add_line(slide, m, rule_y, SLIDE_W - m, rule_y, color=t.secondary, width=0.5)

    # Table caption
    add_textbox(slide, "Table 1: Key Facts & Figures", m, rule_y + 80000,
                SLIDE_W - 2 * m, 350000, font_size=FONT_CAPTION, italic=True, color=sub_color)

    metrics = c.get("key_facts", [
        ("$850M", "Annual Revenue"), ("2,500+", "Team Members"), ("98%", "Client Retention"),
        ("12", "Global Offices"), ("150+", "Enterprise Clients"), ("4.8/5", "Customer Rating"),
    ])

    # Table header rule
    table_y = rule_y + 500000
    content_w = SLIDE_W - 2 * m
    add_line(slide, m, table_y, SLIDE_W - m, table_y, color=text_color, width=1.0)

    # Column headers
    col_val_w = content_w // 3
    col_lbl_w = content_w - col_val_w
    add_textbox(slide, "Metric", m + 50000, table_y + 50000, col_val_w, 300000,
                font_size=FONT_CAPTION, bold=True, color=text_color)
    add_textbox(slide, "Description", m + col_val_w, table_y + 50000, col_lbl_w, 300000,
                font_size=FONT_CAPTION, bold=True, color=text_color)

    # Header bottom rule
    row_start = table_y + 350000
    add_line(slide, m, row_start, SLIDE_W - m, row_start, color=text_color, width=0.75)

    # Data rows
    row_h = 420000
    for i, (value, label) in enumerate(metrics[:6]):
        ry = row_start + 50000 + i * row_h
        add_textbox(slide, value, m + 50000, ry, col_val_w - 100000, row_h - 50000,
                    font_size=FONT_BODY, bold=True, color=text_color)
        add_textbox(slide, label, m + col_val_w, ry, col_lbl_w - 100000, row_h - 50000,
                    font_size=FONT_BODY, color=sub_color)
        # Row divider
        add_line(slide, m, ry + row_h - 50000, SLIDE_W - m, ry + row_h - 50000,
                 color=t.secondary, width=0.5)


# ── Laboratory variant ────────────────────────────────────────────────
def _laboratory(slide, t) -> None:
    """Key Findings — metric cards with large numbers, colored accent lines."""
    c = t.content
    add_background(slide, t.primary)
    content_y = add_slide_title(slide, c.get("key_facts_title", "Key Findings"), theme=t)

    text_color = "#FFFFFF"
    sub_color = tint(t.primary, 0.6)
    dark_fill = tint(t.primary, 0.15)
    p = t.palette

    metrics = c.get("key_facts", [
        ("$850M", "Annual Revenue"), ("2,500+", "Team Members"), ("98%", "Client Retention"),
        ("12", "Global Offices"), ("150+", "Enterprise Clients"), ("4.8/5", "Customer Rating"),
    ])

    m = MARGIN
    gap = 150000
    cols, rows = 3, 2
    card_w = (SLIDE_W - 2 * m - (cols - 1) * gap) // cols
    card_h = (SLIDE_H - content_y - 600000 - (rows - 1) * gap) // rows

    for i, (value, label) in enumerate(metrics[:6]):
        row, col = divmod(i, cols)
        cx = m + col * (card_w + gap)
        cy = content_y + 150000 + row * (card_h + gap)
        accent = p[i % len(p)]

        add_rect(slide, cx, cy, card_w, card_h, fill_color=dark_fill,
                 line_color=tint(t.primary, 0.25), line_width=1)
        # Colored accent line at top of card
        add_rect(slide, cx, cy, card_w, 50000, fill_color=accent)
        add_textbox(slide, value, cx + 100000, cy + 250000, card_w - 200000, 700000,
                    font_size=36, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, label, cx + 100000, cy + 950000, card_w - 200000, 400000,
                    font_size=FONT_CAPTION, color=sub_color, alignment=PP_ALIGN.CENTER)


# ── Dashboard variant ─────────────────────────────────────────────────
def _dashboard(slide, t) -> None:
    """Header band + 2x3 grid of mini stat tiles (dense)."""
    c = t.content
    s = t.ux_style
    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("key_facts_title", "Key Facts & Figures"), theme=t)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    # Accent header band
    add_rect(slide, 0, content_y - 200000, SLIDE_W, 400000, fill_color=t.accent)

    metrics = c.get("key_facts", [
        ("$850M", "Annual Revenue"), ("2,500+", "Team Members"), ("98%", "Client Retention"),
        ("12", "Global Offices"), ("150+", "Enterprise Clients"), ("4.8/5", "Customer Rating"),
    ])

    m = int(MARGIN * 0.6)
    grid_y = content_y + 350000
    gap = 100000
    cols, rows = 3, 2
    card_w = (SLIDE_W - 2 * m - (cols - 1) * gap) // cols
    card_h = (SLIDE_H - grid_y - 400000 - (rows - 1) * gap) // rows
    p = t.palette

    for i, (value, label) in enumerate(metrics[:6]):
        row, col = divmod(i, cols)
        cx = m + col * (card_w + gap)
        cy = grid_y + row * (card_h + gap)
        accent = p[i % len(p)]

        add_styled_card(slide, cx, cy, card_w, card_h, theme=t, accent_color=accent)
        add_textbox(slide, value, cx + 60000, cy + card_h // 5, card_w - 120000, card_h // 3,
                    font_size=32, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, label, cx + 60000, cy + card_h * 3 // 5, card_w - 120000, card_h // 4,
                    font_size=11, color=sub_color, alignment=PP_ALIGN.CENTER)


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    s = t.ux_style

    # ── Style dispatch ────────────────────────────────────────────────
    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("key_facts_title", "Key Facts & Figures"), theme=t)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    metrics = c.get("key_facts", [
        ("$850M", "Annual Revenue"), ("2,500+", "Team Members"), ("98%", "Client Retention"),
        ("12", "Global Offices"), ("150+", "Enterprise Clients"), ("4.8/5", "Customer Rating"),
    ])

    gap = int(200000 * s.gap_factor)
    positions = card_positions(3, 2, top=content_y + 100000, gap=gap)
    p = t.palette

    for i, (value, label) in enumerate(metrics[:6]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        add_styled_card(slide, left, top, width, height, theme=t, accent_color=accent)
        add_textbox(slide, value, left + 100000, top + 300000, width - 200000, 700000,
                    font_size=38, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, label, left + 100000, top + 1000000, width - 200000, 500000,
                    font_size=14, color=sub_color, alignment=PP_ALIGN.CENTER)
