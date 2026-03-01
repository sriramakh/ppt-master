"""Slide 19 — Comparison: Left vs Right columns with VS badge."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import SLIDE_W, SLIDE_H, MARGIN, FONT_CARD_HEADING, FONT_CAPTION, CONTENT_TOP
from pptmaster.builder.helpers import add_textbox, add_rect, add_line, add_circle, add_gold_accent_line, add_styled_card, add_slide_title, add_dark_bg, add_background
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    p = t.palette

    s = t.ux_style

    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)

    content_top = add_slide_title(slide, c.get("comparison_title", "Comparison"), theme=t)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    total_w = SLIDE_W - 2 * MARGIN
    vs_badge_w = 500000
    col_w = (total_w - vs_badge_w) // 2
    col_top = content_top + 100000

    _headers = c.get("comparison_headers", ("Option A", "Option B"))
    opt_a = _headers[0] if len(_headers) > 0 else "Option A"
    opt_b = _headers[1] if len(_headers) > 1 else "Option B"

    # Column header cards
    add_styled_card(slide, MARGIN, col_top, col_w, 500000, theme=t, accent_color=p[0])
    add_rect(slide, MARGIN, col_top, col_w, 500000, fill_color=p[0], corner_radius=40000)
    add_textbox(slide, opt_a, MARGIN, col_top, col_w, 500000,
                font_size=FONT_CARD_HEADING, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    right_x = MARGIN + col_w + vs_badge_w
    add_styled_card(slide, right_x, col_top, col_w, 500000, theme=t, accent_color=p[1])
    add_rect(slide, right_x, col_top, col_w, 500000, fill_color=p[1], corner_radius=40000)
    add_textbox(slide, opt_b, right_x, col_top, col_w, 500000,
                font_size=FONT_CARD_HEADING, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    vs_cx = MARGIN + col_w + vs_badge_w // 2
    add_circle(slide, vs_cx, col_top + 250000, 200000, fill_color=t.accent)
    add_textbox(slide, "VS", vs_cx - 200000, col_top + 50000, 400000, 400000,
                font_size=16, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    rows = c.get("comparison_rows", [
        ("Cost", "$500K", "$350K"), ("Timeline", "6 months", "4 months"),
        ("Scale", "Enterprise", "Mid-market"), ("Support", "24/7", "Business hours"),
        ("Integration", "200+", "50+"), ("ROI", "12 months", "8 months"),
    ])

    row_h = 500000
    rows_top = col_top + 550000

    for i, (feature, opt_a_val, opt_b_val) in enumerate(rows):
        y = rows_top + i * row_h
        if s.dark_mode:
            bg = tint(t.primary, 0.15) if i % 2 == 0 else tint(t.primary, 0.10)
        else:
            bg = t.primary_tint(0.95) if i % 2 == 0 else "#FFFFFF"

        add_rect(slide, MARGIN, y, total_w, row_h, fill_color=bg)
        add_textbox(slide, opt_a_val, MARGIN + 100000, y, col_w - 200000, row_h,
                    font_size=13, color=text_color, alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, feature, MARGIN + col_w, y, vs_badge_w, row_h,
                    font_size=11, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, opt_b_val, right_x + 100000, y, col_w - 200000, row_h,
                    font_size=13, color=text_color, alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant: scholarly ────────────────────────────────────────────────

def _scholarly(slide, t):
    """White bg, 'Table 2:' caption, two-column layout with thin rules."""
    c = t.content
    p = t.palette
    m = MARGIN

    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, c.get("comparison_title", "Comparison"), theme=t)

    # Table caption
    add_textbox(slide, "Table 2: Comparison", m, content_y + 20000,
                SLIDE_W - 2 * m, 300000,
                font_size=FONT_CAPTION, italic=True, color=t.secondary)

    _headers = c.get("comparison_headers", ("Option A", "Option B"))
    opt_a = _headers[0] if len(_headers) > 0 else "Option A"
    opt_b = _headers[1] if len(_headers) > 1 else "Option B"
    rows = c.get("comparison_rows", [
        ("Cost", "$500K", "$350K"), ("Timeline", "6 months", "4 months"),
        ("Scale", "Enterprise", "Mid-market"), ("Support", "24/7", "Business hours"),
        ("Integration", "200+", "50+"), ("ROI", "12 months", "8 months"),
    ])

    rule_color = tint(t.secondary, 0.5)
    total_w = SLIDE_W - 2 * m
    feature_w = int(total_w * 0.25)
    val_w = int(total_w * 0.375)
    table_top = content_y + 400000
    row_h = 450000

    # Column headers with thin rule below
    add_textbox(slide, "Criterion", m, table_top, feature_w, row_h,
                font_size=12, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT,
                vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, opt_a, m + feature_w, table_top, val_w, row_h,
                font_size=12, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, opt_b, m + feature_w + val_w, table_top, val_w, row_h,
                font_size=12, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Double rule below header
    hdr_bottom = table_top + row_h
    add_line(slide, m, hdr_bottom, m + total_w, hdr_bottom, color=t.primary, width=1.0)
    add_line(slide, m, hdr_bottom + 25000, m + total_w, hdr_bottom + 25000,
             color=t.primary, width=0.5)

    for i, (feature, a_val, b_val) in enumerate(rows):
        y = hdr_bottom + 50000 + i * row_h
        add_textbox(slide, feature, m, y, feature_w, row_h,
                    font_size=11, bold=True, color=t.secondary, alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, a_val, m + feature_w, y, val_w, row_h,
                    font_size=11, color=t.primary, alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, b_val, m + feature_w + val_w, y, val_w, row_h,
                    font_size=11, color=t.primary, alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)
        # Thin rule below row
        rule_y = y + row_h
        add_line(slide, m, rule_y, m + total_w, rule_y,
                 color=rule_color, width=0.25)


# ── Variant: laboratory ──────────────────────────────────────────────

def _laboratory(slide, t):
    """Dark bg, comparison items in dark cards with colored accent top borders."""
    c = t.content
    p = t.palette
    m = MARGIN

    add_background(slide, t.primary)
    content_y = add_slide_title(slide, c.get("comparison_title", "Comparison"), theme=t)

    _headers = c.get("comparison_headers", ("Option A", "Option B"))
    opt_a = _headers[0] if len(_headers) > 0 else "Option A"
    opt_b = _headers[1] if len(_headers) > 1 else "Option B"
    rows = c.get("comparison_rows", [
        ("Cost", "$500K", "$350K"), ("Timeline", "6 months", "4 months"),
        ("Scale", "Enterprise", "Mid-market"), ("Support", "24/7", "Business hours"),
        ("Integration", "200+", "50+"), ("ROI", "12 months", "8 months"),
    ])

    total_w = SLIDE_W - 2 * m
    col_w = (total_w - 200000) // 2
    card_fill = tint(t.primary, 0.15)
    col_a_color = p[0] if len(p) > 0 else t.accent
    col_b_color = p[1] if len(p) > 1 else t.accent

    # Column headers
    header_y = content_y + 100000
    add_rect(slide, m, header_y, col_w, 450000, fill_color=col_a_color, corner_radius=30000)
    add_textbox(slide, opt_a, m, header_y, col_w, 450000,
                font_size=FONT_CARD_HEADING, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    right_x = m + col_w + 200000
    add_rect(slide, right_x, header_y, col_w, 450000, fill_color=col_b_color, corner_radius=30000)
    add_textbox(slide, opt_b, right_x, header_y, col_w, 450000,
                font_size=FONT_CARD_HEADING, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Data rows as dark cards
    row_h = 500000
    rows_top = header_y + 520000

    for i, (feature, a_val, b_val) in enumerate(rows):
        y = rows_top + i * (row_h + 50000)
        # Left card
        add_rect(slide, m, y, col_w, row_h, fill_color=card_fill, corner_radius=30000)
        add_rect(slide, m, y, col_w, 40000, fill_color=col_a_color)
        add_textbox(slide, feature, m + 100000, y + 50000, col_w - 200000, 180000,
                    font_size=9, color=tint(t.primary, 0.5))
        add_textbox(slide, a_val, m + 100000, y + 220000, col_w - 200000, 250000,
                    font_size=14, bold=True, color=col_a_color)
        # Right card
        add_rect(slide, right_x, y, col_w, row_h, fill_color=card_fill, corner_radius=30000)
        add_rect(slide, right_x, y, col_w, 40000, fill_color=col_b_color)
        add_textbox(slide, feature, right_x + 100000, y + 50000, col_w - 200000, 180000,
                    font_size=9, color=tint(t.primary, 0.5))
        add_textbox(slide, b_val, right_x + 100000, y + 220000, col_w - 200000, 250000,
                    font_size=14, bold=True, color=col_b_color)


# ── Variant: dashboard ───────────────────────────────────────────────

def _dashboard(slide, t):
    """White bg with accent header band, compact horizontal comparison."""
    c = t.content
    p = t.palette
    m = int(MARGIN * 0.85)

    add_background(slide, "#FFFFFF")

    # Accent header band
    add_rect(slide, 0, 0, SLIDE_W, 750000, fill_color=t.primary)
    add_textbox(slide, c.get("comparison_title", "Comparison"), m, 150000,
                SLIDE_W - 2 * m, 450000,
                font_size=22, bold=True, color="#FFFFFF")
    add_gold_accent_line(slide, m, 700000, 600000, color=t.accent)

    _headers = c.get("comparison_headers", ("Option A", "Option B"))
    opt_a = _headers[0] if len(_headers) > 0 else "Option A"
    opt_b = _headers[1] if len(_headers) > 1 else "Option B"
    rows = c.get("comparison_rows", [
        ("Cost", "$500K", "$350K"), ("Timeline", "6 months", "4 months"),
        ("Scale", "Enterprise", "Mid-market"), ("Support", "24/7", "Business hours"),
        ("Integration", "200+", "50+"), ("ROI", "12 months", "8 months"),
    ])

    total_w = SLIDE_W - 2 * m
    feature_w = int(total_w * 0.20)
    val_w = int(total_w * 0.40)
    col_a_color = p[0] if len(p) > 0 else t.accent
    col_b_color = p[1] if len(p) > 1 else t.accent

    # Column headers
    hdr_y = 850000
    add_rect(slide, m + feature_w, hdr_y, val_w, 400000, fill_color=tint(col_a_color, 0.9))
    add_textbox(slide, opt_a, m + feature_w, hdr_y, val_w, 400000,
                font_size=14, bold=True, color=col_a_color,
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, m + feature_w + val_w, hdr_y, val_w, 400000, fill_color=tint(col_b_color, 0.9))
    add_textbox(slide, opt_b, m + feature_w + val_w, hdr_y, val_w, 400000,
                font_size=14, bold=True, color=col_b_color,
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Compact data rows
    row_h = 400000
    rows_top = hdr_y + 450000

    for i, (feature, a_val, b_val) in enumerate(rows):
        y = rows_top + i * row_h
        bg = tint(t.primary, 0.95) if i % 2 == 0 else "#FFFFFF"
        add_rect(slide, m, y, total_w, row_h, fill_color=bg)
        add_textbox(slide, feature, m + 80000, y, feature_w - 100000, row_h,
                    font_size=11, bold=True, color=t.secondary,
                    alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, a_val, m + feature_w, y, val_w, row_h,
                    font_size=12, color=t.primary,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, b_val, m + feature_w + val_w, y, val_w, row_h,
                    font_size=12, color=t.primary,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
