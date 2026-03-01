"""Slide 9 — Executive Summary: Bullets left + dark sidebar with 3 metrics."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from pptmaster.builder.design_system import SLIDE_W, MARGIN, FONT_SLIDE_TITLE, FONT_CAPTION, CONTENT_TOP, FOOTER_TOP
from pptmaster.builder.helpers import add_textbox, add_rect, add_bullet_list, add_gold_accent_line, add_line, add_styled_card, add_slide_title, add_dark_bg, add_background
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    s = t.ux_style

    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("exec_title", "Executive Summary"), theme=t)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary
    m = int(MARGIN * s.margin_factor)

    left_w = int((SLIDE_W - 2 * m) * 0.58)
    bullets = c.get("exec_bullets", [])

    add_bullet_list(slide, bullets, m, content_y + 100000, left_w, 4500000,
                    font_size=14, color=text_color, bullet_color=t.accent, spacing=8)

    # Sidebar — in dark mode use a lighter shade instead of dark primary
    right_x = m + left_w + int(300000 * s.gap_factor)
    right_w = SLIDE_W - right_x - m
    sidebar_top = content_y + 100000
    sidebar_h = 4500000

    if s.dark_mode:
        sidebar_fill = tint(t.primary, 0.2)
        sidebar_text = "#FFFFFF"
        sidebar_label_color = tint(t.primary, 0.5)
        sidebar_line_color = tint(t.primary, 0.35)
    else:
        sidebar_fill = t.primary
        sidebar_text = t.accent
        sidebar_label_color = tint(t.primary, 0.6)
        sidebar_line_color = tint(t.primary, 0.3)

    add_rect(slide, right_x, sidebar_top, right_w, sidebar_h,
             fill_color=sidebar_fill, corner_radius=int(80000 * s.margin_factor))

    metrics = c.get("exec_metrics", [("$850M", "Revenue"), ("+23%", "Growth"), ("98%", "Retention")])
    metric_h = sidebar_h // max(len(metrics), 1)

    for i, (value, label) in enumerate(metrics):
        my = sidebar_top + i * metric_h
        add_textbox(slide, value, right_x + 150000, my + 200000, right_w - 300000, 500000,
                    font_size=34, bold=True, color=sidebar_text, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, label, right_x + 150000, my + 700000, right_w - 300000, 350000,
                    font_size=FONT_CAPTION, color=sidebar_label_color, alignment=PP_ALIGN.CENTER)
        if i < len(metrics) - 1:
            ly = my + metric_h - 50000
            add_line(slide, right_x + 200000, ly, right_x + right_w - 200000, ly,
                     color=sidebar_line_color, width=0.5)


# ── Variant: SCHOLARLY — white bg, thin rules, academic feel ─────────────

def _scholarly(slide, t):
    c = t.content
    m = MARGIN

    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, c.get("exec_title", "Abstract"), theme=t)

    # Thin top rule spanning full content width
    rule_left = m
    rule_right = SLIDE_W - m
    add_line(slide, rule_left, content_y + 50000, rule_right, content_y + 50000,
             color=tint(t.secondary, 0.6), width=0.5)

    # Body text as indented paragraphs
    bullets = c.get("exec_bullets", [])
    text_left = m + 200000
    text_w = int((SLIDE_W - 2 * m) * 0.55)
    para_top = content_y + 200000

    for i, bullet in enumerate(bullets):
        add_textbox(slide, bullet, text_left, para_top + i * 550000,
                    text_w, 500000,
                    font_size=12, color=t.primary, alignment=PP_ALIGN.LEFT)

    # Sidebar metrics as "Table N:" style entries
    metrics = c.get("exec_metrics", [("$850M", "Revenue"), ("+23%", "Growth"), ("98%", "Retention")])
    sidebar_left = m + text_w + 500000
    sidebar_w = SLIDE_W - sidebar_left - m
    table_top = content_y + 200000

    add_textbox(slide, "Key Figures", sidebar_left, table_top, sidebar_w, 300000,
                font_size=11, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT)

    # Thin rule under heading
    add_line(slide, sidebar_left, table_top + 300000,
             sidebar_left + sidebar_w, table_top + 300000,
             color=t.primary, width=0.75)

    row_h = 600000
    for i, (value, label) in enumerate(metrics):
        ry = table_top + 400000 + i * row_h

        add_textbox(slide, f"Table {i + 1}:", sidebar_left, ry, sidebar_w * 0.35, 250000,
                    font_size=9, bold=True, color=t.secondary, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, label, sidebar_left + int(sidebar_w * 0.35), ry,
                    int(sidebar_w * 0.35), 250000,
                    font_size=10, color=t.primary, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, value, sidebar_left + int(sidebar_w * 0.7), ry,
                    int(sidebar_w * 0.3), 250000,
                    font_size=12, bold=True, color=t.primary, alignment=PP_ALIGN.RIGHT)

        # Thin row separator
        add_line(slide, sidebar_left, ry + 350000,
                 sidebar_left + sidebar_w, ry + 350000,
                 color=tint(t.secondary, 0.7), width=0.5)


# ── Variant: LABORATORY — dark bg, accent-bordered cards ─────────────────

def _laboratory(slide, t):
    c = t.content
    m = MARGIN

    add_background(slide, t.primary)
    content_y = add_slide_title(slide, c.get("exec_title", "Executive Brief"), theme=t)

    # Left: accent-bordered card for text
    left_w = int((SLIDE_W - 2 * m) * 0.58)
    card_top = content_y + 150000
    card_h = FOOTER_TOP - card_top - 300000
    card_fill = tint(t.primary, 0.08)

    add_rect(slide, m, card_top, left_w, card_h,
             fill_color=card_fill, line_color=tint(t.primary, 0.2),
             line_width=0.75, corner_radius=60000)

    # Left accent border on card
    add_rect(slide, m, card_top + 80000, 40000, card_h - 160000,
             fill_color=t.accent)

    bullets = c.get("exec_bullets", [])
    para_top = card_top + 200000
    for i, bullet in enumerate(bullets):
        add_textbox(slide, f"\u2022  {bullet}", m + 150000, para_top + i * 500000,
                    left_w - 300000, 450000,
                    font_size=12, color="#FFFFFF", alignment=PP_ALIGN.LEFT)

    # Right sidebar: metric tiles with colored top borders
    right_x = m + left_w + 250000
    right_w = SLIDE_W - right_x - m
    metrics = c.get("exec_metrics", [("$850M", "Revenue"), ("+23%", "Growth"), ("98%", "Retention")])
    tile_h = (card_h - (len(metrics) - 1) * 120000) // len(metrics)
    p = t.palette

    for i, (value, label) in enumerate(metrics):
        accent = p[i % len(p)]
        tile_top = card_top + i * (tile_h + 120000)
        tile_fill = tint(t.primary, 0.1)

        add_rect(slide, right_x, tile_top, right_w, tile_h,
                 fill_color=tile_fill, corner_radius=40000)

        # Colored top border
        add_rect(slide, right_x + 1, tile_top, right_w - 2, 40000,
                 fill_color=accent)

        add_textbox(slide, value, right_x + 80000, tile_top + 120000,
                    right_w - 160000, 500000,
                    font_size=30, bold=True, color=accent, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, label, right_x + 80000, tile_top + 620000,
                    right_w - 160000, 300000,
                    font_size=FONT_CAPTION, color=tint(t.primary, 0.5),
                    alignment=PP_ALIGN.CENTER)


# ── Variant: DASHBOARD — compact, dense layout, header band ─────────────

def _dashboard(slide, t):
    c = t.content
    m = int(MARGIN * 0.7)

    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, c.get("exec_title", "Executive Summary"), theme=t)

    # Header accent band
    band_h = 50000
    add_rect(slide, 0, content_y, SLIDE_W, band_h, fill_color=t.accent)

    # Compact bullet points in left area
    bullets = c.get("exec_bullets", [])
    left_w = int((SLIDE_W - 2 * m) * 0.62)
    bullet_top = content_y + band_h + 150000

    for i, bullet in enumerate(bullets):
        add_textbox(slide, f"\u2022  {bullet}", m, bullet_top + i * 400000,
                    left_w, 380000,
                    font_size=11, color=t.primary, alignment=PP_ALIGN.LEFT)

    # Right sidebar — narrow, metrics stacked vertically
    metrics = c.get("exec_metrics", [("$850M", "Revenue"), ("+23%", "Growth"), ("98%", "Retention")])
    right_x = m + left_w + 200000
    right_w = SLIDE_W - right_x - m
    metric_top = content_y + band_h + 100000
    available_h = FOOTER_TOP - metric_top - 200000
    metric_h = available_h // max(len(metrics), 1)
    p = t.palette

    for i, (value, label) in enumerate(metrics):
        accent = p[i % len(p)]
        my = metric_top + i * metric_h

        # Small tile
        add_rect(slide, right_x, my, right_w, metric_h - 80000,
                 fill_color=tint(accent, 0.9), corner_radius=40000)

        # Left accent bar
        add_rect(slide, right_x, my + 30000, 30000, metric_h - 140000,
                 fill_color=accent)

        add_textbox(slide, value, right_x + 60000, my + 50000,
                    right_w - 120000, 350000,
                    font_size=22, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT)
        add_textbox(slide, label, right_x + 60000, my + 400000,
                    right_w - 120000, 250000,
                    font_size=9, color=t.secondary, alignment=PP_ALIGN.LEFT)
