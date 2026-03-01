"""Slide 11 — Process Linear: 11 unique visual variants dispatched by UX style."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_circle, add_line, add_chevron,
    add_gold_accent_line, add_slide_title, add_dark_bg, add_background,
    add_hexagon, add_diamond, add_styled_card,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    s = t.ux_style
    c = t.content

    dispatch = {
        "chevron": _chevron,
        "numbered-line": _numbered_line,
        "block-steps": _block_steps,
        "floating-circles": _floating_circles,
        "progress-bar": _progress_bar_variant,
        "alternating-sides": _alternating_sides,
        "hexagonal": _hexagonal,
        "editorial-flow": _editorial_flow,
        "gradient-circles": _gradient_circles,
        "retro-badges": _retro_badges,
        "creative-bubbles": _creative_bubbles,
        "scholarly-numbered": _scholarly_numbered,
        "laboratory-flow": _laboratory_flow,
        "dashboard-steps": _dashboard_steps,
    }
    builder = dispatch.get(s.process, _chevron)
    builder(slide, t, c)


def _get_steps(c):
    """Extract process steps from content dict with sensible defaults."""
    return c.get("process_steps", [
        ("Discovery", "Research & analysis"), ("Strategy", "Planning & design"),
        ("Develop", "Build & iterate"), ("Deploy", "Launch & integrate"),
        ("Optimize", "Measure & improve"),
    ])


# ── Variant 1: CLASSIC — 5 chevrons in progressive gradient ─────────────

def _chevron(slide, t, c):
    content_top = add_slide_title(slide, c.get("process_title", "Our Process"), theme=t)
    steps = _get_steps(c)

    n = len(steps)
    total_w = SLIDE_W - 2 * MARGIN
    overlap = 120000
    chevron_w = (total_w + overlap * (n - 1)) // n
    chevron_h = 1200000
    chevron_top = content_top + 800000

    for i, (title, desc) in enumerate(steps):
        factor = 0.7 - (i * 0.7 / max(n - 1, 1))
        color = tint(t.primary, factor) if factor > 0 else t.primary
        text_color = t.primary if factor > 0.4 else "#FFFFFF"

        left = MARGIN + i * (chevron_w - overlap)
        add_chevron(slide, left, chevron_top, chevron_w, chevron_h,
                    fill_color=color, text=title, text_color=text_color, font_size=14)
        add_textbox(slide, f"Step {i + 1}", left + chevron_w // 4, chevron_top - 350000,
                    chevron_w // 2, 300000,
                    font_size=10, bold=True, color=t.accent, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, left + 50000, chevron_top + chevron_h + 200000,
                    chevron_w - 100000, 400000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.CENTER)


# ── Variant 2: MINIMAL — numbered circles on a thin line, labels below ──

def _numbered_line(slide, t, c):
    content_top = add_slide_title(slide, c.get("process_title", "Our Process"), theme=t)
    steps = _get_steps(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(steps)
    total_w = SLIDE_W - 2 * m
    step_w = total_w // n
    line_y = content_top + 1200000
    circle_r = 200000

    # Thin horizontal line
    add_line(slide, m + step_w // 2, line_y,
             m + step_w * (n - 1) + step_w // 2, line_y,
             color=tint(t.secondary, 0.7), width=0.75)

    for i, (title, desc) in enumerate(steps):
        cx = m + i * step_w + step_w // 2

        # Small circle on line
        add_circle(slide, cx, line_y, circle_r,
                   fill_color="#FFFFFF", line_color=t.accent, line_width=1.5)

        # Number inside circle
        add_textbox(slide, str(i + 1), cx - circle_r, line_y - circle_r,
                    circle_r * 2, circle_r * 2,
                    font_size=14, bold=True, color=t.accent,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Title below line
        add_textbox(slide, title, cx - step_w // 2 + 40000, line_y + circle_r + 200000,
                    step_w - 80000, 350000,
                    font_size=13, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, desc, cx - step_w // 2 + 40000, line_y + circle_r + 600000,
                    step_w - 80000, 400000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.CENTER)


# ── Variant 3: BOLD — large colored blocks side by side, number on top ───

def _block_steps(slide, t, c):
    content_top = add_slide_title(slide, c.get("process_title", "Our Process").upper(), theme=t)
    steps = _get_steps(c)
    p = t.palette

    n = len(steps)
    gap = 80000
    total_w = SLIDE_W - 2 * MARGIN
    block_w = (total_w - gap * (n - 1)) // n
    block_h = FOOTER_TOP - content_top - 400000
    block_top = content_top + 200000

    for i, (title, desc) in enumerate(steps):
        accent = p[i % len(p)]
        left = MARGIN + i * (block_w + gap)

        # Full-height colored block
        add_rect(slide, left, block_top, block_w, block_h, fill_color=accent)

        # Large number at top
        add_textbox(slide, f"{i + 1:02d}", left + 60000, block_top + 150000,
                    block_w - 120000, 700000,
                    font_size=40, bold=True, color="#FFFFFF", alignment=PP_ALIGN.LEFT)

        # Thick white underline
        add_rect(slide, left + 60000, block_top + 900000, block_w // 2, 50000,
                 fill_color="#FFFFFF")

        # Title
        add_textbox(slide, title.upper(), left + 60000, block_top + 1100000,
                    block_w - 120000, 450000,
                    font_size=14, bold=True, color="#FFFFFF", alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, left + 60000, block_top + 1600000,
                    block_w - 120000, 500000,
                    font_size=11, color=tint(accent, 0.7), alignment=PP_ALIGN.LEFT)


# ── Variant 4: ELEVATED — floating circles connected by lines ────────────

def _floating_circles(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("process_title", "Our Process"), theme=t)
    steps = _get_steps(c)
    p = t.palette

    n = len(steps)
    total_w = SLIDE_W - 2 * MARGIN
    step_w = total_w // n
    circle_y = content_top + 1000000
    circle_r = 380000

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if t.ux_style.dark_mode else t.secondary

    for i, (title, desc) in enumerate(steps):
        accent = p[i % len(p)]
        cx = MARGIN + i * step_w + step_w // 2

        # Connector line to next circle
        if i < n - 1:
            next_cx = MARGIN + (i + 1) * step_w + step_w // 2
            add_line(slide, cx + circle_r, circle_y, next_cx - circle_r, circle_y,
                     color=tint(t.accent, 0.5 if not t.ux_style.dark_mode else 0.2),
                     width=1.5)

        # Shadow circle
        shadow_off = t.ux_style.card_shadow_offset
        add_circle(slide, cx + shadow_off, circle_y + shadow_off, circle_r,
                   fill_color=shade(t.light_bg if not t.ux_style.dark_mode else t.primary, 0.1))

        # Main circle
        fill = "#FFFFFF" if not t.ux_style.dark_mode else tint(t.primary, 0.12)
        add_circle(slide, cx, circle_y, circle_r, fill_color=fill)

        # Accent ring (top of circle)
        add_circle(slide, cx, circle_y - circle_r + 50000, 50000, fill_color=accent)

        # Number
        add_textbox(slide, str(i + 1), cx - 150000, circle_y - 250000, 300000, 250000,
                    font_size=20, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER)

        # Title inside circle
        add_textbox(slide, title, cx - circle_r + 60000, circle_y + 0, circle_r * 2 - 120000, 300000,
                    font_size=12, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)

        # Description below circle
        add_textbox(slide, desc, cx - step_w // 2 + 20000, circle_y + circle_r + 200000,
                    step_w - 40000, 400000,
                    font_size=11, color=sub_color, alignment=PP_ALIGN.CENTER)


# ── Variant 5: DARK — horizontal segmented progress bar ─────────────────

def _progress_bar_variant(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("process_title", "Our Process").upper(), theme=t)
    steps = _get_steps(c)
    p = t.palette

    n = len(steps)
    total_w = SLIDE_W - 2 * MARGIN
    seg_w = total_w // n
    seg_gap = 30000
    actual_seg_w = seg_w - seg_gap
    bar_h = 180000
    bar_top = content_top + 800000

    # Track background
    add_rect(slide, MARGIN, bar_top, total_w, bar_h,
             fill_color=tint(t.primary, 0.1), corner_radius=50000)

    for i, (title, desc) in enumerate(steps):
        accent = p[i % len(p)]
        seg_left = MARGIN + i * seg_w

        # Colored segment
        add_rect(slide, seg_left + seg_gap // 2, bar_top, actual_seg_w, bar_h,
                 fill_color=accent)

        # Step number above
        add_textbox(slide, f"0{i + 1}", seg_left, bar_top - 400000,
                    seg_w, 300000,
                    font_size=22, bold=True, color=accent, alignment=PP_ALIGN.CENTER)

        # Title below bar
        add_textbox(slide, title.upper(), seg_left, bar_top + bar_h + 250000,
                    seg_w, 350000,
                    font_size=13, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, desc, seg_left + 20000, bar_top + bar_h + 700000,
                    seg_w - 40000, 400000,
                    font_size=11, color=tint(t.primary, 0.45), alignment=PP_ALIGN.CENTER)

        # Connector dot below description
        dot_y = bar_top + bar_h + 1250000
        add_circle(slide, seg_left + seg_w // 2, dot_y, 30000, fill_color=accent)

    # Bottom decorative line
    add_line(slide, MARGIN, bar_top + bar_h + 1400000,
             SLIDE_W - MARGIN, bar_top + bar_h + 1400000,
             color=tint(t.primary, 0.15), width=0.5)


# ── Variant 6: SPLIT — steps alternating left/right of a center line ────

def _alternating_sides(slide, t, c):
    content_top = add_slide_title(slide, c.get("process_title", "Our Process"), theme=t)
    steps = _get_steps(c)
    p = t.palette

    n = len(steps)
    mid_x = SLIDE_W // 2
    available_h = FOOTER_TOP - content_top - 300000
    step_h = available_h // n
    card_w = int(SLIDE_W * 0.35)
    card_h = min(step_h - 100000, 800000)

    # Center vertical line
    line_top = content_top + 200000
    line_bottom = content_top + 200000 + n * step_h
    add_line(slide, mid_x, line_top, mid_x, line_bottom,
             color=t.accent, width=2)

    for i, (title, desc) in enumerate(steps):
        accent = p[i % len(p)]
        is_left = i % 2 == 0
        step_top = content_top + 200000 + i * step_h
        node_y = step_top + card_h // 2

        # Node circle on center line
        add_circle(slide, mid_x, node_y, 100000, fill_color=accent)
        add_textbox(slide, str(i + 1), mid_x - 80000, node_y - 80000, 160000, 160000,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        if is_left:
            card_left = mid_x - card_w - 250000
            # Connector line from card to node
            add_line(slide, card_left + card_w, node_y, mid_x - 100000, node_y,
                     color=tint(accent, 0.5), width=1)
        else:
            card_left = mid_x + 250000
            add_line(slide, mid_x + 100000, node_y, card_left, node_y,
                     color=tint(accent, 0.5), width=1)

        # Step card
        add_styled_card(slide, card_left, step_top, card_w, card_h,
                        theme=t, accent_color=accent)

        # Title
        add_textbox(slide, title, card_left + 100000, step_top + 80000,
                    card_w - 200000, 300000,
                    font_size=13, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, card_left + 100000, step_top + 380000,
                    card_w - 200000, 300000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.LEFT)


# ── Variant 7: GEO — hexagonal nodes connected by lines ─────────────────

def _hexagonal(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("process_title", "Our Process"), theme=t)
    steps = _get_steps(c)
    p = t.palette

    n = len(steps)
    total_w = SLIDE_W - 2 * MARGIN
    step_w = total_w // n
    hex_y = content_top + 1100000
    hex_size = 350000

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if t.ux_style.dark_mode else t.secondary

    for i, (title, desc) in enumerate(steps):
        accent = p[i % len(p)]
        cx = MARGIN + i * step_w + step_w // 2

        # Connector line to next hex
        if i < n - 1:
            next_cx = MARGIN + (i + 1) * step_w + step_w // 2
            add_line(slide, cx + hex_size, hex_y, next_cx - hex_size, hex_y,
                     color=tint(t.accent, 0.3 if t.ux_style.dark_mode else 0.6), width=1.5)

        # Hexagon node
        add_hexagon(slide, cx, hex_y, hex_size, fill_color=accent)

        # Number inside hexagon
        add_textbox(slide, str(i + 1), cx - hex_size, hex_y - hex_size,
                    hex_size * 2, hex_size * 2,
                    font_size=18, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Title below hex
        add_textbox(slide, title, cx - step_w // 2 + 20000, hex_y + hex_size + 200000,
                    step_w - 40000, 350000,
                    font_size=13, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, desc, cx - step_w // 2 + 20000, hex_y + hex_size + 600000,
                    step_w - 40000, 400000,
                    font_size=11, color=sub_color, alignment=PP_ALIGN.CENTER)

        # Decorative diamond below description
        add_diamond(slide, cx, hex_y + hex_size + 1150000, 50000, fill_color=accent)


# ── Variant 8: EDITORIAL — clean numbered steps with thin rules ──────────

def _editorial_flow(slide, t, c):
    content_top = add_slide_title(slide, c.get("process_title", "Our Process"), theme=t)
    steps = _get_steps(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(steps)
    available_h = FOOTER_TOP - content_top - 400000
    step_h = available_h // n

    for i, (title, desc) in enumerate(steps):
        step_top = content_top + 200000 + i * step_h

        # Thin horizontal rule
        add_line(slide, m, step_top, SLIDE_W - m, step_top,
                 color=tint(t.secondary, 0.8), width=0.5)

        # Short accent mark at start of rule
        add_line(slide, m, step_top, m + 500000, step_top,
                 color=t.accent, width=1.5)

        # Number — left column
        num_w = 800000
        add_textbox(slide, f"0{i + 1}", m, step_top + 100000, num_w, 400000,
                    font_size=24, bold=True, color=t.accent, alignment=PP_ALIGN.LEFT)

        # Title — middle column
        title_left = m + num_w + 100000
        title_w = 3500000
        add_textbox(slide, title, title_left, step_top + 100000, title_w, 400000,
                    font_size=14, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT)

        # Description — right column
        desc_left = title_left + title_w + 200000
        desc_w = SLIDE_W - m - desc_left
        add_textbox(slide, desc, desc_left, step_top + 100000, desc_w, 400000,
                    font_size=12, color=t.secondary, alignment=PP_ALIGN.LEFT)

    # Bottom rule
    add_line(slide, m, FOOTER_TOP - 300000, SLIDE_W - m, FOOTER_TOP - 300000,
             color=tint(t.secondary, 0.8), width=0.5)


# ── Variant 9: GRADIENT — circles with gradient-like colors fading ───────

def _gradient_circles(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("process_title", "Our Process"), theme=t)
    steps = _get_steps(c)

    n = len(steps)
    total_w = SLIDE_W - 2 * MARGIN
    step_w = total_w // n
    circle_y = content_top + 1100000
    circle_r = 360000

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if t.ux_style.dark_mode else t.secondary

    for i, (title, desc) in enumerate(steps):
        cx = MARGIN + i * step_w + step_w // 2
        # Progressive tint from accent to lighter
        factor = i * 0.15
        circle_color = tint(t.accent, factor)

        # Connector line
        if i < n - 1:
            next_cx = MARGIN + (i + 1) * step_w + step_w // 2
            add_line(slide, cx + circle_r, circle_y, next_cx - circle_r, circle_y,
                     color=tint(t.accent, 0.6), width=1)

        # Outer glow circle
        add_circle(slide, cx, circle_y, circle_r + 60000,
                   fill_color=tint(circle_color, 0.85))

        # Main circle
        add_circle(slide, cx, circle_y, circle_r, fill_color=circle_color)

        # Number
        add_textbox(slide, str(i + 1), cx - 150000, circle_y - 200000, 300000, 250000,
                    font_size=18, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER)

        # Title inside circle
        add_textbox(slide, title, cx - circle_r + 40000, circle_y + 50000,
                    circle_r * 2 - 80000, 250000,
                    font_size=11, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER)

        # Description below
        add_textbox(slide, desc, cx - step_w // 2 + 20000, circle_y + circle_r + 200000,
                    step_w - 40000, 400000,
                    font_size=11, color=sub_color, alignment=PP_ALIGN.CENTER)


# ── Variant 10: RETRO — rounded badge nodes with dotted connectors ───────

def _retro_badges(slide, t, c):
    add_background(slide, t.light_bg)
    content_top = add_slide_title(slide, c.get("process_title", "Our Process"), theme=t)
    steps = _get_steps(c)
    p = t.palette

    n = len(steps)
    total_w = SLIDE_W - 2 * MARGIN
    step_w = total_w // n
    badge_y = content_top + 1000000
    badge_r = 320000

    for i, (title, desc) in enumerate(steps):
        accent = p[i % len(p)]
        cx = MARGIN + i * step_w + step_w // 2

        # Dotted connector to next badge
        if i < n - 1:
            next_cx = MARGIN + (i + 1) * step_w + step_w // 2
            add_line(slide, cx + badge_r + 40000, badge_y,
                     next_cx - badge_r - 40000, badge_y,
                     color=t.accent, width=1.5, dash="dash")

        # Outer border circle
        add_circle(slide, cx, badge_y, badge_r,
                   fill_color="#FFFFFF", line_color=accent, line_width=2.5)

        # Inner border circle
        add_circle(slide, cx, badge_y, badge_r - 50000,
                   fill_color="#FFFFFF", line_color=accent, line_width=0.75)

        # Number
        add_textbox(slide, str(i + 1), cx - 150000, badge_y - 200000, 300000, 250000,
                    font_size=18, bold=True, color=accent, alignment=PP_ALIGN.CENTER)

        # Title inside badge
        add_textbox(slide, title, cx - badge_r + 40000, badge_y + 50000,
                    badge_r * 2 - 80000, 250000,
                    font_size=11, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)

        # Description below badge
        add_textbox(slide, desc, cx - step_w // 2 + 20000, badge_y + badge_r + 200000,
                    step_w - 40000, 400000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.CENTER)

        # Decorative dots below description
        dot_y = badge_y + badge_r + 700000
        for d in range(3):
            add_circle(slide, cx - 60000 + d * 60000, dot_y, 15000, fill_color=accent)


# ── Variant 11: MAGAZINE — overlapping circles of different sizes ────────

def _creative_bubbles(slide, t, c):
    content_top = add_slide_title(slide, c.get("process_title", "Our Process"), theme=t)
    steps = _get_steps(c)
    p = t.palette

    n = len(steps)
    total_w = SLIDE_W - 2 * MARGIN
    step_w = total_w // n

    # Varying circle sizes for visual interest
    radii = [420000, 350000, 400000, 360000, 380000]
    y_offsets = [0, 150000, -100000, 200000, -50000]
    base_y = content_top + 1200000

    for i, (title, desc) in enumerate(steps):
        accent = p[i % len(p)]
        cx = MARGIN + i * step_w + step_w // 2
        r = radii[i % len(radii)]
        cy = base_y + y_offsets[i % len(y_offsets)]

        # Large background bubble (overlaps with neighbors)
        add_circle(slide, cx, cy, r + 80000, fill_color=tint(accent, 0.85))

        # Main bubble
        add_circle(slide, cx, cy, r, fill_color=accent)

        # Number — top area of circle
        add_textbox(slide, f"{i + 1}", cx - r + 40000, cy - r + 100000,
                    r * 2 - 80000, 350000,
                    font_size=24, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER)

        # Title — center of circle
        add_textbox(slide, title, cx - r + 40000, cy - 80000,
                    r * 2 - 80000, 300000,
                    font_size=12, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER)

        # Description below circle
        add_textbox(slide, desc, cx - step_w // 2 + 20000, cy + r + 150000,
                    step_w - 40000, 400000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.CENTER)


# ── Variant 12: SCHOLARLY — white bg, numbered paragraph list ────────────

def _scholarly_numbered(slide, t, c):
    add_background(slide, "#FFFFFF")
    content_top = add_slide_title(slide, c.get("process_title", "Methodology"), theme=t)
    steps = _get_steps(c)

    n = len(steps)
    m = MARGIN
    available_h = FOOTER_TOP - content_top - 400000
    step_h = available_h // n

    for i, (title, desc) in enumerate(steps):
        step_top = content_top + 200000 + i * step_h

        # Thin horizontal rule between steps
        add_line(slide, m, step_top, SLIDE_W - m, step_top,
                 color=tint(t.secondary, 0.75), width=0.5)

        # Step number and title as paragraph text
        add_textbox(slide, f"Step {i + 1}.", m + 100000, step_top + 80000,
                    1200000, 350000,
                    font_size=13, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT)

        add_textbox(slide, title, m + 1300000, step_top + 80000,
                    3500000, 350000,
                    font_size=13, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT)

        # Description as indented paragraph body
        add_textbox(slide, desc, m + 1300000, step_top + 420000,
                    SLIDE_W - m - 1300000 - m, 350000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.LEFT)

    # Bottom rule
    add_line(slide, m, FOOTER_TOP - 300000, SLIDE_W - m, FOOTER_TOP - 300000,
             color=tint(t.secondary, 0.75), width=0.5)


# ── Variant 13: LABORATORY — dark bg, rectangles with accent arrows ──────

def _laboratory_flow(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("process_title", "Process Flow").upper(), theme=t)
    steps = _get_steps(c)
    p = t.palette

    n = len(steps)
    m = MARGIN
    total_w = SLIDE_W - 2 * m
    gap = 150000
    box_w = (total_w - (n - 1) * gap) // n
    box_h = 1600000
    box_top = content_top + 700000

    for i, (title, desc) in enumerate(steps):
        accent = p[i % len(p)]
        left = m + i * (box_w + gap)
        box_fill = tint(t.primary, 0.1)

        # Dark rectangle
        add_rect(slide, left, box_top, box_w, box_h,
                 fill_color=box_fill, line_color=tint(t.primary, 0.2),
                 line_width=0.5, corner_radius=40000)

        # Accent top bar
        add_rect(slide, left + 1, box_top, box_w - 2, 40000, fill_color=accent)

        # Monospace-style step number
        add_textbox(slide, f"[{i + 1:02d}]", left + 60000, box_top + 150000,
                    box_w - 120000, 300000,
                    font_size=16, bold=True, color=accent, alignment=PP_ALIGN.LEFT)

        # Title
        add_textbox(slide, title.upper(), left + 60000, box_top + 500000,
                    box_w - 120000, 350000,
                    font_size=12, bold=True, color="#FFFFFF", alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, left + 60000, box_top + 900000,
                    box_w - 120000, 400000,
                    font_size=10, color=tint(t.primary, 0.45), alignment=PP_ALIGN.LEFT)

        # Accent-colored arrow to next step
        if i < n - 1:
            arrow_left = left + box_w + 10000
            arrow_w = gap - 20000
            arrow_y = box_top + box_h // 2
            add_line(slide, arrow_left, arrow_y, arrow_left + arrow_w, arrow_y,
                     color=accent, width=2)
            # Arrowhead triangle (small right-pointing indicator)
            add_textbox(slide, "\u25B6", arrow_left + arrow_w - 120000, arrow_y - 100000,
                        120000, 200000,
                        font_size=10, color=accent, alignment=PP_ALIGN.CENTER,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Decorative bottom line
    add_line(slide, m, box_top + box_h + 300000, SLIDE_W - m, box_top + box_h + 300000,
             color=tint(t.primary, 0.15), width=0.5)


# ── Variant 14: DASHBOARD — compact circles on line, description below ───

def _dashboard_steps(slide, t, c):
    add_background(slide, "#FFFFFF")
    content_top = add_slide_title(slide, c.get("process_title", "Our Process"), theme=t)
    steps = _get_steps(c)
    p = t.palette

    # Header accent band
    band_h = 45000
    add_rect(slide, 0, content_top, SLIDE_W, band_h, fill_color=t.accent)

    n = len(steps)
    m = int(MARGIN * 0.8)
    total_w = SLIDE_W - 2 * m
    step_w = total_w // n
    circle_r = 180000
    line_y = content_top + band_h + 700000

    # Horizontal connecting line
    first_cx = m + step_w // 2
    last_cx = m + (n - 1) * step_w + step_w // 2
    add_line(slide, first_cx, line_y, last_cx, line_y,
             color=tint(t.secondary, 0.6), width=1.5)

    for i, (title, desc) in enumerate(steps):
        accent = p[i % len(p)]
        cx = m + i * step_w + step_w // 2

        # Filled numbered circle
        add_circle(slide, cx, line_y, circle_r, fill_color=accent)
        add_textbox(slide, str(i + 1), cx - circle_r, line_y - circle_r,
                    circle_r * 2, circle_r * 2,
                    font_size=16, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Title below circle
        add_textbox(slide, title, cx - step_w // 2 + 20000, line_y + circle_r + 150000,
                    step_w - 40000, 300000,
                    font_size=12, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)

        # Description below title
        add_textbox(slide, desc, cx - step_w // 2 + 20000, line_y + circle_r + 480000,
                    step_w - 40000, 350000,
                    font_size=10, color=t.secondary, alignment=PP_ALIGN.CENTER)
