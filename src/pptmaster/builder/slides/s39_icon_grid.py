"""Slide 39 — Icon Grid: 14 unique visual variants dispatched by UX style.

Displays 4-6 capability/feature items in a grid, each with real PNG icons
(from the generated icon set) or colored circle fallback placeholders.
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP,
    col_span, card_positions,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_circle, add_line, add_hexagon,
    add_icon_image, add_gold_accent_line, add_slide_title,
    add_dark_bg, add_background, add_styled_card,
    draw_icon_or_placeholder, resolve_icon_path,
)
from pptmaster.assets.color_utils import tint, shade


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    s = t.ux_style
    c = t.content

    dispatch = {
        "card-icons": _card_icons,
        "list-icons": _list_icons,
        "bold-cards": _bold_cards,
        "floating-icons": _floating_icons,
        "dark-icons": _dark_icons,
        "split-icons": _split_icons,
        "hex-icons": _hex_icons,
        "editorial-icons": _editorial_icons,
        "gradient-icons": _gradient_icons,
        "retro-icons": _retro_icons,
        "creative-icons": _creative_icons,
        "scholarly-icons": _scholarly_icons,
        "laboratory-icons": _laboratory_icons,
        "dashboard-icons": _dashboard_icons,
    }
    builder = dispatch.get(s.icon_grid, _card_icons)
    builder(slide, t, c)


# ---------------------------------------------------------------------------
# Data helpers
# ---------------------------------------------------------------------------

_DEFAULT_ITEMS = [
    ("chart", "Analytics", "Real-time data insights and reporting"),
    ("shield", "Security", "Enterprise-grade protection"),
    ("globe", "Global Reach", "Operations in 40+ countries"),
    ("lightning", "Performance", "Sub-50ms response times"),
    ("users", "Team", "2,500+ professionals worldwide"),
    ("trophy", "Awards", "Industry recognition and accolades"),
]


def _normalize_item(item):
    """Safely convert any item shape to (icon, title, desc) tuple."""
    if isinstance(item, (list, tuple)):
        parts = list(item)
        return (
            parts[0] if len(parts) > 0 else "star",
            str(parts[1]) if len(parts) > 1 else "",
            str(parts[2]) if len(parts) > 2 else "",
        )
    return ("star", str(item), "")


def _get_items(c):
    """Return icon-grid items from content dict (4-6 items), normalized to (icon, title, desc)."""
    raw = c.get("icon_grid_items", _DEFAULT_ITEMS)
    return [_normalize_item(it) for it in raw]


def _get_title(c):
    return c.get("icon_grid_title", "Key Capabilities")


def _draw_icon_placeholder(slide, cx, cy, radius, color, icon_keyword):
    """Place a real PNG icon if available, else draw colored circle + letter."""
    draw_icon_or_placeholder(slide, icon_keyword, cx, cy, radius, color)


def _grid_layout(n):
    """Return (n_cols, n_rows) for *n* items — prefer wider grids."""
    if n <= 3:
        return n, 1
    if n <= 4:
        return 2, 2
    if n <= 6:
        return 3, 2
    return 3, 3  # safety cap


# ===================================================================
# Variant 1  CLASSIC — card-icons
# Cards arranged in a 2x3 or 3x2 grid with colored circle + title + desc
# ===================================================================

def _card_icons(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, _get_title(c), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    n = min(len(items), 6)
    n_cols, n_rows = _grid_layout(n)
    positions = card_positions(n_cols, n_rows, top=content_top + 80000,
                               gap=200000, row_gap=200000)

    for i, (icon, title, desc) in enumerate(items[:n]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        add_styled_card(slide, left, top, width, height, theme=t, accent_color=accent)

        # Icon circle
        circle_r = 200000
        cx = left + width // 2
        cy = top + 350000
        _draw_icon_placeholder(slide, cx, cy, circle_r, accent, icon)

        # Title
        add_textbox(slide, title, left + 80000, cy + circle_r + 120000,
                    width - 160000, 350000,
                    font_size=15, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, desc, left + 100000, cy + circle_r + 500000,
                    width - 200000, height - (cy - top) - circle_r - 600000,
                    font_size=11, color=sub_color,
                    alignment=PP_ALIGN.CENTER)


# ===================================================================
# Variant 2  MINIMAL — list-icons
# Clean horizontal list: icon circles left-aligned with text beside them
# ===================================================================

def _list_icons(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, _get_title(c), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    n = min(len(items), 6)
    available_h = FOOTER_TOP - content_top - 200000
    row_h = available_h // n
    circle_r = 160000

    for i, (icon, title, desc) in enumerate(items[:n]):
        accent = p[i % len(p)]
        row_top = content_top + 100000 + i * row_h
        cx = m + circle_r + 40000
        cy = row_top + row_h // 2

        _draw_icon_placeholder(slide, cx, cy, circle_r, accent, icon)

        text_left = cx + circle_r + 200000
        text_w = SLIDE_W - text_left - m

        add_textbox(slide, title, text_left, cy - 220000, text_w, 280000,
                    font_size=14, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT)
        add_textbox(slide, desc, text_left, cy + 60000, text_w, 250000,
                    font_size=11, color=sub_color,
                    alignment=PP_ALIGN.LEFT)

        # Thin separator below (except last)
        if i < n - 1:
            sep_y = row_top + row_h - 10000
            add_line(slide, m, sep_y, SLIDE_W - m, sep_y,
                     color=tint(t.secondary, 0.8), width=0.5)


# ===================================================================
# Variant 3  BOLD — bold-cards
# Bold colored cards, large icon circles, uppercase titles
# ===================================================================

def _bold_cards(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, _get_title(c).upper(), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    n = min(len(items), 6)
    n_cols, n_rows = _grid_layout(n)
    positions = card_positions(n_cols, n_rows, top=content_top + 80000,
                               gap=120000, row_gap=120000)

    for i, (icon, title, desc) in enumerate(items[:n]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        # Bold color fill card
        add_rect(slide, left, top, width, height, fill_color=accent)

        # Large icon (real PNG or white circle fallback)
        circle_r = 240000
        cx = left + width // 2
        cy = top + 420000
        _path = resolve_icon_path(icon)
        if _path:
            isize = int(circle_r * 1.6)
            add_icon_image(slide, _path, cx - isize // 2, cy - isize // 2, isize, isize)
        else:
            add_circle(slide, cx, cy, circle_r, fill_color="#FFFFFF")
            add_textbox(slide, icon[0].upper(), cx - circle_r, cy - circle_r,
                        circle_r * 2, circle_r * 2,
                        font_size=24, bold=True, color=accent,
                        alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Uppercase title
        add_textbox(slide, title.upper(), left + 60000, cy + circle_r + 150000,
                    width - 120000, 380000,
                    font_size=16, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, desc, left + 80000, cy + circle_r + 560000,
                    width - 160000, height - (cy - top) - circle_r - 650000,
                    font_size=11, color=tint("#FFFFFF", 0.15),
                    alignment=PP_ALIGN.CENTER)


# ===================================================================
# Variant 4  ELEVATED — floating-icons
# Floating cards with shadow, icon circles, pill-shaped cards
# ===================================================================

def _floating_icons(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, _get_title(c), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    n = min(len(items), 6)
    n_cols, n_rows = _grid_layout(n)
    positions = card_positions(n_cols, n_rows, top=content_top + 100000,
                               gap=250000, row_gap=250000)

    for i, (icon, title, desc) in enumerate(items[:n]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        # Floating card (heavy shadow via styled card)
        add_styled_card(slide, left, top, width, height, theme=t, accent_color=accent)

        # Double-ring icon
        circle_r = 180000
        cx = left + width // 2
        cy = top + 380000
        add_circle(slide, cx, cy, circle_r + 50000, fill_color=tint(accent, 0.85))
        _draw_icon_placeholder(slide, cx, cy, circle_r, accent, icon)

        # Title
        add_textbox(slide, title, left + 60000, cy + circle_r + 140000,
                    width - 120000, 350000,
                    font_size=14, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

        # Description
        desc_top = cy + circle_r + 520000
        add_textbox(slide, desc, left + 80000, desc_top,
                    width - 160000, height - (desc_top - top) - 100000,
                    font_size=11, color=sub_color,
                    alignment=PP_ALIGN.CENTER)


# ===================================================================
# Variant 5  DARK — dark-icons
# Dark background with glowing icon circles
# ===================================================================

def _dark_icons(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, _get_title(c).upper(), theme=t)
    items = _get_items(c)
    p = t.palette

    n = min(len(items), 6)
    n_cols, n_rows = _grid_layout(n)
    positions = card_positions(n_cols, n_rows, top=content_top + 100000,
                               gap=200000, row_gap=200000)

    for i, (icon, title, desc) in enumerate(items[:n]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]
        card_fill = tint(t.primary, 0.08)

        # Dark card with thin border
        add_rect(slide, left, top, width, height, fill_color=card_fill,
                 line_color=tint(t.primary, 0.2), line_width=0.75,
                 corner_radius=60000)

        # Neon accent bar top
        add_rect(slide, left + 1, top, width - 2, 50000, fill_color=accent)

        # Glow ring + icon circle
        circle_r = 190000
        cx = left + width // 2
        cy = top + 420000
        add_circle(slide, cx, cy, circle_r + 70000, fill_color=tint(accent, 0.12))
        add_circle(slide, cx, cy, circle_r + 35000, fill_color=tint(accent, 0.06))
        _draw_icon_placeholder(slide, cx, cy, circle_r, accent, icon)

        # Title (neon-bright)
        add_textbox(slide, title, left + 60000, cy + circle_r + 160000,
                    width - 120000, 350000,
                    font_size=14, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER)

        # Description
        desc_top = cy + circle_r + 540000
        add_textbox(slide, desc, left + 80000, desc_top,
                    width - 160000, height - (desc_top - top) - 100000,
                    font_size=11, color=tint(t.primary, 0.5),
                    alignment=PP_ALIGN.CENTER)


# ===================================================================
# Variant 6  SPLIT — split-icons
# Icons grid on left (2x3), summary text panel on right
# ===================================================================

def _split_icons(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, _get_title(c), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    # Layout: 60% left grid, 40% right summary
    split_x = MARGIN + int((SLIDE_W - 2 * MARGIN) * 0.58)
    left_w = split_x - MARGIN - 100000
    right_left = split_x + 100000
    right_w = SLIDE_W - MARGIN - right_left

    # Vertical divider
    add_line(slide, split_x, content_top + 100000, split_x, FOOTER_TOP - 200000,
             color=t.accent, width=2)

    # Left: icon grid (2 cols)
    n = min(len(items), 6)
    icon_cols = 2
    icon_rows = (n + 1) // 2
    available_h = FOOTER_TOP - content_top - 300000
    cell_w = left_w // icon_cols
    cell_h = available_h // icon_rows
    circle_r = 140000

    for i, (icon, title, desc) in enumerate(items[:n]):
        col = i % icon_cols
        row = i // icon_cols
        accent = p[i % len(p)]

        cell_left = MARGIN + col * cell_w
        cell_top = content_top + 150000 + row * cell_h
        cx = cell_left + cell_w // 2
        cy = cell_top + 200000

        _draw_icon_placeholder(slide, cx, cy, circle_r, accent, icon)

        add_textbox(slide, title, cell_left + 20000, cy + circle_r + 80000,
                    cell_w - 40000, 280000,
                    font_size=12, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, cell_left + 30000, cy + circle_r + 380000,
                    cell_w - 60000, 350000,
                    font_size=10, color=sub_color,
                    alignment=PP_ALIGN.CENTER)

    # Right: summary panel
    summary_title = c.get("icon_grid_summary_title", "Our Strengths")
    summary_text = c.get("icon_grid_summary_text",
                         "We combine deep expertise with cutting-edge technology "
                         "to deliver exceptional results for our clients.")

    add_textbox(slide, summary_title, right_left, content_top + 400000,
                right_w, 450000,
                font_size=22, bold=True, color=text_color,
                alignment=PP_ALIGN.LEFT)
    add_gold_accent_line(slide, right_left, content_top + 900000,
                         600000, color=t.accent)
    add_textbox(slide, summary_text, right_left, content_top + 1100000,
                right_w, 2000000,
                font_size=13, color=sub_color,
                alignment=PP_ALIGN.LEFT)

    # Stats below summary
    stats = c.get("icon_grid_stats", [("40+", "Countries"), ("2.5K+", "Team")])
    stat_top = content_top + 3400000
    stat_w = right_w // len(stats)
    for j, (val, lbl) in enumerate(stats[:3]):
        sx = right_left + j * stat_w
        add_textbox(slide, val, sx, stat_top, stat_w, 500000,
                    font_size=28, bold=True, color=t.accent,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, sx, stat_top + 480000, stat_w, 300000,
                    font_size=10, color=sub_color,
                    alignment=PP_ALIGN.CENTER)


# ===================================================================
# Variant 7  GEO — hex-icons
# Hexagonal icon placeholders in a grid
# ===================================================================

def _hex_icons(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, _get_title(c), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if s.dark_mode else t.secondary

    n = min(len(items), 6)
    n_cols, n_rows = _grid_layout(n)
    positions = card_positions(n_cols, n_rows, top=content_top + 80000,
                               gap=180000, row_gap=180000)

    for i, (icon, title, desc) in enumerate(items[:n]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        # Card with bottom accent (GEO style)
        add_styled_card(slide, left, top, width, height, theme=t, accent_color=accent)

        # Hexagonal icon or real PNG icon
        hex_size = 200000
        cx = left + width // 2
        cy = top + 380000
        _path = resolve_icon_path(icon)
        if _path:
            isize = int(hex_size * 1.6)
            add_icon_image(slide, _path, cx - isize // 2, cy - isize // 2, isize, isize)
        else:
            add_hexagon(slide, cx, cy, hex_size, fill_color=accent)
            add_textbox(slide, icon[0].upper(), cx - hex_size, cy - hex_size,
                        hex_size * 2, hex_size * 2,
                        font_size=18, bold=True, color="#FFFFFF",
                        alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Title
        add_textbox(slide, title.upper(), left + 60000, cy + hex_size + 160000,
                    width - 120000, 350000,
                    font_size=13, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

        # Description
        desc_top = cy + hex_size + 540000
        add_textbox(slide, desc, left + 80000, desc_top,
                    width - 160000, height - (desc_top - top) - 80000,
                    font_size=11, color=sub_color,
                    alignment=PP_ALIGN.CENTER)

        # Decorative bottom line
        add_line(slide, left + 100000, top + height - 120000,
                 left + width - 100000, top + height - 120000,
                 color=accent, width=0.75)


# ===================================================================
# Variant 8  EDITORIAL — editorial-icons
# Elegant two-column layout with thin rules and serif-feel spacing
# ===================================================================

def _editorial_icons(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, _get_title(c), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    n = min(len(items), 6)
    col_w = (SLIDE_W - 2 * m - 400000) // 2
    row_h = 1350000
    circle_r = 140000

    for i, (icon, title, desc) in enumerate(items[:n]):
        col = i % 2
        row = i // 2
        accent = p[i % len(p)]

        cell_left = m + col * (col_w + 400000)
        cell_top = content_top + 200000 + row * (row_h + 100000)

        # Thin top rule
        add_line(slide, cell_left, cell_top, cell_left + col_w, cell_top,
                 color=tint(t.secondary, 0.75), width=0.5)
        # Short accent mark on left
        add_line(slide, cell_left, cell_top, cell_left + 400000, cell_top,
                 color=t.accent, width=1.5)

        # Icon circle, left-aligned
        cx = cell_left + circle_r + 40000
        cy = cell_top + row_h // 2
        _draw_icon_placeholder(slide, cx, cy, circle_r, accent, icon)

        # Title and desc beside icon
        text_left = cx + circle_r + 180000
        text_w = col_w - (text_left - cell_left) - 40000

        add_textbox(slide, title, text_left, cy - 240000, text_w, 280000,
                    font_size=14, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT)
        add_textbox(slide, desc, text_left, cy + 60000, text_w, 400000,
                    font_size=11, color=sub_color,
                    alignment=PP_ALIGN.LEFT)

    # Bottom rule
    add_line(slide, m, FOOTER_TOP - 300000, SLIDE_W - m, FOOTER_TOP - 300000,
             color=tint(t.secondary, 0.8), width=0.5)


# ===================================================================
# Variant 9  GRADIENT — gradient-icons
# Very rounded cards with tinted background circles behind icons
# ===================================================================

def _gradient_icons(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, _get_title(c), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    n = min(len(items), 6)
    n_cols, n_rows = _grid_layout(n)
    positions = card_positions(n_cols, n_rows, top=content_top + 100000,
                               gap=250000, row_gap=250000)

    for i, (icon, title, desc) in enumerate(items[:n]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        # Soft rounded card (no accent, per gradient style)
        add_styled_card(slide, left, top, width, height, theme=t)

        # Large tinted bg circle behind icon
        cx = left + width // 2
        cy = top + 450000
        add_circle(slide, cx, cy, 320000, fill_color=tint(accent, 0.88))

        # Icon circle on top
        circle_r = 180000
        _draw_icon_placeholder(slide, cx, cy, circle_r, accent, icon)

        # Title
        add_textbox(slide, title, left + 60000, cy + 350000,
                    width - 120000, 350000,
                    font_size=14, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

        # Description
        desc_top = cy + 730000
        add_textbox(slide, desc, left + 80000, desc_top,
                    width - 160000, height - (desc_top - top) - 100000,
                    font_size=11, color=sub_color,
                    alignment=PP_ALIGN.CENTER)

        # Subtle accent dot below description
        dot_y = top + height - 200000
        add_circle(slide, cx, dot_y, 30000, fill_color=tint(accent, 0.5))


# ===================================================================
# Variant 10  RETRO — retro-icons
# Vintage styled icon badges with double borders and decorative elements
# ===================================================================

def _retro_icons(slide, t, c):
    add_background(slide, t.light_bg)
    content_top = add_slide_title(slide, _get_title(c), theme=t)
    items = _get_items(c)
    p = t.palette

    n = min(len(items), 6)
    n_cols, n_rows = _grid_layout(n)
    positions = card_positions(n_cols, n_rows, top=content_top + 100000,
                               gap=220000, row_gap=220000)

    for i, (icon, title, desc) in enumerate(items[:n]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        # Double-bordered card
        add_rect(slide, left, top, width, height, fill_color="#FFFFFF",
                 line_color=accent, line_width=2.0, corner_radius=100000)
        add_rect(slide, left + 40000, top + 40000, width - 80000, height - 80000,
                 line_color=accent, line_width=0.75, corner_radius=80000)

        # Decorative corner dots
        add_circle(slide, left + 120000, top + 120000, 35000, fill_color=accent)
        add_circle(slide, left + width - 120000, top + 120000, 35000, fill_color=accent)

        # Badge-style icon circle with ring
        cx = left + width // 2
        cy = top + 420000
        add_circle(slide, cx, cy, 220000, fill_color=tint(accent, 0.85),
                   line_color=accent, line_width=2.0)
        _draw_icon_placeholder(slide, cx, cy, 160000, accent, icon)

        # Title with decorative dashes
        add_textbox(slide, f"\u2014 {title} \u2014", left + 60000, cy + 260000,
                    width - 120000, 350000,
                    font_size=13, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER)

        # Decorative rule below title
        rule_left = left + width // 4
        rule_y = cy + 620000
        add_line(slide, rule_left, rule_y, rule_left + width // 2, rule_y,
                 color=accent, width=1.0)

        # Description
        add_textbox(slide, desc, left + 80000, rule_y + 80000,
                    width - 160000, height - (rule_y - top) - 150000,
                    font_size=11, color=t.secondary,
                    alignment=PP_ALIGN.CENTER)


# ===================================================================
# Variant 11  MAGAZINE — creative-icons
# Creative asymmetric layout: alternating large/small icon panels
# ===================================================================

def _creative_icons(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, _get_title(c), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    n = min(len(items), 6)
    available_h = FOOTER_TOP - content_top - 200000
    content_w = SLIDE_W - 2 * MARGIN

    # Asymmetric: first item large (left 40%), rest stacked right (60%)
    large_w = int(content_w * 0.38)
    small_left = MARGIN + large_w + 150000
    small_w = SLIDE_W - MARGIN - small_left
    large_h = available_h - 50000
    small_count = n - 1
    small_gap = 80000
    small_h = (available_h - small_gap * max(small_count - 1, 0)) // max(small_count, 1)

    # Large feature item
    if n > 0:
        icon0, title0, desc0 = items[0]
        accent0 = p[0]
        card_top = content_top + 100000

        add_rect(slide, MARGIN, card_top, large_w, large_h, fill_color=accent0)

        # Large icon
        cx = MARGIN + large_w // 2
        cy = card_top + large_h // 3
        large_r = 280000
        add_circle(slide, cx, cy, large_r, fill_color="#FFFFFF")
        add_textbox(slide, icon0[0].upper(), cx - large_r, cy - large_r,
                    large_r * 2, large_r * 2,
                    font_size=32, bold=True, color=accent0,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        add_textbox(slide, title0, MARGIN + 100000, cy + large_r + 150000,
                    large_w - 200000, 400000,
                    font_size=18, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc0, MARGIN + 120000, cy + large_r + 580000,
                    large_w - 240000, 600000,
                    font_size=12, color=tint("#FFFFFF", 0.15),
                    alignment=PP_ALIGN.CENTER)

    # Small items stacked on right
    for j, (icon, title, desc) in enumerate(items[1:n]):
        accent = p[(j + 1) % len(p)]
        ry = content_top + 100000 + j * (small_h + small_gap)

        add_rect(slide, small_left, ry, small_w, small_h, fill_color=accent)

        # Small icon (real PNG or white circle fallback)
        sr = 120000
        scx = small_left + 200000
        scy = ry + small_h // 2
        _path = resolve_icon_path(icon)
        if _path:
            isize = int(sr * 1.6)
            add_icon_image(slide, _path, scx - isize // 2, scy - isize // 2, isize, isize)
        else:
            add_circle(slide, scx, scy, sr, fill_color="#FFFFFF")
            add_textbox(slide, icon[0].upper(), scx - sr, scy - sr, sr * 2, sr * 2,
                        font_size=14, bold=True, color=accent,
                        alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Title and desc beside
        t_left = scx + sr + 150000
        t_w = small_left + small_w - t_left - 60000
        add_textbox(slide, title, t_left, scy - 180000, t_w, 220000,
                    font_size=13, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT)
        add_textbox(slide, desc, t_left, scy + 50000, t_w, 220000,
                    font_size=10, color=tint("#FFFFFF", 0.15),
                    alignment=PP_ALIGN.LEFT)


# ===================================================================
# Variant 12  SCHOLARLY — scholarly-icons
# Clean numbered icon list with figure captions and thin rules
# ===================================================================

def _scholarly_icons(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, _get_title(c), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    n = min(len(items), 6)
    col_w = (SLIDE_W - 2 * m - 300000) // 3
    row_h = (FOOTER_TOP - content_top - 600000) // 2
    circle_r = 150000

    for i, (icon, title, desc) in enumerate(items[:n]):
        col = i % 3
        row = i // 3
        accent = p[i % len(p)]

        cell_left = m + col * (col_w + 150000)
        cell_top = content_top + 250000 + row * (row_h + 100000)

        # Numbered circle (scholarly uses numbers)
        cx = cell_left + col_w // 2
        cy = cell_top + 200000
        add_circle(slide, cx, cy, circle_r, fill_color=accent,
                   line_color=shade(accent, 0.2), line_width=1.0)
        add_textbox(slide, str(i + 1), cx - circle_r, cy - circle_r,
                    circle_r * 2, circle_r * 2,
                    font_size=18, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Title centered below
        add_textbox(slide, title, cell_left, cy + circle_r + 80000,
                    col_w, 280000,
                    font_size=13, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

        # Thin rule
        rule_y = cy + circle_r + 400000
        add_line(slide, cell_left + 60000, rule_y, cell_left + col_w - 60000, rule_y,
                 color=tint(t.secondary, 0.75), width=0.5)

        # Description (caption style)
        add_textbox(slide, desc, cell_left + 40000, rule_y + 60000,
                    col_w - 80000, row_h - (rule_y - cell_top) - 120000,
                    font_size=10, italic=True, color=sub_color,
                    alignment=PP_ALIGN.CENTER)

    # Figure label at bottom
    add_textbox(slide, f"Fig. {c.get('icon_grid_fig_num', '5')}. {_get_title(c)}",
                m, FOOTER_TOP - 350000, SLIDE_W - 2 * m, 250000,
                font_size=9, italic=True, color=sub_color,
                alignment=PP_ALIGN.LEFT)


# ===================================================================
# Variant 13  LABORATORY — laboratory-icons
# Dark bg, data-coded icon circles with monospace-feel numbering
# ===================================================================

def _laboratory_icons(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, _get_title(c), theme=t)
    items = _get_items(c)
    p = t.palette

    n = min(len(items), 6)
    n_cols, n_rows = _grid_layout(n)
    positions = card_positions(n_cols, n_rows, top=content_top + 100000,
                               gap=160000, row_gap=160000)

    for i, (icon, title, desc) in enumerate(items[:n]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]
        card_fill = tint(t.primary, 0.06)

        # Card with color-coded left border (lab style)
        add_rect(slide, left, top, width, height, fill_color=card_fill,
                 line_color=tint(t.primary, 0.15), line_width=0.75,
                 corner_radius=30000)
        # Left accent bar
        add_rect(slide, left, top + 40000, 50000, height - 80000,
                 fill_color=accent)

        # Specimen ID label at top-right
        add_textbox(slide, f"CAP-{i + 1:02d}", left + width - 500000, top + 60000,
                    440000, 200000,
                    font_size=8, bold=True, color=tint(t.primary, 0.35),
                    alignment=PP_ALIGN.RIGHT)

        # Icon circle with data-ring
        circle_r = 170000
        cx = left + width // 2
        cy = top + 400000
        # Outer data ring
        add_circle(slide, cx, cy, circle_r + 40000, fill_color=tint(accent, 0.08),
                   line_color=accent, line_width=1.5)
        _draw_icon_placeholder(slide, cx, cy, circle_r, accent, icon)

        # Title
        add_textbox(slide, title, left + 80000, cy + circle_r + 120000,
                    width - 160000, 300000,
                    font_size=13, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER)

        # Thin separator
        sep_y = cy + circle_r + 450000
        add_line(slide, left + 120000, sep_y, left + width - 120000, sep_y,
                 color=tint(t.primary, 0.15), width=0.5)

        # Description
        add_textbox(slide, desc, left + 100000, sep_y + 60000,
                    width - 200000, height - (sep_y - top) - 120000,
                    font_size=10, color=tint(t.primary, 0.45),
                    alignment=PP_ALIGN.CENTER)

    # Grid overlay label
    add_textbox(slide, f"[CAPABILITY MATRIX  |  n={n}]",
                MARGIN, FOOTER_TOP - 280000,
                3000000, 200000,
                font_size=8, color=tint(t.primary, 0.3),
                alignment=PP_ALIGN.LEFT)


# ===================================================================
# Variant 14  DASHBOARD — dashboard-icons
# Dense compact icon tiles in a tight grid, minimal padding
# ===================================================================

def _dashboard_icons(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, _get_title(c), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    n = min(len(items), 6)
    n_cols, n_rows = _grid_layout(n)
    # Tight gaps for dashboard density
    positions = card_positions(n_cols, n_rows, top=content_top + 50000,
                               gap=100000, row_gap=100000)

    for i, (icon, title, desc) in enumerate(items[:n]):
        left, top, width, height = positions[i]
        accent = p[i % len(p)]

        # Compact tile card
        card_fill = "#FFFFFF" if not s.dark_mode else tint(t.primary, 0.1)
        add_rect(slide, left, top, width, height, fill_color=card_fill,
                 line_color=tint(t.secondary, 0.8) if not s.dark_mode
                 else tint(t.primary, 0.2),
                 line_width=0.5, corner_radius=50000)

        # Accent header strip
        strip_h = 55000
        add_rect(slide, left, top, width, strip_h, fill_color=accent)

        # Small icon circle, top-left corner
        circle_r = 130000
        cx = left + 200000
        cy = top + strip_h + 180000
        _draw_icon_placeholder(slide, cx, cy, circle_r, accent, icon)

        # Title to right of icon
        title_left = cx + circle_r + 100000
        title_w = left + width - title_left - 60000
        add_textbox(slide, title, title_left, cy - 140000, title_w, 280000,
                    font_size=12, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT)

        # Description below icon row, full width
        desc_top = top + strip_h + circle_r + 280000
        add_textbox(slide, desc, left + 80000, desc_top,
                    width - 160000, height - (desc_top - top) - 80000,
                    font_size=10, color=sub_color,
                    alignment=PP_ALIGN.LEFT)

        # Tiny status dot (bottom-right of tile)
        add_circle(slide, left + width - 130000, top + height - 130000,
                   25000, fill_color=accent)
