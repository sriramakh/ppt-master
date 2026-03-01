"""Slide 37 — Matrix Quadrant: 2x2 strategic matrix with axis labels, 14 visual variants."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_circle, add_line, add_diamond,
    add_gold_accent_line, add_slide_title, add_dark_bg, add_background,
    add_styled_card, add_color_cell,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    p = t.palette
    s = t.ux_style

    dispatch = {
        "quadrant-grid": _quadrant_grid,
        "labeled-grid": _labeled_grid,
        "color-blocks": _color_blocks,
        "floating-quad": _floating_quad,
        "dark-grid": _dark_grid,
        "split-quad": _split_quad,
        "angular-grid": _angular_grid,
        "editorial-grid": _editorial_grid,
        "gradient-grid": _gradient_grid,
        "retro-grid": _retro_grid,
        "creative-grid": _creative_grid,
        "scholarly-grid": _scholarly_grid,
        "laboratory-grid": _laboratory_grid,
        "dashboard-grid": _dashboard_grid,
    }
    builder = dispatch.get(s.matrix, _quadrant_grid)
    builder(slide, t, c)


# ── Data helpers ─────────────────────────────────────────────────────────

def _get_data(c):
    """Extract matrix content from content dict with defaults."""
    title = c.get("matrix_title", "Strategic Matrix")
    x_axis = c.get("matrix_x_axis", "Impact")
    y_axis = c.get("matrix_y_axis", "Effort")
    quadrants = c.get("matrix_quadrants", [
        ("Quick Wins", "High impact, low effort \u2014 prioritize these"),
        ("Major Projects", "High impact, high effort \u2014 plan carefully"),
        ("Fill-Ins", "Low impact, low effort \u2014 delegate or automate"),
        ("Thankless Tasks", "Low impact, high effort \u2014 reconsider"),
    ])
    return title, x_axis, y_axis, quadrants


def _quad_colors(t):
    """Return 4 accent colors from the theme palette."""
    p = t.palette
    return [p[i % len(p)] for i in range(4)]


# ── Variant 1: CLASSIC — 2x2 colored quadrants with axis labels ─────────

def _quadrant_grid(slide, t, c):
    add_dark_bg(slide, t)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title, theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    gap = int(120000 * s.gap_factor)
    axis_margin = 450000  # space for axis labels

    grid_left = m + axis_margin
    grid_top = content_y + 100000
    grid_w = SLIDE_W - 2 * m - axis_margin
    grid_h = FOOTER_TOP - grid_top - 350000
    quad_w = (grid_w - gap) // 2
    quad_h = (grid_h - gap) // 2

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if s.dark_mode else t.secondary

    # Draw 4 colored quadrants
    positions = [
        (0, 0), (1, 0),  # top-left, top-right
        (0, 1), (1, 1),  # bottom-left, bottom-right
    ]
    for i, (col, row) in enumerate(positions):
        qx = grid_left + col * (quad_w + gap)
        qy = grid_top + row * (quad_h + gap)
        bg = tint(colors[i], 0.25 if s.dark_mode else 0.82)

        add_rect(slide, qx, qy, quad_w, quad_h, fill_color=bg, corner_radius=60000)
        # Accent strip at top
        add_rect(slide, qx, qy, quad_w, 55000, fill_color=colors[i])
        # Label
        label, desc = quads[i] if i < len(quads) else ("", "")
        add_textbox(slide, label, qx + 150000, qy + 120000, quad_w - 300000, 350000,
                    font_size=16, bold=True, color=colors[i])
        add_textbox(slide, desc, qx + 150000, qy + 500000, quad_w - 300000, quad_h - 650000,
                    font_size=11, color=text_color)

    # Y-axis label (vertical text simulated as short textbox)
    y_label_x = m
    y_label_y = grid_top + grid_h // 2 - 200000
    add_textbox(slide, y_axis, y_label_x, y_label_y, axis_margin - 80000, 400000,
                font_size=12, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)
    # Y-axis arrow line
    add_line(slide, m + axis_margin - 120000, grid_top + grid_h,
             m + axis_margin - 120000, grid_top,
             color=sub_color, width=1.0)

    # X-axis label
    x_label_y = grid_top + grid_h + 60000
    add_textbox(slide, x_axis, grid_left, x_label_y, grid_w, 280000,
                font_size=12, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER)
    # X-axis arrow line
    add_line(slide, grid_left, grid_top + grid_h + 40000,
             grid_left + grid_w, grid_top + grid_h + 40000,
             color=sub_color, width=1.0)


# ── Variant 2: MINIMAL — clean thin-bordered grid with labels ───────────

def _labeled_grid(slide, t, c):
    add_dark_bg(slide, t)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title, theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    axis_margin = 500000

    grid_left = m + axis_margin
    grid_top = content_y + 150000
    grid_w = SLIDE_W - 2 * m - axis_margin
    grid_h = FOOTER_TOP - grid_top - 400000
    quad_w = grid_w // 2
    quad_h = grid_h // 2

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary
    line_color = tint(t.secondary, 0.7) if not s.dark_mode else tint(t.primary, 0.25)

    # Outer border
    add_rect(slide, grid_left, grid_top, grid_w, grid_h,
             line_color=line_color, line_width=0.5)
    # Cross lines
    mid_x = grid_left + grid_w // 2
    mid_y = grid_top + grid_h // 2
    add_line(slide, mid_x, grid_top, mid_x, grid_top + grid_h,
             color=line_color, width=0.5)
    add_line(slide, grid_left, mid_y, grid_left + grid_w, mid_y,
             color=line_color, width=0.5)

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, (col, row) in enumerate(positions):
        qx = grid_left + col * quad_w
        qy = grid_top + row * quad_h
        label, desc = quads[i] if i < len(quads) else ("", "")
        # Small color dot
        add_circle(slide, qx + 200000, qy + 160000, 50000, fill_color=colors[i])
        add_textbox(slide, label, qx + 300000, qy + 100000, quad_w - 400000, 300000,
                    font_size=13, bold=True, color=text_color)
        add_textbox(slide, desc, qx + 150000, qy + 420000, quad_w - 300000, quad_h - 520000,
                    font_size=10, color=sub_color)

    # Axis labels
    add_textbox(slide, y_axis, m, grid_top + grid_h // 2 - 180000,
                axis_margin - 120000, 360000,
                font_size=11, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, x_axis, grid_left, grid_top + grid_h + 80000,
                grid_w, 260000,
                font_size=11, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER)

    # Thin axis arrows
    add_line(slide, grid_left - 60000, grid_top + grid_h,
             grid_left - 60000, grid_top,
             color=sub_color, width=0.5)
    add_line(slide, grid_left, grid_top + grid_h + 60000,
             grid_left + grid_w, grid_top + grid_h + 60000,
             color=sub_color, width=0.5)


# ── Variant 3: BOLD — bold solid colored blocks with large text ──────────

def _color_blocks(slide, t, c):
    add_dark_bg(slide, t)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title.upper(), theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    axis_margin = 500000
    gap = 80000

    grid_left = m + axis_margin
    grid_top = content_y + 80000
    grid_w = SLIDE_W - 2 * m - axis_margin
    grid_h = FOOTER_TOP - grid_top - 350000
    quad_w = (grid_w - gap) // 2
    quad_h = (grid_h - gap) // 2

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, (col, row) in enumerate(positions):
        qx = grid_left + col * (quad_w + gap)
        qy = grid_top + row * (quad_h + gap)
        label, desc = quads[i] if i < len(quads) else ("", "")

        # Solid color block, sharp corners
        add_rect(slide, qx, qy, quad_w, quad_h, fill_color=colors[i])
        # Large bold label centered
        add_textbox(slide, label.upper(), qx + 100000, qy + quad_h // 2 - 400000,
                    quad_w - 200000, 500000,
                    font_size=22, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)
        # Description below
        add_textbox(slide, desc, qx + 120000, qy + quad_h // 2 + 100000,
                    quad_w - 240000, 400000,
                    font_size=10, color=tint("#FFFFFF", 0.15), alignment=PP_ALIGN.CENTER)

    # Thick axis labels
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    add_textbox(slide, y_axis.upper(), m, grid_top + grid_h // 2 - 200000,
                axis_margin - 100000, 400000,
                font_size=13, bold=True, color=text_color, alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)
    # Thick left line
    add_line(slide, grid_left - 100000, grid_top + grid_h,
             grid_left - 100000, grid_top,
             color=t.accent, width=3.0)

    add_textbox(slide, x_axis.upper(), grid_left, grid_top + grid_h + 60000,
                grid_w, 280000,
                font_size=13, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
    # Thick bottom line
    add_line(slide, grid_left, grid_top + grid_h + 40000,
             grid_left + grid_w, grid_top + grid_h + 40000,
             color=t.accent, width=3.0)


# ── Variant 4: ELEVATED — floating card quadrants with shadow ────────────

def _floating_quad(slide, t, c):
    add_dark_bg(slide, t)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title, theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    axis_margin = 500000
    gap = int(250000 * s.gap_factor)

    grid_left = m + axis_margin
    grid_top = content_y + 150000
    grid_w = SLIDE_W - 2 * m - axis_margin
    grid_h = FOOTER_TOP - grid_top - 400000
    quad_w = (grid_w - gap) // 2
    quad_h = (grid_h - gap) // 2

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, (col, row) in enumerate(positions):
        qx = grid_left + col * (quad_w + gap)
        qy = grid_top + row * (quad_h + gap)
        label, desc = quads[i] if i < len(quads) else ("", "")

        add_styled_card(slide, qx, qy, quad_w, quad_h,
                        theme=t, accent_color=colors[i])
        # Quadrant number badge
        badge_cx = qx + 200000
        badge_cy = qy + 200000
        add_circle(slide, badge_cx, badge_cy, 100000, fill_color=colors[i])
        add_textbox(slide, str(i + 1), badge_cx - 100000, badge_cy - 100000,
                    200000, 200000,
                    font_size=11, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        # Label
        add_textbox(slide, label, qx + 360000, qy + 130000, quad_w - 460000, 300000,
                    font_size=15, bold=True, color=text_color)
        # Description
        add_textbox(slide, desc, qx + 150000, qy + 500000, quad_w - 300000, quad_h - 650000,
                    font_size=11, color=sub_color)

    # Axis labels with accent dots
    add_circle(slide, m + 120000, grid_top + grid_h // 2, 40000, fill_color=t.accent)
    add_textbox(slide, y_axis, m, grid_top + grid_h // 2 + 80000,
                axis_margin - 80000, 300000,
                font_size=11, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER)
    add_line(slide, m + axis_margin - 100000, grid_top + grid_h,
             m + axis_margin - 100000, grid_top,
             color=tint(sub_color, 0.5), width=0.75)

    add_circle(slide, grid_left + grid_w // 2, grid_top + grid_h + 120000,
               40000, fill_color=t.accent)
    add_textbox(slide, x_axis, grid_left, grid_top + grid_h + 180000,
                grid_w, 260000,
                font_size=11, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER)
    add_line(slide, grid_left, grid_top + grid_h + 80000,
             grid_left + grid_w, grid_top + grid_h + 80000,
             color=tint(sub_color, 0.5), width=0.75)


# ── Variant 5: DARK — dark bg with glowing borders ──────────────────────

def _dark_grid(slide, t, c):
    add_background(slide, t.primary)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title.upper(), theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    axis_margin = 500000
    gap = 140000

    grid_left = m + axis_margin
    grid_top = content_y + 100000
    grid_w = SLIDE_W - 2 * m - axis_margin
    grid_h = FOOTER_TOP - grid_top - 380000
    quad_w = (grid_w - gap) // 2
    quad_h = (grid_h - gap) // 2

    card_fill = tint(t.primary, 0.08)

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, (col, row) in enumerate(positions):
        qx = grid_left + col * (quad_w + gap)
        qy = grid_top + row * (quad_h + gap)
        label, desc = quads[i] if i < len(quads) else ("", "")

        # Dark card with glowing border
        add_rect(slide, qx, qy, quad_w, quad_h, fill_color=card_fill,
                 line_color=colors[i], line_width=1.5, corner_radius=60000)
        # Glowing top accent
        add_rect(slide, qx + 1, qy, quad_w - 2, 50000, fill_color=colors[i])

        add_textbox(slide, label.upper(), qx + 150000, qy + 150000,
                    quad_w - 300000, 350000,
                    font_size=15, bold=True, color=colors[i])
        add_textbox(slide, desc, qx + 150000, qy + 530000,
                    quad_w - 300000, quad_h - 680000,
                    font_size=10, color=tint(t.primary, 0.5))

    # Neon axis lines
    add_line(slide, grid_left - 100000, grid_top + grid_h,
             grid_left - 100000, grid_top,
             color=t.accent, width=1.5)
    add_textbox(slide, y_axis.upper(), m, grid_top + grid_h // 2 - 180000,
                axis_margin - 100000, 360000,
                font_size=11, bold=True, color=tint(t.primary, 0.45),
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    add_line(slide, grid_left, grid_top + grid_h + 50000,
             grid_left + grid_w, grid_top + grid_h + 50000,
             color=t.accent, width=1.5)
    add_textbox(slide, x_axis.upper(), grid_left, grid_top + grid_h + 80000,
                grid_w, 260000,
                font_size=11, bold=True, color=tint(t.primary, 0.45),
                alignment=PP_ALIGN.CENTER)


# ── Variant 6: SPLIT — matrix left, descriptions right ──────────────────

def _split_quad(slide, t, c):
    add_dark_bg(slide, t)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title, theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    # Left: compact 2x2 matrix
    matrix_left = m
    matrix_top = content_y + 100000
    matrix_w = (SLIDE_W - 2 * m) * 55 // 100
    matrix_h = FOOTER_TOP - matrix_top - 350000
    gap = 100000
    mq_w = (matrix_w - gap) // 2
    mq_h = (matrix_h - gap) // 2

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, (col, row) in enumerate(positions):
        qx = matrix_left + col * (mq_w + gap)
        qy = matrix_top + row * (mq_h + gap)
        label = quads[i][0] if i < len(quads) else ""

        add_color_cell(slide, qx, qy, mq_w, mq_h,
                       fill_color=tint(colors[i], 0.25 if s.dark_mode else 0.75),
                       text=label, text_color=text_color, font_size=14, bold=True,
                       line_color=colors[i], line_width=1.0)

    # Axis labels on left grid
    add_textbox(slide, f"\u2190 {y_axis} \u2192", matrix_left, matrix_top - 200000,
                matrix_w, 180000,
                font_size=9, color=sub_color, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, f"\u2190 {x_axis} \u2192", matrix_left, matrix_top + matrix_h + 30000,
                matrix_w, 180000,
                font_size=9, color=sub_color, alignment=PP_ALIGN.CENTER)

    # Right: description list
    desc_left = matrix_left + matrix_w + 300000
    desc_w = SLIDE_W - m - desc_left
    item_h = matrix_h // 4

    # Divider line
    div_x = desc_left - 150000
    add_line(slide, div_x, content_y + 200000, div_x, FOOTER_TOP - 400000,
             color=t.accent, width=2.0)

    for i in range(min(4, len(quads))):
        label, desc = quads[i]
        iy = matrix_top + i * item_h
        # Color dot
        add_circle(slide, desc_left + 80000, iy + 120000, 50000, fill_color=colors[i])
        add_textbox(slide, label, desc_left + 180000, iy + 60000,
                    desc_w - 200000, 280000,
                    font_size=13, bold=True, color=text_color)
        add_textbox(slide, desc, desc_left + 180000, iy + 350000,
                    desc_w - 200000, item_h - 420000,
                    font_size=10, color=sub_color)


# ── Variant 7: GEO — angular geometric quadrants with borders ───────────

def _angular_grid(slide, t, c):
    add_dark_bg(slide, t)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title, theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if s.dark_mode else t.secondary
    axis_margin = 480000
    gap = 100000

    grid_left = m + axis_margin
    grid_top = content_y + 100000
    grid_w = SLIDE_W - 2 * m - axis_margin
    grid_h = FOOTER_TOP - grid_top - 380000
    quad_w = (grid_w - gap) // 2
    quad_h = (grid_h - gap) // 2

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, (col, row) in enumerate(positions):
        qx = grid_left + col * (quad_w + gap)
        qy = grid_top + row * (quad_h + gap)
        label, desc = quads[i] if i < len(quads) else ("", "")

        # Sharp card with thick border and bottom accent
        card_fill = tint(colors[i], 0.15 if s.dark_mode else 0.9)
        add_rect(slide, qx, qy, quad_w, quad_h, fill_color=card_fill,
                 line_color=colors[i], line_width=1.5)
        # Bottom accent bar
        add_rect(slide, qx, qy + quad_h - 60000, quad_w, 60000, fill_color=colors[i])

        # Diamond marker
        add_diamond(slide, qx + 200000, qy + 200000, 80000, fill_color=colors[i])

        add_textbox(slide, label, qx + 340000, qy + 130000, quad_w - 440000, 300000,
                    font_size=14, bold=True, color=text_color)
        add_textbox(slide, desc, qx + 150000, qy + 480000, quad_w - 300000, quad_h - 650000,
                    font_size=10, color=sub_color)

    # Angular axis markers
    add_line(slide, grid_left - 80000, grid_top + grid_h,
             grid_left - 80000, grid_top,
             color=colors[0], width=1.5)
    add_textbox(slide, y_axis, m, grid_top + grid_h // 2 - 180000,
                axis_margin - 100000, 360000,
                font_size=11, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)

    add_line(slide, grid_left, grid_top + grid_h + 50000,
             grid_left + grid_w, grid_top + grid_h + 50000,
             color=colors[1], width=1.5)
    add_textbox(slide, x_axis, grid_left, grid_top + grid_h + 80000,
                grid_w, 260000,
                font_size=11, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER)


# ── Variant 8: EDITORIAL — elegant thin-ruled grid with typography ───────

def _editorial_grid(slide, t, c):
    add_dark_bg(slide, t)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title, theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary
    rule_color = tint(t.secondary, 0.7) if not s.dark_mode else tint(t.primary, 0.2)
    axis_margin = 520000

    grid_left = m + axis_margin
    grid_top = content_y + 200000
    grid_w = SLIDE_W - 2 * m - axis_margin
    grid_h = FOOTER_TOP - grid_top - 450000
    quad_w = grid_w // 2
    quad_h = grid_h // 2

    # Thin cross rules
    mid_x = grid_left + grid_w // 2
    mid_y = grid_top + grid_h // 2
    add_line(slide, mid_x, grid_top, mid_x, grid_top + grid_h,
             color=rule_color, width=0.5)
    add_line(slide, grid_left, mid_y, grid_left + grid_w, mid_y,
             color=rule_color, width=0.5)

    # Short accent on intersection
    accent_len = 200000
    add_line(slide, mid_x - accent_len, mid_y, mid_x + accent_len, mid_y,
             color=t.accent, width=1.5)
    add_line(slide, mid_x, mid_y - accent_len, mid_x, mid_y + accent_len,
             color=t.accent, width=1.5)

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, (col, row) in enumerate(positions):
        qx = grid_left + col * quad_w
        qy = grid_top + row * quad_h
        label, desc = quads[i] if i < len(quads) else ("", "")

        # Short accent mark above label
        accent_x = qx + 150000
        accent_y = qy + 100000
        add_line(slide, accent_x, accent_y, accent_x + 400000, accent_y,
                 color=colors[i], width=1.5)

        add_textbox(slide, label, qx + 150000, qy + 140000, quad_w - 300000, 320000,
                    font_size=14, bold=True, color=text_color, italic=True)
        add_textbox(slide, desc, qx + 150000, qy + 480000, quad_w - 300000, quad_h - 580000,
                    font_size=10, color=sub_color)

    # Elegant axis labels
    add_textbox(slide, y_axis, m, grid_top + grid_h // 2 - 160000,
                axis_margin - 120000, 320000,
                font_size=11, italic=True, color=sub_color,
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_line(slide, grid_left - 80000, grid_top + grid_h,
             grid_left - 80000, grid_top,
             color=rule_color, width=0.5)

    add_textbox(slide, x_axis, grid_left, grid_top + grid_h + 80000,
                grid_w, 260000,
                font_size=11, italic=True, color=sub_color, alignment=PP_ALIGN.CENTER)
    add_line(slide, grid_left, grid_top + grid_h + 50000,
             grid_left + grid_w, grid_top + grid_h + 50000,
             color=rule_color, width=0.5)

    # Bottom editorial rule
    add_line(slide, m, FOOTER_TOP - 200000, SLIDE_W - m, FOOTER_TOP - 200000,
             color=rule_color, width=0.5)


# ── Variant 9: GRADIENT — gradient-tinted quadrants ──────────────────────

def _gradient_grid(slide, t, c):
    add_dark_bg(slide, t)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title, theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary
    axis_margin = 480000
    gap = int(200000 * s.gap_factor)

    grid_left = m + axis_margin
    grid_top = content_y + 150000
    grid_w = SLIDE_W - 2 * m - axis_margin
    grid_h = FOOTER_TOP - grid_top - 400000
    quad_w = (grid_w - gap) // 2
    quad_h = (grid_h - gap) // 2

    # Soft background circle behind grid
    cx = grid_left + grid_w // 2
    cy = grid_top + grid_h // 2
    add_circle(slide, cx, cy, min(grid_w, grid_h) // 2 + 100000,
               fill_color=tint(t.accent, 0.92 if not s.dark_mode else 0.08))

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    tint_factors = [0.7, 0.6, 0.8, 0.5] if not s.dark_mode else [0.15, 0.2, 0.1, 0.25]
    for i, (col, row) in enumerate(positions):
        qx = grid_left + col * (quad_w + gap)
        qy = grid_top + row * (quad_h + gap)
        label, desc = quads[i] if i < len(quads) else ("", "")

        bg = tint(colors[i], tint_factors[i])
        add_rect(slide, qx, qy, quad_w, quad_h, fill_color=bg,
                 corner_radius=120000)

        add_textbox(slide, label, qx + 150000, qy + quad_h // 2 - 350000,
                    quad_w - 300000, 350000,
                    font_size=15, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, desc, qx + 120000, qy + quad_h // 2 + 50000,
                    quad_w - 240000, 350000,
                    font_size=10, color=sub_color, alignment=PP_ALIGN.CENTER)

    # Axis labels with dots
    add_circle(slide, m + 140000, grid_top + grid_h // 2, 35000, fill_color=t.accent)
    add_textbox(slide, y_axis, m, grid_top + grid_h // 2 + 60000,
                axis_margin - 80000, 280000,
                font_size=11, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER)

    add_circle(slide, grid_left + grid_w // 2, grid_top + grid_h + 130000,
               35000, fill_color=t.accent)
    add_textbox(slide, x_axis, grid_left, grid_top + grid_h + 180000,
                grid_w, 260000,
                font_size=11, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER)


# ── Variant 10: RETRO — vintage styled grid with decorative borders ─────

def _retro_grid(slide, t, c):
    add_background(slide, t.light_bg)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title, theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    axis_margin = 500000
    gap = 160000

    grid_left = m + axis_margin
    grid_top = content_y + 150000
    grid_w = SLIDE_W - 2 * m - axis_margin
    grid_h = FOOTER_TOP - grid_top - 400000
    quad_w = (grid_w - gap) // 2
    quad_h = (grid_h - gap) // 2

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, (col, row) in enumerate(positions):
        qx = grid_left + col * (quad_w + gap)
        qy = grid_top + row * (quad_h + gap)
        label, desc = quads[i] if i < len(quads) else ("", "")

        # Double-bordered retro card
        add_rect(slide, qx, qy, quad_w, quad_h, fill_color="#FFFFFF",
                 line_color=colors[i], line_width=2.0, corner_radius=100000)
        add_rect(slide, qx + 40000, qy + 40000, quad_w - 80000, quad_h - 80000,
                 line_color=colors[i], line_width=0.75, corner_radius=80000)

        # Decorative corner dots
        add_circle(slide, qx + 100000, qy + 100000, 30000, fill_color=colors[i])
        add_circle(slide, qx + quad_w - 100000, qy + 100000, 30000, fill_color=colors[i])

        # Label with dashes
        add_textbox(slide, f"\u2014 {label} \u2014", qx + 100000, qy + 200000,
                    quad_w - 200000, 320000,
                    font_size=14, bold=True, color=colors[i], alignment=PP_ALIGN.CENTER)

        # Decorative rule
        rule_left = qx + quad_w // 4
        add_line(slide, rule_left, qy + 540000, rule_left + quad_w // 2, qy + 540000,
                 color=colors[i], width=1.0)

        add_textbox(slide, desc, qx + 120000, qy + 600000,
                    quad_w - 240000, quad_h - 750000,
                    font_size=10, color=t.secondary, alignment=PP_ALIGN.CENTER)

    # Retro axis labels with decorations
    add_textbox(slide, f"\u2500 {y_axis} \u2500", m, grid_top + grid_h // 2 - 180000,
                axis_margin - 80000, 360000,
                font_size=11, bold=True, color=t.accent, alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, f"\u2500 {x_axis} \u2500", grid_left, grid_top + grid_h + 80000,
                grid_w, 260000,
                font_size=11, bold=True, color=t.accent, alignment=PP_ALIGN.CENTER)


# ── Variant 11: MAGAZINE — creative asymmetric quadrants ─────────────────

def _creative_grid(slide, t, c):
    add_dark_bg(slide, t)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title, theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    # Asymmetric: top-left large, top-right narrow, bottom-left narrow, bottom-right large
    total_w = SLIDE_W - 2 * m
    total_h = FOOTER_TOP - content_y - 500000
    gap = 80000
    start_y = content_y + 100000

    large_w = int(total_w * 0.6)
    small_w = total_w - large_w - gap
    large_h = int(total_h * 0.55)
    small_h = total_h - large_h - gap

    layouts = [
        (m, start_y, large_w, large_h),                               # Q1 top-left (large)
        (m + large_w + gap, start_y, small_w, large_h),               # Q2 top-right (tall narrow)
        (m, start_y + large_h + gap, small_w, small_h),               # Q3 bottom-left (short narrow)
        (m + small_w + gap, start_y + large_h + gap, large_w, small_h),  # Q4 bottom-right (large)
    ]

    for i, (qx, qy, qw, qh) in enumerate(layouts):
        label, desc = quads[i] if i < len(quads) else ("", "")

        add_rect(slide, qx, qy, qw, qh, fill_color=colors[i])

        # Large number watermark
        add_textbox(slide, str(i + 1), qx + qw - 450000, qy + 50000,
                    400000, 400000,
                    font_size=48, bold=True, color=tint("#FFFFFF", 0.2),
                    alignment=PP_ALIGN.RIGHT)
        # Label
        label_fs = 18 if qw > total_w // 2 else 14
        add_textbox(slide, label, qx + 150000, qy + 120000, qw - 300000, 350000,
                    font_size=label_fs, bold=True, color="#FFFFFF")
        # Desc
        add_textbox(slide, desc, qx + 150000, qy + 480000, qw - 300000, qh - 580000,
                    font_size=10, color=tint("#FFFFFF", 0.15))

    # Axis annotations overlaid
    add_textbox(slide, f"{x_axis} \u2192", m, start_y + total_h + 40000,
                total_w, 220000,
                font_size=10, bold=True, color=sub_color, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, f"{y_axis} \u2191", m, start_y - 200000,
                total_w // 3, 180000,
                font_size=10, bold=True, color=sub_color, alignment=PP_ALIGN.LEFT)


# ── Variant 12: SCHOLARLY — figure-captioned, thin rules, centered ───────

def _scholarly_grid(slide, t, c):
    add_dark_bg(slide, t)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title, theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary
    rule_color = tint(t.secondary, 0.7) if not s.dark_mode else tint(t.primary, 0.2)
    axis_margin = 550000

    grid_left = m + axis_margin
    grid_top = content_y + 250000
    grid_w = SLIDE_W - 2 * m - axis_margin - 200000
    grid_h = FOOTER_TOP - grid_top - 650000
    quad_w = grid_w // 2
    quad_h = grid_h // 2

    # Outer thin rule
    add_rect(slide, grid_left, grid_top, grid_w, grid_h,
             line_color=rule_color, line_width=0.5)

    # Cross rules
    mid_x = grid_left + grid_w // 2
    mid_y = grid_top + grid_h // 2
    add_line(slide, mid_x, grid_top, mid_x, grid_top + grid_h,
             color=rule_color, width=0.5)
    add_line(slide, grid_left, mid_y, grid_left + grid_w, mid_y,
             color=rule_color, width=0.5)

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, (col, row) in enumerate(positions):
        qx = grid_left + col * quad_w
        qy = grid_top + row * quad_h
        label, desc = quads[i] if i < len(quads) else ("", "")

        # Small numbered reference
        ref = chr(ord("a") + i)
        add_textbox(slide, f"({ref})", qx + 80000, qy + 60000, 200000, 200000,
                    font_size=9, italic=True, color=sub_color)
        # Thin short color accent
        add_line(slide, qx + 80000, qy + 280000, qx + 280000, qy + 280000,
                 color=colors[i], width=1.0)
        add_textbox(slide, label, qx + 80000, qy + 310000, quad_w - 160000, 280000,
                    font_size=13, bold=True, color=text_color)
        add_textbox(slide, desc, qx + 80000, qy + 600000, quad_w - 160000, quad_h - 700000,
                    font_size=9, color=sub_color)

    # Axis labels (scholarly style)
    add_textbox(slide, y_axis, m, grid_top + grid_h // 2 - 150000,
                axis_margin - 150000, 300000,
                font_size=10, italic=True, color=sub_color,
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_line(slide, grid_left - 60000, grid_top + grid_h,
             grid_left - 60000, grid_top,
             color=rule_color, width=0.5)
    # Small arrow heads
    add_textbox(slide, "\u25B2", grid_left - 110000, grid_top - 100000,
                100000, 100000, font_size=7, color=rule_color, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x_axis, grid_left, grid_top + grid_h + 60000,
                grid_w, 220000,
                font_size=10, italic=True, color=sub_color, alignment=PP_ALIGN.CENTER)
    add_line(slide, grid_left, grid_top + grid_h + 40000,
             grid_left + grid_w, grid_top + grid_h + 40000,
             color=rule_color, width=0.5)
    add_textbox(slide, "\u25B6", grid_left + grid_w + 20000, grid_top + grid_h + 10000,
                100000, 100000, font_size=7, color=rule_color, alignment=PP_ALIGN.CENTER)

    # Figure caption below
    caption_y = grid_top + grid_h + 300000
    add_textbox(slide, f"Figure 1. {title} \u2014 {y_axis} vs. {x_axis}",
                grid_left, caption_y, grid_w, 220000,
                font_size=9, italic=True, color=sub_color, alignment=PP_ALIGN.CENTER)


# ── Variant 13: LABORATORY — dark bg, data-coded, grid overlay ──────────

def _laboratory_grid(slide, t, c):
    add_background(slide, t.primary)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title, theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    axis_margin = 520000
    gap = 100000

    grid_left = m + axis_margin
    grid_top = content_y + 100000
    grid_w = SLIDE_W - 2 * m - axis_margin
    grid_h = FOOTER_TOP - grid_top - 400000
    quad_w = (grid_w - gap) // 2
    quad_h = (grid_h - gap) // 2

    # Subtle grid overlay lines
    grid_spacing = 300000
    grid_line_color = tint(t.primary, 0.06)
    x = grid_left
    while x <= grid_left + grid_w:
        add_line(slide, x, grid_top, x, grid_top + grid_h,
                 color=grid_line_color, width=0.25)
        x += grid_spacing
    y = grid_top
    while y <= grid_top + grid_h:
        add_line(slide, grid_left, y, grid_left + grid_w, y,
                 color=grid_line_color, width=0.25)
        y += grid_spacing

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, (col, row) in enumerate(positions):
        qx = grid_left + col * (quad_w + gap)
        qy = grid_top + row * (quad_h + gap)
        label, desc = quads[i] if i < len(quads) else ("", "")

        card_bg = tint(t.primary, 0.06)
        add_rect(slide, qx, qy, quad_w, quad_h, fill_color=card_bg,
                 line_color=colors[i], line_width=0.75, corner_radius=30000)
        # Color-coded left border strip
        add_rect(slide, qx, qy, 50000, quad_h, fill_color=colors[i])

        # Data code label
        code = f"Q{i + 1}"
        add_textbox(slide, code, qx + 100000, qy + 60000, 200000, 200000,
                    font_size=9, bold=True, color=colors[i])
        # Thin separator
        add_line(slide, qx + 100000, qy + 260000, qx + quad_w - 100000, qy + 260000,
                 color=tint(t.primary, 0.12), width=0.5)

        add_textbox(slide, label, qx + 100000, qy + 300000, quad_w - 200000, 300000,
                    font_size=13, bold=True, color="#FFFFFF")
        add_textbox(slide, desc, qx + 100000, qy + 620000, quad_w - 200000, quad_h - 750000,
                    font_size=9, color=tint(t.primary, 0.45))

    # Lab-style axis
    axis_line_color = tint(t.primary, 0.25)
    add_line(slide, grid_left - 80000, grid_top + grid_h,
             grid_left - 80000, grid_top,
             color=axis_line_color, width=1.0)
    add_textbox(slide, y_axis, m, grid_top + grid_h // 2 - 150000,
                axis_margin - 120000, 300000,
                font_size=10, bold=True, color=tint(t.primary, 0.4),
                alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    add_line(slide, grid_left, grid_top + grid_h + 50000,
             grid_left + grid_w, grid_top + grid_h + 50000,
             color=axis_line_color, width=1.0)
    add_textbox(slide, x_axis, grid_left, grid_top + grid_h + 80000,
                grid_w, 240000,
                font_size=10, bold=True, color=tint(t.primary, 0.4),
                alignment=PP_ALIGN.CENTER)

    # Tick marks on axes
    for tick in range(5):
        frac = tick / 4
        tx = grid_left + int(grid_w * frac)
        add_line(slide, tx, grid_top + grid_h + 50000,
                 tx, grid_top + grid_h + 100000,
                 color=axis_line_color, width=0.5)
        ty = grid_top + int(grid_h * (1 - frac))
        add_line(slide, grid_left - 80000, ty,
                 grid_left - 130000, ty,
                 color=axis_line_color, width=0.5)


# ── Variant 14: DASHBOARD — dense compact matrix with metrics ────────────

def _dashboard_grid(slide, t, c):
    add_dark_bg(slide, t)
    title, x_axis, y_axis, quads = _get_data(c)
    content_y = add_slide_title(slide, title, theme=t)
    s = t.ux_style
    colors = _quad_colors(t)

    m = int(MARGIN * s.margin_factor)
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if s.dark_mode else t.secondary
    axis_margin = 420000
    gap = 60000

    grid_left = m + axis_margin
    grid_top = content_y + 80000
    grid_w = SLIDE_W - 2 * m - axis_margin
    grid_h = FOOTER_TOP - grid_top - 350000
    quad_w = (grid_w - gap) // 2
    quad_h = (grid_h - gap) // 2

    # Header bar for axis labels
    header_h = 260000
    add_rect(slide, grid_left, grid_top - header_h - 20000, grid_w, header_h,
             fill_color=tint(t.accent, 0.15 if s.dark_mode else 0.85))
    # High/Low labels along x-axis in header
    half_w = grid_w // 2
    add_textbox(slide, f"Low {x_axis}", grid_left + 80000,
                grid_top - header_h - 10000, half_w - 160000, header_h,
                font_size=9, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, f"High {x_axis}", grid_left + half_w + 80000,
                grid_top - header_h - 10000, half_w - 160000, header_h,
                font_size=9, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Side labels for y-axis
    side_w = axis_margin - 60000
    side_half_h = grid_h // 2
    add_rect(slide, m, grid_top, side_w, side_half_h - gap // 2,
             fill_color=tint(t.accent, 0.15 if s.dark_mode else 0.85))
    add_textbox(slide, f"Low\n{y_axis}", m + 20000, grid_top + 20000,
                side_w - 40000, side_half_h - gap // 2 - 40000,
                font_size=9, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, m, grid_top + side_half_h + gap // 2, side_w, side_half_h - gap // 2,
             fill_color=tint(t.accent, 0.1 if s.dark_mode else 0.9))
    add_textbox(slide, f"High\n{y_axis}", m + 20000,
                grid_top + side_half_h + gap // 2 + 20000,
                side_w - 40000, side_half_h - gap // 2 - 40000,
                font_size=9, bold=True, color=sub_color, alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, (col, row) in enumerate(positions):
        qx = grid_left + col * (quad_w + gap)
        qy = grid_top + row * (quad_h + gap)
        label, desc = quads[i] if i < len(quads) else ("", "")

        card_bg = tint(colors[i], 0.12 if s.dark_mode else 0.88)
        add_rect(slide, qx, qy, quad_w, quad_h, fill_color=card_bg,
                 line_color=tint(colors[i], 0.4 if s.dark_mode else 0.5),
                 line_width=0.5, corner_radius=50000)

        # Compact colored header bar
        bar_h = 200000
        add_rect(slide, qx, qy, quad_w, bar_h, fill_color=colors[i],
                 corner_radius=50000)
        # Fix bottom corners of header bar
        add_rect(slide, qx, qy + bar_h - 50000, quad_w, 50000, fill_color=colors[i])

        add_textbox(slide, label, qx + 100000, qy + 20000, quad_w - 200000, bar_h - 40000,
                    font_size=11, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Metric-style large number
        add_textbox(slide, f"Q{i + 1}", qx + quad_w - 350000, qy + bar_h + 60000,
                    280000, 250000,
                    font_size=28, bold=True, color=tint(colors[i], 0.6 if s.dark_mode else 0.3),
                    alignment=PP_ALIGN.RIGHT)

        add_textbox(slide, desc, qx + 100000, qy + bar_h + 80000,
                    quad_w - 450000, quad_h - bar_h - 150000,
                    font_size=9, color=text_color)
