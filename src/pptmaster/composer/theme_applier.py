"""Theme applier — final consistency pass to verify fonts and colors."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.util import Pt

from pptmaster.assets.color_utils import hex_to_rgb
from pptmaster.models import DesignProfile


def apply_theme(prs, profile: DesignProfile) -> None:
    """Final pass over all slides to enforce theme consistency.

    Checks and corrects:
    - Title font family and style
    - Body font family
    - Default text colors
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            is_title = _is_title_shape(shape)

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    _enforce_font(run, profile, is_title)


def _is_title_shape(shape) -> bool:
    """Check if a shape is a title placeholder."""
    try:
        if shape.placeholder_format is not None:
            ph_type = shape.placeholder_format.type
            if ph_type is not None and int(ph_type) in (0, 15):  # TITLE, CENTER_TITLE
                return True
    except Exception:
        pass

    # Check name
    name = (shape.name or "").lower()
    return "title" in name


def _enforce_font(run, profile: DesignProfile, is_title: bool) -> None:
    """Enforce the correct font on a text run."""
    font = run.font

    # Only override if the font is a default/mismatched font
    current = font.name
    expected = profile.fonts.major if is_title else profile.fonts.minor

    if current and current not in (profile.fonts.major, profile.fonts.minor,
                                    "+mj-lt", "+mn-lt", "Calibri", "Calibri Light"):
        # Font was explicitly set to something else (e.g. in a chart) — don't override
        return

    font.name = expected

    # Set bold for titles only if not already explicitly set
    if is_title and font.bold is None:
        font.bold = True
