"""Slide 5 — Our Values: 4 equal columns with accent circles and centered text."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import SLIDE_W, SLIDE_H, MARGIN, FONT_SLIDE_TITLE, FONT_BODY, FONT_CARD_HEADING, FONT_CAPTION, CONTENT_TOP, col_span
from pptmaster.builder.helpers import add_textbox, add_circle, add_gold_accent_line, add_styled_card, add_slide_title, add_dark_bg, add_background, add_rect, add_line
from pptmaster.assets.color_utils import tint, shade


# ── Scholarly variant ─────────────────────────────────────────────────
def _scholarly(slide, t) -> None:
    """Core Principles — numbered vertical list with thin dividers."""
    c = t.content
    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, c.get("values_title", "Core Principles"), theme=t)

    text_color = t.primary
    sub_color = t.secondary
    m = MARGIN

    # Thin top rule
    rule_y = content_y + 50000
    add_line(slide, m, rule_y, SLIDE_W - m, rule_y, color=t.secondary, width=0.5)

    values = c.get("values", [
        ("Integrity", "We act with honesty and ethical leadership."),
        ("Innovation", "We embrace creative thinking and improvement."),
        ("Excellence", "We pursue the highest standards."),
        ("Collaboration", "We achieve more together."),
    ])

    y = rule_y + 200000
    content_w = SLIDE_W - 2 * m
    for i, (title, description) in enumerate(values[:4]):
        num = f"{i + 1}."
        add_textbox(slide, num, m, y, 200000, 400000,
                    font_size=FONT_CARD_HEADING, bold=True, color=t.accent)
        add_textbox(slide, title, m + 200000, y, content_w - 200000, 350000,
                    font_size=FONT_CARD_HEADING, bold=True, color=text_color)
        add_textbox(slide, description, m + 200000, y + 350000, content_w - 400000, 400000,
                    font_size=FONT_BODY, color=sub_color)
        y += 850000
        # Thin divider between values
        if i < len(values[:4]) - 1:
            add_line(slide, m + 100000, y, SLIDE_W - m - 100000, y,
                     color=t.secondary, width=0.5)
            y += 150000


# ── Laboratory variant ────────────────────────────────────────────────
def _laboratory(slide, t) -> None:
    """Research Pillars — horizontal dark cards with colored left borders."""
    c = t.content
    add_background(slide, t.primary)
    content_y = add_slide_title(slide, c.get("values_title", "Research Pillars"), theme=t)

    text_color = "#FFFFFF"
    sub_color = tint(t.primary, 0.6)
    dark_fill = tint(t.primary, 0.15)
    p = t.palette

    values = c.get("values", [
        ("Integrity", "We act with honesty and ethical leadership."),
        ("Innovation", "We embrace creative thinking and improvement."),
        ("Excellence", "We pursue the highest standards."),
        ("Collaboration", "We achieve more together."),
    ])

    m = MARGIN
    avail_w = SLIDE_W - 2 * m
    card_gap = 150000
    n = min(len(values), 4)
    card_w = (avail_w - (n - 1) * card_gap) // n
    card_y = content_y + 200000
    card_h = 3800000

    for i, (title, description) in enumerate(values[:n]):
        cx = m + i * (card_w + card_gap)
        accent = p[i % len(p)]
        add_rect(slide, cx, card_y, card_w, card_h, fill_color=dark_fill,
                 line_color=tint(t.primary, 0.25), line_width=1)
        add_rect(slide, cx, card_y, 50000, card_h, fill_color=accent)  # left accent
        add_textbox(slide, title, cx + 130000, card_y + 250000, card_w - 200000, 400000,
                    font_size=FONT_CARD_HEADING, bold=True, color=accent)
        add_textbox(slide, description, cx + 130000, card_y + 750000, card_w - 200000, 2800000,
                    font_size=13, color=sub_color)


# ── Dashboard variant ─────────────────────────────────────────────────
def _dashboard(slide, t) -> None:
    """Header band + compact horizontal card strip."""
    c = t.content
    add_dark_bg(slide, t)
    s = t.ux_style
    content_y = add_slide_title(slide, c.get("values_title", "Our Values"), theme=t)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    # Accent header band
    add_rect(slide, 0, content_y - 200000, SLIDE_W, 400000, fill_color=t.accent)

    values = c.get("values", [
        ("Integrity", "We act with honesty and ethical leadership."),
        ("Innovation", "We embrace creative thinking and improvement."),
        ("Excellence", "We pursue the highest standards."),
        ("Collaboration", "We achieve more together."),
    ])

    m = int(MARGIN * 0.7)
    strip_y = content_y + 400000
    strip_h = SLIDE_H - strip_y - 500000
    n = min(len(values), 4)
    gap = 120000
    card_w = (SLIDE_W - 2 * m - (n - 1) * gap) // n
    p = t.palette

    for i, (title, description) in enumerate(values[:n]):
        cx = m + i * (card_w + gap)
        accent = p[i % len(p)]
        add_styled_card(slide, cx, strip_y, card_w, strip_h, theme=t, accent_color=accent)
        add_textbox(slide, title, cx + 100000, strip_y + 200000, card_w - 200000, 400000,
                    font_size=FONT_CARD_HEADING, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, description, cx + 80000, strip_y + 700000, card_w - 160000, strip_h - 900000,
                    font_size=12, color=sub_color, alignment=PP_ALIGN.CENTER)


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    s = t.ux_style

    # ── Style dispatch ────────────────────────────────────────────────
    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("values_title", "Our Values"), theme=t)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    values = c.get("values", [
        ("Integrity", "We act with honesty and ethical leadership."),
        ("Innovation", "We embrace creative thinking and improvement."),
        ("Excellence", "We pursue the highest standards."),
        ("Collaboration", "We achieve more together."),
    ])

    circle_radius = 250000
    circle_top_y = content_y + 500000
    gap = int(200000 * s.gap_factor)

    for i, (title, description) in enumerate(values[:4]):
        left, width = col_span(4, i, gap=gap)
        center_x = left + width // 2

        add_circle(slide, center_x, circle_top_y, circle_radius, fill_color=t.accent)
        add_textbox(slide, title[0], center_x - circle_radius, circle_top_y - circle_radius,
                    circle_radius * 2, circle_radius * 2,
                    font_size=28, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, title, left, circle_top_y + circle_radius + 250000, width, 400000,
                    font_size=FONT_CARD_HEADING, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, description, left + 50000, circle_top_y + circle_radius + 700000,
                    width - 100000, 1800000, font_size=13, color=sub_color, alignment=PP_ALIGN.CENTER)
