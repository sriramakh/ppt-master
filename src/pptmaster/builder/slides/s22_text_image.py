"""Slide 22 — Text + Image: 60/40 split, bullets left, image placeholder right."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, FONT_SLIDE_TITLE, FONT_BODY, FONT_CAPTION,
    CONTENT_TOP, FOOTER_TOP,
)
from pptmaster.builder.helpers import (
    add_textbox, add_bullet_list, add_image_placeholder,
    add_styled_card, add_slide_title, add_dark_bg, add_background,
    add_rect, add_line,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content

    s = t.ux_style
    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)

    content_top = add_slide_title(slide, c.get("approach_title", "Our Approach"), theme=t)

    s = t.ux_style
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    total_w = SLIDE_W - 2 * MARGIN
    gap = 300000
    left_w = int(total_w * 0.58)
    ct = content_top + 100000

    add_textbox(slide, c.get("approach_intro", ""), MARGIN, ct, left_w, 800000,
                font_size=FONT_BODY, color=sub_color)
    add_bullet_list(slide, c.get("approach_bullets", []), MARGIN, ct + 900000, left_w, 3500000,
                    font_size=14, color=text_color, bullet_color=t.accent, spacing=8)

    right_x = MARGIN + left_w + gap
    right_w = total_w - left_w - gap
    img_bg = tint(t.primary, 0.15) if s.dark_mode else t.light_bg
    add_image_placeholder(slide, right_x, ct, right_w, 4500000, label="Insert Image", bg_color=img_bg)


# ── Scholarly variant: white bg, indented paragraphs, figure caption ──────

def _scholarly(slide, t):
    c = t.content
    add_background(slide, "#FFFFFF")
    content_top = add_slide_title(slide, c.get("approach_title", "Our Approach"), theme=t)

    total_w = SLIDE_W - 2 * MARGIN
    gap = 400000
    left_w = int(total_w * 0.55)
    ct = content_top + 100000

    # Top thin rule
    add_line(slide, MARGIN, ct, MARGIN + total_w, ct,
             color=tint(t.secondary, 0.7), width=0.5)

    # Intro as indented paragraph
    indent = 300000
    add_textbox(slide, c.get("approach_intro", ""), MARGIN + indent, ct + 150000,
                left_w - indent, 800000,
                font_size=FONT_BODY, color=t.primary)

    # Bullets as indented body text
    add_bullet_list(slide, c.get("approach_bullets", []), MARGIN + indent, ct + 1000000,
                    left_w - indent, 3200000,
                    font_size=14, color=t.primary, bullet_color=t.secondary, spacing=10)

    # Image placeholder on right with figure caption
    right_x = MARGIN + left_w + gap
    right_w = total_w - left_w - gap
    img_h = 3600000
    add_image_placeholder(slide, right_x, ct + 150000, right_w, img_h,
                          label="Insert Image", bg_color=tint(t.secondary, 0.9))
    add_textbox(slide, "Figure 1: Reference Image", right_x, ct + 150000 + img_h + 80000,
                right_w, 300000,
                font_size=FONT_CAPTION, italic=True, color=t.secondary,
                alignment=PP_ALIGN.CENTER)

    # Bottom thin rule
    rule_y = FOOTER_TOP - 200000
    add_line(slide, MARGIN, rule_y, MARGIN + total_w, rule_y,
             color=tint(t.secondary, 0.7), width=0.5)


# ── Laboratory variant: dark bg, accent-bordered image placeholder ────────

def _laboratory(slide, t):
    c = t.content
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("approach_title", "Our Approach"), theme=t)

    total_w = SLIDE_W - 2 * MARGIN
    gap = 350000
    left_w = int(total_w * 0.55)
    ct = content_top + 100000

    # Text on left in white
    add_textbox(slide, c.get("approach_intro", ""), MARGIN, ct, left_w, 800000,
                font_size=FONT_BODY, color=tint(t.primary, 0.6))

    add_bullet_list(slide, c.get("approach_bullets", []), MARGIN, ct + 900000,
                    left_w, 3500000,
                    font_size=14, color="#FFFFFF", bullet_color=t.accent, spacing=8)

    # Image placeholder on right with colored border accent
    right_x = MARGIN + left_w + gap
    right_w = total_w - left_w - gap
    img_h = 4300000
    border_w = 40000
    add_rect(slide, right_x - border_w, ct - border_w,
             right_w + 2 * border_w, img_h + 2 * border_w,
             fill_color=t.accent)
    add_image_placeholder(slide, right_x, ct, right_w, img_h,
                          label="Insert Image",
                          bg_color=tint(t.primary, 0.12))


# ── Dashboard variant: header band, compact card layout ──────────────────

def _dashboard(slide, t):
    c = t.content
    add_background(slide, "#FFFFFF")

    # Accent header band
    band_h = 80000
    add_rect(slide, 0, 0, SLIDE_W, band_h, fill_color=t.accent)

    content_top = add_slide_title(slide, c.get("approach_title", "Our Approach"), theme=t)

    total_w = SLIDE_W - 2 * MARGIN
    gap = 250000
    left_w = int(total_w * 0.56)
    ct = content_top + 50000

    # Left compact card with text
    card_h = 4400000
    add_styled_card(slide, MARGIN, ct, left_w, card_h, theme=t, accent_color=t.accent)
    add_textbox(slide, c.get("approach_intro", ""), MARGIN + 150000, ct + 100000,
                left_w - 300000, 600000,
                font_size=13, color=t.secondary)
    add_bullet_list(slide, c.get("approach_bullets", []), MARGIN + 150000, ct + 750000,
                    left_w - 300000, 3400000,
                    font_size=13, color=t.primary, bullet_color=t.accent, spacing=6)

    # Right compact card with image
    right_x = MARGIN + left_w + gap
    right_w = total_w - left_w - gap
    add_styled_card(slide, right_x, ct, right_w, card_h, theme=t, accent_color=t.accent)
    add_image_placeholder(slide, right_x + 80000, ct + 80000,
                          right_w - 160000, card_h - 160000,
                          label="Insert Image", bg_color=t.light_bg)
