"""Slide 4 — Company Overview: Mission card + 4 icon cards (2x2)."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from pptmaster.builder.design_system import SLIDE_W, SLIDE_H, MARGIN, FONT_SLIDE_TITLE, FONT_BODY, FONT_CARD_HEADING, FONT_CAPTION, CONTENT_TOP, col_span
from pptmaster.builder.helpers import add_textbox, add_card, add_gold_accent_line, add_circle, add_styled_card, add_slide_title, add_dark_bg, add_background, add_rect, add_line
from pptmaster.assets.color_utils import tint, shade


# ── Scholarly variant ─────────────────────────────────────────────────
def _scholarly(slide, t, *, company_name: str = "") -> None:
    """Institution Profile — centered text paragraphs with thin rules."""
    c = t.content
    name = company_name or t.company_name

    add_background(slide, "#FFFFFF")
    content_y = add_slide_title(slide, c.get("overview_title", "Institution Profile"), theme=t)

    text_color = t.primary
    sub_color = t.secondary
    m = MARGIN

    # Thin top rule
    rule_y = content_y + 50000
    add_line(slide, m, rule_y, SLIDE_W - m, rule_y, color=t.secondary, width=0.5)

    # Mission text — centered paragraph
    mission_y = rule_y + 150000
    add_textbox(slide, "Mission", m, mission_y, SLIDE_W - 2 * m, 350000,
                font_size=FONT_CARD_HEADING, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, c.get("overview_mission", ""), m + 300000, mission_y + 400000,
                SLIDE_W - 2 * m - 600000, 1200000, font_size=FONT_BODY, color=sub_color,
                alignment=PP_ALIGN.CENTER)

    # Thin divider
    div_y = mission_y + 1700000
    add_line(slide, m + 300000, div_y, SLIDE_W - m - 300000, div_y, color=t.secondary, width=0.5)

    # Facts as simple text list
    facts = c.get("overview_facts", [("Founded", "2005"), ("Employees", "2,500+"),
                                      ("Offices", "12"), ("Revenue", "$850M")])
    y = div_y + 200000
    for label, value in facts[:4]:
        add_textbox(slide, f"{label}:  {value}", m + 400000, y, SLIDE_W - 2 * m - 800000, 400000,
                    font_size=FONT_BODY, color=text_color, alignment=PP_ALIGN.CENTER)
        y += 450000
        add_line(slide, m + 600000, y - 50000, SLIDE_W - m - 600000, y - 50000,
                 color=t.secondary, width=0.5)


# ── Laboratory variant ────────────────────────────────────────────────
def _laboratory(slide, t, *, company_name: str = "") -> None:
    """Dark data-dense overview with accent-bordered cards."""
    c = t.content
    name = company_name or t.company_name

    add_background(slide, t.primary)
    content_y = add_slide_title(slide, c.get("overview_title", "Company Overview"), theme=t)

    text_color = "#FFFFFF"
    sub_color = tint(t.primary, 0.6)

    # Abstract card with accent left border
    m = MARGIN
    abs_x = m
    abs_w = SLIDE_W // 2 - m - 100000
    abs_y = content_y + 100000
    abs_h = 4200000
    dark_fill = tint(t.primary, 0.15)

    add_rect(slide, abs_x, abs_y, abs_w, abs_h, fill_color=dark_fill,
             line_color=tint(t.primary, 0.25), line_width=1)
    add_rect(slide, abs_x, abs_y, 60000, abs_h, fill_color=t.accent)  # left accent border
    add_textbox(slide, "Abstract", abs_x + 180000, abs_y + 200000, abs_w - 280000, 350000,
                font_size=FONT_CARD_HEADING, bold=True, color=t.accent)
    add_textbox(slide, c.get("overview_mission", ""), abs_x + 180000, abs_y + 600000,
                abs_w - 280000, 3400000, font_size=FONT_BODY, color=sub_color)

    # 4 stat cards on right with colored left borders
    facts = c.get("overview_facts", [("Founded", "2005"), ("Employees", "2,500+"),
                                      ("Offices", "12"), ("Revenue", "$850M")])
    right_x = SLIDE_W // 2 + 100000
    right_w = SLIDE_W - m - right_x
    card_gap = 150000
    card_h = (abs_h - 3 * card_gap) // 4
    p = t.palette

    for i, (label, value) in enumerate(facts[:4]):
        cy = abs_y + i * (card_h + card_gap)
        accent = p[i % len(p)]
        add_rect(slide, right_x, cy, right_w, card_h, fill_color=dark_fill,
                 line_color=tint(t.primary, 0.25), line_width=1)
        add_rect(slide, right_x, cy, 50000, card_h, fill_color=accent)  # left accent
        add_textbox(slide, value, right_x + 130000, cy + 80000, right_w - 200000, card_h // 2,
                    font_size=24, bold=True, color=text_color)
        add_textbox(slide, label, right_x + 130000, cy + card_h // 2, right_w - 200000, card_h // 2,
                    font_size=FONT_CAPTION, color=sub_color)


# ── Dashboard variant ─────────────────────────────────────────────────
def _dashboard(slide, t, *, company_name: str = "") -> None:
    """Accent header band, sidebar info panel + KPI tiles."""
    c = t.content
    name = company_name or t.company_name

    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("overview_title", "Company Overview"), theme=t)

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if t.ux_style.dark_mode else t.secondary

    # Accent header band
    add_rect(slide, 0, content_y - 200000, SLIDE_W, 400000, fill_color=t.accent)

    m = int(MARGIN * 0.7)
    panel_y = content_y + 300000
    avail_h = SLIDE_H - panel_y - 500000

    # Left sidebar panel — company mission
    panel_w = SLIDE_W // 3 - m
    panel_fill = tint(t.primary, 0.15) if t.ux_style.dark_mode else "#F1F5F9"
    add_rect(slide, m, panel_y, panel_w, avail_h, fill_color=panel_fill,
             corner_radius=80000, line_color=t.accent, line_width=1)
    add_textbox(slide, name or "Company", m + 150000, panel_y + 150000, panel_w - 300000, 350000,
                font_size=FONT_CARD_HEADING, bold=True, color=text_color)
    add_textbox(slide, c.get("overview_mission", ""), m + 150000, panel_y + 600000,
                panel_w - 300000, avail_h - 800000, font_size=13, color=sub_color)

    # Right: KPI tiles (2x2)
    facts = c.get("overview_facts", [("Founded", "2005"), ("Employees", "2,500+"),
                                      ("Offices", "12"), ("Revenue", "$850M")])
    tile_x0 = m + panel_w + 200000
    tile_area_w = SLIDE_W - tile_x0 - m
    tile_gap = 150000
    tile_w = (tile_area_w - tile_gap) // 2
    tile_h = (avail_h - tile_gap) // 2
    p = t.palette

    for i, (label, value) in enumerate(facts[:4]):
        row, col = divmod(i, 2)
        tx = tile_x0 + col * (tile_w + tile_gap)
        ty = panel_y + row * (tile_h + tile_gap)
        accent = p[i % len(p)]
        add_styled_card(slide, tx, ty, tile_w, tile_h, theme=t, accent_color=accent)
        add_textbox(slide, value, tx + 80000, ty + tile_h // 4, tile_w - 160000, tile_h // 3,
                    font_size=28, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, label, tx + 80000, ty + tile_h * 2 // 3, tile_w - 160000, tile_h // 4,
                    font_size=FONT_CAPTION, color=sub_color, alignment=PP_ALIGN.CENTER)


def build(slide, *, theme=None, company_name: str = "") -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    name = company_name or t.company_name
    c = t.content
    s = t.ux_style

    # ── Style dispatch ────────────────────────────────────────────────
    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t, company_name=name)

    add_dark_bg(slide, t)
    content_y = add_slide_title(slide, c.get("overview_title", "Company Overview"), theme=t)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    # Left: Mission card
    left_x, left_w = col_span(2, 0, gap=int(300000 * s.gap_factor))
    mission_top = content_y + 100000
    mission_h = 4200000

    add_styled_card(slide, left_x, mission_top, left_w, mission_h,
                    theme=t, accent_color=t.accent)
    add_textbox(slide, "Our Mission", left_x + 300000, mission_top + 250000, left_w - 500000, 400000,
                font_size=FONT_CARD_HEADING, bold=True, color=text_color)
    add_textbox(slide, c.get("overview_mission", ""), left_x + 300000, mission_top + 750000,
                left_w - 500000, 3200000, font_size=FONT_BODY, color=sub_color)

    # Right: 4 icon cards (2x2)
    right_x, right_w = col_span(2, 1, gap=int(300000 * s.gap_factor))
    card_gap = int(200000 * s.gap_factor)
    card_w = (right_w - card_gap) // 2
    card_h = (mission_h - card_gap) // 2

    facts = c.get("overview_facts", [("Founded", "2005"), ("Employees", "2,500+"), ("Offices", "12"), ("Revenue", "$850M")])
    p = t.palette

    for i, (label, value) in enumerate(facts[:4]):
        row, col = divmod(i, 2)
        cx = right_x + col * (card_w + card_gap)
        cy = mission_top + row * (card_h + card_gap)
        accent = p[i % len(p)]

        add_styled_card(slide, cx, cy, card_w, card_h, theme=t, accent_color=accent)
        add_circle(slide, cx + card_w // 2, cy + 500000, 180000, fill_color=accent)
        add_textbox(slide, value, cx + 100000, cy + 800000, card_w - 200000, 500000,
                    font_size=24, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, label, cx + 100000, cy + 1250000, card_w - 200000, 300000,
                    font_size=FONT_CAPTION, color=sub_color, alignment=PP_ALIGN.CENTER)
