"""Slide 26 — Infographic Dashboard: 3 mini KPIs + bar chart + progress bars."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP, FONT_CAPTION, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_styled_card, add_slide_title, add_dark_bg,
    add_background, add_line,
)
from pptmaster.builder.slides._chart_helpers import profile_from_theme
from pptmaster.models import ChartSpec
from pptmaster.composer.chart_renderer import add_chart
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    p = t.palette

    s = t.ux_style
    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)

    content_top = add_slide_title(slide, c.get("infographic_title", "Performance Dashboard"), theme=t)

    s = t.ux_style
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    # Top: 3 mini KPI cards
    kpis = c.get("infographic_kpis", [("$850M", "Revenue"), ("2,500", "Employees"), ("98%", "Satisfaction")])
    kpi_top = content_top + 50000
    kpi_h = 900000

    for i, (value, label) in enumerate(kpis[:3]):
        left, width = col_span(3, i, gap=200000)
        accent = p[i % len(p)]
        add_styled_card(slide, left, kpi_top, width, kpi_h, theme=t, accent_color=accent)
        add_textbox(slide, value, left + 50000, kpi_top + 100000, width - 100000, 450000,
                    font_size=28, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, label, left + 50000, kpi_top + 530000, width - 100000, 300000,
                    font_size=11, color=sub_color, alignment=PP_ALIGN.CENTER)

    # Bottom left: mini bar chart
    chart_top = kpi_top + kpi_h + 200000
    chart_w = int((SLIDE_W - 2 * MARGIN) * 0.55)
    chart_h = 3200000

    spec = ChartSpec(
        chart_type="bar",
        title=c.get("infographic_chart_title", "Quarterly Revenue"),
        categories=c.get("infographic_chart_cats", ["Q1", "Q2", "Q3", "Q4"]),
        series=c.get("infographic_chart_series", [{"name": "2025", "values": [180, 200, 195, 210]}]),
        show_legend=True,
    )
    add_chart(slide, spec, profile_from_theme(t), left=MARGIN, top=chart_top, width=chart_w, height=chart_h)

    # Bottom right: progress bars
    prog_left = MARGIN + chart_w + 300000
    prog_w = SLIDE_W - prog_left - MARGIN
    prog_top = chart_top + 200000

    add_textbox(slide, "Project Completion", prog_left, prog_top, prog_w, 350000,
                font_size=14, bold=True, color=text_color)

    progress_items = c.get("infographic_progress", [
        ("Phase 1", 1.0), ("Phase 2", 0.75), ("Phase 3", 0.45), ("Phase 4", 0.15),
    ])
    bar_start_y = prog_top + 450000
    bar_spacing = 600000
    bar_h = 80000

    for i, (label, progress) in enumerate(progress_items[:4]):
        color = p[i % len(p)]
        y = bar_start_y + i * bar_spacing
        add_textbox(slide, f"{label}  ({int(progress * 100)}%)", prog_left, y, prog_w, 250000,
                    font_size=11, color=text_color)
        track_color = tint(color, 0.3) if s.dark_mode else tint(color, 0.8)
        add_rect(slide, prog_left, y + 280000, prog_w, bar_h, fill_color=track_color, corner_radius=40000)
        fill_w = int(prog_w * progress)
        if fill_w > 0:
            add_rect(slide, prog_left, y + 280000, fill_w, bar_h, fill_color=color, corner_radius=40000)


# ── Scholarly variant: white bg, clean centered stats, figure caption ─────

def _scholarly(slide, t):
    c = t.content
    p = t.palette
    add_background(slide, "#FFFFFF")
    content_top = add_slide_title(slide, c.get("infographic_title", "Performance Dashboard"), theme=t)

    kpis = c.get("infographic_kpis", [
        ("$850M", "Revenue"), ("2,500", "Employees"), ("98%", "Satisfaction"),
    ])

    ct = content_top + 100000
    total_w = SLIDE_W - 2 * MARGIN

    # Top thin rule framing data area
    add_line(slide, MARGIN, ct, MARGIN + total_w, ct,
             color=tint(t.secondary, 0.7), width=0.5)

    # Centered stat layout — 3 KPIs in a row
    kpi_top = ct + 200000
    for i, (value, label) in enumerate(kpis[:3]):
        left, width = col_span(3, i, gap=300000)
        add_textbox(slide, value, left, kpi_top, width, 600000,
                    font_size=36, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, label, left, kpi_top + 600000, width, 300000,
                    font_size=12, color=t.secondary, alignment=PP_ALIGN.CENTER)

    # Mid thin rule
    mid_rule_y = kpi_top + 1050000
    add_line(slide, MARGIN + total_w // 4, mid_rule_y,
             MARGIN + 3 * total_w // 4, mid_rule_y,
             color=tint(t.secondary, 0.7), width=0.5)

    # Progress bars below in a clean layout
    progress_items = c.get("infographic_progress", [
        ("Phase 1", 1.0), ("Phase 2", 0.75), ("Phase 3", 0.45), ("Phase 4", 0.15),
    ])
    bar_top = mid_rule_y + 250000
    bar_spacing = 550000
    bar_h = 60000
    bar_w = int(total_w * 0.7)
    bar_left = MARGIN + (total_w - bar_w) // 2

    for i, (label, progress) in enumerate(progress_items[:4]):
        y = bar_top + i * bar_spacing
        add_textbox(slide, f"{label}  ({int(progress * 100)}%)", bar_left, y, bar_w, 220000,
                    font_size=11, color=t.primary)
        add_rect(slide, bar_left, y + 250000, bar_w, bar_h,
                 fill_color=tint(t.secondary, 0.85), corner_radius=30000)
        fill_w = int(bar_w * progress)
        if fill_w > 0:
            add_rect(slide, bar_left, y + 250000, fill_w, bar_h,
                     fill_color=t.primary, corner_radius=30000)

    # Bottom thin rule
    bottom_rule_y = FOOTER_TOP - 250000
    add_line(slide, MARGIN, bottom_rule_y, MARGIN + total_w, bottom_rule_y,
             color=tint(t.secondary, 0.7), width=0.5)

    # Figure caption
    add_textbox(slide, "Figure 1: Key Performance Indicators", MARGIN, bottom_rule_y + 60000,
                total_w, 250000,
                font_size=FONT_CAPTION, italic=True, color=t.secondary,
                alignment=PP_ALIGN.CENTER)


# ── Laboratory variant: dark bg, dark metric tiles with accent numbers ────

def _laboratory(slide, t):
    c = t.content
    p = t.palette
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("infographic_title", "Performance Dashboard"), theme=t)

    kpis = c.get("infographic_kpis", [
        ("$850M", "Revenue"), ("2,500", "Employees"), ("98%", "Satisfaction"),
    ])

    ct = content_top + 100000

    # KPI tiles in dark cards with accent numbers
    kpi_h = 900000
    for i, (value, label) in enumerate(kpis[:3]):
        left, width = col_span(3, i, gap=200000)
        accent = p[i % len(p)]

        # Dark tile
        add_rect(slide, left, ct, width, kpi_h,
                 fill_color=tint(t.primary, 0.08), corner_radius=50000)
        # Value in accent color
        add_textbox(slide, value, left + 50000, ct + 100000, width - 100000, 450000,
                    font_size=30, bold=True, color=accent, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, label, left + 50000, ct + 530000, width - 100000, 300000,
                    font_size=11, color=tint(t.primary, 0.5), alignment=PP_ALIGN.CENTER)

    # Progress bars in dark style
    progress_items = c.get("infographic_progress", [
        ("Phase 1", 1.0), ("Phase 2", 0.75), ("Phase 3", 0.45), ("Phase 4", 0.15),
    ])
    prog_top = ct + kpi_h + 350000
    total_w = SLIDE_W - 2 * MARGIN
    bar_spacing = 550000
    bar_h = 70000

    add_textbox(slide, "Project Completion", MARGIN, prog_top, total_w, 300000,
                font_size=14, bold=True, color="#FFFFFF")

    bar_start_y = prog_top + 400000
    for i, (label, progress) in enumerate(progress_items[:4]):
        accent = p[i % len(p)]
        y = bar_start_y + i * bar_spacing
        add_textbox(slide, f"{label}  ({int(progress * 100)}%)", MARGIN, y,
                    total_w, 220000, font_size=11, color=tint(t.primary, 0.5))
        track_color = tint(t.primary, 0.12)
        add_rect(slide, MARGIN, y + 260000, total_w, bar_h,
                 fill_color=track_color, corner_radius=35000)
        fill_w = int(total_w * progress)
        if fill_w > 0:
            add_rect(slide, MARGIN, y + 260000, fill_w, bar_h,
                     fill_color=accent, corner_radius=35000)


# ── Dashboard variant: header band, dense 2x3 metric tile grid ───────────

def _dashboard(slide, t):
    c = t.content
    p = t.palette
    add_background(slide, "#FFFFFF")

    # Accent header band
    add_rect(slide, 0, 0, SLIDE_W, 80000, fill_color=t.accent)

    content_top = add_slide_title(slide, c.get("infographic_title", "Performance Dashboard"), theme=t)

    kpis = c.get("infographic_kpis", [
        ("$850M", "Revenue"), ("2,500", "Employees"), ("98%", "Satisfaction"),
    ])

    # Extend KPIs with progress data for dense grid
    progress_items = c.get("infographic_progress", [
        ("Phase 1", 1.0), ("Phase 2", 0.75), ("Phase 3", 0.45),
    ])

    ct = content_top + 50000

    # Row 1: 3 KPI tiles
    tile_h = 900000
    for i, (value, label) in enumerate(kpis[:3]):
        left, width = col_span(3, i, gap=200000)
        accent = p[i % len(p)]
        add_styled_card(slide, left, ct, width, tile_h, theme=t, accent_color=accent)
        add_textbox(slide, value, left + 40000, ct + 80000, width - 80000, 420000,
                    font_size=26, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, label, left + 40000, ct + 480000, width - 80000, 280000,
                    font_size=10, color=t.secondary, alignment=PP_ALIGN.CENTER)

    # Row 2: progress tiles as mini cards
    row2_top = ct + tile_h + 200000
    tile2_h = 750000
    for i, (label, progress) in enumerate(progress_items[:3]):
        left, width = col_span(3, i, gap=200000)
        accent = p[i % len(p)]
        add_styled_card(slide, left, row2_top, width, tile2_h, theme=t, accent_color=accent)
        add_textbox(slide, f"{int(progress * 100)}%", left + 40000, row2_top + 80000,
                    width - 80000, 350000,
                    font_size=24, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, label, left + 40000, row2_top + 420000, width - 80000, 250000,
                    font_size=10, color=t.secondary, alignment=PP_ALIGN.CENTER)

    # Bottom: mini bar chart
    chart_top = row2_top + tile2_h + 200000
    chart_w = SLIDE_W - 2 * MARGIN
    chart_h = FOOTER_TOP - chart_top - 100000
    if chart_h > 500000:
        spec = ChartSpec(
            chart_type="bar",
            title=c.get("infographic_chart_title", "Quarterly Revenue"),
            categories=c.get("infographic_chart_cats", ["Q1", "Q2", "Q3", "Q4"]),
            series=c.get("infographic_chart_series",
                         [{"name": "2025", "values": [180, 200, 195, 210]}]),
            show_legend=True,
        )
        add_chart(slide, spec, profile_from_theme(t),
                  left=MARGIN, top=chart_top, width=chart_w, height=chart_h)
