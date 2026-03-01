"""Slide 35 — Milestone Roadmap: 14 unique visual variants dispatched by UX style."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_circle, add_line, add_diamond,
    add_notched_arrow, add_round_rect_2, add_gold_accent_line,
    add_slide_title, add_dark_bg, add_background, add_styled_card,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    s = t.ux_style
    c = t.content

    dispatch = {
        "arrow-path": _arrow_path,
        "dot-timeline": _dot_timeline,
        "block-arrows": _block_arrows,
        "floating-path": _floating_path,
        "glow-path": _glow_path,
        "alternating-milestone": _alternating_milestone,
        "diamond-path": _diamond_path,
        "editorial-path": _editorial_path,
        "gradient-path": _gradient_path,
        "retro-path": _retro_path,
        "mosaic-path": _mosaic_path,
        "scholarly-path": _scholarly_path,
        "laboratory-path": _laboratory_path,
        "dashboard-path": _dashboard_path,
    }
    builder = dispatch.get(s.milestone, _arrow_path)
    builder(slide, t, c)


def _get_items(c):
    """Extract milestone items from content dict with sensible defaults."""
    return c.get("milestone_items", [
        ("Jan 2026", "Project Kickoff", "Assemble team and define scope"),
        ("Mar 2026", "Alpha Release", "Core features complete"),
        ("May 2026", "Beta Testing", "User acceptance testing"),
        ("Jul 2026", "Launch", "Public release and marketing"),
        ("Sep 2026", "Scale", "Performance optimization"),
        ("Nov 2026", "Review", "Post-launch assessment"),
    ])


# ── Variant 1: CLASSIC — Horizontal path with notched arrows at each milestone ──

def _arrow_path(slide, t, c):
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones"), theme=t)
    items = _get_items(c)
    p = t.palette

    n = len(items)
    gap = 60000
    total_w = SLIDE_W - 2 * MARGIN
    arrow_w = (total_w - gap * (n - 1)) // n
    arrow_h = 700000
    arrow_top = content_top + 1200000

    # Horizontal guide line behind the arrows
    add_line(slide, MARGIN, arrow_top + arrow_h // 2,
             SLIDE_W - MARGIN, arrow_top + arrow_h // 2,
             color=tint(t.secondary, 0.7), width=1)

    for i, (date, title, desc) in enumerate(items):
        accent = p[i % len(p)]
        left = MARGIN + i * (arrow_w + gap)

        # Notched arrow shape with date inside
        add_notched_arrow(slide, left, arrow_top, arrow_w, arrow_h,
                          fill_color=accent, text=date,
                          text_color="#FFFFFF", font_size=11)

        # Title above arrow
        add_textbox(slide, title, left, arrow_top - 500000,
                    arrow_w, 400000,
                    font_size=13, bold=True, color=t.primary,
                    alignment=PP_ALIGN.CENTER)

        # Description below arrow
        add_textbox(slide, desc, left, arrow_top + arrow_h + 200000,
                    arrow_w, 500000,
                    font_size=10, color=t.secondary,
                    alignment=PP_ALIGN.CENTER)

        # Small circle connector between arrow and title
        add_circle(slide, left + arrow_w // 2, arrow_top - 100000, 40000,
                   fill_color=accent)


# ── Variant 2: MINIMAL — Dots on a thin horizontal line, labels alternate ──

def _dot_timeline(slide, t, c):
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones"), theme=t)
    items = _get_items(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(items)
    total_w = SLIDE_W - 2 * m
    step_w = total_w // n
    line_y = content_top + 2000000
    dot_r = 55000

    # Thin horizontal line spanning first to last dot
    first_cx = m + step_w // 2
    last_cx = m + (n - 1) * step_w + step_w // 2
    add_line(slide, first_cx, line_y, last_cx, line_y,
             color=tint(t.secondary, 0.7), width=0.5)

    for i, (date, title, desc) in enumerate(items):
        cx = m + i * step_w + step_w // 2
        is_above = i % 2 == 0

        # Small dot on line
        add_circle(slide, cx, line_y, dot_r, fill_color=t.accent)

        if is_above:
            # Date, title, desc above the line
            add_textbox(slide, date, cx - step_w // 2 + 20000,
                        line_y - dot_r - 950000,
                        step_w - 40000, 250000,
                        font_size=9, bold=True, color=t.accent,
                        alignment=PP_ALIGN.CENTER)
            add_textbox(slide, title, cx - step_w // 2 + 20000,
                        line_y - dot_r - 650000,
                        step_w - 40000, 300000,
                        font_size=12, bold=True, color=t.primary,
                        alignment=PP_ALIGN.CENTER)
            add_textbox(slide, desc, cx - step_w // 2 + 20000,
                        line_y - dot_r - 350000,
                        step_w - 40000, 300000,
                        font_size=10, color=t.secondary,
                        alignment=PP_ALIGN.CENTER)
            # Thin connector from dot up
            add_line(slide, cx, line_y - dot_r, cx, line_y - dot_r - 200000,
                     color=tint(t.accent, 0.5), width=0.5)
        else:
            # Date, title, desc below the line
            add_line(slide, cx, line_y + dot_r, cx, line_y + dot_r + 200000,
                     color=tint(t.accent, 0.5), width=0.5)
            add_textbox(slide, date, cx - step_w // 2 + 20000,
                        line_y + dot_r + 250000,
                        step_w - 40000, 250000,
                        font_size=9, bold=True, color=t.accent,
                        alignment=PP_ALIGN.CENTER)
            add_textbox(slide, title, cx - step_w // 2 + 20000,
                        line_y + dot_r + 500000,
                        step_w - 40000, 300000,
                        font_size=12, bold=True, color=t.primary,
                        alignment=PP_ALIGN.CENTER)
            add_textbox(slide, desc, cx - step_w // 2 + 20000,
                        line_y + dot_r + 800000,
                        step_w - 40000, 300000,
                        font_size=10, color=t.secondary,
                        alignment=PP_ALIGN.CENTER)


# ── Variant 3: BOLD — Bold colored block arrows with large dates ──

def _block_arrows(slide, t, c):
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones").upper(), theme=t)
    items = _get_items(c)
    p = t.palette

    n = len(items)
    gap = 50000
    total_w = SLIDE_W - 2 * MARGIN
    block_w = (total_w - gap * (n - 1)) // n
    block_h = FOOTER_TOP - content_top - 400000
    block_top = content_top + 200000

    for i, (date, title, desc) in enumerate(items):
        accent = p[i % len(p)]
        left = MARGIN + i * (block_w + gap)

        # Tall notched arrow block
        add_notched_arrow(slide, left, block_top, block_w, block_h,
                          fill_color=accent)

        # Large date at top of block
        add_textbox(slide, date, left + 60000, block_top + 200000,
                    block_w - 120000, 600000,
                    font_size=24, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT)

        # Thick white separator line
        add_rect(slide, left + 60000, block_top + 850000,
                 block_w * 2 // 3, 50000, fill_color="#FFFFFF")

        # Title
        add_textbox(slide, title.upper(), left + 60000, block_top + 1050000,
                    block_w - 120000, 500000,
                    font_size=13, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, left + 60000, block_top + 1600000,
                    block_w - 120000, 600000,
                    font_size=10, color=tint(accent, 0.7),
                    alignment=PP_ALIGN.LEFT)


# ── Variant 4: ELEVATED — Floating milestone cards connected by a curved path ──

def _floating_path(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones"), theme=t)
    items = _get_items(c)
    p = t.palette

    n = len(items)
    total_w = SLIDE_W - 2 * MARGIN
    card_w = int(total_w * 0.20)
    card_h = 1400000
    x_step = (total_w - card_w) // max(n - 1, 1)
    available_h = FOOTER_TOP - content_top - card_h - 500000
    y_step = available_h // max(n - 1, 1)

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if t.ux_style.dark_mode else t.secondary

    for i, (date, title, desc) in enumerate(items):
        accent = p[i % len(p)]
        card_left = MARGIN + i * x_step
        card_top = content_top + 200000 + i * y_step

        # Dashed connector line from previous to this card
        if i > 0:
            prev_left = MARGIN + (i - 1) * x_step
            prev_top = content_top + 200000 + (i - 1) * y_step
            add_line(slide, prev_left + card_w, prev_top + card_h // 2,
                     card_left, card_top + card_h // 2,
                     color=tint(t.accent, 0.3), width=1.5, dash="dash")

        # Floating card with rounded corners
        add_styled_card(slide, card_left, card_top, card_w, card_h,
                        theme=t, accent_color=accent)

        # Round-rect milestone badge at top-right
        badge_w = 400000
        badge_h = 250000
        add_round_rect_2(slide, card_left + card_w - badge_w - 50000,
                         card_top - badge_h // 2,
                         badge_w, badge_h,
                         fill_color=accent, text=f"M{i + 1}",
                         text_color="#FFFFFF", font_size=11)

        # Date
        add_textbox(slide, date, card_left + 80000, card_top + 120000,
                    card_w - 160000, 280000,
                    font_size=10, bold=True, color=accent,
                    alignment=PP_ALIGN.LEFT)

        # Title
        add_textbox(slide, title, card_left + 80000, card_top + 420000,
                    card_w - 160000, 400000,
                    font_size=13, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, card_left + 80000, card_top + 860000,
                    card_w - 160000, 450000,
                    font_size=10, color=sub_color,
                    alignment=PP_ALIGN.LEFT)


# ── Variant 5: DARK — Dark bg with glowing milestone nodes ──

def _glow_path(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones").upper(), theme=t)
    items = _get_items(c)
    p = t.palette

    n = len(items)
    total_w = SLIDE_W - 2 * MARGIN
    step_w = total_w // n
    node_y = content_top + 1400000
    node_r = 120000
    glow_r = 220000

    # Horizontal path
    first_cx = MARGIN + step_w // 2
    last_cx = MARGIN + (n - 1) * step_w + step_w // 2
    add_line(slide, first_cx, node_y, last_cx, node_y,
             color=tint(t.primary, 0.15), width=2)

    for i, (date, title, desc) in enumerate(items):
        accent = p[i % len(p)]
        cx = MARGIN + i * step_w + step_w // 2

        # Outer glow ring
        add_circle(slide, cx, node_y, glow_r, fill_color=tint(accent, 0.12))

        # Inner glow ring
        add_circle(slide, cx, node_y, int(glow_r * 0.7),
                   fill_color=tint(accent, 0.2))

        # Core node
        add_circle(slide, cx, node_y, node_r, fill_color=accent)

        # Milestone number inside node
        add_textbox(slide, str(i + 1), cx - node_r, node_y - node_r,
                    node_r * 2, node_r * 2,
                    font_size=12, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Date above
        add_textbox(slide, date, cx - step_w // 2 + 20000,
                    node_y - glow_r - 400000,
                    step_w - 40000, 280000,
                    font_size=10, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER)

        # Title below
        add_textbox(slide, title, cx - step_w // 2 + 20000,
                    node_y + glow_r + 150000,
                    step_w - 40000, 350000,
                    font_size=13, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, desc, cx - step_w // 2 + 20000,
                    node_y + glow_r + 550000,
                    step_w - 40000, 400000,
                    font_size=10, color=tint(t.primary, 0.45),
                    alignment=PP_ALIGN.CENTER)


# ── Variant 6: SPLIT — Milestones alternating left/right of center line ──

def _alternating_milestone(slide, t, c):
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones"), theme=t)
    items = _get_items(c)
    p = t.palette

    n = len(items)
    mid_x = SLIDE_W // 2
    available_h = FOOTER_TOP - content_top - 300000
    step_h = available_h // n
    card_w = int(SLIDE_W * 0.34)
    card_h = min(step_h - 80000, 800000)

    # Central vertical guide
    line_top = content_top + 200000
    line_bottom = line_top + n * step_h
    add_line(slide, mid_x, line_top, mid_x, line_bottom,
             color=t.accent, width=2)

    for i, (date, title, desc) in enumerate(items):
        accent = p[i % len(p)]
        is_left = i % 2 == 0
        step_top = content_top + 200000 + i * step_h
        node_y = step_top + card_h // 2

        # Circle node at center line
        add_circle(slide, mid_x, node_y, 90000, fill_color=accent)
        add_textbox(slide, str(i + 1), mid_x - 70000, node_y - 70000,
                    140000, 140000,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        if is_left:
            card_left = mid_x - card_w - 200000
            add_line(slide, card_left + card_w, node_y, mid_x - 90000, node_y,
                     color=tint(accent, 0.5), width=1)
        else:
            card_left = mid_x + 200000
            add_line(slide, mid_x + 90000, node_y, card_left, node_y,
                     color=tint(accent, 0.5), width=1)

        add_styled_card(slide, card_left, step_top, card_w, card_h,
                        theme=t, accent_color=accent)

        # Date
        add_textbox(slide, date, card_left + 100000, step_top + 50000,
                    card_w - 200000, 220000,
                    font_size=10, bold=True, color=accent,
                    alignment=PP_ALIGN.LEFT)

        # Title
        add_textbox(slide, title, card_left + 100000, step_top + 280000,
                    card_w - 200000, 280000,
                    font_size=13, bold=True, color=t.primary,
                    alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, card_left + 100000, step_top + 550000,
                    card_w - 200000, 250000,
                    font_size=11, color=t.secondary,
                    alignment=PP_ALIGN.LEFT)


# ── Variant 7: GEO — Diamond shapes on a path with geometric styling ──

def _diamond_path(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones"), theme=t)
    items = _get_items(c)
    p = t.palette

    n = len(items)
    total_w = SLIDE_W - 2 * MARGIN
    step_w = total_w // n
    line_y = content_top + 1500000
    diamond_size = 180000

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if t.ux_style.dark_mode else t.secondary

    # Horizontal spine
    first_cx = MARGIN + step_w // 2
    last_cx = MARGIN + (n - 1) * step_w + step_w // 2
    add_line(slide, first_cx, line_y, last_cx, line_y,
             color=tint(t.accent, 0.3 if t.ux_style.dark_mode else 0.6),
             width=2)

    for i, (date, title, desc) in enumerate(items):
        accent = p[i % len(p)]
        cx = MARGIN + i * step_w + step_w // 2

        # Diamond marker
        add_diamond(slide, cx, line_y, diamond_size, fill_color=accent)

        # Milestone number inside diamond
        add_textbox(slide, str(i + 1), cx - diamond_size, line_y - diamond_size,
                    diamond_size * 2, diamond_size * 2,
                    font_size=12, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Date above diamond
        add_textbox(slide, date, cx - step_w // 2 + 20000,
                    line_y - diamond_size - 420000,
                    step_w - 40000, 280000,
                    font_size=10, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER)

        # Title below diamond
        add_textbox(slide, title, cx - step_w // 2 + 20000,
                    line_y + diamond_size + 200000,
                    step_w - 40000, 300000,
                    font_size=13, bold=True, color=text_color,
                    alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, desc, cx - step_w // 2 + 20000,
                    line_y + diamond_size + 550000,
                    step_w - 40000, 400000,
                    font_size=10, color=sub_color,
                    alignment=PP_ALIGN.CENTER)

        # Small decorative diamond below text
        add_diamond(slide, cx, line_y + diamond_size + 1050000, 35000,
                    fill_color=accent)


# ── Variant 8: EDITORIAL — Thin rules, elegant numbered milestones ──

def _editorial_path(slide, t, c):
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones"), theme=t)
    items = _get_items(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(items)
    available_h = FOOTER_TOP - content_top - 400000
    step_h = available_h // n

    date_col_w = 1100000
    num_col_w = 600000
    line_x = m + date_col_w + num_col_w + 200000
    content_left = line_x + 250000
    content_w = SLIDE_W - m - content_left

    for i, (date, title, desc) in enumerate(items):
        step_top = content_top + 200000 + i * step_h

        # Thin horizontal rule across full width
        add_line(slide, m, step_top, SLIDE_W - m, step_top,
                 color=tint(t.secondary, 0.8), width=0.5)

        # Short accent mark at line intersection
        add_line(slide, line_x - 40000, step_top, line_x + 40000, step_top,
                 color=t.accent, width=2)

        # Number in small circle
        num_cx = m + date_col_w + num_col_w // 2
        add_textbox(slide, f"{i + 1:02d}", num_cx - num_col_w // 2,
                    step_top + 100000, num_col_w, 350000,
                    font_size=20, bold=True, color=t.accent,
                    alignment=PP_ALIGN.CENTER)

        # Date in left column
        add_textbox(slide, date, m, step_top + 100000, date_col_w, 350000,
                    font_size=11, bold=True, color=t.secondary,
                    alignment=PP_ALIGN.RIGHT)

        # Title in content column
        add_textbox(slide, title, content_left, step_top + 100000,
                    content_w, 350000,
                    font_size=14, bold=True, color=t.primary,
                    alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, content_left, step_top + 470000,
                    content_w, 300000,
                    font_size=12, color=t.secondary,
                    alignment=PP_ALIGN.LEFT)

    # Vertical connecting line through the number column
    first_top = content_top + 200000
    last_top = content_top + 200000 + (n - 1) * step_h
    add_line(slide, line_x, first_top, line_x, last_top,
             color=tint(t.secondary, 0.6), width=0.75)

    # Bottom rule
    add_line(slide, m, FOOTER_TOP - 300000, SLIDE_W - m, FOOTER_TOP - 300000,
             color=tint(t.secondary, 0.8), width=0.5)


# ── Variant 9: GRADIENT — Gradient-colored path with round nodes ──

def _gradient_path(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones"), theme=t)
    items = _get_items(c)

    n = len(items)
    total_w = SLIDE_W - 2 * MARGIN
    step_w = total_w // n
    node_y = content_top + 1100000
    node_r = 250000

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if t.ux_style.dark_mode else t.secondary

    for i, (date, title, desc) in enumerate(items):
        cx = MARGIN + i * step_w + step_w // 2
        # Progressive gradient: lighten accent more for later milestones
        factor = i * 0.10
        node_color = tint(t.accent, factor)

        # Connector line to next node
        if i < n - 1:
            next_cx = MARGIN + (i + 1) * step_w + step_w // 2
            add_line(slide, cx + node_r, node_y, next_cx - node_r, node_y,
                     color=tint(t.accent, 0.4), width=1.5)

        # Soft outer halo
        add_circle(slide, cx, node_y, node_r + 70000,
                   fill_color=tint(node_color, 0.85))

        # Main round node
        add_circle(slide, cx, node_y, node_r, fill_color=node_color)

        # Date inside top of node
        add_textbox(slide, date, cx - node_r + 30000, node_y - node_r + 50000,
                    node_r * 2 - 60000, 200000,
                    font_size=8, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

        # Title centered in node
        add_textbox(slide, title, cx - node_r + 30000, node_y - 70000,
                    node_r * 2 - 60000, 250000,
                    font_size=11, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

        # Description card below node
        card_left = cx - step_w // 2 + 30000
        card_w = step_w - 60000
        card_top = node_y + node_r + 200000
        card_h = 1200000

        add_styled_card(slide, card_left, card_top, card_w, card_h,
                        theme=t, accent_color=node_color)

        add_textbox(slide, desc, card_left + 60000, card_top + 100000,
                    card_w - 120000, 500000,
                    font_size=10, color=sub_color,
                    alignment=PP_ALIGN.CENTER)


# ── Variant 10: RETRO — Vintage styled with decorative milestone markers ──

def _retro_path(slide, t, c):
    add_background(slide, t.light_bg)
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones"), theme=t)
    items = _get_items(c)
    p = t.palette

    n = len(items)
    badge_x = MARGIN + 600000
    available_h = FOOTER_TOP - content_top - 400000
    step_h = available_h // n
    badge_r = 180000
    card_left = badge_x + 500000
    card_w = SLIDE_W - MARGIN - card_left
    card_h = min(step_h - 80000, 900000)

    for i, (date, title, desc) in enumerate(items):
        accent = p[i % len(p)]
        node_y = content_top + 200000 + i * step_h + card_h // 2

        # Dotted vertical connector to next badge
        if i < n - 1:
            next_y = content_top + 200000 + (i + 1) * step_h + card_h // 2
            add_line(slide, badge_x, node_y + badge_r + 20000,
                     badge_x, next_y - badge_r - 20000,
                     color=t.accent, width=1.5, dash="dash")

        # Outer badge border ring
        add_circle(slide, badge_x, node_y, badge_r,
                   fill_color="#FFFFFF", line_color=accent, line_width=2.5)

        # Inner badge border ring
        add_circle(slide, badge_x, node_y, badge_r - 40000,
                   fill_color="#FFFFFF", line_color=accent, line_width=0.75)

        # Milestone number inside badge
        add_textbox(slide, f"M{i + 1}", badge_x - badge_r, node_y - badge_r,
                    badge_r * 2, badge_r * 2,
                    font_size=13, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Dotted horizontal connector to card
        card_top = content_top + 200000 + i * step_h
        add_line(slide, badge_x + badge_r + 20000, node_y, card_left, node_y,
                 color=tint(accent, 0.5), width=1, dash="dash")

        # Card
        add_styled_card(slide, card_left, card_top, card_w, card_h,
                        theme=t, accent_color=accent)

        # Date
        add_textbox(slide, date, card_left + 80000, card_top + 60000,
                    card_w - 160000, 220000,
                    font_size=10, bold=True, color=accent,
                    alignment=PP_ALIGN.LEFT)

        # Title
        add_textbox(slide, title, card_left + 80000, card_top + 300000,
                    card_w - 160000, 280000,
                    font_size=13, bold=True, color=t.primary,
                    alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, card_left + 80000, card_top + 590000,
                    card_w - 160000, 280000,
                    font_size=11, color=t.secondary,
                    alignment=PP_ALIGN.LEFT)

        # Decorative triple dots
        dot_y = card_top + card_h - 90000
        for d in range(3):
            add_circle(slide, card_left + card_w - 200000 + d * 60000,
                       dot_y, 15000, fill_color=accent)


# ── Variant 11: MAGAZINE — Creative staggered cards with oversized dates ──

def _mosaic_path(slide, t, c):
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones"), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(items)
    gap = int(80000 * s.gap_factor)
    total_w = SLIDE_W - 2 * m

    # Row 1: up to 3 cards
    row1_count = min(n, 3)
    row1_w = (total_w - gap * (row1_count - 1)) // row1_count
    row1_h = 1800000
    row1_top = content_top + 200000

    # Row 2: remaining cards
    row2_count = min(n - row1_count, 3)
    row2_w = ((total_w - gap * max(row2_count - 1, 0))
              // max(row2_count, 1)) if row2_count else 0
    row2_h = 1800000
    row2_top = row1_top + row1_h + gap

    for i, (date, title, desc) in enumerate(items):
        accent = p[i % len(p)]

        if i < row1_count:
            card_left = m + i * (row1_w + gap)
            card_top = row1_top
            card_w = row1_w
            card_h = row1_h
        elif i < row1_count + row2_count:
            j = i - row1_count
            row2_total = row2_count * row2_w + (row2_count - 1) * gap
            row2_start = m + (total_w - row2_total) // 2
            card_left = row2_start + j * (row2_w + gap)
            card_top = row2_top
            card_w = row2_w
            card_h = row2_h
        else:
            continue

        # Full-color mosaic block
        add_rect(slide, card_left, card_top, card_w, card_h,
                 fill_color=accent, corner_radius=0)

        # Oversized date watermark
        add_textbox(slide, date, card_left + 60000, card_top + 80000,
                    card_w - 120000, 700000,
                    font_size=32, bold=True, color=tint(accent, 0.25),
                    alignment=PP_ALIGN.LEFT)

        # Milestone number
        add_textbox(slide, f"M{i + 1}", card_left + card_w - 500000,
                    card_top + 80000,
                    400000, 350000,
                    font_size=14, bold=True, color=tint(accent, 0.5),
                    alignment=PP_ALIGN.RIGHT)

        # Title
        add_textbox(slide, title.upper(), card_left + 60000, card_top + 800000,
                    card_w - 120000, 400000,
                    font_size=15, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, card_left + 60000, card_top + 1250000,
                    card_w - 120000, 450000,
                    font_size=11, color=tint(accent, 0.7),
                    alignment=PP_ALIGN.LEFT)


# ── Variant 12: SCHOLARLY — Numbered figure, thin timeline, centered ──

def _scholarly_path(slide, t, c):
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones"), theme=t)
    items = _get_items(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(items)
    total_w = SLIDE_W - 2 * m
    step_w = total_w // n
    line_y = content_top + 1600000
    node_r = 90000

    # Thin horizontal timeline
    first_cx = m + step_w // 2
    last_cx = m + (n - 1) * step_w + step_w // 2
    add_line(slide, first_cx, line_y, last_cx, line_y,
             color=tint(t.secondary, 0.6), width=0.75)

    for i, (date, title, desc) in enumerate(items):
        cx = m + i * step_w + step_w // 2

        # Simple circle node on line
        add_circle(slide, cx, line_y, node_r,
                   fill_color="#FFFFFF", line_color=t.accent, line_width=1.0)

        # Number inside node
        add_textbox(slide, str(i + 1), cx - node_r, line_y - node_r,
                    node_r * 2, node_r * 2,
                    font_size=11, bold=True, color=t.accent,
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Date above node
        add_textbox(slide, date, cx - step_w // 2 + 20000,
                    line_y - node_r - 380000,
                    step_w - 40000, 250000,
                    font_size=9, bold=False, color=t.secondary,
                    alignment=PP_ALIGN.CENTER)

        # Title below node
        add_textbox(slide, title, cx - step_w // 2 + 20000,
                    line_y + node_r + 180000,
                    step_w - 40000, 300000,
                    font_size=12, bold=True, color=t.primary,
                    alignment=PP_ALIGN.CENTER)

        # Description further below
        add_textbox(slide, desc, cx - step_w // 2 + 20000,
                    line_y + node_r + 500000,
                    step_w - 40000, 350000,
                    font_size=10, color=t.secondary,
                    alignment=PP_ALIGN.CENTER)

    # Figure caption at bottom
    fig_label = f"Figure. Project milestone timeline ({n} phases)"
    add_textbox(slide, fig_label, m, FOOTER_TOP - 400000,
                total_w, 250000,
                font_size=9, italic=True, color=t.secondary,
                alignment=PP_ALIGN.CENTER)

    # Thin rule above caption
    add_line(slide, m + total_w // 4, FOOTER_TOP - 480000,
             m + total_w * 3 // 4, FOOTER_TOP - 480000,
             color=tint(t.secondary, 0.7), width=0.5)


# ── Variant 13: LABORATORY — Dark bg, data-coded nodes, grid overlay ──

def _laboratory_path(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones"), theme=t)
    items = _get_items(c)
    p = t.palette

    n = len(items)
    total_w = SLIDE_W - 2 * MARGIN
    step_w = total_w // n
    node_y = content_top + 1500000
    node_r = 100000

    # Subtle grid overlay: vertical lines
    grid_spacing = 800000
    grid_color = tint(t.primary, 0.08)
    for gx in range(MARGIN, SLIDE_W - MARGIN + 1, grid_spacing):
        add_line(slide, gx, content_top, gx, FOOTER_TOP - 200000,
                 color=grid_color, width=0.25)

    # Subtle grid overlay: horizontal lines
    for gy in range(content_top, FOOTER_TOP - 200000, grid_spacing):
        add_line(slide, MARGIN, gy, SLIDE_W - MARGIN, gy,
                 color=grid_color, width=0.25)

    # Main horizontal path
    first_cx = MARGIN + step_w // 2
    last_cx = MARGIN + (n - 1) * step_w + step_w // 2
    add_line(slide, first_cx, node_y, last_cx, node_y,
             color=tint(t.primary, 0.2), width=1.5)

    for i, (date, title, desc) in enumerate(items):
        accent = p[i % len(p)]
        cx = MARGIN + i * step_w + step_w // 2

        # Data-coded node: outer square (rect) with inner circle
        sq_size = node_r + 50000
        add_rect(slide, cx - sq_size, node_y - sq_size,
                 sq_size * 2, sq_size * 2,
                 fill_color=tint(accent, 0.15),
                 line_color=accent, line_width=0.75)
        add_circle(slide, cx, node_y, node_r, fill_color=accent)

        # Node label
        add_textbox(slide, f"M{i + 1}", cx - node_r, node_y - node_r,
                    node_r * 2, node_r * 2,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Date above node (monospace feel)
        add_textbox(slide, date, cx - step_w // 2 + 20000,
                    node_y - sq_size - 380000,
                    step_w - 40000, 250000,
                    font_size=9, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER)

        # Evidence box below node
        box_top = node_y + sq_size + 200000
        box_w = step_w - 100000
        box_h = 900000
        box_left = cx - box_w // 2
        add_rect(slide, box_left, box_top, box_w, box_h,
                 fill_color=tint(t.primary, 0.08),
                 line_color=tint(t.primary, 0.2), line_width=0.5)

        # Color-coded left accent on evidence box
        add_rect(slide, box_left, box_top, 35000, box_h,
                 fill_color=accent)

        # Title inside box
        add_textbox(slide, title, box_left + 60000, box_top + 60000,
                    box_w - 120000, 350000,
                    font_size=12, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT)

        # Description inside box
        add_textbox(slide, desc, box_left + 60000, box_top + 420000,
                    box_w - 120000, 400000,
                    font_size=10, color=tint(t.primary, 0.45),
                    alignment=PP_ALIGN.LEFT)


# ── Variant 14: DASHBOARD — Dense compact timeline with status indicators ──

def _dashboard_path(slide, t, c):
    content_top = add_slide_title(
        slide, c.get("milestone_title", "Project Milestones"), theme=t)
    items = _get_items(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(items)
    total_w = SLIDE_W - 2 * m

    # Header band
    header_h = 350000
    header_top = content_top + 100000
    add_rect(slide, m, header_top, total_w, header_h,
             fill_color=t.primary)

    # Header columns: #, Date, Milestone, Description, Status
    col_num_w = 400000
    col_date_w = 1200000
    col_title_w = 2600000
    col_desc_w = total_w - col_num_w - col_date_w - col_title_w - 1200000
    col_status_w = 1200000

    headers = ["#", "Date", "Milestone", "Description", "Status"]
    col_widths = [col_num_w, col_date_w, col_title_w, col_desc_w, col_status_w]
    col_left = m
    for hi, (header, hw) in enumerate(zip(headers, col_widths)):
        add_textbox(slide, header, col_left + 40000, header_top + 60000,
                    hw - 80000, 230000,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT)
        col_left += hw

    # Row entries
    row_h = int((FOOTER_TOP - header_top - header_h - 400000) / n)
    row_h = min(row_h, 650000)

    statuses = ["Complete", "In Progress", "Planned"]

    for i, (date, title, desc) in enumerate(items):
        accent = p[i % len(p)]
        row_top = header_top + header_h + i * row_h

        # Alternating row background
        if i % 2 == 0:
            add_rect(slide, m, row_top, total_w, row_h,
                     fill_color=tint(t.primary, 0.95))

        col_left = m

        # Number column
        add_textbox(slide, str(i + 1), col_left + 40000, row_top + 30000,
                    col_num_w - 80000, row_h - 60000,
                    font_size=11, bold=True, color=t.accent,
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)
        col_left += col_num_w

        # Date column
        add_textbox(slide, date, col_left + 40000, row_top + 30000,
                    col_date_w - 80000, row_h - 60000,
                    font_size=10, bold=True, color=t.primary,
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)
        col_left += col_date_w

        # Title column
        add_textbox(slide, title, col_left + 40000, row_top + 30000,
                    col_title_w - 80000, row_h - 60000,
                    font_size=11, bold=True, color=t.primary,
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)
        col_left += col_title_w

        # Description column
        add_textbox(slide, desc, col_left + 40000, row_top + 30000,
                    col_desc_w - 80000, row_h - 60000,
                    font_size=10, color=t.secondary,
                    alignment=PP_ALIGN.LEFT,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)
        col_left += col_desc_w

        # Status indicator column: colored pill + label
        status = statuses[min(i, len(statuses) - 1)]
        if i == 0:
            status = "Complete"
            status_color = "#10B981"
        elif i < n - 1:
            status = "In Progress" if i < n // 2 else "Planned"
            status_color = accent if i < n // 2 else tint(t.secondary, 0.4)
        else:
            status = "Planned"
            status_color = tint(t.secondary, 0.4)

        pill_w = 900000
        pill_h = 240000
        pill_left = col_left + (col_status_w - pill_w) // 2
        pill_top = row_top + (row_h - pill_h) // 2
        add_rect(slide, pill_left, pill_top, pill_w, pill_h,
                 fill_color=status_color, corner_radius=120000)
        add_textbox(slide, status, pill_left, pill_top,
                    pill_w, pill_h,
                    font_size=9, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER,
                    vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Bottom border line
    bottom_y = header_top + header_h + n * row_h
    add_line(slide, m, bottom_y, m + total_w, bottom_y,
             color=tint(t.secondary, 0.5), width=1)

    # Summary bar under table
    add_textbox(slide, f"{n} milestones tracked", m, bottom_y + 80000,
                total_w, 250000,
                font_size=9, color=t.secondary,
                alignment=PP_ALIGN.RIGHT)
