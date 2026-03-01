"""Slide 24 — Three Column: 3 service pillars with circles and accent lines."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, FONT_SLIDE_TITLE, FONT_CARD_HEADING,
    CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_circle, add_line, add_rect,
    add_styled_card, add_slide_title, add_dark_bg, add_background,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    p = t.palette

    s = t.ux_style
    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)

    content_top = add_slide_title(slide, c.get("pillars_title", "Key Focus Areas"), theme=t)

    s = t.ux_style
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    pillars = c.get("pillars", [
        ("Consulting", "Strategic advisory."), ("Technology", "Enterprise solutions."), ("Analytics", "Data-driven insights."),
    ])

    ct = content_top + 100000
    circle_radius = 280000

    for i, (title, desc) in enumerate(pillars[:3]):
        left, width = col_span(3, i, gap=300000)
        center_x = left + width // 2
        accent = p[i % len(p)]

        add_styled_card(slide, left, ct, width, 4200000, theme=t, accent_color=accent)
        circle_y = ct + 650000
        add_circle(slide, center_x, circle_y, circle_radius, fill_color=t.accent)
        add_textbox(slide, f"0{i+1}", center_x - circle_radius, circle_y - circle_radius,
                    circle_radius * 2, circle_radius * 2,
                    font_size=24, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, title, left + 100000, circle_y + circle_radius + 250000, width - 200000, 400000,
                    font_size=FONT_CARD_HEADING, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, left + 150000, circle_y + circle_radius + 700000, width - 300000, 1200000,
                    font_size=13, color=sub_color, alignment=PP_ALIGN.CENTER)
        add_line(slide, left + width // 4, ct + 4200000 - 60000,
                 left + 3 * width // 4, ct + 4200000 - 60000, color=accent, width=3)


# ── Scholarly variant: white bg, 3 text columns with thin vertical rules ─

def _scholarly(slide, t):
    c = t.content
    add_background(slide, "#FFFFFF")
    content_top = add_slide_title(slide, c.get("pillars_title", "Key Focus Areas"), theme=t)

    pillars = c.get("pillars", [
        ("Consulting", "Strategic advisory."),
        ("Technology", "Enterprise solutions."),
        ("Analytics", "Data-driven insights."),
    ])

    ct = content_top + 100000
    col_h = 4000000

    for i, (title, desc) in enumerate(pillars[:3]):
        left, width = col_span(3, i, gap=400000)

        add_textbox(slide, title, left, ct, width, 400000,
                    font_size=FONT_CARD_HEADING, bold=True, color=t.primary,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, left + 80000, ct + 500000, width - 160000, col_h - 600000,
                    font_size=13, color=t.secondary, alignment=PP_ALIGN.CENTER)

    # Thin vertical rules between columns
    for i in range(2):
        left_col, w_col = col_span(3, i, gap=400000)
        rule_x = left_col + w_col + 200000
        add_line(slide, rule_x, ct, rule_x, ct + col_h,
                 color=tint(t.secondary, 0.7), width=0.5)


# ── Laboratory variant: dark bg, 3 dark cards with colored top accent bars ─

def _laboratory(slide, t):
    c = t.content
    p = t.palette
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("pillars_title", "Key Focus Areas"), theme=t)

    pillars = c.get("pillars", [
        ("Consulting", "Strategic advisory."),
        ("Technology", "Enterprise solutions."),
        ("Analytics", "Data-driven insights."),
    ])

    ct = content_top + 100000

    for i, (title, desc) in enumerate(pillars[:3]):
        left, width = col_span(3, i, gap=300000)
        accent = p[i % len(p)]
        card_h = 4200000

        # Dark card
        add_rect(slide, left, ct, width, card_h,
                 fill_color=tint(t.primary, 0.08), corner_radius=60000)
        # Colored top accent bar
        add_rect(slide, left, ct, width, 40000, fill_color=accent)

        add_textbox(slide, title, left + 100000, ct + 200000, width - 200000, 400000,
                    font_size=FONT_CARD_HEADING, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, left + 120000, ct + 700000, width - 240000, card_h - 900000,
                    font_size=13, color=tint(t.primary, 0.55),
                    alignment=PP_ALIGN.CENTER)


# ── Dashboard variant: header band, 3 compact cards with gradient accents ─

def _dashboard(slide, t):
    c = t.content
    p = t.palette
    add_background(slide, "#FFFFFF")

    # Accent header band
    add_rect(slide, 0, 0, SLIDE_W, 80000, fill_color=t.accent)

    content_top = add_slide_title(slide, c.get("pillars_title", "Key Focus Areas"), theme=t)

    pillars = c.get("pillars", [
        ("Consulting", "Strategic advisory."),
        ("Technology", "Enterprise solutions."),
        ("Analytics", "Data-driven insights."),
    ])

    ct = content_top + 50000

    for i, (title, desc) in enumerate(pillars[:3]):
        left, width = col_span(3, i, gap=250000)
        accent = p[i % len(p)]
        card_h = 4300000

        # Shadow behind card
        add_rect(slide, left + 15000, ct + 15000, width, card_h,
                 fill_color=tint(t.secondary, 0.85), corner_radius=40000)
        add_styled_card(slide, left, ct, width, card_h, theme=t, accent_color=accent)

        # Gradient-style top accent (two-tone thin bars)
        add_rect(slide, left, ct, width, 30000, fill_color=accent)
        add_rect(slide, left, ct + 30000, width, 15000,
                 fill_color=tint(accent, 0.5))

        add_textbox(slide, title, left + 80000, ct + 120000, width - 160000, 380000,
                    font_size=FONT_CARD_HEADING, bold=True, color=t.primary,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, left + 100000, ct + 550000, width - 200000, card_h - 700000,
                    font_size=12, color=t.secondary, alignment=PP_ALIGN.CENTER)
