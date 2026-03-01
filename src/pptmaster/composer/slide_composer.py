"""Slide Composer — orchestrates PPTX generation from outline + design profile.

Pipeline: open template → delete sample slides → for each slide: match layout →
add slide → render content → place icons → set footer → add notes → save.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from pptmaster.analyzer.potx_handler import ensure_openable
from pptmaster.assets.color_utils import hex_to_rgb, tint
from pptmaster.assets.icon_manager import IconManager
from pptmaster.models import (
    DesignProfile,
    LayoutInfo,
    PresentationOutline,
    SlideContent,
    SlideType,
)

from .chart_renderer import add_chart
from .layout_matcher import match_layout
from .shape_renderer import add_separator_line, add_background_rect, clone_icon
from .table_renderer import add_table
from .text_renderer import set_body_text, set_bullets, set_metric_text, set_subtitle_text, set_title
from .theme_applier import apply_theme

_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def compose(
    outline: PresentationOutline,
    profile: DesignProfile,
    output_path: str | Path,
    layout_assignments: list[tuple[SlideContent, LayoutInfo | None]] | None = None,
) -> Path:
    """Compose a PPTX presentation from an outline and design profile."""
    output_path = Path(output_path)

    # 1. Open the template
    openable = ensure_openable(profile.template_path)
    prs = Presentation(openable)

    # 2. Delete existing sample slides
    _delete_sample_slides(prs)

    # 3. Build icon manager
    icon_mgr = IconManager(profile.icons) if profile.icons else None

    # 4. Add slides
    used_layout_indices: list[int] = []

    for i, slide_content in enumerate(outline.slides):
        layout: LayoutInfo | None = None
        if layout_assignments and i < len(layout_assignments):
            _, layout = layout_assignments[i]

        if layout is None:
            layout = match_layout(slide_content, profile, used_layout_indices)

        if layout is None:
            continue

        used_layout_indices.append(layout.index)

        slide_layout = prs.slide_layouts[layout.index]
        slide = prs.slides.add_slide(slide_layout)

        _render_slide(slide, slide_content, layout, profile, icon_mgr)

        # Clean up empty placeholders
        _hide_empty_placeholders(slide)

        if slide_content.speaker_notes:
            _set_speaker_notes(slide, slide_content.speaker_notes)

    # 5. Add slide numbers
    _add_slide_numbers(prs, profile)

    # 6. Final theme consistency pass
    apply_theme(prs, profile)

    # 6. Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    return output_path


# ── Slide deletion ────────────────────────────────────────────────────


def _delete_sample_slides(prs: Presentation) -> None:
    """Delete all existing slides from the presentation via XML manipulation."""
    prs_elem = prs.element
    sldIdLst = prs_elem.find(f"{{{_NS_P}}}sldIdLst")
    if sldIdLst is None:
        return

    rids_to_remove = []
    for sldId in list(sldIdLst):
        rid = sldId.get(f"{{{_NS_R}}}id")
        if rid:
            rids_to_remove.append(rid)
        sldIdLst.remove(sldId)

    for rid in rids_to_remove:
        try:
            prs.part.rels.pop(rid)
        except KeyError:
            pass


# ── Placeholder classification ────────────────────────────────────────


def _classify_placeholders(slide) -> dict[str, list]:
    """Classify slide placeholders into typed groups.

    Returns dict with keys: 'title', 'body', 'pic', 'generic', 'other'
    Each value is a list of placeholder shapes sorted by left position.
    """
    result: dict[str, list] = {
        "title": [], "body": [], "pic": [], "generic": [], "other": []
    }

    for ph in slide.placeholders:
        xml_type = _get_ph_xml_type(ph)

        if xml_type in ("title", "ctrTitle"):
            result["title"].append(ph)
        elif xml_type in ("ftr", "dt", "sldNum"):
            result["other"].append(ph)
        elif xml_type == "pic":
            result["pic"].append(ph)
        elif xml_type in ("body", "subTitle"):
            result["body"].append(ph)
        else:
            # Generic/object placeholder — classify by size heuristic
            # Wide placeholders are likely body, square ones could be icons
            try:
                w = ph.width or 0
                h = ph.height or 0
                if w > 0 and h > 0 and h / w > 0.8 and w < 5000000:
                    result["pic"].append(ph)  # Roughly square and small = image
                else:
                    result["generic"].append(ph)
            except Exception:
                result["generic"].append(ph)

    # Sort body and generic by left position (for multi-column ordering)
    for key in ("body", "generic", "pic"):
        result[key].sort(key=lambda p: (p.left or 0))

    return result


def _get_ph_xml_type(ph) -> str:
    """Get the XML placeholder type string."""
    elem = ph._element
    nvPr = elem.find(f".//{{{_NS_P}}}nvPr")
    if nvPr is not None:
        ph_elem = nvPr.find(f"{{{_NS_P}}}ph")
        if ph_elem is not None:
            return ph_elem.get("type", "")
    return ""


def _get_text_placeholders(slide) -> list:
    """Get body + generic placeholders suitable for text content."""
    classified = _classify_placeholders(slide)
    return classified["body"] + classified["generic"]


# ── Main render dispatch ──────────────────────────────────────────────


def _render_slide(
    slide,
    content: SlideContent,
    layout: LayoutInfo,
    profile: DesignProfile,
    icon_mgr: IconManager | None,
) -> None:
    """Render content onto a slide based on its type."""
    if content.title:
        set_title(slide, content.title, profile)

    renderer = _RENDERERS.get(content.slide_type, _render_content_text)
    renderer(slide, content, layout, profile, icon_mgr)


# ── Slide type renderers ──────────────────────────────────────────────


def _render_title_slide(
    slide, content: SlideContent, layout: LayoutInfo,
    profile: DesignProfile, icon_mgr: IconManager | None,
) -> None:
    """Render a title slide with subtitle."""
    text_phs = _get_text_placeholders(slide)
    subtitle = content.subtitle or content.body or ""
    if text_phs and subtitle:
        set_body_text(text_phs[0], subtitle, profile)


def _render_content_text(
    slide, content: SlideContent, layout: LayoutInfo,
    profile: DesignProfile, icon_mgr: IconManager | None,
) -> None:
    """Render a text content slide with body and/or bullets."""
    text_phs = _get_text_placeholders(slide)
    if not text_phs:
        return

    if content.bullets:
        set_bullets(text_phs[0], content.bullets, profile, font_size=22)
    elif content.body:
        # Split long body text into bullet-like sentences for better readability
        sentences = _split_to_bullets(content.body)
        if len(sentences) > 1:
            set_bullets(text_phs[0], sentences, profile, font_size=22)
        else:
            set_body_text(text_phs[0], content.body, profile)

    # If we have a subtitle, add it as a supporting line
    if content.subtitle:
        sw = profile.slide_size.width
        sh = profile.slide_size.height
        set_subtitle_text(
            slide, content.subtitle, profile,
            left=int(sw * 0.06), top=int(sh * 0.85),
            width=int(sw * 0.88), height=int(sh * 0.08),
        )

    _add_visual_separator(slide, profile)


def _render_content_image(
    slide, content: SlideContent, layout: LayoutInfo,
    profile: DesignProfile, icon_mgr: IconManager | None,
) -> None:
    """Render a content + image slide. Text in body placeholders only."""
    text_phs = _get_text_placeholders(slide)
    if text_phs:
        if content.bullets:
            set_bullets(text_phs[0], content.bullets, profile)
        elif content.body:
            set_body_text(text_phs[0], content.body, profile)

    _add_visual_separator(slide, profile)


def _render_multi_column(
    slide, content: SlideContent, layout: LayoutInfo,
    profile: DesignProfile, icon_mgr: IconManager | None,
) -> None:
    """Render a multi-column slide with heading + bullets per column."""
    text_phs = _get_text_placeholders(slide)
    sw = profile.slide_size.width
    sh = profile.slide_size.height

    if content.columns and text_phs:
        accent_colors = profile.colors.accent_colors
        for i, col in enumerate(content.columns):
            if i < len(text_phs):
                col_color = accent_colors[i % len(accent_colors)]
                _render_column_content(text_phs[i], col, profile, heading_color=col_color)
                # Add accent bar above column placeholder
                ph = text_phs[i]
                bar_height = int(sh * 0.008)
                _add_colored_rect(
                    slide,
                    ph.left or 0,
                    (ph.top or 0) - bar_height - int(sh * 0.005),
                    ph.width or int(sw * 0.3),
                    bar_height,
                    col_color,
                )
    elif content.columns and not text_phs:
        # No placeholders — create text boxes programmatically
        _render_columns_as_textboxes(slide, content, profile, icon_mgr)
    elif content.bullets and text_phs:
        set_bullets(text_phs[0], content.bullets, profile)

    _add_visual_separator(slide, profile)

    # Add colored circle indicators above column headings
    if content.columns and content.slide_type in (SlideType.THREE_COLUMN, SlideType.FOUR_COLUMN):
        _add_column_icons(slide, content, profile, icon_mgr)


def _render_key_metrics(
    slide, content: SlideContent, layout: LayoutInfo,
    profile: DesignProfile, icon_mgr: IconManager | None,
) -> None:
    """Render a key metrics slide with big numbers in colored cards."""
    if not content.metrics:
        _render_content_text(slide, content, layout, profile, icon_mgr)
        return

    sw = profile.slide_size.width
    sh = profile.slide_size.height
    n = len(content.metrics)
    accent_colors = profile.colors.accent_colors

    # Layout: evenly spaced cards with colored top accent bar
    margin_x = int(sw * 0.06)
    gap = int(sw * 0.025)
    start_y = int(sh * 0.30)
    total_width = sw - 2 * margin_x
    card_width = (total_width - (n - 1) * gap) // n
    card_height = int(sh * 0.45)
    accent_bar_height = int(sh * 0.018)

    for i, metric in enumerate(content.metrics):
        color = metric.color or accent_colors[i % len(accent_colors)]
        left = margin_x + i * (card_width + gap)

        # Colored accent bar at top of each card
        _add_colored_rect(slide, left, start_y, card_width, accent_bar_height, color)

        # Light background card
        card_bg_color = tint(color, 0.92)
        _add_colored_rect(
            slide, left, start_y + accent_bar_height,
            card_width, card_height - accent_bar_height, card_bg_color
        )

        # Metric value and label
        set_metric_text(
            slide, metric.value, metric.label,
            left=left + int(card_width * 0.05),
            top=start_y + accent_bar_height + int(card_height * 0.08),
            width=int(card_width * 0.9),
            height=int(card_height * 0.75),
            profile=profile, color=color,
        )

    _add_visual_separator(slide, profile)


def _render_chart(
    slide, content: SlideContent, layout: LayoutInfo,
    profile: DesignProfile, icon_mgr: IconManager | None,
) -> None:
    """Render a chart slide. Chart takes most of the slide area."""
    if not content.chart:
        _render_content_text(slide, content, layout, profile, icon_mgr)
        return

    sw = profile.slide_size.width
    sh = profile.slide_size.height

    # If there's also a subtitle/body, put it at the top
    text_phs = _get_text_placeholders(slide)
    chart_top = int(sh * 0.20)

    if content.body and text_phs:
        set_body_text(text_phs[0], content.body, profile)
        chart_top = int(sh * 0.32)

    add_chart(
        slide, content.chart, profile,
        left=int(sw * 0.06),
        top=chart_top,
        width=int(sw * 0.88),
        height=int(sh * 0.90) - chart_top,
    )


def _render_table(
    slide, content: SlideContent, layout: LayoutInfo,
    profile: DesignProfile, icon_mgr: IconManager | None,
) -> None:
    """Render a table slide."""
    if not content.table:
        _render_content_text(slide, content, layout, profile, icon_mgr)
        return

    sw = profile.slide_size.width
    sh = profile.slide_size.height

    num_rows = len(content.table.rows) + (1 if content.table.headers else 0)
    # Dynamic height based on row count — ensure sufficient row height for text
    row_height = min(int(sh * 0.10), int(sh * 0.65 / max(num_rows, 1)))
    row_height = max(row_height, 500000)  # Minimum ~0.55 inch per row
    table_height = num_rows * row_height

    add_table(
        slide, content.table, profile,
        left=int(sw * 0.06),
        top=int(sh * 0.22),
        width=int(sw * 0.88),
        height=table_height,
    )


def _render_divider(
    slide, content: SlideContent, layout: LayoutInfo,
    profile: DesignProfile, icon_mgr: IconManager | None,
) -> None:
    """Render a section divider slide with subtitle."""
    text_phs = _get_text_placeholders(slide)
    if text_phs and content.subtitle:
        set_body_text(text_phs[0], content.subtitle, profile)
    elif content.subtitle:
        # No text placeholder — add a text box
        sw = profile.slide_size.width
        sh = profile.slide_size.height
        txBox = slide.shapes.add_textbox(
            Emu(int(sw * 0.08)), Emu(int(sh * 0.55)),
            Emu(int(sw * 0.84)), Emu(int(sh * 0.15)),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = content.subtitle
        run.font.name = profile.fonts.minor
        run.font.size = Pt(22)
        r, g, b = hex_to_rgb(profile.colors.dk2)
        run.font.color.rgb = RGBColor(r, g, b)


def _render_thank_you(
    slide, content: SlideContent, layout: LayoutInfo,
    profile: DesignProfile, icon_mgr: IconManager | None,
) -> None:
    """Render a thank you / closing slide."""
    text_phs = _get_text_placeholders(slide)
    if text_phs:
        text = content.body or content.subtitle or ""
        if content.bullets:
            set_bullets(text_phs[0], content.bullets, profile)
        elif text:
            set_body_text(text_phs[0], text, profile)


def _render_timeline(
    slide, content: SlideContent, layout: LayoutInfo,
    profile: DesignProfile, icon_mgr: IconManager | None,
) -> None:
    """Render a timeline slide."""
    _render_multi_column(slide, content, layout, profile, icon_mgr)


# ── Visual helpers ────────────────────────────────────────────────────


def _add_visual_separator(slide, profile: DesignProfile) -> None:
    """Add a thin separator line under the title area."""
    sw = profile.slide_size.width
    sh = profile.slide_size.height

    # Find the title's bottom edge for proper positioning
    title_bottom = int(sh * 0.18)
    classified = _classify_placeholders(slide)
    if classified["title"]:
        t = classified["title"][0]
        if t.top is not None and t.height is not None:
            title_bottom = t.top + t.height + int(sh * 0.01)

    add_separator_line(
        slide, profile,
        left=int(sw * 0.04),
        top=title_bottom,
        width=int(sw * 0.92),
    )


def _add_colored_rect(
    slide, left: int, top: int, width: int, height: int, color: str
) -> None:
    """Add a colored rectangle (for cards, accent bars, etc.)."""
    r, g, b = hex_to_rgb(color)
    shape = slide.shapes.add_shape(
        1,  # RECTANGLE
        Emu(left), Emu(top), Emu(width), Emu(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()  # No outline


def _render_column_content(
    placeholder, col, profile: DesignProfile, heading_color: str = "",
) -> None:
    """Render content for a single column into a placeholder."""
    from pptmaster.models import ColumnContent

    placeholder.text = ""
    tf = placeholder.text_frame
    tf.clear()
    tf.word_wrap = True

    hc = heading_color or profile.colors.accent1
    accent_r, accent_g, accent_b = hex_to_rgb(hc)
    text_r, text_g, text_b = hex_to_rgb(profile.colors.dk2)

    # Column heading (bold, accent color, larger)
    if col.heading:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = col.heading
        run.font.name = profile.fonts.minor
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(accent_r, accent_g, accent_b)
        p.space_after = Pt(10)

    # Bullets with colored bullet character
    for i, bullet in enumerate(col.bullets):
        p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(3)
        # Add colored bullet character
        from .text_renderer import _set_bullet_color, _set_bullet_char
        _set_bullet_char(p, "\u2022")
        _set_bullet_color(p, accent_r, accent_g, accent_b)
        run = p.add_run()
        run.text = bullet
        run.font.name = profile.fonts.minor
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(text_r, text_g, text_b)

    # Body text
    if col.body and not col.bullets:
        p = tf.add_paragraph() if col.heading else tf.paragraphs[0]
        run = p.add_run()
        run.text = col.body
        run.font.name = profile.fonts.minor
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(text_r, text_g, text_b)


def _render_columns_as_textboxes(
    slide, content: SlideContent, profile: DesignProfile,
    icon_mgr: IconManager | None,
) -> None:
    """Create multi-column layout using text boxes when no placeholders exist."""
    sw = profile.slide_size.width
    sh = profile.slide_size.height
    n = len(content.columns)
    if n == 0:
        return

    margin_x = int(sw * 0.06)
    gap = int(sw * 0.02)
    start_y = int(sh * 0.25)
    total_w = sw - 2 * margin_x
    col_w = (total_w - (n - 1) * gap) // n
    col_h = int(sh * 0.60)

    accent_colors = profile.colors.accent_colors

    for i, col in enumerate(content.columns):
        left = margin_x + i * (col_w + gap)
        color = accent_colors[i % len(accent_colors)]

        # Colored top accent bar
        _add_colored_rect(slide, left, start_y, col_w, int(sh * 0.012), color)

        # Text box for column content
        txBox = slide.shapes.add_textbox(
            Emu(left), Emu(start_y + int(sh * 0.02)),
            Emu(col_w), Emu(col_h),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        text_r, text_g, text_b = hex_to_rgb(profile.colors.dk2)
        cr, cg, cb = hex_to_rgb(color)

        # Heading
        if col.heading:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = col.heading
            run.font.name = profile.fonts.minor
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = RGBColor(cr, cg, cb)
            p.space_after = Pt(10)

        # Bullets
        for j, bullet in enumerate(col.bullets):
            p = tf.add_paragraph()
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = f"• {bullet}"
            run.font.name = profile.fonts.minor
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(text_r, text_g, text_b)

        # Body
        if col.body and not col.bullets:
            p = tf.add_paragraph() if col.heading else tf.paragraphs[0]
            run = p.add_run()
            run.text = col.body
            run.font.name = profile.fonts.minor
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(text_r, text_g, text_b)


def _add_column_icons(
    slide, content: SlideContent, profile: DesignProfile,
    icon_mgr: IconManager | None,
) -> None:
    """Add colored circle indicators with icon keyword above column headings."""
    sw = profile.slide_size.width
    sh = profile.slide_size.height
    n = len(content.columns)
    if n == 0:
        return

    accent_colors = profile.colors.accent_colors

    # Detect column positions from placeholders on the slide
    classified = _classify_placeholders(slide)
    text_phs = classified["body"] + classified["generic"]

    if text_phs and len(text_phs) >= n:
        # Position circles centered above each placeholder
        col_positions = []
        for j in range(min(n, len(text_phs))):
            ph = text_phs[j]
            center_x = (ph.left or 0) + (ph.width or 0) // 2
            ph_top = ph.top or int(sh * 0.25)
            col_positions.append((center_x, ph_top))
    else:
        # Fallback: evenly spaced
        margin_x = int(sw * 0.06)
        gap = int(sw * 0.02)
        total_w = sw - 2 * margin_x
        col_w = (total_w - (n - 1) * gap) // n
        col_positions = [
            (margin_x + j * (col_w + gap) + col_w // 2, int(sh * 0.25))
            for j in range(n)
        ]

    circle_size = int(sh * 0.07)  # Circle size

    for i, col in enumerate(content.columns):
        if i >= len(col_positions):
            break
        center_x, ph_top = col_positions[i]
        color = accent_colors[i % len(accent_colors)]

        # Position circle just above the placeholder
        circle_left = center_x - circle_size // 2
        circle_top = ph_top - circle_size - int(sh * 0.015)

        # Add colored circle
        from pptx.enum.shapes import MSO_SHAPE
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(circle_left), Emu(circle_top),
            Emu(circle_size), Emu(circle_size),
        )
        r, g, b = hex_to_rgb(color)
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(r, g, b)
        circle.line.fill.background()  # No outline

        # Try to clone a real icon from the icon manager
        kw = col.icon_keyword or col.heading or ""
        if kw and icon_mgr:
            icons = icon_mgr.search(kw, limit=1)
            if icons:
                try:
                    clone_icon(
                        slide, icons[0],
                        left=circle_left + circle_size // 4,
                        top=circle_top + circle_size // 4,
                        width=circle_size // 2,
                        height=circle_size // 2,
                        color="#FFFFFF",
                    )
                    continue
                except Exception:
                    pass

        # Fallback: put icon keyword abbreviation as text in circle
        _NS_A_local = "http://schemas.openxmlformats.org/drawingml/2006/main"
        try:
            txBody = circle._element.find(f"{{{_NS_A_local}}}txBody")
            if txBody is not None:
                bodyPr = txBody.find(f"{{{_NS_A_local}}}bodyPr")
                if bodyPr is not None:
                    bodyPr.set("anchor", "ctr")
        except Exception:
            pass

        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        # Use first 2 chars of icon keyword or heading
        abbrev = (col.icon_keyword or col.heading or "")[:2].upper()
        run.text = abbrev
        run.font.name = profile.fonts.minor
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)


# ── Other helpers ─────────────────────────────────────────────────────


def _hide_empty_placeholders(slide) -> None:
    """Remove or hide empty placeholders to avoid ugly empty boxes."""
    sp_tree = slide._element.find(f".//{{{_NS_P}}}cSld/{{{_NS_P}}}spTree")
    if sp_tree is None:
        return

    to_remove = []
    for ph in slide.placeholders:
        xml_type = _get_ph_xml_type(ph)

        # Skip title, footer, date, slide number — these are controlled by the layout
        if xml_type in ("title", "ctrTitle", "ftr", "dt", "sldNum"):
            continue

        # Check if this placeholder has meaningful content
        has_content = False
        if ph.has_text_frame:
            text = ph.text_frame.text.strip()
            if text:
                has_content = True
        # Picture placeholders with inserted images have a blipFill
        blip = ph._element.find(f".//{{{_NS_A}}}blipFill")
        if blip is not None:
            has_content = True

        if not has_content:
            to_remove.append(ph._element)

    for elem in to_remove:
        parent = elem.getparent()
        if parent is not None:
            parent.remove(elem)


def _add_slide_numbers(prs: Presentation, profile: DesignProfile) -> None:
    """Add slide numbers to each slide (bottom-right corner)."""
    sw = profile.slide_size.width
    sh = profile.slide_size.height
    total = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        # Skip title and thank you slides
        if i == 0 or i == total - 1:
            continue

        num_box = slide.shapes.add_textbox(
            Emu(int(sw * 0.92)), Emu(int(sh * 0.94)),
            Emu(int(sw * 0.06)), Emu(int(sh * 0.04)),
        )
        tf = num_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(i + 1)
        run.font.name = profile.fonts.minor
        run.font.size = Pt(12)
        r, g, b = hex_to_rgb(profile.colors.dk2)
        run.font.color.rgb = RGBColor(r, g, b)


def _split_to_bullets(text: str) -> list[str]:
    """Split a paragraph into bullet-friendly sentences."""
    import re
    # Split on sentence boundaries
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    # Filter out very short fragments
    bullets = [s.rstrip('.') for s in sentences if len(s.strip()) > 10]
    # Cap at 5 bullets
    return bullets[:5]


def _set_speaker_notes(slide, notes_text: str) -> None:
    """Add speaker notes to a slide."""
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text
    except Exception:
        pass


# ── Renderer dispatch table ───────────────────────────────────────────

_RENDERERS = {
    SlideType.TITLE: _render_title_slide,
    SlideType.CONTENT_TEXT: _render_content_text,
    SlideType.CONTENT_IMAGE_RIGHT: _render_content_image,
    SlideType.CONTENT_IMAGE_LEFT: _render_content_image,
    SlideType.TWO_COLUMN: _render_multi_column,
    SlideType.THREE_COLUMN: _render_multi_column,
    SlideType.FOUR_COLUMN: _render_multi_column,
    SlideType.KEY_METRICS: _render_key_metrics,
    SlideType.CHART_BAR: _render_chart,
    SlideType.CHART_LINE: _render_chart,
    SlideType.CHART_PIE: _render_chart,
    SlideType.CHART_DONUT: _render_chart,
    SlideType.TABLE: _render_table,
    SlideType.TIMELINE: _render_timeline,
    SlideType.PROFILE_GRID: _render_multi_column,
    SlideType.DIVIDER: _render_divider,
    SlideType.FULL_IMAGE: _render_content_image,
    SlideType.MISSION_VISION: _render_multi_column,
    SlideType.THANK_YOU: _render_thank_you,
}
