"""Slide 31 — Funnel Diagram: 14 unique visual variants dispatched by UX style."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_circle, add_line, add_funnel_tier,
    add_gold_accent_line, add_slide_title, add_dark_bg, add_background,
    add_styled_card, add_color_cell,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content
    p = t.palette
    s = t.ux_style

    dispatch = {
        "stacked": _stacked,
        "outline": _outline,
        "bold-gradient": _bold_gradient,
        "floating-funnel": _floating_funnel,
        "dark-funnel": _dark_funnel,
        "split-funnel": _split_funnel,
        "angular-funnel": _angular_funnel,
        "editorial-funnel": _editorial_funnel,
        "gradient-funnel": _gradient_funnel,
        "retro-funnel": _retro_funnel,
        "creative-funnel": _creative_funnel,
        "scholarly-funnel": _scholarly_funnel,
        "laboratory-funnel": _laboratory_funnel,
        "dashboard-funnel": _dashboard_funnel,
    }
    builder = dispatch.get(s.funnel, _stacked)
    builder(slide, t, c)


def _get_stages(c):
    """Extract funnel stages from content dict with sensible defaults."""
    return c.get("funnel_stages", [
        ("Awareness", "10,000", "Total market reach"),
        ("Interest", "5,200", "Engaged prospects"),
        ("Consideration", "2,800", "Qualified leads"),
        ("Intent", "1,400", "Sales pipeline"),
        ("Purchase", "680", "Converted customers"),
    ])


# ── Variant 1: CLASSIC — stacked trapezoids narrowing downward ────────

def _stacked(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("funnel_title", "Conversion Funnel"), theme=t)
    stages = _get_stages(c)
    p = t.palette

    n = len(stages)
    available_h = FOOTER_TOP - content_top - 300000
    tier_h = min(available_h // n - 40000, 700000)
    tier_gap = 50000
    # Leave 2200000 EMU for right-side value/desc labels
    max_w = SLIDE_W - 2 * MARGIN - 2200000
    min_w = max_w * 35 // 100

    text_color = "#FFFFFF"  # always on colored fill

    for i, (label, value, desc) in enumerate(stages):
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.65
        tier_w = int(max_w * fraction)
        tier_left = (SLIDE_W - tier_w) // 2 - 600000  # shift left to balance
        tier_top = content_top + 100000 + i * (tier_h + tier_gap)
        accent = p[i % len(p)]

        add_funnel_tier(slide, tier_left, tier_top, tier_w, tier_h,
                        fill_color=accent, text=label, text_color="#FFFFFF",
                        font_size=14)

        # Value on the right side — clamp to slide bounds
        val_left = tier_left + tier_w + 120000
        val_w = min(1200000, SLIDE_W - MARGIN - val_left)
        add_textbox(slide, value, val_left, tier_top + 20000,
                    val_w, tier_h // 2,
                    font_size=18, bold=True,
                    color=accent if not t.ux_style.dark_mode else accent,
                    alignment=PP_ALIGN.LEFT)

        # Description below value
        desc_w = min(1800000, SLIDE_W - MARGIN - val_left)
        add_textbox(slide, desc, val_left, tier_top + tier_h // 2 - 20000,
                    desc_w, tier_h // 2,
                    font_size=10,
                    color=tint(t.primary, 0.5) if t.ux_style.dark_mode else t.secondary,
                    alignment=PP_ALIGN.LEFT)


# ── Variant 2: MINIMAL — outlined trapezoids, no fill ─────────────────

def _outline(slide, t, c):
    content_top = add_slide_title(slide, c.get("funnel_title", "Conversion Funnel"), theme=t)
    stages = _get_stages(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(stages)
    available_h = FOOTER_TOP - content_top - 400000
    tier_h = min(available_h // n - 30000, 600000)
    tier_gap = 30000
    max_w = SLIDE_W - 2 * m - 600000
    center_x = SLIDE_W // 2

    for i, (label, value, desc) in enumerate(stages):
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.6
        tier_w = int(max_w * fraction)
        tier_left = center_x - tier_w // 2
        tier_top = content_top + 150000 + i * (tier_h + tier_gap)

        # Outlined rectangle (no fill, thin border) simulating funnel tier
        add_rect(slide, tier_left, tier_top, tier_w, tier_h,
                 line_color=tint(t.secondary, 0.5), line_width=1.0)

        # Label — left inside
        add_textbox(slide, label, tier_left + 150000, tier_top + 60000,
                    tier_w // 2 - 200000, tier_h - 120000,
                    font_size=13, bold=True, color=t.primary,
                    alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Value — right inside
        add_textbox(slide, value, tier_left + tier_w // 2, tier_top + 40000,
                    tier_w // 2 - 150000, tier_h // 2,
                    font_size=16, bold=True, color=t.accent,
                    alignment=PP_ALIGN.RIGHT)

        # Description — right inside below value
        add_textbox(slide, desc, tier_left + tier_w // 2, tier_top + tier_h // 2 - 20000,
                    tier_w // 2 - 150000, tier_h // 2,
                    font_size=10, color=t.secondary,
                    alignment=PP_ALIGN.RIGHT)

        # Small accent dot at left
        add_circle(slide, tier_left + 60000, tier_top + tier_h // 2, 25000,
                   fill_color=t.accent)


# ── Variant 3: BOLD — full-width bold colored blocks narrowing ─────────

def _bold_gradient(slide, t, c):
    content_top = add_slide_title(
        slide, c.get("funnel_title", "Conversion Funnel").upper(), theme=t)
    stages = _get_stages(c)
    p = t.palette

    n = len(stages)
    available_h = FOOTER_TOP - content_top - 200000
    tier_h = available_h // n
    max_w = SLIDE_W - 2 * MARGIN

    for i, (label, value, desc) in enumerate(stages):
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.55
        tier_w = int(max_w * fraction)
        tier_left = (SLIDE_W - tier_w) // 2
        tier_top = content_top + 100000 + i * tier_h
        accent = p[i % len(p)]

        # Full colored block
        add_rect(slide, tier_left, tier_top, tier_w, tier_h - 20000,
                 fill_color=accent)

        # Large value on the left
        add_textbox(slide, value, tier_left + 80000, tier_top + 30000,
                    1500000, tier_h - 80000,
                    font_size=28, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Label in center
        add_textbox(slide, label.upper(), tier_left + 1800000, tier_top + 20000,
                    tier_w - 3600000, tier_h - 60000,
                    font_size=16, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Description on the right
        add_textbox(slide, desc, tier_left + tier_w - 2200000, tier_top + 30000,
                    2100000, tier_h - 80000,
                    font_size=11, color=tint(accent, 0.7),
                    alignment=PP_ALIGN.RIGHT, vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant 4: ELEVATED — floating card tiers with shadow ─────────────

def _floating_funnel(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("funnel_title", "Conversion Funnel"), theme=t)
    stages = _get_stages(c)
    p = t.palette
    s = t.ux_style

    n = len(stages)
    available_h = FOOTER_TOP - content_top - 400000
    tier_h = min(available_h // n - 50000, 650000)
    tier_gap = 50000
    max_w = SLIDE_W - 2 * MARGIN - 600000

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    for i, (label, value, desc) in enumerate(stages):
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.55
        tier_w = int(max_w * fraction)
        tier_left = (SLIDE_W - tier_w) // 2
        tier_top = content_top + 200000 + i * (tier_h + tier_gap)
        accent = p[i % len(p)]

        # Shadow card behind
        shadow_off = s.card_shadow_offset
        add_rect(slide, tier_left + shadow_off, tier_top + shadow_off,
                 tier_w, tier_h,
                 fill_color=shade(t.light_bg if not s.dark_mode else t.primary, 0.1),
                 corner_radius=s.card_radius)

        # Main card
        fill = "#FFFFFF" if not s.dark_mode else tint(t.primary, 0.12)
        add_rect(slide, tier_left, tier_top, tier_w, tier_h,
                 fill_color=fill, corner_radius=s.card_radius)

        # Left accent bar
        add_rect(slide, tier_left, tier_top, 60000, tier_h,
                 fill_color=accent)

        # Label
        add_textbox(slide, label, tier_left + 150000, tier_top + 60000,
                    tier_w // 2 - 200000, tier_h - 120000,
                    font_size=14, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Value — right section
        add_textbox(slide, value, tier_left + tier_w - 1800000, tier_top + 30000,
                    1200000, tier_h // 2,
                    font_size=20, bold=True, color=accent,
                    alignment=PP_ALIGN.RIGHT)

        # Description — right section below value
        add_textbox(slide, desc, tier_left + tier_w - 1800000, tier_top + tier_h // 2,
                    1600000, tier_h // 2 - 30000,
                    font_size=10, color=sub_color,
                    alignment=PP_ALIGN.RIGHT)


# ── Variant 5: DARK — dark bg with glowing colored tiers ──────────────

def _dark_funnel(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(
        slide, c.get("funnel_title", "Conversion Funnel").upper(), theme=t)
    stages = _get_stages(c)
    p = t.palette

    n = len(stages)
    available_h = FOOTER_TOP - content_top - 300000
    tier_h = min(available_h // n - 40000, 650000)
    tier_gap = 40000
    max_w = SLIDE_W - 2 * MARGIN - 1000000

    for i, (label, value, desc) in enumerate(stages):
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.6
        tier_w = int(max_w * fraction)
        tier_left = (SLIDE_W - tier_w) // 2
        tier_top = content_top + 150000 + i * (tier_h + tier_gap)
        accent = p[i % len(p)]

        # Glow background (wider, slightly transparent effect)
        glow_w = tier_w + 100000
        glow_left = (SLIDE_W - glow_w) // 2
        add_rect(slide, glow_left, tier_top - 15000, glow_w, tier_h + 30000,
                 fill_color=tint(accent, 0.15), corner_radius=60000)

        # Main tier
        add_funnel_tier(slide, tier_left, tier_top, tier_w, tier_h,
                        fill_color=accent, text="", font_size=14)

        # Label inside (left-ish)
        add_textbox(slide, label.upper(), tier_left + 200000, tier_top + 20000,
                    tier_w - 400000, tier_h // 2,
                    font_size=12, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

        # Value inside (centered)
        add_textbox(slide, value, tier_left + 200000, tier_top + tier_h // 2 - 30000,
                    tier_w - 400000, tier_h // 2,
                    font_size=18, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

        # Description on far right — clamp width to slide bounds
        desc_left = (SLIDE_W + tier_w) // 2 + 120000
        desc_w = min(1800000, SLIDE_W - MARGIN - desc_left)
        if desc_w > 0:
            add_textbox(slide, desc, desc_left, tier_top + 20000,
                        desc_w, tier_h - 40000,
                        font_size=10, color=tint(t.primary, 0.45),
                        alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Bottom decorative line
    add_line(slide, MARGIN, FOOTER_TOP - 200000,
             SLIDE_W - MARGIN, FOOTER_TOP - 200000,
             color=tint(t.primary, 0.15), width=0.5)


# ── Variant 6: SPLIT — funnel on left, details list on right ──────────

def _split_funnel(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("funnel_title", "Conversion Funnel"), theme=t)
    stages = _get_stages(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    n = len(stages)
    # Left half: funnel
    funnel_area_w = (SLIDE_W - 2 * MARGIN) * 50 // 100
    funnel_left_base = MARGIN
    # Right half: details
    detail_left = MARGIN + funnel_area_w + 200000
    detail_w = SLIDE_W - MARGIN - detail_left

    available_h = FOOTER_TOP - content_top - 400000
    tier_h = min(available_h // n - 30000, 600000)
    tier_gap = 30000
    max_funnel_w = funnel_area_w - 200000

    # Vertical divider line
    div_x = MARGIN + funnel_area_w + 100000
    add_line(slide, div_x, content_top + 100000, div_x, FOOTER_TOP - 200000,
             color=tint(t.accent, 0.3 if s.dark_mode else 0.6), width=1)

    for i, (label, value, desc) in enumerate(stages):
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.6
        tier_w = int(max_funnel_w * fraction)
        tier_left = funnel_left_base + (funnel_area_w - tier_w) // 2
        tier_top = content_top + 200000 + i * (tier_h + tier_gap)
        accent = p[i % len(p)]

        # Funnel tier on the left
        add_funnel_tier(slide, tier_left, tier_top, tier_w, tier_h,
                        fill_color=accent, text=label, text_color="#FFFFFF",
                        font_size=12)

        # Detail row on the right
        row_top = tier_top

        # Accent dot
        add_circle(slide, detail_left + 60000, row_top + tier_h // 2, 40000,
                   fill_color=accent)

        # Value
        add_textbox(slide, value, detail_left + 160000, row_top + 30000,
                    detail_w - 200000, tier_h // 2,
                    font_size=18, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, detail_left + 160000, row_top + tier_h // 2 - 10000,
                    detail_w - 200000, tier_h // 2,
                    font_size=10, color=sub_color,
                    alignment=PP_ALIGN.LEFT)


# ── Variant 7: GEO — angular/geometric tiers with borders ─────────────

def _angular_funnel(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("funnel_title", "Conversion Funnel"), theme=t)
    stages = _get_stages(c)
    p = t.palette
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if s.dark_mode else t.secondary

    n = len(stages)
    available_h = FOOTER_TOP - content_top - 400000
    tier_h = min(available_h // n - 30000, 600000)
    tier_gap = 30000
    max_w = SLIDE_W - 2 * MARGIN - 800000

    for i, (label, value, desc) in enumerate(stages):
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.6
        tier_w = int(max_w * fraction)
        tier_left = (SLIDE_W - tier_w) // 2
        tier_top = content_top + 200000 + i * (tier_h + tier_gap)
        accent = p[i % len(p)]

        # Sharp-cornered bordered rect (geometric feel)
        add_rect(slide, tier_left, tier_top, tier_w, tier_h,
                 fill_color=tint(accent, 0.15 if s.dark_mode else 0.85),
                 line_color=accent, line_width=1.5)

        # Thick bottom accent bar
        add_rect(slide, tier_left, tier_top + tier_h - 40000, tier_w, 40000,
                 fill_color=accent)

        # Step number — angular box at far left
        num_w = 250000
        add_rect(slide, tier_left - num_w - 40000, tier_top, num_w, tier_h,
                 fill_color=accent)
        add_textbox(slide, f"{i + 1:02d}",
                    tier_left - num_w - 40000, tier_top, num_w, tier_h,
                    font_size=16, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Label
        add_textbox(slide, label, tier_left + 100000, tier_top + 40000,
                    tier_w // 2 - 150000, tier_h - 100000,
                    font_size=13, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Value
        add_textbox(slide, value, tier_left + tier_w // 2, tier_top + 20000,
                    tier_w // 2 - 120000, tier_h // 2,
                    font_size=18, bold=True, color=accent,
                    alignment=PP_ALIGN.RIGHT)

        # Description
        add_textbox(slide, desc, tier_left + tier_w // 2, tier_top + tier_h // 2 - 10000,
                    tier_w // 2 - 120000, tier_h // 2 - 40000,
                    font_size=10, color=sub_color,
                    alignment=PP_ALIGN.RIGHT)


# ── Variant 8: EDITORIAL — thin rules separating tiers, elegant ───────

def _editorial_funnel(slide, t, c):
    content_top = add_slide_title(slide, c.get("funnel_title", "Conversion Funnel"), theme=t)
    stages = _get_stages(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(stages)
    available_h = FOOTER_TOP - content_top - 400000
    row_h = available_h // n

    for i, (label, value, desc) in enumerate(stages):
        row_top = content_top + 200000 + i * row_h

        # Thin horizontal rule at top of each row
        add_line(slide, m, row_top, SLIDE_W - m, row_top,
                 color=tint(t.secondary, 0.8), width=0.5)

        # Short accent mark at start of rule
        add_line(slide, m, row_top, m + 400000, row_top,
                 color=t.accent, width=1.5)

        # Elegant step number — left column
        num_w = 600000
        add_textbox(slide, f"0{i + 1}", m, row_top + 100000,
                    num_w, 400000,
                    font_size=28, bold=True, color=t.accent,
                    alignment=PP_ALIGN.LEFT)

        # Label — second column
        label_left = m + num_w + 100000
        label_w = 2800000
        add_textbox(slide, label, label_left, row_top + 100000,
                    label_w, 350000,
                    font_size=15, bold=True, color=t.primary,
                    alignment=PP_ALIGN.LEFT)

        # Value — third column
        val_left = label_left + label_w + 200000
        val_w = 1500000
        add_textbox(slide, value, val_left, row_top + 100000,
                    val_w, 350000,
                    font_size=20, bold=True, color=t.accent,
                    alignment=PP_ALIGN.CENTER)

        # Description — right column
        desc_left = val_left + val_w + 200000
        desc_w = SLIDE_W - m - desc_left
        add_textbox(slide, desc, desc_left, row_top + 120000,
                    desc_w, 350000,
                    font_size=11, color=t.secondary,
                    alignment=PP_ALIGN.LEFT)

    # Bottom rule
    add_line(slide, m, FOOTER_TOP - 300000, SLIDE_W - m, FOOTER_TOP - 300000,
             color=tint(t.secondary, 0.8), width=0.5)


# ── Variant 9: GRADIENT — rounded tiers with progressive tinting ──────

def _gradient_funnel(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("funnel_title", "Conversion Funnel"), theme=t)
    stages = _get_stages(c)
    s = t.ux_style

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if s.dark_mode else t.secondary

    n = len(stages)
    available_h = FOOTER_TOP - content_top - 400000
    tier_h = min(available_h // n - 40000, 650000)
    tier_gap = 40000
    max_w = SLIDE_W - 2 * MARGIN - 600000

    for i, (label, value, desc) in enumerate(stages):
        # Progressive tint from accent
        factor = i * 0.15
        tier_color = tint(t.accent, factor)
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.55
        tier_w = int(max_w * fraction)
        tier_left = (SLIDE_W - tier_w) // 2
        tier_top = content_top + 200000 + i * (tier_h + tier_gap)

        # Outer glow rounded rect
        add_rect(slide, tier_left - 30000, tier_top - 15000,
                 tier_w + 60000, tier_h + 30000,
                 fill_color=tint(tier_color, 0.85),
                 corner_radius=120000)

        # Main rounded tier
        add_rect(slide, tier_left, tier_top, tier_w, tier_h,
                 fill_color=tier_color, corner_radius=100000)

        # Label centered
        add_textbox(slide, label, tier_left + 100000, tier_top + 20000,
                    tier_w - 200000, tier_h // 2,
                    font_size=13, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

        # Value centered below label
        add_textbox(slide, value, tier_left + 100000, tier_top + tier_h // 2 - 30000,
                    tier_w - 200000, tier_h // 2,
                    font_size=18, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

        # Description to the right
        desc_left = tier_left + tier_w + 100000
        add_textbox(slide, desc, desc_left, tier_top + 20000,
                    1800000, tier_h - 40000,
                    font_size=10, color=sub_color,
                    alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE)


# ── Variant 10: RETRO — vintage styled with decorative borders ────────

def _retro_funnel(slide, t, c):
    add_background(slide, t.light_bg)
    content_top = add_slide_title(slide, c.get("funnel_title", "Conversion Funnel"), theme=t)
    stages = _get_stages(c)
    p = t.palette

    n = len(stages)
    available_h = FOOTER_TOP - content_top - 400000
    tier_h = min(available_h // n - 40000, 600000)
    tier_gap = 40000
    max_w = SLIDE_W - 2 * MARGIN - 800000

    for i, (label, value, desc) in enumerate(stages):
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.55
        tier_w = int(max_w * fraction)
        tier_left = (SLIDE_W - tier_w) // 2
        tier_top = content_top + 200000 + i * (tier_h + tier_gap)
        accent = p[i % len(p)]

        # Outer decorative border
        add_rect(slide, tier_left - 20000, tier_top - 20000,
                 tier_w + 40000, tier_h + 40000,
                 line_color=accent, line_width=2.5,
                 corner_radius=100000)

        # Inner border
        add_rect(slide, tier_left, tier_top, tier_w, tier_h,
                 fill_color="#FFFFFF",
                 line_color=accent, line_width=0.75,
                 corner_radius=80000)

        # Step number badge at left
        badge_cx = tier_left + 60000
        badge_cy = tier_top + tier_h // 2
        add_circle(slide, badge_cx, badge_cy, 120000,
                   fill_color=accent)
        add_textbox(slide, str(i + 1), badge_cx - 100000, badge_cy - 100000,
                    200000, 200000,
                    font_size=14, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Label
        add_textbox(slide, label, tier_left + 280000, tier_top + 40000,
                    tier_w // 2 - 300000, tier_h - 80000,
                    font_size=13, bold=True, color=t.primary,
                    alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Value
        add_textbox(slide, value, tier_left + tier_w // 2, tier_top + 20000,
                    tier_w // 2 - 100000, tier_h // 2,
                    font_size=18, bold=True, color=accent,
                    alignment=PP_ALIGN.RIGHT)

        # Description
        add_textbox(slide, desc, tier_left + tier_w // 2, tier_top + tier_h // 2,
                    tier_w // 2 - 100000, tier_h // 2 - 20000,
                    font_size=10, color=t.secondary,
                    alignment=PP_ALIGN.RIGHT)

        # Decorative dots below each tier
        dot_y = tier_top + tier_h + 10000
        for d in range(3):
            add_circle(slide, (SLIDE_W // 2) - 60000 + d * 60000, dot_y, 10000,
                       fill_color=accent)


# ── Variant 11: MAGAZINE — creative/asymmetric with oversized numbers ──

def _creative_funnel(slide, t, c):
    content_top = add_slide_title(slide, c.get("funnel_title", "Conversion Funnel"), theme=t)
    stages = _get_stages(c)
    p = t.palette

    n = len(stages)
    available_h = FOOTER_TOP - content_top - 300000
    tier_h = available_h // n
    max_w = SLIDE_W - 2 * MARGIN

    # Alternating offsets for asymmetric feel
    x_offsets = [0, 200000, -150000, 250000, -100000]
    y_offsets = [0, 30000, -20000, 40000, -30000]

    for i, (label, value, desc) in enumerate(stages):
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.5
        tier_w = int(max_w * fraction)
        base_left = (SLIDE_W - tier_w) // 2 + x_offsets[i % len(x_offsets)]
        tier_top = content_top + 150000 + i * tier_h + y_offsets[i % len(y_offsets)]
        accent = p[i % len(p)]

        # Background tint block
        add_rect(slide, base_left, tier_top, tier_w, tier_h - 30000,
                 fill_color=tint(accent, 0.88))

        # Oversized number — bleeds partially over the edge
        num_left = base_left - 100000
        add_textbox(slide, str(i + 1), num_left, tier_top - 50000,
                    600000, tier_h + 20000,
                    font_size=52, bold=True, color=tint(accent, 0.6),
                    alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Label — bold and large
        add_textbox(slide, label, base_left + 500000, tier_top + 30000,
                    tier_w - 1200000, tier_h // 2 - 30000,
                    font_size=16, bold=True, color=t.primary,
                    alignment=PP_ALIGN.LEFT)

        # Value badge
        val_left = base_left + tier_w - 1500000
        val_top = tier_top + 40000
        add_rect(slide, val_left, val_top, 1200000, 350000,
                 fill_color=accent, corner_radius=60000)
        add_textbox(slide, value, val_left + 40000, val_top + 20000,
                    1120000, 310000,
                    font_size=16, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Description
        add_textbox(slide, desc, base_left + 500000, tier_top + tier_h // 2 + 10000,
                    tier_w - 1200000, tier_h // 2 - 60000,
                    font_size=10, color=t.secondary,
                    alignment=PP_ALIGN.LEFT)


# ── Variant 12: SCHOLARLY — figure-captioned, thin rules, centered ─────

def _scholarly_funnel(slide, t, c):
    """Academic figure style: centered funnel with numbered figure caption,
    thin horizontal rules between tiers, serif-like spacing, footnote references."""
    content_top = add_slide_title(slide, c.get("funnel_title", "Conversion Funnel"), theme=t)
    stages = _get_stages(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(stages)
    # Reserve space for figure caption at top and footnote at bottom
    caption_h = 350000
    footnote_h = 300000
    available_h = FOOTER_TOP - content_top - caption_h - footnote_h - 400000
    tier_h = min(available_h // n, 550000)
    funnel_top = content_top + caption_h + 100000

    # Figure caption (e.g., "Figure 1: Conversion pipeline analysis")
    add_textbox(slide, f"Figure 1: {c.get('funnel_title', 'Conversion Funnel')} Analysis",
                m, content_top + 80000,
                SLIDE_W - 2 * m, 280000,
                font_size=11, italic=True, color=t.secondary,
                alignment=PP_ALIGN.CENTER)

    # Thin top rule for the figure box
    add_line(slide, m + 1500000, funnel_top - 40000,
             SLIDE_W - m - 1500000, funnel_top - 40000,
             color=tint(t.secondary, 0.7), width=0.5)

    max_w = SLIDE_W - 2 * m - 2400000  # generous side margins for scholarly feel

    for i, (label, value, desc) in enumerate(stages):
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.5
        tier_w = int(max_w * fraction)
        tier_left = (SLIDE_W - tier_w) // 2
        tier_top = funnel_top + i * tier_h

        # Thin rule above each tier (except the first, handled above)
        if i > 0:
            rule_w = tier_w + 200000
            rule_left = (SLIDE_W - rule_w) // 2
            add_line(slide, rule_left, tier_top,
                     rule_left + rule_w, tier_top,
                     color=tint(t.secondary, 0.75), width=0.5)

        # Subtle fill — very light tint
        add_rect(slide, tier_left, tier_top + 10000, tier_w, tier_h - 20000,
                 fill_color=tint(t.accent, 0.92))

        # Centered label
        add_textbox(slide, label, tier_left, tier_top + 20000,
                    tier_w, tier_h * 40 // 100,
                    font_size=13, bold=True, color=t.primary,
                    alignment=PP_ALIGN.CENTER)

        # Centered value with parenthetical notation: "n = 10,000"
        add_textbox(slide, f"n = {value}", tier_left, tier_top + tier_h * 38 // 100,
                    tier_w, tier_h * 35 // 100,
                    font_size=11, color=t.accent,
                    alignment=PP_ALIGN.CENTER)

        # Description as small italic note
        add_textbox(slide, desc, tier_left, tier_top + tier_h * 68 // 100,
                    tier_w, tier_h * 28 // 100,
                    font_size=9, italic=True, color=t.secondary,
                    alignment=PP_ALIGN.CENTER)

    # Thin bottom rule for the figure box
    bottom_rule_y = funnel_top + n * tier_h + 20000
    add_line(slide, m + 1500000, bottom_rule_y,
             SLIDE_W - m - 1500000, bottom_rule_y,
             color=tint(t.secondary, 0.7), width=0.5)

    # Footnote / source line
    add_textbox(slide, "Source: Internal analytics pipeline. Conversion rates calculated quarterly.",
                m, bottom_rule_y + 60000,
                SLIDE_W - 2 * m, 250000,
                font_size=8, italic=True, color=tint(t.secondary, 0.3),
                alignment=PP_ALIGN.CENTER)


# ── Variant 13: LABORATORY — dark bg, grid overlay, data-coded ─────────

def _laboratory_funnel(slide, t, c):
    """Data laboratory style: dark background, grid overlay lines, data labels
    with monospace-style formatting, evidence boxes alongside tiers."""
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("funnel_title", "Conversion Funnel"), theme=t)
    stages = _get_stages(c)
    p = t.palette
    m = MARGIN

    n = len(stages)
    available_h = FOOTER_TOP - content_top - 400000
    tier_h = min(available_h // n - 30000, 580000)
    tier_gap = 30000

    # Grid overlay — faint horizontal lines across the slide
    grid_start_y = content_top + 100000
    grid_end_y = FOOTER_TOP - 200000
    grid_step = (grid_end_y - grid_start_y) // 12
    for g in range(13):
        gy = grid_start_y + g * grid_step
        add_line(slide, m, gy, SLIDE_W - m, gy,
                 color=tint(t.primary, 0.08), width=0.25)

    # Grid overlay — faint vertical lines
    v_step = (SLIDE_W - 2 * m) // 10
    for v in range(11):
        vx = m + v * v_step
        add_line(slide, vx, grid_start_y, vx, grid_end_y,
                 color=tint(t.primary, 0.08), width=0.25)

    # Funnel area on the left 55%
    funnel_area_w = int((SLIDE_W - 2 * m) * 0.55)
    max_funnel_w = funnel_area_w - 200000
    funnel_center_x = m + funnel_area_w // 2

    # Evidence panel area on the right 40%
    evidence_left = m + funnel_area_w + 200000
    evidence_w = SLIDE_W - m - evidence_left

    # Panel header
    add_rect(slide, evidence_left, content_top + 100000, evidence_w, 300000,
             fill_color=tint(t.primary, 0.12))
    add_textbox(slide, "DATA EVIDENCE LOG", evidence_left + 60000, content_top + 130000,
                evidence_w - 120000, 240000,
                font_size=9, bold=True, color=tint(t.accent, 0.7),
                alignment=PP_ALIGN.LEFT)

    for i, (label, value, desc) in enumerate(stages):
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.6
        tier_w = int(max_funnel_w * fraction)
        tier_left = funnel_center_x - tier_w // 2
        tier_top = content_top + 250000 + i * (tier_h + tier_gap)
        accent = p[i % len(p)]

        # Tier with accent color
        add_funnel_tier(slide, tier_left, tier_top, tier_w, tier_h,
                        fill_color=accent, text="", font_size=12)

        # Data label inside — formatted like code: "STAGE_01"
        stage_code = f"STG_{i + 1:02d}"
        add_textbox(slide, stage_code, tier_left + 80000, tier_top + 15000,
                    tier_w - 160000, tier_h * 40 // 100,
                    font_size=8, bold=True, color=tint(accent, 0.7),
                    alignment=PP_ALIGN.LEFT)

        # Label inside tier
        add_textbox(slide, label, tier_left + 80000, tier_top + tier_h * 30 // 100,
                    tier_w - 160000, tier_h * 40 // 100,
                    font_size=12, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

        # Value inside tier
        add_textbox(slide, value, tier_left + 80000, tier_top + tier_h * 60 // 100,
                    tier_w - 160000, tier_h * 35 // 100,
                    font_size=14, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

        # Evidence box on the right panel
        ebox_top = content_top + 450000 + i * (tier_h + tier_gap)
        ebox_h = tier_h - 20000

        # Evidence box background
        add_rect(slide, evidence_left, ebox_top, evidence_w, ebox_h,
                 fill_color=tint(t.primary, 0.08),
                 line_color=tint(accent, 0.3), line_width=0.5)

        # Left accent stripe on evidence box
        add_rect(slide, evidence_left, ebox_top, 30000, ebox_h,
                 fill_color=accent)

        # Evidence entry: stage + value
        add_textbox(slide, f"{label}  |  {value}", evidence_left + 60000, ebox_top + 20000,
                    evidence_w - 100000, ebox_h * 45 // 100,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.LEFT)

        # Evidence entry: description + conversion note
        conv_pct = ""
        if i > 0:
            conv_pct = f"  [{label} conv.]"
        add_textbox(slide, f"{desc}{conv_pct}", evidence_left + 60000,
                    ebox_top + ebox_h * 45 // 100,
                    evidence_w - 100000, ebox_h * 50 // 100,
                    font_size=8, color=tint(t.primary, 0.45),
                    alignment=PP_ALIGN.LEFT)

    # Bottom data timestamp
    add_textbox(slide, "DATA SNAPSHOT: Q4 ANALYSIS  |  PIPELINE REV 2.1",
                m, FOOTER_TOP - 250000,
                SLIDE_W - 2 * m, 200000,
                font_size=7, bold=True, color=tint(t.primary, 0.25),
                alignment=PP_ALIGN.LEFT)


# ── Variant 14: DASHBOARD — dense, compact funnel with stats sidebar ──

def _dashboard_funnel(slide, t, c):
    """Dashboard style: dense layout with compact funnel tiers on the left,
    metric tiles on the right with conversion rates and delta indicators."""
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("funnel_title", "Conversion Funnel"), theme=t)
    stages = _get_stages(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if s.dark_mode else t.secondary
    card_fill = tint(t.primary, 0.12) if s.dark_mode else "#FFFFFF"
    tile_border = tint(t.primary, 0.25) if s.dark_mode else tint(t.secondary, 0.7)

    n = len(stages)
    available_h = FOOTER_TOP - content_top - 300000

    # Left area: compact funnel (55%)
    funnel_area_w = int((SLIDE_W - 2 * m) * 0.52)
    max_funnel_w = funnel_area_w - 100000
    funnel_center_x = m + funnel_area_w // 2

    # Right area: metric tiles (42%)
    tiles_left = m + funnel_area_w + 150000
    tiles_w = SLIDE_W - m - tiles_left

    # Compact tier sizing
    tier_h = min(available_h // n - 20000, 500000)
    tier_gap = 20000

    # Section header for metrics panel
    header_top = content_top + 50000
    add_rect(slide, tiles_left, header_top, tiles_w, 260000,
             fill_color=t.accent, corner_radius=s.card_radius)
    add_textbox(slide, "CONVERSION METRICS", tiles_left + 60000, header_top + 30000,
                tiles_w - 120000, 200000,
                font_size=10, bold=True, color="#FFFFFF",
                alignment=PP_ALIGN.CENTER)

    for i, (label, value, desc) in enumerate(stages):
        fraction = 1.0 - (i / max(n - 1, 1)) * 0.55
        tier_w = int(max_funnel_w * fraction)
        tier_left = funnel_center_x - tier_w // 2
        tier_top = content_top + 150000 + i * (tier_h + tier_gap)
        accent = p[i % len(p)]

        # Compact funnel tier
        add_rect(slide, tier_left, tier_top, tier_w, tier_h - 10000,
                 fill_color=accent, corner_radius=s.card_radius // 2)

        # Label inside
        add_textbox(slide, label, tier_left + 50000, tier_top + 10000,
                    tier_w - 100000, tier_h * 50 // 100,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

        # Value inside
        add_textbox(slide, value, tier_left + 50000, tier_top + tier_h * 45 // 100,
                    tier_w - 100000, tier_h * 45 // 100,
                    font_size=13, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER)

        # ── Metric tile on the right ──
        tile_top = content_top + 350000 + i * (tier_h + tier_gap)
        tile_h = tier_h - 30000

        # Tile background
        add_rect(slide, tiles_left, tile_top, tiles_w, tile_h,
                 fill_color=card_fill,
                 line_color=tile_border, line_width=0.5,
                 corner_radius=s.card_radius)

        # Tile: accent dot at left
        add_circle(slide, tiles_left + 80000, tile_top + tile_h // 2, 30000,
                   fill_color=accent)

        # Tile: label
        add_textbox(slide, label, tiles_left + 150000, tile_top + 15000,
                    tiles_w * 45 // 100, tile_h * 45 // 100,
                    font_size=9, bold=True, color=text_color,
                    alignment=PP_ALIGN.LEFT)

        # Tile: description
        add_textbox(slide, desc, tiles_left + 150000, tile_top + tile_h * 45 // 100,
                    tiles_w * 55 // 100, tile_h * 45 // 100,
                    font_size=7, color=sub_color,
                    alignment=PP_ALIGN.LEFT)

        # Tile: value (right-aligned, large)
        val_left = tiles_left + tiles_w * 55 // 100
        add_textbox(slide, value, val_left, tile_top + 10000,
                    tiles_w * 40 // 100, tile_h * 50 // 100,
                    font_size=14, bold=True, color=accent,
                    alignment=PP_ALIGN.RIGHT)

        # Tile: conversion rate indicator (percentage drop from previous)
        if i > 0:
            conv_label = f"Stage {i} \u2192 {i + 1}"
            add_textbox(slide, conv_label, val_left, tile_top + tile_h * 50 // 100,
                        tiles_w * 40 // 100, tile_h * 40 // 100,
                        font_size=7, color=sub_color,
                        alignment=PP_ALIGN.RIGHT)
        else:
            add_textbox(slide, "Top of funnel", val_left, tile_top + tile_h * 50 // 100,
                        tiles_w * 40 // 100, tile_h * 40 // 100,
                        font_size=7, color=sub_color,
                        alignment=PP_ALIGN.RIGHT)

    # Summary bar at bottom
    summary_top = FOOTER_TOP - 350000
    add_rect(slide, m, summary_top, SLIDE_W - 2 * m, 250000,
             fill_color=tint(t.primary, 0.08) if s.dark_mode else tint(t.secondary, 0.9),
             corner_radius=s.card_radius)
    add_textbox(slide, f"{n} stages  |  {stages[0][1]} \u2192 {stages[-1][1]}  |  Full pipeline view",
                m + 80000, summary_top + 30000,
                SLIDE_W - 2 * m - 160000, 190000,
                font_size=8, color=sub_color,
                alignment=PP_ALIGN.CENTER)
