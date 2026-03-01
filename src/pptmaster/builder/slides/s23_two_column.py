"""Slide 23 — Two Column: 2 equal columns with accent bars, 4 bullets each."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, FONT_SLIDE_TITLE, FONT_CARD_HEADING,
    FONT_BODY, CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_bullet_list, add_rect, add_line,
    add_styled_card, add_slide_title, add_dark_bg, add_background,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    c = t.content

    s = t.ux_style
    _STYLE_DISPATCH = {"scholarly": _scholarly, "laboratory": _laboratory, "dashboard": _dashboard}
    _fn = _STYLE_DISPATCH.get(s.name)
    if _fn:
        return _fn(slide, t)

    add_dark_bg(slide, t)

    content_top = add_slide_title(slide, c.get("two_col_title", "Strategic Priorities"), theme=t)

    s = t.ux_style
    text_color = "#FFFFFF" if s.dark_mode else t.primary
    sub_color = tint(t.primary, 0.6) if s.dark_mode else t.secondary

    columns = c.get("col2", [
        {"heading": "Short-Term Goals", "bullets": ["Goal 1", "Goal 2", "Goal 3", "Goal 4"]},
        {"heading": "Long-Term Vision", "bullets": ["Vision 1", "Vision 2", "Vision 3", "Vision 4"]},
    ])

    ct = content_top + 100000

    for i, col_data in enumerate(columns[:2]):
        left, width = col_span(2, i, gap=400000)
        add_styled_card(slide, left, ct, width, 4500000, theme=t, accent_color=t.accent)
        add_textbox(slide, col_data["heading"], left + 150000, ct + 100000, width - 300000, 500000,
                    font_size=FONT_CARD_HEADING, bold=True, color=text_color)
        add_bullet_list(slide, col_data["bullets"], left + 150000, ct + 700000,
                        width - 300000, 3600000, font_size=14, color=text_color,
                        bullet_color=t.accent, spacing=8)


# ── Scholarly variant: white bg, thin vertical rule between columns ───────

def _scholarly(slide, t):
    c = t.content
    add_background(slide, "#FFFFFF")
    content_top = add_slide_title(slide, c.get("two_col_title", "Strategic Priorities"), theme=t)

    columns = c.get("col2", [
        {"heading": "Short-Term Goals", "bullets": ["Goal 1", "Goal 2", "Goal 3", "Goal 4"]},
        {"heading": "Long-Term Vision", "bullets": ["Vision 1", "Vision 2", "Vision 3", "Vision 4"]},
    ])

    ct = content_top + 100000
    total_w = SLIDE_W - 2 * MARGIN
    col_w = (total_w - 400000) // 2  # gap = 400000
    col_h = 4200000

    for i, col_data in enumerate(columns[:2]):
        left = MARGIN + i * (col_w + 400000)
        add_textbox(slide, col_data["heading"], left, ct, col_w, 450000,
                    font_size=FONT_CARD_HEADING, bold=True, color=t.primary,
                    alignment=PP_ALIGN.CENTER)
        add_bullet_list(slide, col_data["bullets"], left + 100000, ct + 550000,
                        col_w - 200000, col_h - 550000,
                        font_size=14, color=t.primary, bullet_color=t.secondary, spacing=10)

    # Thin vertical rule between columns
    mid_x = MARGIN + col_w + 200000
    add_line(slide, mid_x, ct, mid_x, ct + col_h,
             color=tint(t.secondary, 0.7), width=0.5)


# ── Laboratory variant: dark bg, accent-bordered cards ────────────────────

def _laboratory(slide, t):
    c = t.content
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("two_col_title", "Strategic Priorities"), theme=t)

    columns = c.get("col2", [
        {"heading": "Short-Term Goals", "bullets": ["Goal 1", "Goal 2", "Goal 3", "Goal 4"]},
        {"heading": "Long-Term Vision", "bullets": ["Vision 1", "Vision 2", "Vision 3", "Vision 4"]},
    ])

    ct = content_top + 100000
    p = t.palette

    for i, col_data in enumerate(columns[:2]):
        left, width = col_span(2, i, gap=400000)
        accent = p[i % len(p)]
        card_h = 4400000

        # Dark card background
        add_rect(slide, left, ct, width, card_h,
                 fill_color=tint(t.primary, 0.08), corner_radius=60000)
        # Colored left accent border
        add_rect(slide, left, ct, 40000, card_h, fill_color=accent)

        add_textbox(slide, col_data["heading"], left + 200000, ct + 150000,
                    width - 350000, 450000,
                    font_size=FONT_CARD_HEADING, bold=True, color="#FFFFFF")
        add_bullet_list(slide, col_data["bullets"], left + 200000, ct + 700000,
                        width - 350000, card_h - 850000,
                        font_size=14, color=tint(t.primary, 0.6),
                        bullet_color=accent, spacing=8)


# ── Dashboard variant: header band, compact shadow cards ──────────────────

def _dashboard(slide, t):
    c = t.content
    add_background(slide, "#FFFFFF")

    # Accent header band
    add_rect(slide, 0, 0, SLIDE_W, 80000, fill_color=t.accent)

    content_top = add_slide_title(slide, c.get("two_col_title", "Strategic Priorities"), theme=t)

    columns = c.get("col2", [
        {"heading": "Short-Term Goals", "bullets": ["Goal 1", "Goal 2", "Goal 3", "Goal 4"]},
        {"heading": "Long-Term Vision", "bullets": ["Vision 1", "Vision 2", "Vision 3", "Vision 4"]},
    ])

    ct = content_top + 50000

    for i, col_data in enumerate(columns[:2]):
        left, width = col_span(2, i, gap=350000)
        card_h = 4400000

        # Subtle shadow (offset rect behind card)
        add_rect(slide, left + 20000, ct + 20000, width, card_h,
                 fill_color=tint(t.secondary, 0.85), corner_radius=40000)
        add_styled_card(slide, left, ct, width, card_h, theme=t, accent_color=t.accent)

        add_textbox(slide, col_data["heading"], left + 130000, ct + 80000,
                    width - 260000, 400000,
                    font_size=FONT_CARD_HEADING, bold=True, color=t.primary)
        add_bullet_list(slide, col_data["bullets"], left + 130000, ct + 550000,
                        width - 260000, card_h - 700000,
                        font_size=13, color=t.primary, bullet_color=t.accent, spacing=6)
