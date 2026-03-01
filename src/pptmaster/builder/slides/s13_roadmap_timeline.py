"""Slide 13 — Roadmap Timeline: 11 unique visual variants dispatched by UX style."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptmaster.builder.design_system import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, FOOTER_TOP, col_span,
)
from pptmaster.builder.helpers import (
    add_textbox, add_rect, add_circle, add_line, add_card,
    add_gold_accent_line, add_slide_title, add_dark_bg, add_background,
    add_hexagon, add_diamond, add_styled_card,
)
from pptmaster.assets.color_utils import tint, shade


def build(slide, *, theme=None) -> None:
    from pptmaster.builder.themes import DEFAULT_THEME
    t = theme or DEFAULT_THEME
    s = t.ux_style
    c = t.content

    dispatch = {
        "horizontal": _horizontal,
        "dot-line": _dot_line,
        "bold-blocks": _bold_blocks,
        "cascading-cards": _cascading_cards,
        "vertical-glow": _vertical_glow,
        "zigzag": _zigzag,
        "diamond": _diamond,
        "editorial-vertical": _editorial_vertical,
        "gradient-flow": _gradient_flow,
        "retro-vertical": _retro_vertical,
        "mosaic": _mosaic,
        "scholarly-timeline": _scholarly_timeline,
        "laboratory-timeline": _laboratory_timeline,
        "dashboard-timeline": _dashboard_timeline,
    }
    builder = dispatch.get(s.timeline, _horizontal)
    builder(slide, t, c)


def _get_milestones(c):
    """Extract milestone data from content dict with sensible defaults."""
    return c.get("milestones", [
        ("Q1 2026", "Foundation", "Core launch"), ("Q2 2026", "Growth", "Market entry"),
        ("Q3 2026", "Scale", "Enterprise"), ("Q4 2026", "Optimize", "Refine"),
        ("Q1 2027", "Expand", "International"),
    ])


# ── Variant 1: CLASSIC — horizontal line, gold markers, alternating cards ──

def _horizontal(slide, t, c):
    content_top = add_slide_title(slide, c.get("roadmap_title", "Strategic Roadmap"), theme=t)
    milestones = _get_milestones(c)

    n = len(milestones)
    timeline_y = content_top + 2000000
    total_w = SLIDE_W - 2 * MARGIN - 400000
    start_x = MARGIN + 200000
    spacing = total_w // (n - 1) if n > 1 else total_w
    marker_radius = 100000

    # Horizontal timeline spine
    add_line(slide, start_x, timeline_y, start_x + total_w, timeline_y,
             color=t.primary, width=2.5)

    for i, (date, title, desc) in enumerate(milestones):
        mx = start_x + i * spacing
        add_circle(slide, mx, timeline_y, marker_radius, fill_color=t.accent)

        is_above = i % 2 == 0
        card_w = min(spacing - 200000, 1800000) if spacing > 200000 else 800000
        card_h = 1200000

        if is_above:
            card_top = timeline_y - marker_radius - 150000 - card_h
            add_line(slide, mx, card_top + card_h, mx, timeline_y - marker_radius,
                     color=t.accent, width=1)
        else:
            card_top = timeline_y + marker_radius + 150000
            add_line(slide, mx, timeline_y + marker_radius, mx, card_top,
                     color=t.accent, width=1)

        card_left = mx - card_w // 2
        add_styled_card(slide, card_left, card_top, card_w, card_h,
                        theme=t, accent_color=t.accent)

        add_textbox(slide, date, card_left + 50000, card_top + 80000,
                    card_w - 100000, 250000,
                    font_size=10, bold=True, color=t.accent, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, title, card_left + 50000, card_top + 350000,
                    card_w - 100000, 350000,
                    font_size=14, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, card_left + 50000, card_top + 700000,
                    card_w - 100000, 450000,
                    font_size=10, color=t.secondary, alignment=PP_ALIGN.CENTER)


# ── Variant 2: MINIMAL — small dots on thin line, labels below ─────────

def _dot_line(slide, t, c):
    content_top = add_slide_title(slide, c.get("roadmap_title", "Strategic Roadmap"), theme=t)
    milestones = _get_milestones(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(milestones)
    total_w = SLIDE_W - 2 * m
    step_w = total_w // n
    line_y = content_top + 1400000
    dot_r = 70000

    # Thin horizontal line spanning first to last dot
    add_line(slide, m + step_w // 2, line_y,
             m + step_w * (n - 1) + step_w // 2, line_y,
             color=tint(t.secondary, 0.7), width=0.5)

    for i, (date, title, desc) in enumerate(milestones):
        cx = m + i * step_w + step_w // 2

        # Small dot
        add_circle(slide, cx, line_y, dot_r, fill_color=t.accent)

        # Date above dot
        add_textbox(slide, date, cx - step_w // 2 + 30000, line_y - dot_r - 400000,
                    step_w - 60000, 300000,
                    font_size=9, bold=True, color=t.accent, alignment=PP_ALIGN.CENTER)

        # Title below
        add_textbox(slide, title, cx - step_w // 2 + 30000, line_y + dot_r + 200000,
                    step_w - 60000, 350000,
                    font_size=13, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, desc, cx - step_w // 2 + 30000, line_y + dot_r + 600000,
                    step_w - 60000, 400000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.CENTER)


# ── Variant 3: BOLD — large colored blocks in a row, bold numbers ──────

def _bold_blocks(slide, t, c):
    content_top = add_slide_title(slide, c.get("roadmap_title", "Strategic Roadmap").upper(), theme=t)
    milestones = _get_milestones(c)
    p = t.palette

    n = len(milestones)
    gap = 80000
    total_w = SLIDE_W - 2 * MARGIN
    block_w = (total_w - gap * (n - 1)) // n
    block_h = FOOTER_TOP - content_top - 400000
    block_top = content_top + 200000

    for i, (date, title, desc) in enumerate(milestones):
        accent = p[i % len(p)]
        left = MARGIN + i * (block_w + gap)

        # Full-height colored block
        add_rect(slide, left, block_top, block_w, block_h, fill_color=accent)

        # Date in large text at top
        add_textbox(slide, date, left + 60000, block_top + 150000,
                    block_w - 120000, 500000,
                    font_size=28, bold=True, color="#FFFFFF", alignment=PP_ALIGN.LEFT)

        # Thick white accent bar
        add_rect(slide, left + 60000, block_top + 700000, block_w // 2, 50000,
                 fill_color="#FFFFFF")

        # Title
        add_textbox(slide, title.upper(), left + 60000, block_top + 900000,
                    block_w - 120000, 500000,
                    font_size=14, bold=True, color="#FFFFFF", alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, left + 60000, block_top + 1450000,
                    block_w - 120000, 500000,
                    font_size=11, color=tint(accent, 0.7), alignment=PP_ALIGN.LEFT)


# ── Variant 4: ELEVATED — cascading staircase of cards, left to right ──

def _cascading_cards(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("roadmap_title", "Strategic Roadmap"), theme=t)
    milestones = _get_milestones(c)
    p = t.palette

    n = len(milestones)
    total_w = SLIDE_W - 2 * MARGIN
    card_w = int(total_w * 0.22)
    card_h = 1300000
    x_step = (total_w - card_w) // max(n - 1, 1)
    available_h = FOOTER_TOP - content_top - card_h - 400000
    y_step = available_h // max(n - 1, 1)

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if t.ux_style.dark_mode else t.secondary

    for i, (date, title, desc) in enumerate(milestones):
        accent = p[i % len(p)]
        card_left = MARGIN + i * x_step
        card_top = content_top + 200000 + i * y_step

        # Connector line from previous card to this one
        if i > 0:
            prev_left = MARGIN + (i - 1) * x_step
            prev_top = content_top + 200000 + (i - 1) * y_step
            add_line(slide, prev_left + card_w, prev_top + card_h // 2,
                     card_left, card_top + card_h // 2,
                     color=tint(t.accent, 0.4 if not t.ux_style.dark_mode else 0.2),
                     width=1.5, dash="dash")

        # Card with shadow
        add_styled_card(slide, card_left, card_top, card_w, card_h,
                        theme=t, accent_color=accent)

        # Milestone number badge
        badge_r = 100000
        badge_x = card_left + card_w - 120000
        badge_y = card_top - 60000
        add_circle(slide, badge_x, badge_y, badge_r, fill_color=accent)
        add_textbox(slide, str(i + 1), badge_x - badge_r, badge_y - badge_r,
                    badge_r * 2, badge_r * 2,
                    font_size=12, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Date
        add_textbox(slide, date, card_left + 80000, card_top + 100000,
                    card_w - 160000, 280000,
                    font_size=10, bold=True, color=accent, alignment=PP_ALIGN.LEFT)

        # Title
        add_textbox(slide, title, card_left + 80000, card_top + 400000,
                    card_w - 160000, 400000,
                    font_size=13, bold=True, color=text_color, alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, card_left + 80000, card_top + 800000,
                    card_w - 160000, 400000,
                    font_size=10, color=sub_color, alignment=PP_ALIGN.LEFT)


# ── Variant 5: DARK — vertical timeline, glowing dots, cards to right ──

def _vertical_glow(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("roadmap_title", "Strategic Roadmap").upper(), theme=t)
    milestones = _get_milestones(c)
    p = t.palette

    n = len(milestones)
    line_x = MARGIN + 800000
    available_h = FOOTER_TOP - content_top - 400000
    step_h = available_h // n
    card_w = SLIDE_W - line_x - MARGIN - 400000
    card_h = min(step_h - 120000, 900000)
    dot_r = 80000
    glow_r = 160000

    for i, (date, title, desc) in enumerate(milestones):
        accent = p[i % len(p)]
        node_y = content_top + 200000 + i * step_h + card_h // 2

        # Vertical line segment
        if i < n - 1:
            next_y = content_top + 200000 + (i + 1) * step_h + card_h // 2
            add_line(slide, line_x, node_y + dot_r, line_x, next_y - dot_r,
                     color=tint(t.primary, 0.2), width=1.5)

        # Glowing outer ring
        add_circle(slide, line_x, node_y, glow_r, fill_color=tint(accent, 0.15))

        # Inner dot
        add_circle(slide, line_x, node_y, dot_r, fill_color=accent)

        # Date label to the left of the line
        add_textbox(slide, date, MARGIN, node_y - 150000, line_x - MARGIN - 200000, 300000,
                    font_size=10, bold=True, color=accent, alignment=PP_ALIGN.RIGHT)

        # Connector line to card
        card_left = line_x + 300000
        card_top = content_top + 200000 + i * step_h
        add_line(slide, line_x + glow_r, node_y, card_left, node_y,
                 color=tint(accent, 0.3), width=1)

        # Card (dark style: subtle border, low fill)
        card_fill = tint(t.primary, 0.12)
        add_rect(slide, card_left, card_top, card_w, card_h,
                 fill_color=card_fill, line_color=tint(t.primary, 0.25), line_width=0.5,
                 corner_radius=60000)

        # Accent bar on left edge of card
        add_rect(slide, card_left, card_top + 80000, 40000, card_h - 160000,
                 fill_color=accent)

        # Title inside card
        add_textbox(slide, title, card_left + 120000, card_top + 100000,
                    card_w - 200000, 380000,
                    font_size=14, bold=True, color="#FFFFFF", alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, card_left + 120000, card_top + 500000,
                    card_w - 200000, 350000,
                    font_size=11, color=tint(t.primary, 0.5), alignment=PP_ALIGN.LEFT)


# ── Variant 6: SPLIT — zigzag path, cards alternating left and right ───

def _zigzag(slide, t, c):
    content_top = add_slide_title(slide, c.get("roadmap_title", "Strategic Roadmap"), theme=t)
    milestones = _get_milestones(c)
    p = t.palette

    n = len(milestones)
    mid_x = SLIDE_W // 2
    available_h = FOOTER_TOP - content_top - 300000
    step_h = available_h // n
    card_w = int(SLIDE_W * 0.34)
    card_h = min(step_h - 100000, 800000)

    # Central vertical guide
    line_top = content_top + 200000
    line_bottom = line_top + n * step_h
    add_line(slide, mid_x, line_top, mid_x, line_bottom,
             color=tint(t.secondary, 0.5), width=1.5, dash="dash")

    for i, (date, title, desc) in enumerate(milestones):
        accent = p[i % len(p)]
        is_left = i % 2 == 0
        step_top = content_top + 200000 + i * step_h
        node_y = step_top + card_h // 2

        # Circle node at center
        add_circle(slide, mid_x, node_y, 90000, fill_color=accent)
        add_textbox(slide, str(i + 1), mid_x - 70000, node_y - 70000, 140000, 140000,
                    font_size=10, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        if is_left:
            card_left = mid_x - card_w - 200000
            add_line(slide, card_left + card_w, node_y, mid_x - 90000, node_y,
                     color=tint(accent, 0.5), width=1)
        else:
            card_left = mid_x + 200000
            add_line(slide, mid_x + 90000, node_y, card_left, node_y,
                     color=tint(accent, 0.5), width=1)

        add_styled_card(slide, card_left, step_top, card_w, card_h,
                        theme=t, accent_color=accent)

        # Date
        add_textbox(slide, date, card_left + 100000, step_top + 60000,
                    card_w - 200000, 250000,
                    font_size=10, bold=True, color=accent, alignment=PP_ALIGN.LEFT)

        # Title
        add_textbox(slide, title, card_left + 100000, step_top + 300000,
                    card_w - 200000, 300000,
                    font_size=13, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, card_left + 100000, step_top + 580000,
                    card_w - 200000, 250000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.LEFT)


# ── Variant 7: GEO — diamond-shaped markers on horizontal line ─────────

def _diamond(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("roadmap_title", "Strategic Roadmap"), theme=t)
    milestones = _get_milestones(c)
    p = t.palette

    n = len(milestones)
    total_w = SLIDE_W - 2 * MARGIN
    step_w = total_w // n
    line_y = content_top + 1400000
    diamond_size = 200000

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.5) if t.ux_style.dark_mode else t.secondary

    # Horizontal spine
    add_line(slide, MARGIN + step_w // 2, line_y,
             MARGIN + step_w * (n - 1) + step_w // 2, line_y,
             color=tint(t.accent, 0.3 if t.ux_style.dark_mode else 0.6), width=2)

    for i, (date, title, desc) in enumerate(milestones):
        accent = p[i % len(p)]
        cx = MARGIN + i * step_w + step_w // 2

        # Diamond marker
        add_diamond(slide, cx, line_y, diamond_size, fill_color=accent)

        # Number inside diamond
        add_textbox(slide, str(i + 1), cx - diamond_size, line_y - diamond_size,
                    diamond_size * 2, diamond_size * 2,
                    font_size=12, bold=True, color="#FFFFFF",
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Date above
        add_textbox(slide, date, cx - step_w // 2 + 20000, line_y - diamond_size - 450000,
                    step_w - 40000, 300000,
                    font_size=10, bold=True, color=accent, alignment=PP_ALIGN.CENTER)

        # Title below
        add_textbox(slide, title, cx - step_w // 2 + 20000, line_y + diamond_size + 200000,
                    step_w - 40000, 350000,
                    font_size=13, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, desc, cx - step_w // 2 + 20000, line_y + diamond_size + 600000,
                    step_w - 40000, 400000,
                    font_size=11, color=sub_color, alignment=PP_ALIGN.CENTER)

        # Small decorative diamond below description
        add_diamond(slide, cx, line_y + diamond_size + 1150000, 40000, fill_color=accent)


# ── Variant 8: EDITORIAL — vertical timeline, thin rules, dates left ───

def _editorial_vertical(slide, t, c):
    content_top = add_slide_title(slide, c.get("roadmap_title", "Strategic Roadmap"), theme=t)
    milestones = _get_milestones(c)
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(milestones)
    available_h = FOOTER_TOP - content_top - 400000
    step_h = available_h // n

    date_col_w = 1200000
    line_x = m + date_col_w + 200000
    content_left = line_x + 300000
    content_w = SLIDE_W - m - content_left

    for i, (date, title, desc) in enumerate(milestones):
        step_top = content_top + 200000 + i * step_h

        # Thin horizontal rule across full width
        add_line(slide, m, step_top, SLIDE_W - m, step_top,
                 color=tint(t.secondary, 0.8), width=0.5)

        # Short accent mark at line intersection
        add_line(slide, line_x - 50000, step_top, line_x + 50000, step_top,
                 color=t.accent, width=2)

        # Date in left column
        add_textbox(slide, date, m, step_top + 100000, date_col_w, 350000,
                    font_size=11, bold=True, color=t.accent, alignment=PP_ALIGN.RIGHT)

        # Title
        add_textbox(slide, title, content_left, step_top + 100000,
                    content_w, 350000,
                    font_size=14, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, content_left, step_top + 480000,
                    content_w, 350000,
                    font_size=12, color=t.secondary, alignment=PP_ALIGN.LEFT)

    # Vertical rule connecting intersections
    first_top = content_top + 200000
    last_top = content_top + 200000 + (n - 1) * step_h
    add_line(slide, line_x, first_top, line_x, last_top,
             color=tint(t.secondary, 0.6), width=0.75)

    # Bottom rule
    add_line(slide, m, FOOTER_TOP - 300000, SLIDE_W - m, FOOTER_TOP - 300000,
             color=tint(t.secondary, 0.8), width=0.5)


# ── Variant 9: GRADIENT — horizontal flow, gradient circles, soft cards ─

def _gradient_flow(slide, t, c):
    add_dark_bg(slide, t)
    content_top = add_slide_title(slide, c.get("roadmap_title", "Strategic Roadmap"), theme=t)
    milestones = _get_milestones(c)

    n = len(milestones)
    total_w = SLIDE_W - 2 * MARGIN
    step_w = total_w // n
    circle_y = content_top + 1000000
    circle_r = 280000

    text_color = "#FFFFFF" if t.ux_style.dark_mode else t.primary
    sub_color = tint(t.primary, 0.55) if t.ux_style.dark_mode else t.secondary

    for i, (date, title, desc) in enumerate(milestones):
        cx = MARGIN + i * step_w + step_w // 2
        # Progressive tint from accent
        factor = i * 0.12
        circle_color = tint(t.accent, factor)

        # Connector line
        if i < n - 1:
            next_cx = MARGIN + (i + 1) * step_w + step_w // 2
            add_line(slide, cx + circle_r, circle_y, next_cx - circle_r, circle_y,
                     color=tint(t.accent, 0.5), width=1)

        # Soft glow halo
        add_circle(slide, cx, circle_y, circle_r + 80000,
                   fill_color=tint(circle_color, 0.85))

        # Main circle
        add_circle(slide, cx, circle_y, circle_r, fill_color=circle_color)

        # Date inside top of circle
        add_textbox(slide, date, cx - circle_r + 30000, circle_y - circle_r + 60000,
                    circle_r * 2 - 60000, 200000,
                    font_size=8, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER)

        # Title centered in circle
        add_textbox(slide, title, cx - circle_r + 30000, circle_y - 80000,
                    circle_r * 2 - 60000, 250000,
                    font_size=12, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER)

        # Card below circle
        card_left = cx - step_w // 2 + 30000
        card_w = step_w - 60000
        card_top = circle_y + circle_r + 200000
        card_h = 1300000

        add_styled_card(slide, card_left, card_top, card_w, card_h,
                        theme=t, accent_color=circle_color)

        card_text_color = text_color
        add_textbox(slide, desc, card_left + 60000, card_top + 100000,
                    card_w - 120000, 500000,
                    font_size=11, color=sub_color, alignment=PP_ALIGN.CENTER)


# ── Variant 10: RETRO — vertical, decorated badges, dotted connections ──

def _retro_vertical(slide, t, c):
    add_background(slide, t.light_bg)
    content_top = add_slide_title(slide, c.get("roadmap_title", "Strategic Roadmap"), theme=t)
    milestones = _get_milestones(c)
    p = t.palette

    n = len(milestones)
    badge_x = MARGIN + 600000
    available_h = FOOTER_TOP - content_top - 400000
    step_h = available_h // n
    badge_r = 200000
    card_left = badge_x + 500000
    card_w = SLIDE_W - MARGIN - card_left
    card_h = min(step_h - 100000, 900000)

    for i, (date, title, desc) in enumerate(milestones):
        accent = p[i % len(p)]
        node_y = content_top + 200000 + i * step_h + card_h // 2

        # Dotted connector to next badge
        if i < n - 1:
            next_y = content_top + 200000 + (i + 1) * step_h + card_h // 2
            add_line(slide, badge_x, node_y + badge_r + 20000,
                     badge_x, next_y - badge_r - 20000,
                     color=t.accent, width=1.5, dash="dash")

        # Outer badge border
        add_circle(slide, badge_x, node_y, badge_r,
                   fill_color="#FFFFFF", line_color=accent, line_width=2.5)

        # Inner badge border
        add_circle(slide, badge_x, node_y, badge_r - 40000,
                   fill_color="#FFFFFF", line_color=accent, line_width=0.75)

        # Number inside badge
        add_textbox(slide, str(i + 1), badge_x - badge_r, node_y - badge_r,
                    badge_r * 2, badge_r * 2,
                    font_size=14, bold=True, color=accent,
                    alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Connector line to card
        card_top = content_top + 200000 + i * step_h
        add_line(slide, badge_x + badge_r + 20000, node_y, card_left, node_y,
                 color=tint(accent, 0.5), width=1, dash="dash")

        # Card
        add_styled_card(slide, card_left, card_top, card_w, card_h,
                        theme=t, accent_color=accent)

        # Date
        add_textbox(slide, date, card_left + 80000, card_top + 60000,
                    card_w - 160000, 250000,
                    font_size=10, bold=True, color=accent, alignment=PP_ALIGN.LEFT)

        # Title
        add_textbox(slide, title, card_left + 80000, card_top + 320000,
                    card_w - 160000, 300000,
                    font_size=13, bold=True, color=t.primary, alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, card_left + 80000, card_top + 620000,
                    card_w - 160000, 280000,
                    font_size=11, color=t.secondary, alignment=PP_ALIGN.LEFT)

        # Decorative triple dots after card
        dot_y = card_top + card_h - 100000
        for d in range(3):
            add_circle(slide, card_left + card_w - 200000 + d * 60000, dot_y, 15000,
                       fill_color=accent)


# ── Variant 11: MAGAZINE — mosaic grid layout, 3+2 cards, no timeline ──

def _mosaic(slide, t, c):
    content_top = add_slide_title(slide, c.get("roadmap_title", "Strategic Roadmap"), theme=t)
    milestones = _get_milestones(c)
    p = t.palette
    s = t.ux_style
    m = int(MARGIN * s.margin_factor)

    n = len(milestones)
    gap = int(100000 * s.gap_factor)
    total_w = SLIDE_W - 2 * m

    # Row 1: up to 3 cards
    row1_count = min(n, 3)
    row1_w = (total_w - gap * (row1_count - 1)) // row1_count
    row1_h = 1800000
    row1_top = content_top + 200000

    # Row 2: remaining cards
    row2_count = min(n - row1_count, 3)
    row2_w = (total_w - gap * max(row2_count - 1, 0)) // max(row2_count, 1) if row2_count else 0
    row2_h = 1800000
    row2_top = row1_top + row1_h + gap

    for i, (date, title, desc) in enumerate(milestones):
        accent = p[i % len(p)]

        if i < row1_count:
            card_left = m + i * (row1_w + gap)
            card_top = row1_top
            card_w = row1_w
            card_h = row1_h
        elif i < row1_count + row2_count:
            j = i - row1_count
            # Center row2 if fewer cards
            row2_total = row2_count * row2_w + (row2_count - 1) * gap
            row2_start = m + (total_w - row2_total) // 2
            card_left = row2_start + j * (row2_w + gap)
            card_top = row2_top
            card_w = row2_w
            card_h = row2_h
        else:
            continue

        # Full-color card (no standard card — pure mosaic block)
        add_rect(slide, card_left, card_top, card_w, card_h,
                 fill_color=accent, corner_radius=0)

        # Large milestone number
        add_textbox(slide, f"{i + 1:02d}", card_left + 80000, card_top + 100000,
                    card_w - 160000, 600000,
                    font_size=36, bold=True, color=tint(accent, 0.3), alignment=PP_ALIGN.LEFT)

        # Date
        add_textbox(slide, date, card_left + 80000, card_top + 650000,
                    card_w - 160000, 280000,
                    font_size=10, bold=True, color="#FFFFFF", alignment=PP_ALIGN.LEFT)

        # Title
        add_textbox(slide, title.upper(), card_left + 80000, card_top + 950000,
                    card_w - 160000, 400000,
                    font_size=15, bold=True, color="#FFFFFF", alignment=PP_ALIGN.LEFT)

        # Description
        add_textbox(slide, desc, card_left + 80000, card_top + 1350000,
                    card_w - 160000, 400000,
                    font_size=11, color=tint(accent, 0.7), alignment=PP_ALIGN.LEFT)


# ── Variant 12: SCHOLARLY — white bg, horizontal thin line, minimal dots ─

def _scholarly_timeline(slide, t, c):
    add_background(slide, "#FFFFFF")
    content_top = add_slide_title(slide, c.get("roadmap_title", "Timeline"), theme=t)
    milestones = _get_milestones(c)

    n = len(milestones)
    m = MARGIN
    total_w = SLIDE_W - 2 * m - 400000
    start_x = m + 200000
    spacing = total_w // (n - 1) if n > 1 else total_w
    line_y = content_top + 2000000
    dot_r = 50000

    # Thin horizontal timeline
    add_line(slide, start_x, line_y, start_x + total_w, line_y,
             color=tint(t.secondary, 0.5), width=0.75)

    for i, (date, title, desc) in enumerate(milestones):
        mx = start_x + i * spacing

        # Small dot marker
        add_circle(slide, mx, line_y, dot_r,
                   fill_color=t.primary)

        # Date above the line
        add_textbox(slide, date, mx - 600000, line_y - dot_r - 500000,
                    1200000, 350000,
                    font_size=10, bold=True, color=t.accent, alignment=PP_ALIGN.CENTER)

        # Title below the line
        add_textbox(slide, title, mx - 600000, line_y + dot_r + 200000,
                    1200000, 350000,
                    font_size=12, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)

        # Description further below
        add_textbox(slide, desc, mx - 600000, line_y + dot_r + 580000,
                    1200000, 350000,
                    font_size=10, color=t.secondary, alignment=PP_ALIGN.CENTER)

    # Bottom thin rule
    add_line(slide, m, FOOTER_TOP - 350000, SLIDE_W - m, FOOTER_TOP - 350000,
             color=tint(t.secondary, 0.6), width=0.5)


# ── Variant 13: LABORATORY — dark bg, glowing dots, dark cards below ─────

def _laboratory_timeline(slide, t, c):
    add_background(slide, t.primary)
    content_top = add_slide_title(slide, c.get("roadmap_title", "Research Timeline").upper(), theme=t)
    milestones = _get_milestones(c)
    p = t.palette

    n = len(milestones)
    total_w = SLIDE_W - 2 * MARGIN - 400000
    start_x = MARGIN + 200000
    spacing = total_w // (n - 1) if n > 1 else total_w
    track_y = content_top + 1000000
    dot_r = 80000
    glow_r = 160000

    # Horizontal track line
    add_line(slide, start_x, track_y, start_x + total_w, track_y,
             color=tint(t.primary, 0.2), width=2)

    for i, (date, title, desc) in enumerate(milestones):
        accent = p[i % len(p)]
        mx = start_x + i * spacing

        # Glowing outer ring
        add_circle(slide, mx, track_y, glow_r, fill_color=tint(accent, 0.15))

        # Inner accent dot
        add_circle(slide, mx, track_y, dot_r, fill_color=accent)

        # Date above dot
        add_textbox(slide, date, mx - 700000, track_y - glow_r - 400000,
                    1400000, 300000,
                    font_size=10, bold=True, color=accent, alignment=PP_ALIGN.CENTER)

        # Dark card below milestone
        card_w = min(spacing - 100000, 1800000) if spacing > 100000 else 800000
        card_h = 1400000
        card_left = mx - card_w // 2
        card_top = track_y + glow_r + 200000
        card_fill = tint(t.primary, 0.1)

        add_rect(slide, card_left, card_top, card_w, card_h,
                 fill_color=card_fill, line_color=tint(t.primary, 0.2),
                 line_width=0.5, corner_radius=50000)

        # Accent bar on top of card
        add_rect(slide, card_left + 1, card_top, card_w - 2, 35000,
                 fill_color=accent)

        # Title in card
        add_textbox(slide, title, card_left + 60000, card_top + 120000,
                    card_w - 120000, 400000,
                    font_size=13, bold=True, color="#FFFFFF", alignment=PP_ALIGN.CENTER)

        # Description in card
        add_textbox(slide, desc, card_left + 60000, card_top + 550000,
                    card_w - 120000, 400000,
                    font_size=10, color=tint(t.primary, 0.45), alignment=PP_ALIGN.CENTER)

        # Connector line from dot to card
        add_line(slide, mx, track_y + glow_r, mx, card_top,
                 color=tint(accent, 0.25), width=1)


# ── Variant 14: DASHBOARD — dense timeline, compact cards, header band ───

def _dashboard_timeline(slide, t, c):
    add_background(slide, "#FFFFFF")
    content_top = add_slide_title(slide, c.get("roadmap_title", "Strategic Roadmap"), theme=t)
    milestones = _get_milestones(c)
    p = t.palette

    # Header accent band
    band_h = 45000
    add_rect(slide, 0, content_top, SLIDE_W, band_h, fill_color=t.accent)

    n = len(milestones)
    m = int(MARGIN * 0.7)
    total_w = SLIDE_W - 2 * m - 200000
    start_x = m + 100000
    spacing = total_w // (n - 1) if n > 1 else total_w
    line_y = content_top + band_h + 600000
    dot_r = 70000

    # Timeline spine
    add_line(slide, start_x, line_y, start_x + total_w, line_y,
             color=tint(t.secondary, 0.5), width=1.5)

    for i, (date, title, desc) in enumerate(milestones):
        accent = p[i % len(p)]
        mx = start_x + i * spacing

        # Colored dot
        add_circle(slide, mx, line_y, dot_r, fill_color=accent)

        # Date above
        add_textbox(slide, date, mx - 500000, line_y - dot_r - 350000,
                    1000000, 250000,
                    font_size=9, bold=True, color=accent, alignment=PP_ALIGN.CENTER)

        # Compact card below — clamp to slide bounds
        card_w = min(spacing - 60000, 1600000) if spacing > 60000 else 700000
        card_h = 1000000
        card_left = mx - card_w // 2
        # Ensure card doesn't extend past right edge
        if card_left + card_w > SLIDE_W - m:
            card_left = SLIDE_W - m - card_w
        # Ensure card doesn't extend past left edge
        if card_left < m:
            card_left = m
        card_top = line_y + dot_r + 150000

        add_rect(slide, card_left, card_top, card_w, card_h,
                 fill_color=tint(accent, 0.92), corner_radius=40000)

        # Left accent bar
        add_rect(slide, card_left, card_top + 20000, 25000, card_h - 40000,
                 fill_color=accent)

        # Connector from dot to card
        add_line(slide, mx, line_y + dot_r, mx, card_top,
                 color=tint(accent, 0.5), width=1)

        # Title
        add_textbox(slide, title, card_left + 50000, card_top + 60000,
                    card_w - 100000, 350000,
                    font_size=11, bold=True, color=t.primary, alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, desc, card_left + 50000, card_top + 420000,
                    card_w - 100000, 350000,
                    font_size=9, color=t.secondary, alignment=PP_ALIGN.CENTER)
